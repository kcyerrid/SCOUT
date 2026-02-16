import json
import copy
import os
import re
import sys
import subprocess
import shutil
import tkinter as tk
from tkinter import filedialog, messagebox, simpledialog, ttk
from pathlib import Path
from datetime import date, datetime, timedelta
import webbrowser
from urllib.parse import quote, unquote, urlparse
from PIL import Image, ImageTk
import sqlite3
import hashlib
from dataclasses import dataclass, asdict
try:
    import requests  # type: ignore
    from readability import Document  # type: ignore
    from bs4 import BeautifulSoup  # type: ignore
except Exception:
    requests = None
    Document = None
    BeautifulSoup = None
try:
    from tkinterdnd2 import DND_FILES, DND_TEXT, TkinterDnD  # type: ignore
except Exception:
    DND_FILES = None
    DND_TEXT = None
    TkinterDnD = None
try:
    from tkcalendar import DateEntry  # type: ignore
except ImportError:
    DateEntry = None

BaseTk = TkinterDnD.Tk if (TkinterDnD and hasattr(TkinterDnD, "Tk")) else tk.Tk

# ----------------------------
# Helpers: config IO + keywords schema
# ----------------------------
_LIST_NAME_RE = re.compile(r"^[a-z0-9_]+$")

# ----------------------------
# PPTX generation engine
# ----------------------------
PPTX_PRESENTATION_TYPES = {
    "threat_intel": {
        "label": "Threat Intel",
        "entity_types": {"threat_actor", "malware", "campaign"},
        "include_title_slide": True,
        "include_unmapped": True,
        "max_bullets": 6,
        "notes_mode": "overflow_only",
        "title_keys": [
            "campaign_name",
            "malware_name",
            "common_name",
            "actor_name",
            "technique_name",
            "report_title",
            "title",
        ],
        "sections": [
            {"title": "Executive Summary", "match": ["executive summary", "campaign overview"]},
            {"title": "Attribution", "match": ["attribution", "attribution notes"]},
            {"title": "Background & History", "match": ["background", "history"]},
            {"title": "Campaign Timeline", "match": ["campaign timeline", "timeline"]},
            {"title": "Targeting Profile", "match": ["targeting profile", "targets", "target sectors", "target regions"]},
            {"title": "TTPs & ATT&CK Mapping", "match": ["tactics", "techniques", "procedures", "ttp", "att&ck"]},
            {"title": "Malware & Tools", "match": ["malware & tools", "malware", "tools used", "tooling"]},
            {"title": "Infrastructure", "match": ["infrastructure", "c2", "command & control"]},
            {"title": "Defensive Guidance", "match": ["defensive recommendations", "defensive guidance", "response"]},
            {"title": "Analyst Notes", "match": ["analyst notes", "notes"]},
            {"title": "References", "match": ["references"]},
        ],
    }
}


def _pptx_parse_frontmatter(text: str):
    stripped = text.lstrip()
    if not stripped.startswith("---"):
        return {}, text
    lines = text.splitlines()
    if not lines or lines[0].strip() != "---":
        return {}, text
    fm_lines = []
    body_start = None
    for i in range(1, len(lines)):
        if lines[i].strip() == "---":
            body_start = i + 1
            break
        fm_lines.append(lines[i])
    if body_start is None:
        return {}, text
    fm = {}
    for line in fm_lines:
        if ":" not in line:
            continue
        key, val = line.split(":", 1)
        key = key.strip()
        if not key:
            continue
        val = val.strip().strip('"').strip("'")
        if val != "":
            fm[key] = val
    return fm, "\n".join(lines[body_start:])


def _pptx_normalize_heading(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", (value or "").lower()).strip()


def _pptx_split_sections(body: str, fallback_title: str):
    lines = body.splitlines()
    sections = []
    current_title = ""
    current_level = 1
    current_body = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("#"):
            heading = stripped.lstrip("#").strip()
            level = len(stripped) - len(stripped.lstrip("#"))
            if current_title or current_body:
                sections.append(
                    {"title": current_title or fallback_title, "level": current_level, "lines": current_body}
                )
                current_body = []
            current_title = heading
            current_level = level
        else:
            current_body.append(line)
    if current_title or current_body:
        sections.append({"title": current_title or fallback_title, "level": current_level, "lines": current_body})
    if not sections:
        sections = [{"title": fallback_title, "level": 1, "lines": lines}]
    return sections


def _pptx_body_to_bullets(lines):
    bullets = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith(("-", "*")):
            bullets.append(stripped[1:].strip())
        else:
            m = re.match(r"^\d+\.\s+(.*)$", stripped)
            bullets.append(m.group(1) if m else stripped)
    return bullets


def _pptx_find_layout(prs, prefer_names=None, require_title=False, require_body=False, exclude_names=None):
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
    except Exception:
        return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    prefer_names = [n.lower() for n in (prefer_names or [])]
    exclude_names = [n.lower() for n in (exclude_names or [])]

    for layout in prs.slide_layouts:
        name = (getattr(layout, "name", "") or "").strip().lower()
        if name:
            if any(ex in name for ex in exclude_names):
                continue
            if prefer_names and any(pref in name for pref in prefer_names):
                return layout

    for layout in prs.slide_layouts:
        has_title = False
        has_body = False
        for ph in layout.placeholders:
            pht = ph.placeholder_format.type
            if pht == PP_PLACEHOLDER.TITLE:
                has_title = True
            if pht == PP_PLACEHOLDER.BODY:
                has_body = True
        if (not require_title or has_title) and (not require_body or has_body):
            name = (getattr(layout, "name", "") or "").strip().lower()
            if name and any(ex in name for ex in exclude_names):
                continue
            return layout

    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]


def _pptx_infer_title(frontmatter: dict, fallback: str, type_cfg: dict):
    for key in type_cfg.get("title_keys", []):
        val = (frontmatter.get(key) or "").strip()
        if val:
            return val
    return fallback


def _pptx_detect_type(note_path: Path):
    try:
        text = note_path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return None
    frontmatter, _ = _pptx_parse_frontmatter(text)
    entity_type = (frontmatter.get("entity_type") or "").strip().lower()
    for key, cfg in PPTX_PRESENTATION_TYPES.items():
        if entity_type and entity_type in cfg.get("entity_types", set()):
            return key
    return None


def _pptx_add_slide(prs, layout, title, body_lines, max_bullets, notes_text, theme_color=None):
    from pptx.util import Inches
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.enum.text import MSO_AUTO_SIZE

    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title

    body_placeholder = None
    for ph in slide.placeholders:
        pht = ph.placeholder_format.type
        if pht == PP_PLACEHOLDER.BODY:
            body_placeholder = ph
            break

    if body_placeholder is None:
        left = Inches(1.0)
        top = Inches(1.8)
        width = Inches(8.3)
        height = Inches(4.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
    else:
        tf = body_placeholder.text_frame
        tf.clear()

    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    bullets = _pptx_body_to_bullets(body_lines)
    overflow = []
    if max_bullets and len(bullets) > max_bullets:
        overflow = bullets[max_bullets:]
        bullets = bullets[:max_bullets]

    if bullets:
        tf.text = bullets[0]
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
    else:
        tf.text = ""

    if notes_text:
        if overflow:
            notes_text = notes_text + "\n\nAdditional bullets:\n" + "\n".join(overflow)
        try:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.clear()
            notes_tf.text = notes_text
        except Exception:
            pass
    elif overflow:
        try:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.clear()
            notes_tf.text = "Additional bullets:\n" + "\n".join(overflow)
        except Exception:
            pass

    return slide


def _pptx_build_presentation(note_path: Path, template_path: Path, output_path: Path, presentation_type: str):
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError("python-pptx is required. Install it with: pip install python-pptx") from e

    text = note_path.read_text(encoding="utf-8", errors="ignore")
    frontmatter, body = _pptx_parse_frontmatter(text)
    sections = _pptx_split_sections(body, note_path.stem)
    normalized_sections = [
        {"key": _pptx_normalize_heading(sec["title"]), "section": sec}
        for sec in sections
    ]

    type_cfg = PPTX_PRESENTATION_TYPES.get(presentation_type) or PPTX_PRESENTATION_TYPES["threat_intel"]
    title_text = _pptx_infer_title(frontmatter, note_path.stem, type_cfg)

    prs = Presentation(str(template_path))
    title_layout = _pptx_find_layout(
        prs,
        prefer_names=["title slide"],
        require_title=True,
        exclude_names=["section"],
    )
    content_layout = _pptx_find_layout(
        prs,
        prefer_names=["title and content", "title content"],
        require_title=True,
        require_body=True,
        exclude_names=["section"],
    )

    if type_cfg.get("include_title_slide", True):
        _pptx_add_slide(
            prs,
            title_layout,
            title_text,
            body_lines=[],
            max_bullets=0,
            notes_text="",
        )

    used_indexes = set()
    for rule in type_cfg.get("sections", []):
        matchers = [ _pptx_normalize_heading(m) for m in rule.get("match", []) ]
        found = None
        for idx, item in enumerate(normalized_sections):
            if idx in used_indexes:
                continue
            key = item["key"]
            if any(m in key for m in matchers):
                found = (idx, item["section"])
                break
        if found:
            idx, sec = found
            used_indexes.add(idx)
            notes_mode = type_cfg.get("notes_mode", "full_body")
            notes_text = ""
            if notes_mode == "full_body":
                notes_text = "\n".join([line for line in sec["lines"] if line.strip()])
            _pptx_add_slide(
                prs,
                content_layout,
                rule.get("title") or sec["title"],
                sec["lines"],
                max_bullets=type_cfg.get("max_bullets", 6),
                notes_text=notes_text,
            )

    if type_cfg.get("include_unmapped", True):
        for idx, item in enumerate(normalized_sections):
            if idx in used_indexes:
                continue
            sec = item["section"]
            notes_mode = type_cfg.get("notes_mode", "full_body")
            notes_text = ""
            if notes_mode == "full_body":
                notes_text = "\n".join([line for line in sec["lines"] if line.strip()])
            _pptx_add_slide(
                prs,
                content_layout,
                sec["title"],
                sec["lines"],
                max_bullets=type_cfg.get("max_bullets", 6),
                notes_text=notes_text,
            )

    temp_path = output_path.with_suffix(output_path.suffix + ".tmp")
    prs.save(str(temp_path))
    try:
        os.replace(temp_path, output_path)
        return output_path
    except PermissionError:
        timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        fallback = output_path.with_name(f"{output_path.stem} - Export {timestamp}{output_path.suffix}")
        os.replace(temp_path, fallback)
        return fallback
    except Exception:
        try:
            if temp_path.exists():
                temp_path.unlink()
        except Exception:
            pass
        raise


def _config_path_for_launcher(app_dir: Path) -> Path:
    return app_dir / "config.json"


def _read_config_file(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def _write_config_file_with_backup(path: Path, cfg: dict) -> None:
    # Backup first
    if path.exists():
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        bak = path.with_name(f"{path.stem}.bak.{ts}{path.suffix}")
        try:
            bak.write_text(path.read_text(encoding="utf-8"), encoding="utf-8")
        except Exception:
            # best effort; do not block save
            pass
    path.write_text(json.dumps(cfg, indent=2), encoding="utf-8")


def _ensure_keywords_schema(cfg: dict) -> dict:
    """
    Ensures cfg["rss"]["keywords"] exists with safe defaults.
    Returns the keywords dict.
    """
    cfg.setdefault("rss", {})
    rss = cfg["rss"]
    rss.setdefault("keywords", {})
    kw = rss["keywords"]

    # lists is dict[str, list[str]]
    if "lists" not in kw or not isinstance(kw.get("lists"), dict):
        kw["lists"] = {}

    # global behavior
    if kw.get("match_mode") not in ("any", "all"):
        kw["match_mode"] = "any"
    if not isinstance(kw.get("case_sensitive"), bool):
        kw["case_sensitive"] = False
    if not isinstance(kw.get("min_hits"), int) or kw.get("min_hits", 1) < 1:
        kw["min_hits"] = 1

    # optional list->tags mapping
    if "keyword_tag_map" not in rss or not isinstance(rss.get("keyword_tag_map"), dict):
        rss["keyword_tag_map"] = {}

    return kw


# ----------------------------
# Shift log rules schema
# ----------------------------
# RULE: Every time a new entity is created, an entry MUST be appended to the daily note's
# Shift Log. When adding a new entity type: 1) Add it to SHIFT_LOG_ENTITY_REGISTRY below,
# 2) Call _log_shift_entry_from_launcher() in the entity creation flow.
SHIFT_LOG_ENTITY_REGISTRY = {
    "incident": {
        "label": "Incident",
        "date_field": "created",
        "id_keys": ["incident_id"],
        "title_keys": ["title"],
        "header_template": "### {time} — {incident_id} - {title}",
        "body_template": "\n".join(
            [
                "Context:",
                "",
                "- What triggered this work?",
                "- Why did you look into it?",
                "",
                "Actions:",
                "",
                "- What did you do?",
                "- Queries executed?",
                "- Tools used?",
                "",
                "Findings:",
                "",
                "- What did you learn?",
                "- Any anomalies or notable patterns?",
                "",
                "Links:",
                "- [Alert link]",
                "- [Incident link]",
                "- [Investigation link]",
                "- [Observation link]",
            ]
        ),
    },
    "meeting": {
        "label": "Meeting",
        "date_field": "start_time",
        "id_keys": [],
        "title_keys": ["title"],
        "header_template": "### {created_time} - New Meeting Created - {title} - {start_datetime} - {end_datetime}",
        "body_template": "",
    },
    "playbook": {
        "label": "Playbook",
        "date_field": "created",
        "id_keys": ["playbook_id"],
        "title_keys": ["title"],
        "header_template": "### {time} - New Playbook Created - {playbook_id} - {title}",
        "body_template": "",
    },
    "procedure": {
        "label": "Procedure",
        "date_field": "created",
        "id_keys": ["procedure_id"],
        "title_keys": ["title"],
        "header_template": "### {time} - New Procedure Created - {procedure_full_id} - {title}",
        "body_template": "",
    },
    "project": {
        "label": "Project",
        "date_field": "created",
        "id_keys": ["project_id"],
        "title_keys": ["project_name", "title"],
        "header_template": "### {time} - New Project Created - {project_id} - {title}",
        "body_template": "",
    },
    "goal": {
        "label": "Goal",
        "date_field": "created",
        "id_keys": ["goal_id"],
        "title_keys": ["title"],
        "header_template": "### {time} - New Goal Created - {goal_id} - {title}",
        "body_template": "",
    },
    "itid": {
        "label": "ITID",
        "date_field": "created",
        "id_keys": ["itid_id"],
        "title_keys": ["itid_name", "title"],
        "header_template": "### {time} - New ITID Created - {itid_id} - {title}",
        "body_template": "",
    },
    "faq": {
        "label": "FAQ",
        "date_field": "created",
        "id_keys": ["faq_id", "faq_ID"],
        "title_keys": ["faq_title", "title"],
        "header_template": "### {time} - New FAQ Created - {faq_id} - {title}",
        "body_template": "",
    },
    "how_to": {
        "label": "How-To",
        "date_field": "created",
        "id_keys": ["how_to_id"],
        "title_keys": ["title"],
        "header_template": "### {time} - New How-To Created - {how_to_id} - {title}",
        "body_template": "",
    },
    "threat_actor": {
        "label": "Threat Actor",
        "date_field": "created",
        "id_keys": ["actor_id"],
        "title_keys": ["actor_name", "common_name"],
        "header_template": "### {time} - New Threat Actor Created - {actor_id} - {title}",
        "body_template": "",
    },
    "campaign": {
        "label": "Campaign",
        "date_field": "created",
        "id_keys": ["campaign_id"],
        "title_keys": ["campaign_name", "title"],
        "header_template": "### {time} - New Campaign Created - {campaign_id} - {title}",
        "body_template": "",
    },
    "mitre_ttp": {
        "label": "MITRE TTP",
        "date_field": "created",
        "id_keys": ["subtechnique_id", "technique_id"],
        "title_keys": ["technique_name", "title"],
        "header_template": "### {time} - New MITRE TTP Created - {ttp_id} - {title}",
        "body_template": "",
    },
    "malware": {
        "label": "Malware",
        "date_field": "created",
        "id_keys": ["malware_id"],
        "title_keys": ["malware_name", "title"],
        "header_template": "### {time} - New Malware Created - {malware_id} - {title}",
        "body_template": "",
    },
    "ioc": {
        "label": "IOC",
        "date_field": "created",
        "id_keys": ["ioc_type"],
        "title_keys": ["value", "title"],
        "header_template": "### {time} - New IOC Created - {ioc_type} - {title}",
        "body_template": "",
    },
    "sla": {
        "label": "SLA",
        "date_field": "created",
        "id_keys": ["sla_id"],
        "title_keys": ["sla_title", "title"],
        "header_template": "### {time} - New SLA Created - {sla_id} - {title}",
        "body_template": "",
    },
}


def _default_shift_log_rules() -> dict:
    entities = {}
    for key, meta in SHIFT_LOG_ENTITY_REGISTRY.items():
        entities[key] = {
            "enabled": True,
            "date_field": meta.get("date_field", "created"),
            "header_template": meta.get("header_template", "### {time} - New Note Created - {title}"),
            "body_template": meta.get("body_template", ""),
        }
    return {"enabled": True, "entities": entities}


def _ensure_shift_log_rules_schema(cfg: dict) -> dict:
    if not isinstance(cfg, dict):
        return _default_shift_log_rules()
    rules = cfg.setdefault("shift_log_rules", {})
    defaults = _default_shift_log_rules()
    if not isinstance(rules, dict) or not rules:
        cfg["shift_log_rules"] = copy.deepcopy(defaults)
        return cfg["shift_log_rules"]
    rules.setdefault("enabled", defaults["enabled"])
    rules.setdefault("entities", {})
    for entity_key, def_rule in defaults["entities"].items():
        ent = rules["entities"].setdefault(entity_key, {})
        if not isinstance(ent, dict):
            rules["entities"][entity_key] = copy.deepcopy(def_rule)
            continue
        for k, v in def_rule.items():
            ent.setdefault(k, v)
    return rules


def _split_datetime(value: str) -> tuple[str, str]:
    raw = (value or "").strip()
    if not raw:
        return "", ""
    parts = raw.split()
    date_str = parts[0] if parts else ""
    time_str = ""
    if len(parts) > 1:
        time_str = parts[1][:5]
    return date_str, time_str


def _first_present(payload: dict, keys: list[str]) -> str:
    for key in keys:
        val = (payload or {}).get(key)
        if val is None:
            continue
        s = str(val).strip()
        if s:
            return s
    return ""


def _build_shift_log_tokens(
    entity_type: str,
    payload: dict,
    *,
    now_dt: datetime,
    vault_root: Path | None = None,
    note_path: Path | None = None,
) -> dict:
    payload = payload or {}
    tokens = {}
    tokens.update({k: ("" if v is None else str(v)) for k, v in payload.items()})
    tokens["entity_type"] = entity_type

    created_date, created_time = _split_datetime(tokens.get("created", ""))
    start_date, start_time = _split_datetime(tokens.get("start_time", ""))
    end_date, end_time = _split_datetime(tokens.get("end_time", ""))

    tokens["time"] = now_dt.strftime("%H:%M")
    tokens["date"] = now_dt.strftime("%Y-%m-%d")
    tokens["created_time"] = created_time or tokens["time"]
    tokens["created_date"] = created_date or tokens["date"]
    tokens["created_datetime"] = f"{tokens['created_date']} {tokens['created_time']}"

    tokens["start_date"] = start_date or tokens["date"]
    tokens["start_time"] = start_time or "00:00"
    tokens["start_datetime"] = f"{tokens['start_date']} {tokens['start_time']}"

    tokens["end_date"] = end_date or tokens["start_date"] or tokens["date"]
    tokens["end_time"] = end_time or "00:00"
    tokens["end_datetime"] = f"{tokens['end_date']} {tokens['end_time']}"

    tokens["meeting_start_date"] = tokens["start_datetime"]
    tokens["meeting_end_date"] = tokens["end_datetime"]

    meta = SHIFT_LOG_ENTITY_REGISTRY.get(entity_type, {})
    title_keys = meta.get("title_keys", [])
    if not title_keys:
        title_keys = ["title", "name", "title_plain"]
    title_plain = tokens.get("title", "").strip() or _first_present(payload, title_keys)
    tokens["title_plain"] = title_plain
    tokens["title"] = title_plain
    tokens["entity_id"] = _first_present(payload, meta.get("id_keys", []))
    tokens["label"] = meta.get("label", entity_type.replace("_", " ").title())
    if entity_type == "mitre_ttp":
        tokens["ttp_id"] = tokens.get("subtechnique_id") or tokens.get("technique_id") or ""
    if entity_type == "ioc":
        if not tokens.get("title"):
            tokens["title"] = tokens.get("value", "")
    if entity_type == "faq":
        tokens["faq_id"] = tokens.get("faq_id") or tokens.get("faq_ID") or ""
    if entity_type == "procedure":
        proc_id = tokens.get("procedure_id", "").strip()
        year = now_dt.strftime("%Y")
        tokens["procedure_full_id"] = f"PRC-{year}-{proc_id}" if proc_id else ""

    if note_path and vault_root:
        try:
            rel = note_path.relative_to(vault_root).as_posix()
            if rel.lower().endswith(".md"):
                rel = rel[:-3]
            link_text = title_plain or Path(rel).stem
            tokens["title_link"] = f"[[{rel}|{link_text}]]"
            tokens["note_link"] = f"[[{rel}]]"
            if tokens.get("title"):
                tokens["title"] = tokens["title_link"]
        except Exception:
            tokens["title_link"] = title_plain
            tokens["note_link"] = ""
    return tokens


class _SafeFormatDict(dict):
    def __missing__(self, key: str) -> str:
        return ""


def _format_shift_log_template(template: str, tokens: dict) -> str:
    try:
        return template.format_map(_SafeFormatDict(tokens))
    except Exception:
        return template


def _append_shift_log_entry(cfg: dict, vault_root: Path, entity_type: str, payload: dict, note_path: Path | None) -> None:
    """
    Append a shift log entry to the daily note. REQUIRED: Every time a new entity is created,
    an entry MUST be added to the daily note's Shift Log as a separate entry.
    """
    rules = _ensure_shift_log_rules_schema(cfg)
    if not rules.get("enabled", True):
        return
    entity_rules = (rules.get("entities") or {}).get(entity_type, {})
    # If entity type has explicit rules and is disabled, respect that
    if entity_rules and not entity_rules.get("enabled", True):
        return
    # Use fallback for new/unregistered entity types so we always log
    use_fallback = not entity_rules
    if use_fallback:
        entity_rules = {
            "date_field": "created",
            "header_template": f"### {{time}} - New {{label}} Created - {{title}}",
            "body_template": "",
        }

    now_dt = datetime.now()
    date_field = (entity_rules.get("date_field") or "").strip() or "created"
    date_source = (payload or {}).get(date_field, "") or now_dt.strftime("%Y-%m-%d")
    primary_date_str, _ = _split_datetime(str(date_source))
    if not primary_date_str:
        primary_date_str = now_dt.strftime("%Y-%m-%d")

    # For meetings: REQUIRED - write to daily shift log on BOTH (1) day created and
    # (2) day held. Creation entry on today; Meeting Scheduled entry on meeting start date.
    # When same day, both entries go in the same daily note.
    # Other entities: append only to primary date.
    today_str = now_dt.strftime("%Y-%m-%d")
    if entity_type == "meeting":
        entries_to_append = [
            (today_str, False),   # Creation entry
            (primary_date_str, True),  # Meeting Scheduled entry
        ]
    else:
        entries_to_append = [(primary_date_str, False)]

    # Resolve daily note config
    try:
        app_dir = get_app_dir()
        tokens = build_tokens(app_dir, cfg)
    except Exception:
        tokens = {}

    paths_cfg = cfg.get("paths", {})
    daily_cfg = cfg.get("daily_note", {})
    daily_root_rel = expand_value(paths_cfg.get("DAILY_NOTES_REL", "10_Operations/09_Journals/01_Daily"), tokens)
    daily_root_rel = daily_root_rel.strip().replace("\\", "/")
    file_pattern = daily_cfg.get("file_pattern", "%Y-%m-%d")
    ext = daily_cfg.get("extension", ".md")
    apply_when = str(daily_cfg.get("apply_template_when", "missing")).strip()
    # If daily note exists: append to it. If not: create it for ALL entities using the standard template.
    create_if_missing = True

    template_rel = str(daily_cfg.get("template_file_rel", "")).strip().replace("\\", "/")
    template_rel = expand_value(template_rel, tokens)
    template_abs = str(Path(vault_root) / Path(*template_rel.split("/"))) if template_rel else ""

    tokens = _build_shift_log_tokens(entity_type, payload or {}, now_dt=now_dt, vault_root=vault_root, note_path=note_path)
    body = _format_shift_log_template(entity_rules.get("body_template", ""), tokens).strip()

    marker = "<!-- SHIFT_NOTES_LOG -->"
    shift_heading = "## 2. Shift Notes (Chronological Log)"

    meeting_scheduled_header = "### {meeting_start_date} - Meeting Scheduled - {title} - {meeting_start_date} - {meeting_end_date}"

    for date_str, use_scheduled_header in entries_to_append:
        if entity_type == "meeting" and use_scheduled_header:
            header = _format_shift_log_template(meeting_scheduled_header, tokens).strip()
        else:
            header = _format_shift_log_template(entity_rules.get("header_template", ""), tokens).strip()
        block = header if not body else f"{header}\n{body}"

        try:
            dt_obj = datetime.strptime(date_str, "%Y-%m-%d")
        except Exception:
            dt_obj = now_dt
        leaf = dt_obj.strftime(file_pattern) + ext
        file_within_vault = f"{daily_root_rel.rstrip('/')}/{leaf}".replace("\\", "/")
        daily_path = Path(vault_root) / Path(*file_within_vault.split("/"))

        try:
            _ensure_daily_note_exists(str(daily_path), template_abs, apply_when, create_if_missing)
        except Exception:
            continue
        if not daily_path.exists():
            continue

        content = daily_path.read_text(encoding="utf-8")
        updated = content

        if marker in content:
            marker_idx = content.index(marker)
            section_start = marker_idx + len(marker)
            next_heading_idx = content.find("\n## ", section_start)
            if next_heading_idx == -1:
                next_heading_idx = len(content)
            section_body = content[section_start:next_heading_idx].rstrip("\n")
            insert = f"{section_body}\n{block}\n" if section_body.strip() else f"\n{block}\n"
            updated = content[:section_start] + insert + content[next_heading_idx:]
        elif shift_heading in content:
            idx = content.index(shift_heading)
            next_heading_idx = content.find("\n## ", idx + len(shift_heading))
            if next_heading_idx == -1:
                next_heading_idx = len(content)
            section_body = content[idx + len(shift_heading):next_heading_idx].rstrip("\n")
            insert = f"{section_body}\n{block}\n" if section_body.strip() else f"\n{block}\n"
            updated = content[: idx + len(shift_heading)] + insert + content[next_heading_idx:]
        else:
            updated = content + f"\n\n{block}\n"

        if updated != content:
            daily_path.write_text(updated, encoding="utf-8")


def _log_shift_entry_from_launcher(master, vault_root: Path, entity_type: str, payload: dict, note_path: Path | None) -> None:
    """
    Append a shift log entry to the Daily Note. REQUIRED: Every entity creation MUST call this
    so an entry is added to the daily note's Shift Log as a separate entry.
    Requires vault_root; note_path is optional (for links).
    """
    if not vault_root:
        return
    cfg = getattr(master, "cfg", {}) or {}
    _append_shift_log_entry(cfg, vault_root, entity_type, payload, note_path)


def _normalize_keyword(s: str) -> str:
    s = (s or "").strip()
    s = re.sub(r"\s+", " ", s)
    return s


def _split_bulk(text: str) -> list:
    """
    Split bulk input that may be newline and/or comma separated.
    Preserves phrases; smart-splits commas unless many commas inside.
    """
    raw = (text or "").strip()
    if not raw:
        return []
    # First split into lines
    lines = [ln.strip() for ln in raw.splitlines() if ln.strip()]
    out = []
    for ln in lines:
        # If it looks CSV-ish, split by commas
        if "," in ln:
            parts = [p.strip() for p in ln.split(",")]
            out.extend([p for p in parts if p])
        else:
            out.append(ln)
    return [_normalize_keyword(x) for x in out if _normalize_keyword(x)]


def _dedupe_preserve_order(items: list, case_sensitive: bool) -> list:
    seen = set()
    out = []
    for s in items:
        k = s if case_sensitive else s.lower()
        if k in seen:
            continue
        seen.add(k)
        out.append(s)
    return out


def _keywords_match(text: str, keywords: list, case_sensitive: bool) -> list:
    """
    Returns matched keywords (substring match).
    """
    if not case_sensitive:
        hay = (text or "").lower()
        return [k for k in keywords if k and k.lower() in hay]
    hay = text or ""
    return [k for k in keywords if k and k in hay]


# ----------------------------
# UI Window
# ----------------------------
class ManageKeywordsWindow(tk.Toplevel):
    """
    A self-contained, drop-in window for managing rss keywords in config.
    """

    def __init__(self, master, cfg_ref: dict, app_dir: Path, theme: dict = None):
        super().__init__(master)
        self.cfg_ref = cfg_ref                  # reference to launcher cfg in memory
        self.app_dir = Path(app_dir)
        self.cfg_path = _config_path_for_launcher(self.app_dir)
        self.theme = theme or {}
        self._dirty = False

        # Ensure schema
        self.kw = _ensure_keywords_schema(self.cfg_ref)
        self.rss = self.cfg_ref["rss"]

        # UI state
        self._list_names = []
        self._current_list = None
        self._filtered_keywords = []
        self._filter_job = None

        self.title("Manage Keywords")
        self.configure(bg="#111111")
        self.geometry("1200x720")
        self.minsize(980, 600)

        self.protocol("WM_DELETE_WINDOW", self._on_close)

        self._build_ui()
        self._load_lists_into_ui()
        self._bind_keys()

    # ----------------------------
    # Theme helpers
    # ----------------------------
    def _c(self, key: str, default: str) -> str:
        return (self.theme.get(key) or default)

    # ----------------------------
    # Build UI
    # ----------------------------
    def _build_ui(self):
        bg = "#111111"
        panel = "#141824"
        input_bg = "#1e1e1e"
        btn_bg = "#1a1a1a"
        btn_hover = "#2a2a2a"
        fg = "#ffffff"
        fg2 = "#cfcfcf"
        meta = "#9aa7c0"

        # Top bar
        top = tk.Frame(self, bg=bg, padx=12, pady=10)
        top.pack(fill="x")

        tk.Label(top, text="Manage Keywords", bg=bg, fg=fg, font=("Segoe UI", 14, "bold")).pack(side="left")

        self._status_var = tk.StringVar(value="")
        tk.Label(top, textvariable=self._status_var, bg=bg, fg=meta, font=("Segoe UI", 9, "normal")).pack(
            side="left", padx=(12, 0)
        )

        top_btns = tk.Frame(top, bg=bg)
        top_btns.pack(side="right")

        def mk_btn(text, cmd):
            b = tk.Button(
                top_btns,
                text=text,
                command=cmd,
                bg=btn_bg,
                fg=fg,
                relief="flat",
                padx=10,
                pady=6,
                font=("Segoe UI", 10, "bold"),
            )
            b.pack(side="left", padx=4)
            # simple hover
            b.bind("<Enter>", lambda e: b.configure(bg=btn_hover))
            b.bind("<Leave>", lambda e: b.configure(bg=btn_bg))
            return b

        mk_btn("Reload", self._reload_from_disk)
        mk_btn("Save", self._save_to_disk)
        mk_btn("Close", self._on_close)

        # Main split (3 columns)
        body = tk.PanedWindow(self, orient="horizontal", bg=bg, sashwidth=6, sashrelief="flat")
        body.pack(fill="both", expand=True, padx=12, pady=(0, 12))

        self.left = tk.Frame(body, bg=panel, padx=10, pady=10)
        self.center = tk.Frame(body, bg=panel, padx=10, pady=10)
        self.right = tk.Frame(body, bg=panel, padx=10, pady=10)

        body.add(self.left, minsize=260)
        body.add(self.center, minsize=460)
        body.add(self.right, minsize=360)

        # LEFT: Lists
        tk.Label(self.left, text="Lists", bg=panel, fg=fg, font=("Segoe UI", 12, "bold")).pack(anchor="w")

        self.listbox = tk.Listbox(
            self.left,
            bg="#1a1a1a",
            fg=fg,
            selectbackground="#2a2a2a",
            selectforeground=fg,
            relief="flat",
            highlightthickness=0,
            font=("Segoe UI", 11, "normal"),
            height=18,
        )
        self.listbox.pack(fill="both", expand=True, pady=(8, 8))
        self.listbox.bind("<<ListboxSelect>>", lambda e: self._on_list_selected())

        list_btns = tk.Frame(self.left, bg=panel)
        list_btns.pack(fill="x", pady=(4, 0))

        self._btn_new_list = tk.Button(list_btns, text="+ New", command=self._new_list, bg="#1a1a1a", fg=fg, relief="flat")
        self._btn_rename_list = tk.Button(list_btns, text="Rename", command=self._rename_list, bg="#1a1a1a", fg=fg, relief="flat")
        self._btn_dup_list = tk.Button(list_btns, text="Duplicate", command=self._duplicate_list, bg="#1a1a1a", fg=fg, relief="flat")
        self._btn_del_list = tk.Button(list_btns, text="Delete", command=self._delete_list, bg="#1a1a1a", fg=fg, relief="flat")

        for b in (self._btn_new_list, self._btn_rename_list, self._btn_dup_list, self._btn_del_list):
            b.pack(side="left", padx=3, pady=4, ipadx=8, ipady=4)
            b.bind("<Enter>", lambda e, bb=b: bb.configure(bg="#2a2a2a"))
            b.bind("<Leave>", lambda e, bb=b: bb.configure(bg="#1a1a1a"))

        self._list_meta = tk.Label(self.left, text="", bg=panel, fg=meta, font=("Segoe UI", 9, "normal"))
        self._list_meta.pack(anchor="w", pady=(8, 0))

        # CENTER: Keywords
        tk.Label(self.center, text="Keywords", bg=panel, fg=fg, font=("Segoe UI", 12, "bold")).pack(anchor="w")

        # Filter row
        filt = tk.Frame(self.center, bg=panel)
        filt.pack(fill="x", pady=(8, 8))

        tk.Label(filt, text="Filter:", bg=panel, fg=fg2, font=("Segoe UI", 10, "normal")).pack(side="left")
        self.kw_filter_var = tk.StringVar(value="")
        self.kw_filter = tk.Entry(
            filt, textvariable=self.kw_filter_var, bg=input_bg, fg=fg, insertbackground=fg, relief="flat",
            font=("Segoe UI", 11, "normal")
        )
        self.kw_filter.pack(side="left", fill="x", expand=True, padx=(8, 0))
        self.kw_filter.bind("<KeyRelease>", lambda e: self._debounced_apply_keyword_filter())

        # Keywords list
        self.kw_list = tk.Listbox(
            self.center,
            bg="#1a1a1a",
            fg=fg,
            selectbackground="#2a2a2a",
            selectforeground=fg,
            relief="flat",
            highlightthickness=0,
            font=("Segoe UI", 11, "normal"),
        )
        self.kw_list.pack(fill="both", expand=True)

        # Add row
        addrow = tk.Frame(self.center, bg=panel)
        addrow.pack(fill="x", pady=(8, 4))

        self.kw_add_var = tk.StringVar(value="")
        self.kw_add = tk.Entry(
            addrow, textvariable=self.kw_add_var, bg=input_bg, fg=fg, insertbackground=fg, relief="flat",
            font=("Segoe UI", 11, "normal")
        )
        self.kw_add.pack(side="left", fill="x", expand=True)

        tk.Button(
            addrow, text="+ Add", command=self._add_keyword, bg="#1a1a1a", fg=fg, relief="flat",
            padx=10, pady=6, font=("Segoe UI", 10, "bold")
        ).pack(side="left", padx=(8, 0))

        # Keyword actions
        kw_btns = tk.Frame(self.center, bg=panel)
        kw_btns.pack(fill="x", pady=(4, 0))

        def mk_kw_btn(text, cmd):
            b = tk.Button(kw_btns, text=text, command=cmd, bg="#1a1a1a", fg=fg, relief="flat",
                          padx=10, pady=6, font=("Segoe UI", 10, "bold"))
            b.pack(side="left", padx=4, pady=4)
            b.bind("<Enter>", lambda e: b.configure(bg="#2a2a2a"))
            b.bind("<Leave>", lambda e: b.configure(bg="#1a1a1a"))
            return b

        mk_kw_btn("Remove", self._remove_keywords)
        mk_kw_btn("Bulk Add", self._bulk_add_dialog)
        mk_kw_btn("Dedupe", self._dedupe_keywords)
        mk_kw_btn("Normalize", self._normalize_keywords)

        # RIGHT: Settings + Preview
        tk.Label(self.right, text="Settings", bg=panel, fg=fg, font=("Segoe UI", 12, "bold")).pack(anchor="w")

        # Global settings frame
        gs = tk.Frame(self.right, bg=panel)
        gs.pack(fill="x", pady=(8, 10))

        # Match mode
        tk.Label(gs, text="Match mode:", bg=panel, fg=fg2, font=("Segoe UI", 10, "normal")).grid(row=0, column=0, sticky="w")
        self.match_mode_var = tk.StringVar(value=self.kw.get("match_mode", "any"))
        self.match_mode = ttk.Combobox(gs, textvariable=self.match_mode_var, values=("any", "all"), state="readonly", width=10)
        self.match_mode.grid(row=0, column=1, sticky="w", padx=(8, 0))
        self.match_mode.bind("<<ComboboxSelected>>", lambda e: self._set_dirty(True))

        # Case sensitive
        self.case_var = tk.BooleanVar(value=bool(self.kw.get("case_sensitive", False)))
        self.case_chk = tk.Checkbutton(
            gs, text="Case sensitive", variable=self.case_var, bg=panel, fg=fg2, activebackground=panel,
            activeforeground=fg2, selectcolor="#1a1a1a", command=lambda: self._set_dirty(True)
        )
        self.case_chk.grid(row=1, column=0, columnspan=2, sticky="w", pady=(6, 0))

        # Min hits
        tk.Label(gs, text="Min hits:", bg=panel, fg=fg2, font=("Segoe UI", 10, "normal")).grid(row=2, column=0, sticky="w", pady=(6, 0))
        self.min_hits_var = tk.IntVar(value=int(self.kw.get("min_hits", 1)))
        self.min_hits_spin = tk.Spinbox(gs, from_=1, to=999, textvariable=self.min_hits_var, width=8, relief="flat",
                                        bg="#1e1e1e", fg=fg, insertbackground=fg, command=lambda: self._set_dirty(True))
        self.min_hits_spin.grid(row=2, column=1, sticky="w", padx=(8, 0), pady=(6, 0))

        # List-level tag mapping
        tk.Label(self.right, text="Selected list tags", bg=panel, fg=fg, font=("Segoe UI", 12, "bold")).pack(anchor="w", pady=(6, 0))
        tk.Label(self.right, text="Applied to matches (optional).", bg=panel, fg="#9aa7c0", font=("Segoe UI", 9, "normal")).pack(anchor="w")

        self.tags_var = tk.StringVar(value="")
        self.tags_entry = tk.Entry(
            self.right, textvariable=self.tags_var, bg="#1e1e1e", fg=fg, insertbackground=fg,
            relief="flat", font=("Segoe UI", 11, "normal")
        )
        self.tags_entry.pack(fill="x", pady=(6, 10))
        self.tags_entry.bind("<KeyRelease>", lambda e: self._set_dirty(True))

        # Preview
        tk.Label(self.right, text="Match Preview", bg=panel, fg=fg, font=("Segoe UI", 12, "bold")).pack(anchor="w", pady=(6, 0))
        tk.Label(self.right, text="Paste headline/body text and test.", bg=panel, fg="#9aa7c0", font=("Segoe UI", 9, "normal")).pack(anchor="w")

        self.preview_text = tk.Text(self.right, height=10, bg="#1e1e1e", fg=fg, insertbackground=fg, relief="flat", wrap="word")
        self.preview_text.pack(fill="both", expand=False, pady=(6, 6))

        pr_btns = tk.Frame(self.right, bg=panel)
        pr_btns.pack(fill="x", pady=(0, 6))

        tk.Button(
            pr_btns, text="Run Test", command=self._run_preview_test, bg="#1a1a1a", fg=fg, relief="flat",
            padx=10, pady=6, font=("Segoe UI", 10, "bold")
        ).pack(side="left", padx=4)

        tk.Button(
            pr_btns, text="Clear", command=lambda: (self.preview_text.delete("1.0", "end"), None), bg="#1a1a1a", fg=fg,
            relief="flat", padx=10, pady=6, font=("Segoe UI", 10, "bold")
        ).pack(side="left", padx=4)

        self.preview_result = tk.Text(self.right, height=8, bg="#111111", fg="#e6e6e6", relief="flat", wrap="word")
        self.preview_result.pack(fill="both", expand=True)
        self.preview_result.configure(state="disabled")

        # Footer
        foot = tk.Frame(self, bg=bg, padx=12, pady=6)
        foot.pack(fill="x")
        self._dirty_lbl = tk.Label(foot, text="Not saved", bg=bg, fg=meta, font=("Segoe UI", 9, "normal"))
        self._dirty_lbl.pack(side="left")

        tk.Label(
            foot,
            text="Shortcuts: Ctrl+S save • Ctrl+R reload • / filter • Del remove",
            bg=bg, fg=meta, font=("Segoe UI", 9, "normal")
        ).pack(side="right")

    # ----------------------------
    # Keyboard bindings
    # ----------------------------
    def _bind_keys(self):
        self.bind("<Control-s>", lambda e: (self._save_to_disk(), "break"))
        self.bind("<Control-r>", lambda e: (self._reload_from_disk(), "break"))
        self.bind("/", lambda e: (self.kw_filter.focus_set(), "break"))
        self.bind("<Escape>", lambda e: (self._on_close(), "break"))

        self.kw_add.bind("<Return>", lambda e: (self._add_keyword(), "break"))
        self.kw_list.bind("<Delete>", lambda e: (self._remove_keywords(), "break"))
        self.listbox.bind("<Delete>", lambda e: (self._delete_list(), "break"))

        self.preview_text.bind("<Control-Return>", lambda e: (self._run_preview_test(), "break"))

    # ----------------------------
    # Dirty handling
    # ----------------------------
    def _set_dirty(self, is_dirty: bool):
        self._dirty = bool(is_dirty)
        self._dirty_lbl.configure(text=("Unsaved changes" if self._dirty else "All changes saved"))

    def _touch(self):
        self._set_dirty(True)

    # ----------------------------
    # Load + refresh UI
    # ----------------------------
    def _load_lists_into_ui(self):
        lists = self.kw.get("lists", {})
        self._list_names = sorted([k for k in lists.keys()])
        self.listbox.delete(0, tk.END)
        for name in self._list_names:
            self.listbox.insert(tk.END, name)

        if self._list_names:
            # keep current if possible
            if self._current_list in self._list_names:
                idx = self._list_names.index(self._current_list)
            else:
                idx = 0
            self.listbox.selection_set(idx)
            self.listbox.activate(idx)
            self.listbox.see(idx)
            self._on_list_selected()
        else:
            self._current_list = None
            self.kw_list.delete(0, tk.END)
            self._list_meta.configure(text="No lists defined.")
            self.tags_var.set("")
            self._set_dirty(False)

    def _on_list_selected(self):
        sel = self.listbox.curselection()
        if not sel:
            return
        idx = sel[0]
        if idx < 0 or idx >= len(self._list_names):
            return
        name = self._list_names[idx]
        self._current_list = name
        self.kw_filter_var.set("")
        self._render_keywords()
        self._render_list_meta()
        self._render_list_tags()

    def _render_list_meta(self):
        if not self._current_list:
            self._list_meta.configure(text="")
            return
        kws = self.kw["lists"].get(self._current_list, [])
        self._list_meta.configure(text=f"Keywords: {len(kws)}")

    def _render_list_tags(self):
        # tag map stored at rss["keyword_tag_map"][list_name] -> list[str] or comma string
        m = self.rss.get("keyword_tag_map", {}) or {}
        v = m.get(self._current_list, [])
        if isinstance(v, str):
            tags = [t.strip() for t in v.split(",") if t.strip()]
        elif isinstance(v, (list, tuple)):
            tags = [str(t).strip() for t in v if str(t).strip()]
        else:
            tags = []
        self.tags_var.set(", ".join(tags))

    def _render_keywords(self):
        self.kw_list.delete(0, tk.END)
        if not self._current_list:
            return
        kws = list(self.kw["lists"].get(self._current_list, []))
        q = (self.kw_filter_var.get() or "").strip().lower()

        if q:
            out = []
            for k in kws:
                if q in (k or "").lower():
                    out.append(k)
            self._filtered_keywords = out
        else:
            self._filtered_keywords = kws

        for k in self._filtered_keywords:
            self.kw_list.insert(tk.END, k)

    def _debounced_apply_keyword_filter(self):
        if self._filter_job is not None:
            try:
                self.after_cancel(self._filter_job)
            except Exception:
                pass
        self._filter_job = self.after(150, self._apply_keyword_filter)

    def _apply_keyword_filter(self):
        self._filter_job = None
        self._render_keywords()

    # ----------------------------
    # List CRUD
    # ----------------------------
    def _prompt(self, title: str, prompt: str, initial: str = "") -> str:
        """
        Simple modal prompt.
        """
        win = tk.Toplevel(self)
        win.title(title)
        win.configure(bg="#111111")
        win.resizable(False, False)
        win.transient(self)
        win.grab_set()

        # center
        self.update_idletasks()
        w, h = 520, 160
        x = self.winfo_rootx() + max((self.winfo_width() - w) // 2, 0)
        y = self.winfo_rooty() + max((self.winfo_height() - h) // 2, 0)
        win.geometry(f"{w}x{h}+{x}+{y}")

        tk.Label(win, text=prompt, bg="#111111", fg="#cfcfcf", font=("Segoe UI", 10, "normal")).pack(
            anchor="w", padx=12, pady=(12, 6)
        )
        var = tk.StringVar(value=initial)
        ent = tk.Entry(win, textvariable=var, bg="#1e1e1e", fg="#ffffff", insertbackground="#ffffff", relief="flat",
                       font=("Segoe UI", 11, "normal"))
        ent.pack(fill="x", padx=12)
        ent.focus_set()

        out = {"value": None}

        def ok():
            out["value"] = var.get()
            win.destroy()

        def cancel():
            out["value"] = None
            win.destroy()

        btns = tk.Frame(win, bg="#111111")
        btns.pack(fill="x", padx=12, pady=12)

        tk.Button(btns, text="Cancel", command=cancel, bg="#1a1a1a", fg="#ffffff", relief="flat",
                  padx=10, pady=6, font=("Segoe UI", 10, "bold")).pack(side="right", padx=4)
        tk.Button(btns, text="OK", command=ok, bg="#1a1a1a", fg="#ffffff", relief="flat",
                  padx=10, pady=6, font=("Segoe UI", 10, "bold")).pack(side="right", padx=4)

        win.bind("<Return>", lambda e: ok())
        win.bind("<Escape>", lambda e: cancel())

        self.wait_window(win)
        return out["value"]

    def _validate_list_name(self, name: str) -> str:
        name = (name or "").strip()
        if not name:
            return "List name is required."
        if not _LIST_NAME_RE.match(name):
            return "List name must be lowercase letters, numbers, or underscore (e.g., vendor_watch)."
        if name in self.kw["lists"]:
            return "That list name already exists."
        return ""

    def _new_list(self):
        name = self._prompt("New List", "List name (lowercase, numbers, underscore):", "company_watch")
        if name is None:
            return
        name = name.strip()
        err = self._validate_list_name(name)
        if err:
            messagebox.showerror("Invalid list name", err)
            return
        self.kw["lists"][name] = []
        self._touch()
        self._load_lists_into_ui()
        # select new list
        if name in self._list_names:
            idx = self._list_names.index(name)
            self.listbox.selection_clear(0, tk.END)
            self.listbox.selection_set(idx)
            self.listbox.activate(idx)
            self._on_list_selected()

    def _rename_list(self):
        if not self._current_list:
            return
        old = self._current_list
        name = self._prompt("Rename List", "New list name:", old)
        if name is None:
            return
        name = name.strip()
        if name == old:
            return
        if not _LIST_NAME_RE.match(name):
            messagebox.showerror("Invalid list name", "Use lowercase letters, numbers, underscore only.")
            return
        if name in self.kw["lists"]:
            messagebox.showerror("Invalid list name", "A list with that name already exists.")
            return

        # move list
        self.kw["lists"][name] = self.kw["lists"].pop(old, [])

        # migrate tag map if present
        km = self.rss.get("keyword_tag_map", {}) or {}
        if old in km:
            km[name] = km.pop(old)
            self.rss["keyword_tag_map"] = km

        self._current_list = name
        self._touch()
        self._load_lists_into_ui()

    def _duplicate_list(self):
        if not self._current_list:
            return
        base = self._current_list
        name = self._prompt("Duplicate List", "New list name:", f"{base}_copy")
        if name is None:
            return
        name = name.strip()
        err = self._validate_list_name(name)
        if err:
            messagebox.showerror("Invalid list name", err)
            return
        self.kw["lists"][name] = list(self.kw["lists"].get(base, []))
        # copy tag map too
        km = self.rss.get("keyword_tag_map", {}) or {}
        if base in km and name not in km:
            km[name] = km.get(base)
            self.rss["keyword_tag_map"] = km
        self._touch()
        self._load_lists_into_ui()

    def _delete_list(self):
        if not self._current_list:
            return
        name = self._current_list
        count = len(self.kw["lists"].get(name, []))
        if not messagebox.askyesno("Delete list", f"Delete list '{name}'?\n\nKeywords: {count}"):
            return
        try:
            self.kw["lists"].pop(name, None)
            km = self.rss.get("keyword_tag_map", {}) or {}
            km.pop(name, None)
            self.rss["keyword_tag_map"] = km
            self._current_list = None
            self._touch()
            self._load_lists_into_ui()
        except Exception as e:
            messagebox.showerror("Delete failed", str(e))

    # ----------------------------
    # Keyword operations
    # ----------------------------
    def _require_list(self) -> bool:
        if not self._current_list:
            messagebox.showinfo("Select a list", "Select or create a keyword list first.")
            return False
        return True

    def _add_keyword(self):
        if not self._require_list():
            return
        s = _normalize_keyword(self.kw_add_var.get())
        if not s:
            return

        kws = list(self.kw["lists"].get(self._current_list, []))
        case_sensitive = bool(self.case_var.get())

        # prevent duplicates
        kkey = s if case_sensitive else s.lower()
        existing = {x if case_sensitive else (x or "").lower() for x in kws}
        if kkey in existing:
            self._status_var.set("Keyword already exists in list.")
            self.kw_add_var.set("")
            return

        kws.append(s)
        self.kw["lists"][self._current_list] = kws
        self.kw_add_var.set("")
        self._touch()
        self._render_keywords()
        self._render_list_meta()

    def _remove_keywords(self):
        if not self._require_list():
            return
        sel = list(self.kw_list.curselection() or [])
        if not sel:
            return

        # map selection to actual keyword values in filtered view
        to_remove = []
        for i in sel:
            try:
                to_remove.append(self._filtered_keywords[i])
            except Exception:
                pass
        if not to_remove:
            return

        kws = list(self.kw["lists"].get(self._current_list, []))
        kws = [k for k in kws if k not in set(to_remove)]
        self.kw["lists"][self._current_list] = kws
        self._touch()
        self._render_keywords()
        self._render_list_meta()

    def _bulk_add_dialog(self):
        if not self._require_list():
            return

        win = tk.Toplevel(self)
        win.title("Bulk Add Keywords")
        win.configure(bg="#111111")
        win.geometry("700x500")
        win.minsize(560, 420)
        win.transient(self)
        win.grab_set()

        tk.Label(
            win,
            text="Paste keywords (newline or comma separated):",
            bg="#111111",
            fg="#cfcfcf",
            font=("Segoe UI", 10, "normal"),
        ).pack(anchor="w", padx=12, pady=(12, 6))

        txt = tk.Text(win, bg="#1e1e1e", fg="#ffffff", insertbackground="#ffffff", relief="flat", wrap="word")
        txt.pack(fill="both", expand=True, padx=12, pady=(0, 12))
        txt.focus_set()

        status = tk.StringVar(value="")
        tk.Label(win, textvariable=status, bg="#111111", fg="#9aa7c0", font=("Segoe UI", 9, "normal")).pack(
            anchor="w", padx=12
        )

        btns = tk.Frame(win, bg="#111111", padx=12, pady=12)
        btns.pack(fill="x")

        def do_add():
            raw = txt.get("1.0", "end")
            items = _split_bulk(raw)
            if not items:
                status.set("No keywords found.")
                return
            kws = list(self.kw["lists"].get(self._current_list, []))
            case_sensitive = bool(self.case_var.get())
            combined = kws + items
            combined = _dedupe_preserve_order([_normalize_keyword(x) for x in combined if _normalize_keyword(x)], case_sensitive)
            self.kw["lists"][self._current_list] = combined
            self._touch()
            self._render_keywords()
            self._render_list_meta()
            status.set(f"Added/merged {len(items)} items. Total now: {len(combined)}")

        def close():
            win.destroy()

        tk.Button(btns, text="Close", command=close, bg="#1a1a1a", fg="#ffffff", relief="flat",
                  padx=10, pady=6, font=("Segoe UI", 10, "bold")).pack(side="right", padx=4)

        tk.Button(btns, text="Add/Merge", command=do_add, bg="#1a1a1a", fg="#ffffff", relief="flat",
                  padx=10, pady=6, font=("Segoe UI", 10, "bold")).pack(side="right", padx=4)

        win.bind("<Escape>", lambda e: close())
        win.bind("<Control-Return>", lambda e: do_add())

        self.wait_window(win)

    def _dedupe_keywords(self):
        if not self._require_list():
            return
        kws = list(self.kw["lists"].get(self._current_list, []))
        case_sensitive = bool(self.case_var.get())
        before = len(kws)
        kws = _dedupe_preserve_order([_normalize_keyword(x) for x in kws if _normalize_keyword(x)], case_sensitive)
        after = len(kws)
        self.kw["lists"][self._current_list] = kws
        if after != before:
            self._touch()
        self._render_keywords()
        self._render_list_meta()
        self._status_var.set(f"Dedupe complete. Removed {before - after} duplicate(s).")

    def _normalize_keywords(self):
        if not self._require_list():
            return
        kws = list(self.kw["lists"].get(self._current_list, []))
        before = list(kws)
        kws = [_normalize_keyword(x) for x in kws]
        kws = [x for x in kws if x]
        self.kw["lists"][self._current_list] = kws
        if kws != before:
            self._touch()
        self._render_keywords()
        self._render_list_meta()
        self._status_var.set("Normalize complete.")

    # ----------------------------
    # Preview + settings apply
    # ----------------------------
    def _apply_settings_to_cfg(self):
        self.kw["match_mode"] = (self.match_mode_var.get() or "any").strip()
        self.kw["case_sensitive"] = bool(self.case_var.get())
        try:
            mh = int(self.min_hits_var.get())
        except Exception:
            mh = 1
        self.kw["min_hits"] = max(1, mh)

        # per-list tags mapping
        if self._current_list:
            tags_raw = (self.tags_var.get() or "").strip()
            tags = [t.strip() for t in tags_raw.split(",") if t.strip()] if tags_raw else []
            self.rss.setdefault("keyword_tag_map", {})
            self.rss["keyword_tag_map"][self._current_list] = tags

    def _run_preview_test(self):
        if not self._require_list():
            return

        self._apply_settings_to_cfg()
        text = self.preview_text.get("1.0", "end").strip()
        if not text:
            messagebox.showinfo("Preview", "Paste text in the preview box first.")
            return

        keywords = list(self.kw["lists"].get(self._current_list, []))
        case_sensitive = bool(self.kw.get("case_sensitive", False))
        match_mode = self.kw.get("match_mode", "any")
        min_hits = int(self.kw.get("min_hits", 1))

        matched = _keywords_match(text, keywords, case_sensitive)

        would_flag = False
        if match_mode == "any":
            would_flag = len(matched) >= min_hits
        else:
            # all: all keywords must be present; min_hits acts as floor too
            would_flag = (len(matched) >= max(min_hits, len(keywords))) if keywords else False

        tag_map = self.rss.get("keyword_tag_map", {}) or {}
        tags = tag_map.get(self._current_list, [])
        if isinstance(tags, str):
            tags = [t.strip() for t in tags.split(",") if t.strip()]

        lines = []
        lines.append(f"List: {self._current_list}")
        lines.append(f"Keywords in list: {len(keywords)}")
        lines.append(f"Matched: {len(matched)}")
        lines.append(f"Match mode: {match_mode} • Min hits: {min_hits} • Case sensitive: {case_sensitive}")
        lines.append(f"Would flag: {'YES' if would_flag else 'NO'}")
        if tags:
            lines.append(f"Tags to apply: {', '.join(tags)}")
        lines.append("")
        lines.append("Matched keywords:")
        if matched:
            for k in matched[:200]:
                lines.append(f" - {k}")
            if len(matched) > 200:
                lines.append(f" ... and {len(matched) - 200} more")
        else:
            lines.append(" (none)")

        self.preview_result.configure(state="normal")
        self.preview_result.delete("1.0", "end")
        self.preview_result.insert("1.0", "\n".join(lines))
        self.preview_result.configure(state="disabled")

    # ----------------------------
    # Save / Reload
    # ----------------------------
    def _save_to_disk(self):
        try:
            self._apply_settings_to_cfg()

            # Validate all lists
            lists = self.kw.get("lists", {})
            for name, arr in list(lists.items()):
                if not _LIST_NAME_RE.match(name):
                    raise ValueError(f"Invalid list name: {name}")
                if not isinstance(arr, list):
                    lists[name] = []
                # normalize, prune empties
                cleaned = [_normalize_keyword(x) for x in arr]
                cleaned = [x for x in cleaned if x]
                # dedupe based on current global case setting
                cleaned = _dedupe_preserve_order(cleaned, bool(self.kw.get("case_sensitive", False)))
                lists[name] = cleaned

            self.kw["lists"] = lists

            # Persist to config.json next to launcher
            if not self.cfg_path.exists():
                # if config isn't present, still save best-effort
                raise FileNotFoundError(f"config.json not found at: {self.cfg_path}")

            _write_config_file_with_backup(self.cfg_path, self.cfg_ref)
            self._set_dirty(False)
            self._status_var.set("Saved config.json (backup created).")
            self._render_list_meta()
            self._render_keywords()
        except Exception as e:
            messagebox.showerror("Save failed", str(e))

    def _reload_from_disk(self):
        if self._dirty:
            if not messagebox.askyesno("Reload", "Discard unsaved changes and reload from disk?"):
                return
        try:
            if not self.cfg_path.exists():
                raise FileNotFoundError(f"config.json not found at: {self.cfg_path}")

            # Reload and update the referenced cfg in-place so launcher stays consistent
            fresh = _read_config_file(self.cfg_path)

            # Replace contents of cfg_ref in-place
            self.cfg_ref.clear()
            self.cfg_ref.update(fresh)

            # Re-ensure schema and relink pointers
            self.kw = _ensure_keywords_schema(self.cfg_ref)
            self.rss = self.cfg_ref["rss"]

            self._current_list = None
            self._set_dirty(False)
            self._status_var.set("Reloaded from disk.")
            self._load_lists_into_ui()
        except Exception as e:
            messagebox.showerror("Reload failed", str(e))

    # ----------------------------
    # Close
    # ----------------------------
    def _on_close(self):
        if self._dirty:
            resp = messagebox.askyesnocancel("Unsaved changes", "Save changes before closing?")
            if resp is None:
                return
            if resp is True:
                self._save_to_disk()
                if self._dirty:
                    return
        self.destroy()


class ShiftLogRulesWindow(tk.Toplevel):
    """
    UI window for managing per-entity Shift Log rules in config.json.
    """

    def __init__(self, master, cfg_ref: dict, app_dir: Path, theme: dict = None):
        super().__init__(master)
        self.cfg_ref = cfg_ref
        self.app_dir = Path(app_dir)
        self.cfg_path = _config_path_for_launcher(self.app_dir)
        self.theme = theme or {}
        self._dirty = False
        self._current_entity = None

        self.rules = _ensure_shift_log_rules_schema(self.cfg_ref)
        self.entity_keys = list(SHIFT_LOG_ENTITY_REGISTRY.keys())

        self.title("Shift Log Rules")
        self.configure(bg="#111111")
        self.geometry("1100x720")
        self.minsize(960, 640)

        self.protocol("WM_DELETE_WINDOW", self._on_close)
        self._build_ui()
        self._load_entities()

    def _set_dirty(self, dirty: bool = True):
        self._dirty = dirty
        self._status_var.set("Unsaved changes" if dirty else "All changes saved")

    def _build_ui(self):
        bg = "#111111"
        panel = "#141824"
        input_bg = "#1e1e1e"
        btn_bg = "#1a1a1a"
        btn_hover = "#2a2a2a"
        fg = "#ffffff"
        fg2 = "#cfcfcf"
        meta = "#9aa7c0"

        top = tk.Frame(self, bg=bg, padx=12, pady=10)
        top.pack(fill="x")

        tk.Label(top, text="Shift Log Rules", bg=bg, fg=fg, font=("Segoe UI", 14, "bold")).pack(side="left")

        self._status_var = tk.StringVar(value="All changes saved")
        tk.Label(top, textvariable=self._status_var, bg=bg, fg=meta, font=("Segoe UI", 9)).pack(
            side="left", padx=(12, 0)
        )

        top_btns = tk.Frame(top, bg=bg)
        top_btns.pack(side="right")

        def mk_btn(text, cmd):
            b = tk.Button(
                top_btns,
                text=text,
                command=cmd,
                bg=btn_bg,
                fg=fg,
                relief="flat",
                padx=10,
                pady=6,
                font=("Segoe UI", 10, "bold"),
            )
            b.pack(side="left", padx=4)
            b.bind("<Enter>", lambda e: b.configure(bg=btn_hover))
            b.bind("<Leave>", lambda e: b.configure(bg=btn_bg))
            return b

        mk_btn("Reload", self._reload_from_disk)
        mk_btn("Reset Defaults", self._reset_defaults)
        mk_btn("Save", self._save_to_disk)
        mk_btn("Close", self._on_close)

        main = tk.Frame(self, bg=bg, padx=12, pady=10)
        main.pack(fill="both", expand=True)

        left = tk.Frame(main, bg=panel, padx=10, pady=10)
        left.pack(side="left", fill="y")

        tk.Label(left, text="Entity Types", bg=panel, fg=fg2, font=("Segoe UI", 11, "bold")).pack(anchor="w")
        self.entity_list = tk.Listbox(left, bg=input_bg, fg=fg, selectbackground="#2a2a2a", height=20)
        self.entity_list.pack(fill="y", expand=True, pady=(6, 0))
        self.entity_list.bind("<<ListboxSelect>>", self._on_select_entity)

        right = tk.Frame(main, bg=panel, padx=12, pady=12)
        right.pack(side="left", fill="both", expand=True, padx=(12, 0))

        self.enabled_var = tk.BooleanVar(value=True)
        self.date_field_var = tk.StringVar()
        self.header_var = tk.StringVar()

        tk.Checkbutton(
            right,
            text="Enable logging for this entity",
            variable=self.enabled_var,
            bg=panel,
            fg=fg2,
            selectcolor=panel,
            activebackground=panel,
            activeforeground=fg2,
            command=lambda: self._set_dirty(True),
        ).pack(anchor="w", pady=(0, 10))

        def labeled_entry(label, var):
            tk.Label(right, text=label, bg=panel, fg=fg2, font=("Segoe UI", 10, "bold")).pack(anchor="w")
            ent = tk.Entry(right, textvariable=var, bg=input_bg, fg=fg, relief="flat")
            ent.pack(fill="x", pady=(2, 10))
            ent.bind("<KeyRelease>", lambda e: self._set_dirty(True))
            return ent

        labeled_entry("Date field (used to choose daily note)", self.date_field_var)
        labeled_entry("Header template", self.header_var)

        tk.Label(right, text="Body template (optional)", bg=panel, fg=fg2, font=("Segoe UI", 10, "bold")).pack(anchor="w")
        self.body_text = tk.Text(right, bg=input_bg, fg=fg, height=12, wrap="word")
        self.body_text.pack(fill="both", expand=True, pady=(2, 10))
        self.body_text.bind("<KeyRelease>", lambda e: self._set_dirty(True))

        token_help = (
            "Available tokens: {time}, {date}, {created_time}, {created_date}, {created_datetime}, "
            "{start_date}, {start_time}, {start_datetime}, {title}, {entity_id}, {label}, plus any payload fields."
        )
        tk.Label(right, text=token_help, bg=panel, fg=meta, font=("Segoe UI", 9), wraplength=650, justify="left").pack(
            anchor="w", pady=(0, 8)
        )

        tk.Button(
            right,
            text="Apply to Selected Entity",
            command=self._save_current_rule,
            bg=btn_bg,
            fg=fg,
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(anchor="e")

    def _load_entities(self):
        self.entity_list.delete(0, "end")
        for key in self.entity_keys:
            label = SHIFT_LOG_ENTITY_REGISTRY.get(key, {}).get("label", key)
            self.entity_list.insert("end", f"{label} ({key})")
        if self.entity_keys:
            self.entity_list.selection_set(0)
            self._load_entity(self.entity_keys[0])

    def _on_select_entity(self, _event=None):
        sel = list(self.entity_list.curselection())
        if not sel:
            return
        self._save_current_rule()
        idx = sel[0]
        if 0 <= idx < len(self.entity_keys):
            self._load_entity(self.entity_keys[idx])

    def _load_entity(self, entity_key: str):
        self._current_entity = entity_key
        rule = (self.rules.get("entities") or {}).get(entity_key, {})
        self.enabled_var.set(bool(rule.get("enabled", True)))
        self.date_field_var.set(rule.get("date_field", "created"))
        self.header_var.set(rule.get("header_template", "### {time} - New Note Created - {title}"))
        self.body_text.delete("1.0", "end")
        self.body_text.insert("1.0", rule.get("body_template", ""))
        self._set_dirty(False)

    def _save_current_rule(self):
        if not self._current_entity:
            return
        rule = (self.rules.get("entities") or {}).setdefault(self._current_entity, {})
        rule["enabled"] = bool(self.enabled_var.get())
        rule["date_field"] = self.date_field_var.get().strip()
        rule["header_template"] = self.header_var.get().strip()
        rule["body_template"] = self.body_text.get("1.0", "end").strip()
        self._set_dirty(True)

    def _save_to_disk(self):
        try:
            self._save_current_rule()
            if not self.cfg_path.exists():
                raise FileNotFoundError(f"config.json not found at: {self.cfg_path}")
            _write_config_file_with_backup(self.cfg_path, self.cfg_ref)
            self._set_dirty(False)
            self._status_var.set("Saved config.json (backup created).")
        except Exception as e:
            messagebox.showerror("Save failed", str(e))

    def _reload_from_disk(self):
        if self._dirty:
            if not messagebox.askyesno("Reload", "Discard unsaved changes and reload from disk?"):
                return
        try:
            if not self.cfg_path.exists():
                raise FileNotFoundError(f"config.json not found at: {self.cfg_path}")
            fresh = _read_config_file(self.cfg_path)
            self.cfg_ref.clear()
            self.cfg_ref.update(fresh)
            self.rules = _ensure_shift_log_rules_schema(self.cfg_ref)
            self._current_entity = None
            self._set_dirty(False)
            self._status_var.set("Reloaded from disk.")
            self._load_entities()
        except Exception as e:
            messagebox.showerror("Reload failed", str(e))

    def _reset_defaults(self):
        if not messagebox.askyesno("Reset Defaults", "Reset all Shift Log rules to defaults?"):
            return
        self.cfg_ref["shift_log_rules"] = copy.deepcopy(_default_shift_log_rules())
        self.rules = _ensure_shift_log_rules_schema(self.cfg_ref)
        self._current_entity = None
        self._set_dirty(True)
        self._load_entities()

    def _on_close(self):
        if self._dirty:
            if not messagebox.askyesno("Close", "Discard unsaved changes?"):
                return
        self.destroy()

# ============================================================
# Launcher integration: add this method in LauncherApp
# ============================================================
# Add inside class LauncherApp:
#
#   def _open_manage_keywords(self):
#       app_dir = get_app_dir()
#       # optional: pass theme colors from cfg["ui"] if desired
#       ManageKeywordsWindow(self, self.cfg, app_dir)
#
# And add to _dispatch():
#
#   elif action == "open_manage_keywords":
#       self._open_manage_keywords()
#       self._record_recent(item)
#
# Finally, add a menu button in config.json:
#   { "label": "Manage Keywords", "description": "...", "action": "open_manage_keywords", "target": "" }
# ============================================================

# ------------------------------------------------------------
# PyInstaller-friendly path resolution
# ------------------------------------------------------------
def get_app_dir() -> Path:
    if getattr(sys, "frozen", False):
        return Path(sys.executable).resolve().parent
    return Path(__file__).resolve().parent


# ------------------------------------------------------------
# Config loading + token expansion
# ------------------------------------------------------------
_TOKEN_RE = re.compile(r"\{([A-Z0-9_]+)\}")


def load_config(app_dir: Path) -> dict:
    cfg_path = app_dir / "config.json"
    if not cfg_path.exists():
        raise FileNotFoundError(f"Missing config.json next to application: {cfg_path}")
    return json.loads(cfg_path.read_text(encoding="utf-8"))


def expand_value(value: str, tokens: dict) -> str:
    if value is None:
        return ""
    value = os.path.expandvars(value)

    def repl(m):
        key = m.group(1)
        return str(tokens.get(key, m.group(0)))

    return _TOKEN_RE.sub(repl, value)


def build_tokens(app_dir: Path, cfg: dict) -> dict:
    tokens = {
        "APP_DIR": str(app_dir),
        "USER_HOME": str(Path.home()),
        "USERNAME": os.environ.get("USERNAME") or os.environ.get("USER") or "",
        "LOCALAPPDATA": os.environ.get("LOCALAPPDATA", ""),
        "APPDATA": os.environ.get("APPDATA", ""),
    }

    for k, v in cfg.get("paths", {}).items():
        tokens[k] = expand_value(v, tokens)

    return tokens


def resolve_path(path_str: str, tokens: dict) -> str:
    expanded = expand_value(path_str, tokens)
    expanded = expanded.replace("/", os.sep).replace("\\", os.sep)
    return os.path.normpath(expanded)

def _yaml_quote(s: str) -> str:
    """Return a safely double-quoted YAML scalar."""
    if s is None:
        return '""'
    s = str(s)
    s = s.replace("\\", "\\\\").replace('"', '\\"')
    return f"\"{s}\""


def _apply_yaml_updates(text: str, updates: dict) -> str:
    """Update YAML frontmatter in `text` using `updates` (supports scalars + list blocks)."""
    if not updates:
        return text

    raw = text or ""
    bom = raw.startswith("\ufeff")
    work = raw.lstrip("\ufeff")

    def _render_scalar(val):
        if val is None:
            return ""
        if isinstance(val, bool):
            return "true" if val else "false"
        if isinstance(val, (int, float)) and not isinstance(val, bool):
            return str(int(val)) if float(val).is_integer() else str(val)
        return _yaml_quote(str(val))

    def _clean_list(items):
        out = []
        for it in items:
            if it is None:
                continue
            s2 = str(it).strip()
            if s2:
                out.append(s2)
        return out

    def _ensure_frontmatter(t: str) -> str:
        s2 = t.lstrip("\ufeff")
        if s2.startswith("---"):
            return t
        return "---\n---\n\n" + t.lstrip("\n")

    work = _ensure_frontmatter(work)

    lines = work.splitlines()
    end_idx = None
    for i in range(1, len(lines)):
        if lines[i].strip() == "---":
            end_idx = i
            break
    if end_idx is None:
        return raw

    fm_lines = lines[: end_idx + 1]
    body_lines = lines[end_idx + 1 :]

    def _rebuild_key_map():
        km = {}
        for ii in range(1, end_idx):
            ln = fm_lines[ii]
            if not ln.strip() or ln.lstrip().startswith("#"):
                continue
            if ":" in ln:
                k = ln.split(":", 1)[0].strip()
                if k:
                    km[k.lower()] = ii
        return km

    key_map = _rebuild_key_map()

    def _remove_block(k_lower: str):
        nonlocal end_idx, key_map
        if k_lower not in key_map:
            return None
        i0 = key_map[k_lower]
        j = i0 + 1
        while j < end_idx:
            nxt = fm_lines[j]
            if nxt.startswith("  ") or nxt.startswith("\t"):
                j += 1
                continue
            break
        del fm_lines[i0:j]
        end_idx -= (j - i0)
        key_map = _rebuild_key_map()
        return i0

    for k, v in updates.items():
        key = str(k).strip() if k is not None else ""
        if not key:
            continue
        k_l = key.lower()

        if isinstance(v, (list, tuple)):
            insert_at = _remove_block(k_l)
            if insert_at is None:
                insert_at = end_idx
            items = _clean_list(v)
            new_block = [f"{key}:"]
            for item in items:
                new_block.append(f"  - {_yaml_quote(item)}")
            fm_lines[insert_at:insert_at] = new_block
            end_idx += len(new_block)
            key_map = _rebuild_key_map()
        else:
            new_line = f"{key}: {_render_scalar(v)}".rstrip()
            if k_l in key_map:
                fm_lines[key_map[k_l]] = new_line
            else:
                fm_lines.insert(end_idx, new_line)
                end_idx += 1
                key_map = _rebuild_key_map()

    merged = "\n".join(fm_lines + body_lines)
    return ("\ufeff" + merged) if bom else merged

def open_obsidian_file(vault_root: Path, note_path: Path):
    """
    Open a note in Obsidian using the obsidian:// URI scheme.
    """
    try:
        vault_name = vault_root.name
        rel = note_path.relative_to(vault_root).as_posix()
        uri = f"obsidian://open?vault={quote(vault_name)}&file={quote(rel)}"
        webbrowser.open(uri)
    except Exception:
        # Fail silently; do not fall back to OS default editor
        pass

# ------------------------------------------------------------
# Persistent state (Recents + RSS triage)
# ------------------------------------------------------------
def get_state_path(app_name: str = "SCOUT_Launcher") -> Path:
    """
    Distribution-safe state location:
      - Windows: %LOCALAPPDATA%\SCOUT_Launcher\state.json
      - Others:  ~/.scout_launcher/state.json
    """
    if sys.platform.startswith("win"):
        base = Path(os.environ.get("LOCALAPPDATA", str(Path.home())))
        folder = base / app_name
    else:
        folder = Path.home() / ".scout_launcher"
    folder.mkdir(parents=True, exist_ok=True)
    return folder / "state.json"


def load_state(state_path: Path) -> dict:
    if not state_path.exists():
        return {"recent": [], "rss_triage": {}, "rss_article_state": {}, "rss_report_queue": []}
    try:
        st = json.loads(state_path.read_text(encoding="utf-8"))
        if "recent" not in st or not isinstance(st["recent"], list):
            st["recent"] = []
        if "rss_triage" not in st or not isinstance(st["rss_triage"], dict):
            st["rss_triage"] = {}
        if "rss_article_state" not in st or not isinstance(st["rss_article_state"], dict):
            st["rss_article_state"] = {}
        if "rss_report_queue" not in st or not isinstance(st["rss_report_queue"], list):
            st["rss_report_queue"] = []
        return st
    except Exception:
        return {"recent": [], "rss_triage": {}, "rss_article_state": {}, "rss_report_queue": []}


def save_state(state_path: Path, state: dict) -> None:
    try:
        state_path.write_text(json.dumps(state, indent=2), encoding="utf-8")
    except Exception:
        pass


def save_json(path: Path, data: dict) -> None:
    """Write a dict as JSON to disk (UTF-8)."""
    try:
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps(data, indent=2), encoding="utf-8")
    except Exception:
        # Best-effort persistence; callers should not fail UI flow if draft can't be saved.
        pass


def item_signature(item: dict) -> str:
    sec = (item.get("section") or "").strip()
    label = (item.get("label") or "").strip()
    action = (item.get("action") or "").strip()
    target = (item.get("target") or "")
    return f"{sec}||{label}||{action}||{target}"


# ------------------------------------------------------------
# General helpers
# ------------------------------------------------------------
def open_url(url: str) -> None:
    webbrowser.open(url)


def open_obsidian_uri(uri: str) -> None:
    if not uri or not uri.startswith("obsidian://"):
        messagebox.showerror("Obsidian URI Error", f"Invalid Obsidian URI:\n{uri}")
        return

    if sys.platform.startswith("win"):
        try:
            os.startfile(uri)
            return
        except Exception:
            pass

        try:
            subprocess.Popen(["cmd", "/c", "start", "", uri], shell=False)
            return
        except Exception as e:
            messagebox.showerror("Obsidian URI Error", f"Failed to open Obsidian URI:\n{uri}\n\n{e}")
            return

    try:
        ok = webbrowser.open(uri)
        if not ok:
            raise RuntimeError("webbrowser.open returned False")
    except Exception as e:
        messagebox.showerror("Obsidian URI Error", f"Failed to open Obsidian URI:\n{uri}\n\n{e}")



def open_note_in_obsidian(vault_name: str, vault_root: Path, note_path: Path) -> bool:
    """
    Open a note inside Obsidian (not via OS default markdown handler).
    Returns True if an Obsidian URI was attempted, False if insufficient info.
    """
    try:
        if not vault_name or vault_root is None or note_path is None:
            return False
        rel = note_path.relative_to(vault_root).as_posix()
        uri = f"obsidian://open?vault={quote(vault_name)}&file={quote(rel)}"
        open_obsidian_uri(uri)
        return True
    except Exception:
        return False

def open_path(path: str) -> None:
    path = os.path.expandvars(path)
    if not os.path.exists(path):
        messagebox.showerror("Not found", f"Path does not exist:\n{path}")
        return

    if sys.platform.startswith("win"):
        os.startfile(path)
    elif sys.platform == "darwin":
        subprocess.run(["open", path], check=False)
    else:
        subprocess.run(["xdg-open", path], check=False)


def run_command(cmd: str) -> None:
    try:
        subprocess.Popen(cmd, shell=True)
    except Exception as e:
        messagebox.showerror("Launch failed", str(e))


def _safe_filename(s: str, max_len: int = 120) -> str:
    s = (s or "").strip()
    s = re.sub(r"[\\/:*?\"<>|]+", "-", s)
    s = re.sub(r"\s+", " ", s).strip()
    if len(s) > max_len:
        s = s[:max_len].rstrip()
    return s or "Untitled"


def _hash_key(s: str) -> str:
    return hashlib.sha256((s or "").encode("utf-8", errors="ignore")).hexdigest()[:20]


def _fmt_when(value) -> str:
    """
    Formats a timestamp or datetime-like value for UI display.
    Accepts ISO strings, datetime objects, or None.
    """
    if not value:
        return ""
    if isinstance(value, datetime):
        return value.strftime("%Y-%m-%d %H:%M")
    try:
        dt = datetime.fromisoformat(str(value).replace("Z", ""))
        return dt.strftime("%Y-%m-%d %H:%M")
    except Exception:
        return str(value)


# ------------------------------------------------------------
# Markdown builder (FIX: _mk_markdown_from_item not defined)
# ------------------------------------------------------------
def _safe_str(v) -> str:
    if v is None:
        return ""
    try:
        return str(v)
    except Exception:
        return ""


def _first_nonempty(*vals) -> str:
    for v in vals:
        s = _safe_str(v).strip()
        if s:
            return s
    return ""


def _mk_markdown_from_item(item: dict, tags=None, notes: str = "") -> str:
    """
    Build an Obsidian-friendly markdown note from a DB/article row.

    Accepts optional:
      - tags: list[str] or comma string
      - notes: analyst notes string
    """
    item = item or {}

    title = _first_nonempty(item.get("title"), item.get("headline"), "Untitled")
    url = _first_nonempty(item.get("url"), item.get("link"))
    source = _first_nonempty(item.get("source"), item.get("feed"), item.get("publisher"))
    published = _first_nonempty(
        item.get("published"),
        item.get("published_at"),
        item.get("pub_date"),
        item.get("created"),
        item.get("created_at"),
    )
    summary = _first_nonempty(item.get("summary"), item.get("description"), item.get("content"))

    # normalize tags
    norm_tags = []
    if isinstance(tags, (list, tuple)):
        norm_tags = [str(t).strip() for t in tags if str(t).strip()]
    elif isinstance(tags, str):
        norm_tags = [t.strip() for t in tags.split(",") if t.strip()]
    else:
        tags_val = item.get("tags")
        if isinstance(tags_val, (list, tuple)):
            norm_tags = [str(t).strip() for t in tags_val if str(t).strip()]
        elif isinstance(tags_val, str):
            norm_tags = [t.strip() for t in tags_val.split(",") if t.strip()]

    def esc(s: str) -> str:
        return (s or "").replace('"', r"\"")

    yaml_lines = [
        "---",
        "entity_type: cti_article",
        f'title: "{esc(title)}"',
    ]
    if source:
        yaml_lines.append(f'source: "{esc(source)}"')
    if published:
        yaml_lines.append(f'published: "{esc(published)}"')
    if url:
        yaml_lines.append(f'url: "{esc(url)}"')
    if norm_tags:
        yaml_lines.append("tags: [" + ", ".join(f'"{esc(t)}"' for t in norm_tags) + "]")
    yaml_lines.append("---")

    body = []
    body.append(f"# {title}\n")
    if url:
        body.append(f"**Link:** {url}\n")
    if source:
        body.append(f"**Source:** {source}\n")
    if published:
        body.append(f"**Published:** {published}\n")

    if summary:
        body.append("\n## Summary\n")
        body.append(summary.strip() + "\n")

    if notes and notes.strip():
        body.append("\n## Analyst Notes\n")
        body.append(notes.strip() + "\n")

    return "\n".join(yaml_lines) + "\n\n" + "".join(body)


def _rss_fp(item: dict) -> str:
    """
    Stable fingerprint for state tracking.
    Prefer URL; fallback to source||title.
    """
    url = _safe_str(item.get("url")).strip()
    if url:
        return url
    src = _safe_str(item.get("source")).strip()
    title = _safe_str(item.get("title")).strip()
    return f"{src}||{title}"


# ------------------------------------------------------------
# RSS helpers (run_rss_collect support)
# ------------------------------------------------------------
_RSS_FEEDS_RE = re.compile(r"\bFeeds processed:\s*(\d+)\b", re.IGNORECASE)
_RSS_NEW_RE = re.compile(r"\bNew items inserted:\s*(\d+)\b", re.IGNORECASE)


def parse_rss_output(text: str) -> dict:
    feeds = None
    new_items = None

    m1 = _RSS_FEEDS_RE.search(text or "")
    if m1:
        try:
            feeds = int(m1.group(1))
        except Exception:
            feeds = None

    m2 = _RSS_NEW_RE.search(text or "")
    if m2:
        try:
            new_items = int(m2.group(1))
        except Exception:
            new_items = None

    return {"feeds_processed": feeds, "new_inserted": new_items}


def run_rss_collect_and_capture(app_dir: Path, python_cmd: str = "python") -> dict:
    """
    Runs: <python_cmd> scout_rss.py collect
    Assumes scout_rss.py is next to the launcher (APP_DIR).
    """
    script_path = app_dir / "scout_rss.py"
    started = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    if not script_path.exists():
        return {
            "ok": False,
            "started": started,
            "ended": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "exit_code": 127,
            "error": f"Missing scout_rss.py at: {script_path}",
            "feeds_processed": None,
            "new_inserted": None,
            "stdout": "",
            "stderr": "",
            "cmd": f'{python_cmd} "{script_path}" collect',
        }

    cmd = [python_cmd, str(script_path), "collect"]

    try:
        p = subprocess.run(
            cmd,
            cwd=str(app_dir),
            capture_output=True,
            text=True,
            shell=False,
        )
        ended = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        combined = (p.stdout or "") + "\n" + (p.stderr or "")
        parsed = parse_rss_output(combined)

        ok = (p.returncode == 0)
        return {
            "ok": ok,
            "started": started,
            "ended": ended,
            "exit_code": p.returncode,
            "error": "" if ok else (p.stderr or "RSS collect failed"),
            "feeds_processed": parsed["feeds_processed"],
            "new_inserted": parsed["new_inserted"],
            "stdout": p.stdout or "",
            "stderr": p.stderr or "",
            "cmd": " ".join(cmd),
        }
    except FileNotFoundError as e:
        return {
            "ok": False,
            "started": started,
            "ended": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "exit_code": 127,
            "error": f"Python not found: {python_cmd} ({e})",
            "feeds_processed": None,
            "new_inserted": None,
            "stdout": "",
            "stderr": "",
            "cmd": " ".join(cmd),
        }
    except Exception as e:
        return {
            "ok": False,
            "started": started,
            "ended": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "exit_code": 1,
            "error": str(e),
            "feeds_processed": None,
            "new_inserted": None,
            "stdout": "",
            "stderr": "",
            "cmd": " ".join(cmd),
        }


# ------------------------------------------------------------
# RSS Review / Import (SQLite auto-detect + triage store fallback)
# ------------------------------------------------------------
def _guess_rss_db_path(cfg: dict, tokens: dict, app_dir: Path) -> str:
    """
    Finds the most likely RSS SQLite DB.
    Priority:
      1) paths.RSS_DB if provided
      2) common filenames in APP_DIR / SCOUT_ROOT / SCOUT_ROOT\data
      3) score-based scan of *.db in those locations (non-recursive; fast)
    """
    paths = cfg.get("paths", {}) or {}

    # 1) Explicit config wins
    if "RSS_DB" in paths:
        p = resolve_path(paths.get("RSS_DB", ""), tokens)
        if p and os.path.exists(p):
            return p

    scout_root = resolve_path(paths.get("SCOUT_ROOT", ""), tokens)

    # 2) Common guesses
    candidates = []
    if scout_root:
        candidates += [
            os.path.join(scout_root, "rss.db"),
            os.path.join(scout_root, "cti_rss.db"),
            os.path.join(scout_root, "scout_rss.db"),
            os.path.join(scout_root, "data", "rss.db"),
            os.path.join(scout_root, "data", "cti_rss.db"),
            os.path.join(scout_root, "data", "scout_rss.db"),
        ]
    candidates += [
        str(app_dir / "rss.db"),
        str(app_dir / "cti_rss.db"),
        str(app_dir / "scout_rss.db"),
        str(app_dir / "data" / "rss.db"),
        str(app_dir / "data" / "cti_rss.db"),
        str(app_dir / "data" / "scout_rss.db"),
    ]

    for c in candidates:
        if c and os.path.exists(c):
            return c

    # 3) Score-based scan (fast, non-recursive)
    scan_dirs = []
    if scout_root and os.path.isdir(scout_root):
        scan_dirs.append(scout_root)
        d = os.path.join(scout_root, "data")
        if os.path.isdir(d):
            scan_dirs.append(d)
    scan_dirs.append(str(app_dir))
    d2 = str(app_dir / "data")
    if os.path.isdir(d2):
        scan_dirs.append(d2)

    best = ("", -1)
    for d in scan_dirs:
        try:
            for fn in os.listdir(d):
                if not fn.lower().endswith(".db"):
                    continue
                p = os.path.join(d, fn)
                if not os.path.isfile(p):
                    continue
                score = _score_db_for_rss(p)
                if score > best[1]:
                    best = (p, score)
        except Exception:
            continue

    return best[0] if best[1] >= 1 else ""


def _quote_ident(name: str) -> str:
    if name is None:
        return ""
    return '"' + str(name).replace('"', '""') + '"'


def _detect_rss_table_and_columns(db_path: str, return_meta: bool = False):
    """
    Attempts to find a table containing RSS items.
    Returns:
      (table_name, colmap) or (table_name, colmap, meta) if return_meta=True
    """
    if not db_path or not os.path.exists(db_path):
        if return_meta:
            return None, {}, {"db_path": db_path, "reason": "missing db", "score": 0}
        return None, {}

    want = {
        "id": ["id", "item_id", "pk", "uuid"],
        "title": ["title", "headline", "name"],
        "url": ["url", "link", "guid", "permalink"],
        "source": ["source", "feed", "feed_title", "publisher", "site"],
        "published": ["published", "pub_date", "date", "published_at", "created_at", "ts"],
        "summary": ["summary", "description", "excerpt"],
        "content": ["content", "body", "full_text", "text"],
        "is_new": ["is_new", "new", "unread", "status"],
        "triaged": ["triaged", "reviewed", "handled", "imported"],
    }

    def cols_for_table(conn, tname: str):
        cur = conn.execute(f"PRAGMA table_info({_quote_ident(tname)})")
        return [r[1] for r in cur.fetchall()]

    def match_map(colnames: list):
        colset = {c.lower(): c for c in colnames}
        cmap = {}
        score = 0

        for key, aliases in want.items():
            hit = None
            for a in aliases:
                if a in colset:
                    hit = colset[a]
                    break
            cmap[key] = hit

        if cmap.get("title"):
            score += 6
        if cmap.get("url"):
            score += 6
        if cmap.get("published"):
            score += 2
        if cmap.get("source"):
            score += 2
        if cmap.get("summary") or cmap.get("content"):
            score += 2
        if cmap.get("triaged") or cmap.get("is_new"):
            score += 1

        return score, cmap

    meta = {"db_path": db_path, "tables": [], "score": 0, "table": None, "columns": {}}

    try:
        conn = sqlite3.connect(db_path)
        try:
            cur = conn.execute("SELECT name FROM sqlite_master WHERE type='table' AND name NOT LIKE 'sqlite_%'")
            tables = [r[0] for r in cur.fetchall()]
            best = (0, None, {})

            for t in tables:
                try:
                    colnames = cols_for_table(conn, t)
                    sc, cmap = match_map(colnames)

                    t_low = (t or "").lower()
                    if any(k in t_low for k in ("rss", "item", "news", "article", "feed")):
                        sc += 2

                    try:
                        rc = conn.execute(f"SELECT COUNT(1) FROM {_quote_ident(t)}").fetchone()[0]
                    except Exception:
                        rc = 0
                    if rc == 0:
                        sc -= 3
                    elif rc > 0:
                        sc += 1

                    meta["tables"].append({"table": t, "cols": colnames, "score": sc, "rows": rc})

                    if sc > best[0]:
                        best = (sc, t, cmap)
                except Exception:
                    continue

            if best[1] and best[0] >= 8:
                meta["score"] = best[0]
                meta["table"] = best[1]
                meta["columns"] = best[2]
                if return_meta:
                    return best[1], best[2], meta
                return best[1], best[2]

            meta["reason"] = "no sufficiently-scored table"
            if return_meta:
                return None, {}, meta
            return None, {}
        finally:
            conn.close()
    except Exception as e:
        meta["reason"] = f"sqlite error: {e}"
        if return_meta:
            return None, {}, meta
        return None, {}


def _score_db_for_rss(db_path: str) -> int:
    try:
        tname, cmap, meta = _detect_rss_table_and_columns(db_path, return_meta=True)
        if not tname:
            return 0
        score = meta.get("score", 0)
        try:
            mtime = os.path.getmtime(db_path)
            age_hours = (datetime.now().timestamp() - mtime) / 3600.0
            if age_hours < 48:
                score += 2
        except Exception:
            pass
        return int(score)
    except Exception:
        return 0


def _fetch_rss_items(db_path: str, limit: int = 250):
    """
    Returns list of dict items from SQLite if detectable, else [].
    Also returns meta describing what was detected.
    """
    tname, cmap, meta = _detect_rss_table_and_columns(db_path, return_meta=True)
    if not tname:
        return [], meta

    keys = ["id", "title", "url", "source", "published", "summary", "content", "is_new", "triaged"]

    sel_exprs = []
    for k in keys:
        col = cmap.get(k)
        sel_exprs.append(_quote_ident(col) if col else "NULL")

    if cmap.get("published"):
        order = f"{_quote_ident(cmap['published'])} DESC"
    elif cmap.get("id"):
        order = f"{_quote_ident(cmap['id'])} DESC"
    else:
        order = "ROWID DESC"

    sql = f"SELECT {', '.join(sel_exprs)} FROM {_quote_ident(tname)} ORDER BY {order} LIMIT ?"

    out = []
    try:
        conn = sqlite3.connect(db_path)
        try:
            cur = conn.execute(sql, (int(limit),))
            rows = cur.fetchall()
            for r in rows:
                item = dict(zip(keys, r))
                for kk in ("title", "url", "source", "summary", "content"):
                    item[kk] = "" if item.get(kk) is None else str(item.get(kk) or "")
                item["published"] = "" if item.get("published") is None else str(item.get("published") or "")
                out.append(item)
        finally:
            conn.close()
    except Exception as e:
        meta["reason"] = f"query error: {e}"
        meta["sql"] = sql
        return [], meta

    meta["sql"] = sql
    meta["table"] = tname
    meta["columns"] = cmap
    meta["rows_returned"] = len(out)
    return out, meta


def _triage_get(state: dict, url: str) -> dict:
    tri = state.get("rss_triage") or {}
    k = _hash_key(url)
    return tri.get(k, {})


def _triage_set(state: dict, url: str, patch: dict) -> None:
    if "rss_triage" not in state or not isinstance(state["rss_triage"], dict):
        state["rss_triage"] = {}
    tri = state["rss_triage"]
    k = _hash_key(url)
    cur = tri.get(k, {})
    if not isinstance(cur, dict):
        cur = {}
    cur.update(patch or {})
    tri[k] = cur


def _count_new_items(state: dict, items: list) -> int:
    n = 0
    for it in items or []:
        url = (it.get("url") or "").strip()
        if not url:
            continue
        tri = _triage_get(state, url)
        if not tri.get("triaged", False) and not tri.get("imported", False):
            n += 1
    return n

def _cfg_get_rss_schedule(cfg: dict) -> dict:
    rss = cfg.get("rss", {}) or {}
    sch = rss.get("schedule", {}) or {}
    out = {
        "enabled": bool(sch.get("enabled", False)),
        "mode": str(sch.get("mode", "interval")),
        "interval_minutes": int(sch.get("interval_minutes", 60) or 60),
        "run_at_startup": bool(sch.get("run_at_startup", False)),
        "quiet_hours": {
            "enabled": bool((sch.get("quiet_hours", {}) or {}).get("enabled", False)),
            "start": str((sch.get("quiet_hours", {}) or {}).get("start", "22:00")),
            "end": str((sch.get("quiet_hours", {}) or {}).get("end", "06:00")),
        },
    }
    # Guardrails
    if out["interval_minutes"] < 5:
        out["interval_minutes"] = 5
    if out["interval_minutes"] > 1440:
        out["interval_minutes"] = 1440
    return out


def _time_in_quiet_hours(now: datetime, start_hhmm: str, end_hhmm: str) -> bool:
    """
    Quiet hours can wrap over midnight (e.g., 22:00 -> 06:00).
    start/end are 'HH:MM' 24h strings.
    """
    def parse_hm(s: str):
        s = (s or "").strip()
        hh, mm = s.split(":")
        return int(hh), int(mm)

    try:
        sh, sm = parse_hm(start_hhmm)
        eh, em = parse_hm(end_hhmm)
    except Exception:
        return False  # invalid quiet hours -> ignore

    start = now.replace(hour=sh, minute=sm, second=0, microsecond=0)
    end = now.replace(hour=eh, minute=em, second=0, microsecond=0)

    if start <= end:
        return start <= now <= end

    # wraps midnight
    return now >= start or now <= end

def _default_inbox_rel(cfg: dict) -> str:
    rr = cfg.get("rss_review", {}) or {}
    inbox_rel = (rr.get("inbox_rel") or "").strip().replace("\\", "/")
    if inbox_rel:
        return inbox_rel
    return "00_System/99_Inbox/_staging"


def _write_obsidian_note(cfg: dict, tokens: dict, title: str, url: str, source: str, published: str, summary: str, kind: str) -> str:
    """
    Creates a markdown note in the staging inbox and returns absolute path.
    """
    def esc(s: str) -> str:
        return (s or "").replace('"', r'\\"')

    paths = cfg.get("paths", {}) or {}
    scout_root = resolve_path(paths.get("SCOUT_ROOT", ""), tokens)

    if not scout_root or not os.path.isdir(scout_root):
        raise ValueError("SCOUT_ROOT is not set to a valid directory in config.json.")

    inbox_rel = _default_inbox_rel(cfg)
    folder_abs = os.path.join(scout_root, *inbox_rel.split("/"))
    os.makedirs(folder_abs, exist_ok=True)

    today = date.today().strftime("%Y-%m-%d")
    safe = _safe_filename(title)
    fname = f"{today} - {safe}.md"
    abs_path = os.path.join(folder_abs, fname)

    if os.path.exists(abs_path):
        stem = f"{today} - {safe}"
        i = 2
        while True:
            cand = os.path.join(folder_abs, f"{stem} ({i}).md")
            if not os.path.exists(cand):
                abs_path = cand
                break
            i += 1

    frontmatter = [
        "---",
        "entity_type: rss_article",
        f"import_kind: {kind}",
        f"source: \"{esc(source)}\"",
        f"url: \"{esc(url)}\"",
        f"published: \"{esc(published)}\"",
        f"imported_at: \"{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\"",
        "tags: [cti, rss]",
        "---",
        "",
    ]
    body = [
        f"# {title}",
        "",
        f"**Source:** {source}" if source else "**Source:**",
        f"**Published:** {published}" if published else "**Published:**",
        f"**URL:** {url}" if url else "**URL:**",
        "",
        "## Summary",
        "",
        (summary or "").strip() or "_No summary captured._",
        "",
        "## Analyst Notes",
        "",
        "",
    ]

    with open(abs_path, "w", encoding="utf-8") as f:
        f.write("\n".join(frontmatter + body))

    return abs_path


def _open_note_in_obsidian(cfg: dict, tokens: dict, abs_note_path: str) -> None:
    paths = cfg.get("paths", {}) or {}
    vault_name = expand_value(paths.get("OBSIDIAN_VAULT", ""), tokens).strip()
    scout_root = resolve_path(paths.get("SCOUT_ROOT", ""), tokens)

    if not vault_name:
        raise ValueError("OBSIDIAN_VAULT is not set in config.json.")
    if not scout_root:
        raise ValueError("SCOUT_ROOT is not set in config.json.")

    rel_path = os.path.relpath(abs_note_path, scout_root).replace("\\", "/")
    uri = (
        "obsidian://advanced-uri?"
        f"vault={quote(vault_name, safe='')}"
        f"&filepath={quote(rel_path, safe='')}"
        f"&viewmode=source"
    )
    open_obsidian_uri(uri)


# ------------------------------------------------------------
# LLM helpers (OpenAI)
# ------------------------------------------------------------
def _llm_config(cfg: dict, tokens: dict) -> dict:
    llm = cfg.get("llm", {}) or {}
    provider = str(llm.get("provider", "")).strip().lower()
    api_key = expand_value(str(llm.get("api_key", "")), tokens).strip()
    model = str(llm.get("model", "")).strip() or "gpt-4o"
    return {"provider": provider, "api_key": api_key, "model": model}


def _openai_chat_json(api_key: str, model: str, messages: list, timeout: int = 60) -> dict:
    if requests is None:
        raise ValueError("The 'requests' package is not available.")
    if not api_key:
        raise ValueError("OpenAI API key is not set in config.json (llm.api_key).")

    url = "https://api.openai.com/v1/chat/completions"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {
        "model": model or "gpt-4o",
        "messages": messages,
        "temperature": 0.2,
        "response_format": {"type": "json_object"},
    }

    try:
        resp = requests.post(url, headers=headers, json=payload, timeout=timeout)
    except Exception as e:
        raise ValueError(f"OpenAI request failed: {e}")

    if resp.status_code >= 400:
        raise ValueError(f"OpenAI request failed ({resp.status_code}): {resp.text[:2000]}")

    data = resp.json()
    content = (
        (data.get("choices") or [{}])[0]
        .get("message", {})
        .get("content", "")
    )
    if not content:
        raise ValueError("OpenAI response was empty.")

    try:
        return json.loads(content)
    except Exception as e:
        raise ValueError(f"Failed to parse OpenAI JSON response: {e}")


def _openai_chat_text(api_key: str, model: str, messages: list, timeout: int = 120) -> str:
    if requests is None:
        raise ValueError("The 'requests' package is not available.")
    if not api_key:
        raise ValueError("OpenAI API key is not set in config.json (llm.api_key).")

    url = "https://api.openai.com/v1/chat/completions"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {
        "model": model or "gpt-4o",
        "messages": messages,
        "temperature": 0.2,
    }

    try:
        resp = requests.post(url, headers=headers, json=payload, timeout=timeout)
    except Exception as e:
        raise ValueError(f"OpenAI request failed: {e}")

    if resp.status_code >= 400:
        raise ValueError(f"OpenAI request failed ({resp.status_code}): {resp.text[:2000]}")

    data = resp.json()
    content = (
        (data.get("choices") or [{}])[0]
        .get("message", {})
        .get("content", "")
    )
    if not content:
        raise ValueError("OpenAI response was empty.")
    return content.strip()


def _coerce_list(value) -> list[str]:
    if value is None:
        return []
    if isinstance(value, list):
        return [str(v).strip() for v in value if str(v).strip()]
    if isinstance(value, str):
        return [v for v in _split_bulk(value) if v]
    return [str(value).strip()] if str(value).strip() else []


def _write_ioc_enrichment_report(
    cfg: dict,
    tokens: dict,
    title: str,
    sections: dict,
    iocs: list[str],
    context: str,
) -> str:
    paths = cfg.get("paths", {}) or {}
    scout_root = resolve_path(paths.get("SCOUT_ROOT", ""), tokens)
    if not scout_root or not os.path.isdir(scout_root):
        raise ValueError("SCOUT_ROOT is not set to a valid directory in config.json.")

    report_rel = "90_Views/02_Reports/IOC Enrichment"
    report_dir = os.path.join(scout_root, *report_rel.split("/"))
    os.makedirs(report_dir, exist_ok=True)

    today = date.today().strftime("%Y-%m-%d")
    safe_title = _safe_filename(title)
    fname = f"{today} - {safe_title}.md"
    abs_path = os.path.join(report_dir, fname)

    if os.path.exists(abs_path):
        stem = f"{today} - {safe_title}"
        i = 2
        while True:
            cand = os.path.join(report_dir, f"{stem} ({i}).md")
            if not os.path.exists(cand):
                abs_path = cand
                break
            i += 1

    tags = "[audience/analyst, ai/llm]"
    period_start = (date.today() - timedelta(days=7)).strftime("%Y-%m-%d")
    period_end = today

    frontmatter = [
        "---",
        "entity_type: report",
        'report_type: "ioc_summary"',
        f'report_date: "{today}"',
        f'report_period_start: "{period_start}"',
        f'report_period_end: "{period_end}"',
        'report_status: "draft"',
        f"tags: {tags}",
        f"created: {today}",
        f"updated: {today}",
        "banner: 99_Attachments/SCOUT_Obsidian_Banner.png",
        "banner-display: contain",
        "banner-repeat: false",
        "banner-height: 100",
        "content-start: 101",
        "---",
        "",
    ]

    exec_summary = _safe_str(sections.get("executive_summary")).strip() or "_No summary provided._"
    analyst_notes = _safe_str(sections.get("analyst_notes")).strip() or "_No analyst notes provided._"
    key_findings = _coerce_list(sections.get("key_findings"))
    metrics = _coerce_list(sections.get("metrics"))
    actions = _coerce_list(sections.get("actions"))
    references = _coerce_list(sections.get("references"))

    body = [
        f"# {title}",
        "",
        "## Input IOCs",
        *(f"- {ioc}" for ioc in iocs),
        "",
        "## Context",
        context.strip() or "_No additional context provided._",
        "",
        "## 1. Executive Summary",
        exec_summary,
        "",
        "---",
        "## 2. Analyst Notes",
        analyst_notes,
        "",
        "---",
        "## 3. Key Findings",
        *([f"- {v}" for v in key_findings] if key_findings else ["- _None provided._"]),
        "",
        "---",
        "## 4. Metrics and Indicators",
        *([f"- {v}" for v in metrics] if metrics else ["- _None provided._"]),
        "",
        "---",
        "## 5. Actions and Recommendations",
        *([f"- {v}" for v in actions] if actions else ["- _None provided._"]),
        "",
        "---",
        "## 6. References",
        *([f"- {v}" for v in references] if references else ["- _None provided._"]),
        "",
    ]

    with open(abs_path, "w", encoding="utf-8") as f:
        f.write("\n".join(frontmatter + body))

    return abs_path


def _write_project_status_rollup_report(cfg: dict, tokens: dict) -> str:
    paths = cfg.get("paths", {}) or {}
    scout_root = resolve_path(paths.get("SCOUT_ROOT", ""), tokens)
    if not scout_root or not os.path.isdir(scout_root):
        raise ValueError("SCOUT_ROOT is not set to a valid directory in config.json.")

    projects_root = Path(scout_root) / "10_Operations" / "14_Projects"
    if not projects_root.exists():
        raise ValueError("Projects folder does not exist: 10_Operations/14_Projects")

    entries = []
    for p in projects_root.rglob("*.md"):
        try:
            text = p.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            continue
        fm, _ = _pptx_parse_frontmatter(text)
        entity_type = (fm.get("entity_type") or "").strip().lower()
        if entity_type and entity_type != "project":
            continue
        entries.append(
            {
                "path": p,
                "project_id": (fm.get("project_id") or "").strip(),
                "project_name": (fm.get("project_name") or fm.get("title") or p.stem).strip(),
                "project_owner": (fm.get("project_owner") or fm.get("owner") or "").strip(),
                "project_status": (fm.get("project_status") or fm.get("status") or "unknown").strip().lower(),
                "project_priority": (fm.get("project_priority") or fm.get("priority") or "").strip().lower(),
                "start_date": (fm.get("start_date") or "").strip(),
                "end_date": (fm.get("end_date") or "").strip(),
            }
        )

    if not entries:
        raise ValueError("No project notes found in 10_Operations/14_Projects.")

    status_order = ["ideation", "planned", "active", "blocked", "done", "archived", "unknown"]
    priority_order = {"critical": 0, "high": 1, "medium": 2, "low": 3}

    def _sort_key(e: dict):
        p = e.get("project_priority") or ""
        return (priority_order.get(p, 99), e.get("project_name") or "")

    counts = {}
    for e in entries:
        counts[e["project_status"]] = counts.get(e["project_status"], 0) + 1

    today = date.today().strftime("%Y-%m-%d")
    report_rel = "90_Views/02_Reports"
    report_dir = os.path.join(scout_root, *report_rel.split("/"))
    os.makedirs(report_dir, exist_ok=True)
    fname = f"{today} - Project Status Rollup.md"
    abs_path = os.path.join(report_dir, fname)

    if os.path.exists(abs_path):
        stem = f"{today} - Project Status Rollup"
        i = 2
        while True:
            cand = os.path.join(report_dir, f"{stem} ({i}).md")
            if not os.path.exists(cand):
                abs_path = cand
                break
            i += 1

    frontmatter = [
        "---",
        "entity_type: report",
        'report_type: "project_status"',
        f'report_date: "{today}"',
        f"report_count: {len(entries)}",
        'report_status: "draft"',
        'tags: ["projects", "status"]',
        f'created: "{today}"',
        f'updated: "{today}"',
        "banner: 99_Attachments/SCOUT_Obsidian_Banner.png",
        "banner-display: contain",
        "banner-repeat: false",
        "banner-height: 100",
        "content-start: 101",
        "---",
        "",
    ]

    body = [f"# Project Status Rollup — {today}", "", "## Summary", ""]
    for status in status_order:
        if status in counts:
            body.append(f"- {status.title()}: {counts[status]}")
    body.append("")

    for status in status_order:
        items = [e for e in entries if e["project_status"] == status]
        if not items:
            continue
        body.append(f"## {status.title()}")
        body.append("")
        for e in sorted(items, key=_sort_key):
            name = e.get("project_name") or "Untitled Project"
            pid = e.get("project_id")
            owner = e.get("project_owner") or "Unassigned"
            priority = (e.get("project_priority") or "unspecified").title()
            start_date = e.get("start_date") or "unknown"
            end_date = e.get("end_date") or "unknown"
            rel = e["path"].relative_to(Path(scout_root)).as_posix()
            link = f"[[{rel}]]"
            id_part = f"{pid} - " if pid else ""
            body.append(f"- {id_part}{name} | {owner} | {priority} | {start_date} → {end_date} | {link}")
        body.append("")

    with open(abs_path, "w", encoding="utf-8") as f:
        f.write("\n".join(frontmatter + body + [""]))

    return abs_path


def _load_vault_prompt(cfg: dict, tokens: dict, rel_path: str) -> str:
    paths = cfg.get("paths", {}) or {}
    scout_root = resolve_path(paths.get("SCOUT_ROOT", ""), tokens)
    if not scout_root or not os.path.isdir(scout_root):
        raise ValueError("SCOUT_ROOT is not set to a valid directory in config.json.")
    rel = (rel_path or "").strip().replace("\\", "/")
    if not rel:
        raise ValueError("Prompt path is empty.")
    abs_path = os.path.join(scout_root, *rel.split("/"))
    if not os.path.exists(abs_path):
        raise ValueError(f"Prompt file not found: {abs_path}")
    return Path(abs_path).read_text(encoding="utf-8", errors="ignore").strip()


def _write_threat_actor_profile_note(cfg: dict, tokens: dict, label: str, content: str) -> str:
    paths = cfg.get("paths", {}) or {}
    scout_root = resolve_path(paths.get("SCOUT_ROOT", ""), tokens)
    if not scout_root or not os.path.isdir(scout_root):
        raise ValueError("SCOUT_ROOT is not set to a valid directory in config.json.")

    inbox_rel = "00_System/99_Inbox/_staging"
    folder_abs = os.path.join(scout_root, *inbox_rel.split("/"))
    os.makedirs(folder_abs, exist_ok=True)

    today = date.today().strftime("%Y-%m-%d")
    safe_label = _safe_filename(label or "Threat Actor Profile")
    fname = f"{today} - TA-{safe_label}.md"
    abs_path = os.path.join(folder_abs, fname)

    if os.path.exists(abs_path):
        stem = f"{today} - TA-{safe_label}"
        i = 2
        while True:
            cand = os.path.join(folder_abs, f"{stem} ({i}).md")
            if not os.path.exists(cand):
                abs_path = cand
                break
            i += 1

    with open(abs_path, "w", encoding="utf-8") as f:
        f.write((content or "").strip() + "\n")

    return abs_path


def _threat_actor_frontmatter_template(actor_name: str, actor_id: str) -> str:
    today = date.today().strftime("%Y-%m-%d")
    actor_name = actor_name or ""
    actor_id = actor_id or ""
    lines = [
        "---",
        "entity_type: threat_actor",
        f"actor_name: \"{actor_name}\"",
        f"common_name: \"{actor_name}\"",
        f"actor_id: \"{actor_id}\"",
        "actor_type: \"\"",
        "aliases: []",
        "country_of_origin: \"\"",
        "suspected_sponsors: []",
        "attribution_confidence: \"\"",
        "first_seen: \"\"",
        "last_seen: \"\"",
        "status: \"\"",
        "motivations: []",
        "objectives: []",
        "victimology_summary: \"\"",
        "target_sectors: []",
        "target_regions: []",
        "related_groups: []",
        "malware: []",
        "tools: []",
        "infrastructure: []",
        "ttps: []",
        "notable_claims: []",
        "intel_sources: []",
        "tags: []",
        f"created: {today}",
        f"last_modified: {today}",
        "---",
        "",
    ]
    return "\n".join(lines)


# ------------------------------------------------------------
# Daily note helpers
# ------------------------------------------------------------
def _daily_note_compute(cfg: dict, tokens: dict):
    daily_cfg = cfg.get("daily_note", {})
    paths_cfg = cfg.get("paths", {})

    vault_name = expand_value(paths_cfg.get("OBSIDIAN_VAULT", ""), tokens).strip()
    daily_root_rel = expand_value(paths_cfg.get("DAILY_NOTES_REL", ""), tokens).strip().replace("\\", "/")
    scout_root = os.path.normpath(expand_value(paths_cfg.get("SCOUT_ROOT", ""), tokens).strip())

    if not vault_name:
        raise ValueError("OBSIDIAN_VAULT is not set in config.json.")
    if not daily_root_rel:
        raise ValueError("DAILY_NOTES_REL is not set in config.json.")
    if not scout_root or not os.path.isdir(scout_root):
        raise ValueError(f"SCOUT_ROOT is not a valid directory: {scout_root}")

    file_pattern = daily_cfg.get("file_pattern", "%Y-%m-%d")
    ext = daily_cfg.get("extension", ".md")
    create_if_missing = bool(daily_cfg.get("create_if_missing", False))
    template_rel = str(daily_cfg.get("template_file_rel", "")).strip().replace("\\", "/")
    apply_when = str(daily_cfg.get("apply_template_when", "missing")).strip()

    today = date.today()
    leaf = today.strftime(file_pattern) + ext
    file_within_vault = f"{daily_root_rel.rstrip('/')}/{leaf}".replace("\\", "/")

    abs_note_path = os.path.join(scout_root, *file_within_vault.split("/"))

    template_abs = ""
    if template_rel:
        template_abs = os.path.join(scout_root, *template_rel.split("/"))

    return vault_name, scout_root, file_within_vault, abs_note_path, template_abs, create_if_missing, apply_when


def _ensure_daily_note_exists(abs_note_path: str, template_abs: str, apply_when: str, create_if_missing: bool) -> bool:
    note_exists = os.path.exists(abs_note_path)
    note_empty = False
    created = False

    if note_exists:
        try:
            note_empty = os.path.getsize(abs_note_path) == 0
        except Exception:
            note_empty = False

    should_apply_template = (
        (apply_when == "missing" and not note_exists)
        or (apply_when == "missing_or_empty" and (not note_exists or note_empty))
    )

    if not note_exists and not create_if_missing:
        raise FileNotFoundError(f"Daily note does not exist and create_if_missing is false: {abs_note_path}")

    if not note_exists or should_apply_template:
        os.makedirs(os.path.dirname(abs_note_path), exist_ok=True)

    if should_apply_template:
        content = ""
        if template_abs:
            if not os.path.exists(template_abs):
                raise FileNotFoundError(f"Template file not found: {template_abs}")
            with open(template_abs, "r", encoding="utf-8") as f:
                content = f.read()

        if not content.strip():
            today_str = date.today().strftime("%Y-%m-%d")
            content = f"# {today_str}\n\n"

        with open(abs_note_path, "w", encoding="utf-8") as f:
            f.write(content)
        created = True

    elif not note_exists:
        today_str = date.today().strftime("%Y-%m-%d")
        with open(abs_note_path, "w", encoding="utf-8") as f:
            f.write(f"# {today_str}\n\n")
        created = True

    return created


def open_daily_note_live(cfg: dict, tokens: dict) -> None:
    try:
        vault_name, scout_root, file_within_vault, abs_note_path, template_abs, create_if_missing, apply_when = _daily_note_compute(
            cfg, tokens
        )
        created = _ensure_daily_note_exists(abs_note_path, template_abs, apply_when, create_if_missing)

        filepath = file_within_vault
        if not filepath.lower().endswith(".md"):
            filepath += ".md"

        if created:
            opened = open_note_in_obsidian(vault_name, Path(scout_root), Path(abs_note_path))
            if not opened:
                uri = (
                    "obsidian://advanced-uri?"
                    f"vault={quote(vault_name, safe='')}"
                    f"&filepath={quote(filepath, safe='')}"
                    f"&viewmode=live"
                )
                open_obsidian_uri(uri)
        else:
            uri = (
                "obsidian://advanced-uri?"
                f"vault={quote(vault_name, safe='')}"
                f"&filepath={quote(filepath, safe='')}"
                f"&viewmode=live"
            )
            open_obsidian_uri(uri)
    except Exception as e:
        messagebox.showerror("Daily Notes Error", str(e))


# ------------------------------------------------------------
# UI
# ------------------------------------------------------------
class LauncherApp(BaseTk):
    """
    Startup performance improvements applied:
      - Deferred boot: render window fast, then load assets/build UI via after(0).
      - Debounced resize: expensive background resizes happen at most every ~60ms during resize storms.
      - Fast resample during early startup settling, then Lanczos for quality.

    FIXES INCORPORATED:
      1) Maximize-on-launch: repeated, delayed maximize attempts + hard fallback sizing.
      2) Header title/subtitle: true visual centering regardless of left/right logo widths.
    """

    def __init__(self, cfg: dict, tokens: dict, state_path: Path):
        super().__init__()
        self.cfg = cfg
        self.tokens = tokens

        # State (recents + triage + rss_article_state)
        self.state_path = state_path
        self.state = load_state(self.state_path)

        app_cfg = cfg.get("app", {})
        ui_cfg = cfg.get("ui", {})

        self.title(app_cfg.get("title", "Launcher"))

        win_cfg = app_cfg.get("window", {})
        w = int(win_cfg.get("width", 1000))
        h = int(win_cfg.get("height", 600))
        min_w = int(win_cfg.get("min_width", 900))
        min_h = int(win_cfg.get("min_height", 520))
        self.geometry(f"{w}x{h}")
        self.minsize(min_w, min_h)

        # ---- MAXIMIZE FIX (more aggressive + staged) ----
        # Many Windows setups ignore the first zoomed call during early startup.
        # We apply maximize several times across the event loop and after first paint.
        self.after(0, self._start_maximized)
        self.after(80, self._start_maximized)
        self.after(220, self._start_maximized)
        self.after(600, self._start_maximized)
        self.after(0, self._rss_sched_init)

        self.center_offset_y = int(app_cfg.get("center_offset_y", -40))

        # Canvas first: allows immediate paint
        self.canvas = tk.Canvas(self, highlightthickness=0)
        self.canvas.pack(fill="both", expand=True)

        # Background placeholders
        self.bg_original = None
        self.bg_tk = None
        self.bg_id = None

        # Logos placeholders
        self.logo_left_original = None
        self.logo_right_original = None
        self.logo_left_tk = None
        self.logo_right_tk = None
        self.logo_left_label = None
        self.logo_right_label = None

        # Optional margin around centered UI so background shows on all sides
        margin_cfg = (ui_cfg.get("background_margin") or {})
        self.bg_margin_x = int(margin_cfg.get("x", 0))
        self.bg_margin_y = int(margin_cfg.get("y", 0))

        # UI container hosted on canvas (recentered on resize)
        self.ui_frame = tk.Frame(self.canvas, bd=0, highlightthickness=0)
        self.ui_window_id = self.canvas.create_window(0, 0, anchor="nw", window=self.ui_frame, tags=("ui",))

        # Loading label
        self._loading = tk.Label(
            self.ui_frame,
            text="Loading…",
            font=("Segoe UI", 14, "bold"),
            fg="#F2F5FF",
            bg="#0D0F14",
            padx=16,
            pady=10,
        )
        self._loading.pack(pady=40)

        # Debounced resize scheduling
        self._pending_resize = None
        self._resize_job = None
        self._bg_last_size = (0, 0)
        self._startup_fast_resample_until = datetime.now().timestamp() + 0.6

        # Debug
        self.print_resize = bool(ui_cfg.get("debug", {}).get("print_resize", False))

        # Badge registry (for dynamic tag updates)
        self._badge_widgets = {}  # item_signature -> tk.Label
        self._badge_items = {}    # item_signature -> item dict

        # Cached RSS list for badge calc
        self._rss_items_cache = []
        self._rss_db_path = ""
        self._rss_meta = {}

        # Bind resize
        self.bind("<Configure>", self._on_resize)

        # Command Palette toggle (Ctrl+P)
        self.bind("<Control-p>", self._toggle_palette)

        # Deferred boot
        self.after(0, self._boot)

    # ----------------------------
    # deferred boot
    # ----------------------------
    def _boot(self):
        try:
            app_cfg = self.cfg.get("app", {})

            bg_path = resolve_path(app_cfg.get("background", ""), self.tokens)
            logo_left_path = resolve_path(app_cfg.get("logo", ""), self.tokens)
            logo_right_path = resolve_path(app_cfg.get("logo_right", ""), self.tokens)

            self.bg_original = self._load_image(bg_path, required=True)
            self.logo_left_original = self._load_image(logo_left_path, required=False)
            self.logo_right_original = self._load_image(logo_right_path, required=False)

            if self.bg_original:
                self.bg_original = self.bg_original.convert("RGB")
            if self.logo_left_original:
                self.logo_left_original = self.logo_left_original.convert("RGBA")
            if self.logo_right_original:
                self.logo_right_original = self.logo_right_original.convert("RGBA")

            try:
                self._loading.destroy()
            except Exception:
                pass
            self._loading = None

            self._build_header()
            self._build_sections()
            self._apply_app_icon()

            self.after(0, self._initial_render)

            # Defer RSS badge computation (keeps boot fast)
            self.after(30, self._refresh_rss_cache_and_badges)

            # One more maximize after UI is built (common fix for stubborn environments)
            self.after(120, self._start_maximized)

        except Exception as e:
            messagebox.showerror("Startup Error", str(e))

    # ----------------------------
    # Robust maximize logic
    # ----------------------------
    def _start_maximized(self):
        try:
            self.deiconify()
            self.lift()
            self.update_idletasks()
        except Exception:
            pass

        # Tk's zoomed (Windows)
        try:
            self.wm_state("zoomed")
            return
        except Exception:
            pass

        # Some builds expose -zoomed
        try:
            self.attributes("-zoomed", True)
            return
        except Exception:
            pass

        # Windows API hard fallback
        if sys.platform.startswith("win"):
            try:
                import ctypes
                SW_MAXIMIZE = 3
                hwnd = self.winfo_id()
                ctypes.windll.user32.ShowWindow(hwnd, SW_MAXIMIZE)
                return
            except Exception:
                pass

        # Generic fallback: size to screen
        try:
            self.update_idletasks()
            sw = self.winfo_screenwidth()
            sh = self.winfo_screenheight()
            self.geometry(f"{sw}x{sh}+0+0")
        except Exception:
            pass

    # ----------------------------
    # RSS state store (read/flag/import/tags/notes)
    # ----------------------------
    def _rss_state_bucket(self) -> dict:
        self.state.setdefault("rss_article_state", {})
        if not isinstance(self.state["rss_article_state"], dict):
            self.state["rss_article_state"] = {}
        return self.state["rss_article_state"]

    def _rss_state_get(self, fp: str) -> dict:
        bucket = self._rss_state_bucket()
        if fp not in bucket or not isinstance(bucket.get(fp), dict):
            bucket[fp] = {
                "read": False,
                "flagged": False,
                "imported": False,
                "tags": [],
                "notes": "",
                "last_seen": "",
                "obsidian_note_path": "",
                "obsidian_note_created_at": "",
                "reported_on": "",
                "reported_url": "",
                "queued_for_report": False,
                "queued_on": "",
                "printed": False,
            }
        return bucket[fp]

    def _rss_state_set(self, fp: str, **kwargs) -> None:
        st = self._rss_state_get(fp)
        st.update(kwargs)
        save_state(self.state_path, self.state)

    def _rss_counts(self, items) -> dict:
        new_ct = 0
        flag_ct = 0
        imp_ct = 0
        rep_ct = 0
        for it in items or []:
            fp = _rss_fp(it)
            st = self._rss_state_get(fp)
            if st.get("flagged"):
                flag_ct += 1
            if st.get("imported"):
                imp_ct += 1
            if _safe_str(st.get("reported_on")).strip():
                rep_ct += 1
            if not st.get("read"):
                new_ct += 1
        return {"new": new_ct, "flagged": flag_ct, "imported": imp_ct, "reported": rep_ct, "total": len(items or [])}

    def _rss_compute_watchlist_hits(self, item: dict, watch_items: list[dict]) -> tuple[int, list[str], list[str]]:
        """
        Returns (score, hits, categories) for a single RSS item using watchlist items.
        """
        blob = " ".join(
            [
                _safe_str(item.get("title")),
                _safe_str(item.get("source")),
                _safe_str(item.get("summary")),
                _safe_str(item.get("content")),
                _safe_str(item.get("url")),
            ]
        ).lower()
        score = 0
        hits: list[str] = []
        categories: set[str] = set()
        if blob and watch_items:
            for w in watch_items:
                term = _safe_str(w.get("term")).strip().lower()
                if term and term in blob:
                    hits.append(term)
                    try:
                        score += int(w.get("weight", 1))
                    except Exception:
                        score += 1
                    for c in (w.get("categories") or []):
                        if c:
                            categories.add(str(c))
        return score, sorted(set(hits)), sorted(categories)

    def _rss_update_watchlist_freq(self, items: list[dict]) -> None:
        """
        Computes keyword hit frequency and updates UI if present.
        """
        counts: dict[str, int] = {}
        for it in items or []:
            for k in (it.get("_watchlist_hits") or []):
                counts[k] = counts.get(k, 0) + 1

        rows = sorted(counts.items(), key=lambda x: (-x[1], x[0]))
        self.state["rss_watchlist_freq"] = rows
        try:
            save_state(self.state_path, self.state)
        except Exception:
            pass

        top = rows[:10]
        bottom = list(reversed(rows[-10:])) if rows else []

        if getattr(self, "_rss_wl_freq_top", None):
            self._rss_wl_freq_top.delete(0, tk.END)
            for term, cnt in top:
                self._rss_wl_freq_top.insert(tk.END, f"{term} ({cnt})")
        if getattr(self, "_rss_wl_freq_bottom", None):
            self._rss_wl_freq_bottom.delete(0, tk.END)
            for term, cnt in bottom:
                self._rss_wl_freq_bottom.insert(tk.END, f"{term} ({cnt})")

    def _rss_parse_watchlist_text(self, raw: str) -> dict:
        """
        Parse watchlist frontmatter keywords with optional categories/weight.
        Supports:
          keywords:
            - microsoft
            - term: "microsoft"
              categories: ["vendor", "brand"]
              weight: 5
        """
        lines = raw.splitlines()
        if not lines or lines[0].strip() != "---":
            return {"items": [], "fm_start": -1, "fm_end": -1, "kw_start": -1, "kw_end": -1, "defaults": {}}

        fm_end = -1
        for i in range(1, len(lines)):
            if lines[i].strip() == "---":
                fm_end = i
                break
        if fm_end == -1:
            return {"items": [], "fm_start": -1, "fm_end": -1, "kw_start": -1, "kw_end": -1, "defaults": {}}

        def _strip_quotes(s: str) -> str:
            return s.strip().strip('"').strip("'")

        def _parse_inline_list(s: str) -> list[str]:
            s = s.strip()
            if not s:
                return []
            if s.startswith("[") and s.endswith("]"):
                inner = s[1:-1].strip()
                if not inner:
                    return []
                return [_strip_quotes(tok) for tok in inner.split(",") if _strip_quotes(tok)]
            return [_strip_quotes(s)]

        defaults = {"categories": [], "weight": 1}
        items = []
        kw_start = -1
        kw_end = -1

        i = 1
        while i < fm_end:
            s = lines[i].strip()
            if s.lower().startswith(("watchlist_categories:", "categories:", "category:")):
                rest = s.split(":", 1)[1]
                defaults["categories"] = _parse_inline_list(rest)
            elif s.lower().startswith(("watchlist_weight:", "weight:")):
                rest = s.split(":", 1)[1].strip()
                try:
                    defaults["weight"] = int(float(rest))
                except Exception:
                    pass
            elif s.lower().startswith("keywords:"):
                kw_start = i
                i += 1
                current = None
                item_indent = None
                in_categories = False
                while i < fm_end:
                    raw_line = lines[i]
                    if not raw_line.strip():
                        i += 1
                        continue
                    if len(raw_line) - len(raw_line.lstrip()) < 2:
                        kw_end = i
                        break
                    t = raw_line.strip()
                    if t.startswith("- "):
                        in_categories = False
                        item_indent = len(raw_line) - len(raw_line.lstrip())
                        rest = t[2:].strip()
                        current = {"term": "", "categories": [], "weight": None}
                        if rest:
                            if ":" in rest:
                                k, v = rest.split(":", 1)
                                k = k.strip().lower()
                                v = v.strip()
                                if k in ("term", "keyword", "value"):
                                    current["term"] = _strip_quotes(v)
                                elif k in ("categories", "category"):
                                    current["categories"] = _parse_inline_list(v)
                                elif k in ("weight", "score", "priority"):
                                    try:
                                        current["weight"] = int(float(v))
                                    except Exception:
                                        pass
                            else:
                                current["term"] = _strip_quotes(rest)
                        items.append(current)
                        i += 1
                        continue
                    if current and item_indent is not None:
                        ind = len(raw_line) - len(raw_line.lstrip())
                        if ind <= item_indent:
                            i += 1
                            continue
                        if in_categories and t.startswith("- "):
                            val = _strip_quotes(t[2:].strip())
                            if val:
                                current["categories"].append(val)
                            i += 1
                            continue
                        if ":" in t:
                            k, v = t.split(":", 1)
                            k = k.strip().lower()
                            v = v.strip()
                            if k in ("term", "keyword", "value"):
                                current["term"] = _strip_quotes(v)
                            elif k in ("categories", "category"):
                                if v:
                                    current["categories"] = _parse_inline_list(v)
                                    in_categories = False
                                else:
                                    in_categories = True
                            elif k in ("weight", "score", "priority"):
                                try:
                                    current["weight"] = int(float(v))
                                except Exception:
                                    pass
                        i += 1
                        continue
                    i += 1
                if kw_end == -1:
                    kw_end = fm_end
            i += 1

        return {
            "items": items,
            "fm_start": 0,
            "fm_end": fm_end,
            "kw_start": kw_start,
            "kw_end": kw_end,
            "defaults": defaults,
        }

    def _rss_format_watchlist_keywords_block(self, items: list[dict]) -> list[str]:
        lines = ["keywords:"]
        for it in items or []:
            term = _normalize_keyword(it.get("term") or "")
            if not term:
                continue
            cats = [c.strip() for c in (it.get("categories") or []) if c and str(c).strip()]
            weight = it.get("weight")
            try:
                weight = int(weight) if weight is not None else 1
            except Exception:
                weight = 1
            lines.append(f'  - term: "{term}"')
            if cats:
                cats_fmt = ", ".join([f'"{c}"' for c in cats])
                lines.append(f"    categories: [{cats_fmt}]")
            lines.append(f"    weight: {weight}")
        return lines

    def _rss_watchlist_items(self, with_stats: bool = False):
        """
        Returns structured watchlist items.
        """
        rss_cfg = self.cfg.get("rss", {}) if isinstance(self.cfg.get("rss", {}), dict) else {}
        items: list[dict] = []
        stats = {
            "watchlists_dir_selected": "",
            "watchlists_dir_resolved": "",
            "vault_root_token": str(self.tokens.get("VAULT_ROOT", "") or ""),
            "watchlist_files_scanned": 0,
            "watchlist_files_with_keywords": 0,
            "watchlist_files_sample": [],
        }

        # Watchlist files
        wdir_cfg = rss_cfg.get("watchlists_dir") or rss_cfg.get("watchlist_dir") or ""
        candidates = []
        if isinstance(wdir_cfg, str) and wdir_cfg.strip():
            candidates.append(wdir_cfg.strip())
        candidates.append("{VAULT_ROOT}/10_Operations/13_Watchlists")
        candidates.append("{VAULT_ROOT}/30_CIPHER/05_Watchlists")

        resolved_dir = ""
        selected_expr = ""

        def _resolve_candidate(expr: str) -> str:
            try:
                s = resolve_path(expr, self.tokens)
            except Exception:
                s = expr
            s = (s or "").strip()
            if not s:
                return ""
            try:
                p = Path(s)
                if not p.is_absolute() and stats["vault_root_token"]:
                    s = str((Path(stats["vault_root_token"]) / s).resolve())
            except Exception:
                pass
            try:
                s = str(Path(s).resolve())
            except Exception:
                pass
            return s

        for expr in candidates:
            cand = _resolve_candidate(expr)
            if not cand:
                continue
            p = Path(cand)
            if p.exists() and p.is_dir():
                resolved_dir = cand
                selected_expr = expr
                break

        if not resolved_dir and candidates:
            selected_expr = candidates[0]
            resolved_dir = _resolve_candidate(selected_expr)

        stats["watchlists_dir_selected"] = selected_expr
        stats["watchlists_dir_resolved"] = resolved_dir

        if resolved_dir and Path(resolved_dir).exists():
            for fp in sorted(Path(resolved_dir).rglob("*.md")):
                if not fp.is_file():
                    continue
                stats["watchlist_files_scanned"] += 1
                if len(stats["watchlist_files_sample"]) < 5:
                    stats["watchlist_files_sample"].append(fp.name)
                try:
                    raw = fp.read_text(encoding="utf-8", errors="ignore")
                except Exception:
                    continue
                parsed = self._rss_parse_watchlist_text(raw)
                kw_items = parsed.get("items") or []
                defaults = parsed.get("defaults") or {}
                if kw_items:
                    stats["watchlist_files_with_keywords"] += 1
                for it in kw_items:
                    term = _normalize_keyword(it.get("term") or "")
                    if not term:
                        continue
                    cats = it.get("categories") or defaults.get("categories") or []
                    if isinstance(cats, str):
                        cats = [c.strip() for c in cats.split(",") if c.strip()]
                    weight = it.get("weight")
                    if weight is None:
                        weight = defaults.get("weight", 1)
                    items.append(
                        {
                            "term": term,
                            "categories": cats,
                            "weight": weight,
                            "source": str(fp),
                            "watchlist": fp.stem,
                        }
                    )

        # Merge by term
        merged = {}
        for it in items:
            term = _normalize_keyword(it.get("term") or "")
            if not term:
                continue
            cats = [c.strip() for c in (it.get("categories") or []) if c and str(c).strip()]
            try:
                weight = int(it.get("weight", 1))
            except Exception:
                weight = 1
            if term not in merged:
                merged[term] = {"term": term, "categories": set(cats), "weight": weight, "sources": set()}
            else:
                merged[term]["categories"].update(cats)
                merged[term]["weight"] = max(merged[term]["weight"], weight)
            if it.get("source"):
                merged[term]["sources"].add(it.get("source"))
        out_items = []
        for v in merged.values():
            out_items.append(
                {
                    "term": v["term"],
                    "categories": sorted(v["categories"]),
                    "weight": v["weight"],
                    "sources": sorted(v["sources"]),
                }
            )

        if with_stats:
            return out_items, stats
        return out_items

    def _rss_watchlist_keywords(self) -> list[str]:
        """
        Returns a de-duplicated list of lowercase keywords used for auto-flagging RSS items.

        Sources:
          - cfg["rss"]["watchlists_dir"] or cfg["rss"]["watchlist_dir"] = "path" (optional)

        Directory parsing supports:
          - .txt: one keyword per line
          - .md: YAML frontmatter `keywords:` list or one keyword per line
          - .yaml/.yml: `keywords:` list or one keyword per line
        """
        items, stats = self._rss_watchlist_items(with_stats=True)
        out = [it.get("term") for it in items if it.get("term")]
        out = [str(k).strip().lower() for k in out if str(k).strip()]
        try:
            state_stats = self.state.get("rss_autoflag_stats", {}) if isinstance(self.state.get("rss_autoflag_stats", {}), dict) else {}
            state_stats.update(stats or {})
            self.state["rss_autoflag_stats"] = state_stats
        except Exception:
            pass
        return out
    def _rss_auto_flag_items(self, items: list[dict]) -> None:
        """Auto-flag RSS items whose title/summary contains any watchlist keyword."""
        try:
            keywords = self._rss_watchlist_keywords()
            stats = self.state.get("rss_autoflag_stats", {}) if isinstance(self.state.get("rss_autoflag_stats", {}), dict) else {}
            stats["keywords_loaded"] = int(len(keywords))
            stats["auto_flagged"] = 0

            if not keywords:
                self.state["rss_autoflag_stats"] = stats
                try:
                    save_state(self.state_path, self.state)
                except Exception:
                    pass
                return

            bucket = self._rss_state_bucket()
            changed = False

            for it in items or []:
                fp = _rss_fp(it)
                if not fp:
                    continue

                st = bucket.get(fp)
                if not isinstance(st, dict):
                    st = self._rss_state_get(fp)

                if st.get("flagged") is True:
                    continue

                title = _safe_str(it.get("title")).lower()
                summary = _safe_str(it.get("summary") or it.get("description")).lower()
                src = _safe_str(it.get("source") or it.get("feed") or it.get("feed_title")).lower()
                blob = (title + " " + summary + " " + src).strip()

                if not blob:
                    continue

                if any(k in blob for k in keywords):
                    st["flagged"] = True
                    st["last_seen"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    if not st.get("notes"):
                        st["notes"] = "Auto-flagged due to watchlist keyword match."
                    bucket[fp] = st
                    changed = True
                    stats["auto_flagged"] = int(stats.get("auto_flagged", 0)) + 1

            self.state["rss_autoflag_stats"] = stats
            if changed:
                save_state(self.state_path, self.state)
            else:
                try:
                    save_state(self.state_path, self.state)
                except Exception:
                    pass
        except Exception:
            pass


    def _rss_load_items(self):
        # Use cache if present; else attempt DB detect + fetch.
        if self._rss_items_cache:
            items = list(self._rss_items_cache)
        else:
            app_dir = get_app_dir()
            self._rss_db_path = _guess_rss_db_path(self.cfg, self.tokens, app_dir)
            if self._rss_db_path:
                items, meta = _fetch_rss_items(self._rss_db_path, limit=800)
                self._rss_meta = meta or {}
                self._rss_items_cache = items or []
                items = list(self._rss_items_cache)
            else:
                self._rss_meta = {"reason": "no db detected"}
                items = []

        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        for it in items:
            fp = _rss_fp(it)
            # NOTE: this does write state on load; keeping your current behavior intact.
            self._rss_state_set(fp, last_seen=now)

        # Auto-flag based on watchlist keywords
        try:
            self._rss_auto_flag_items(items)
        except Exception:
            pass

        return items

    # ----------------------------
    # Badge refresh for RSS review card
    # ----------------------------
    def _refresh_rss_cache_and_badges(self):
        try:
            items = self._rss_load_items()
            new_count = 0
            for it in items:
                fp = _rss_fp(it)
                st = self._rss_state_get(fp)
                if not st.get("read"):
                    new_count += 1

            for sig, wdg in list(self._badge_widgets.items()):
                it = self._badge_items.get(sig, {}) or {}
                if (it.get("action") or "").strip() == "open_rss_review":
                    try:
                        wdg.configure(text=f"NEW: {new_count}")
                    except Exception:
                        pass
        except Exception as e:
            self._rss_meta = {"reason": str(e)}

    # ----------------------------
    # RSS Scheduler (runs only while app is open)
    # ----------------------------
    def _rss_sched_init(self):
        self._rss_sched_job = None
        self._rss_sched_running = False
        self._rss_sched_next_run = None
        self._rss_schedule = _cfg_get_rss_schedule(self.cfg)

        # Optional run-on-startup
        if self._rss_schedule.get("enabled") and self._rss_schedule.get("run_at_startup"):
            self.after(1200, lambda: self._rss_sched_kick(manual=False))

        # Start loop if enabled
        self._rss_sched_apply()

    def _rss_sched_apply(self):
        # cancel any existing job
        if getattr(self, "_rss_sched_job", None) is not None:
            try:
                self.after_cancel(self._rss_sched_job)
            except Exception:
                pass
            self._rss_sched_job = None

        self._rss_schedule = _cfg_get_rss_schedule(self.cfg)
        if not self._rss_schedule.get("enabled"):
            self._rss_sched_next_run = None
            return

        interval_min = int(self._rss_schedule.get("interval_minutes", 60))
        # always compute next run from "now" when applying changes
        self._rss_sched_next_run = datetime.now().timestamp() + (interval_min * 60)

        # Tick every 5s (cheap), execute only when due
        self._rss_sched_job = self.after(5000, self._rss_sched_tick)

    def _rss_sched_tick(self):
        self._rss_sched_job = None

        sch = _cfg_get_rss_schedule(self.cfg)
        if not sch.get("enabled"):
            self._rss_sched_next_run = None
            return

        now = datetime.now()

        qh = sch.get("quiet_hours", {}) or {}
        if qh.get("enabled"):
            if _time_in_quiet_hours(now, qh.get("start", "22:00"), qh.get("end", "06:00")):
                # defer one interval forward
                self._rss_sched_next_run = now.timestamp() + (int(sch.get("interval_minutes", 60)) * 60)
                self._rss_sched_job = self.after(5000, self._rss_sched_tick)
                return

        if self._rss_sched_next_run is None:
            self._rss_sched_next_run = now.timestamp() + (int(sch.get("interval_minutes", 60)) * 60)

        if now.timestamp() >= float(self._rss_sched_next_run):
            self._rss_sched_kick(manual=False)
            self._rss_sched_next_run = now.timestamp() + (int(sch.get("interval_minutes", 60)) * 60)

        self._rss_sched_job = self.after(5000, self._rss_sched_tick)

    def _rss_sched_kick(self, manual: bool = False):
        # prevent overlapping runs
        if getattr(self, "_rss_sched_running", False):
            return

        self._rss_sched_running = True

        def work():
            try:
                # Use same collector you already wired for manual runs
                app_dir = get_app_dir()
                result = run_rss_collect_and_capture(app_dir, python_cmd="python")

                # Save last run info (mirrors your existing pattern)
                self.state["rss_last_run"] = {
                    "ok": bool(result.get("ok")),
                    "ended": result.get("ended"),
                    "feeds_processed": result.get("feeds_processed"),
                    "new_inserted": result.get("new_inserted"),
                    "exit_code": result.get("exit_code"),
                    "error": result.get("error", ""),
                    "stdout": result.get("stdout", ""),
                    "stderr": result.get("stderr", ""),
                    "scheduled": (not manual),
                }
                save_state(self.state_path, self.state)

                # Update badges after run
                self.after(30, self._refresh_rss_cache_and_badges)

            finally:
                self._rss_sched_running = False

        import threading
        threading.Thread(target=work, daemon=True).start()

    # ----------------------------
    # Scheduler UI
    # ----------------------------
    def _open_rss_schedule(self):
        self._open_rss_window(select_tab="schedule")

    def _build_rss_schedule_tab(self, parent: tk.Frame):
        win = getattr(self, "_rss_win", None)

        # Reload config from disk to ensure schedule persistence reflects the actual config.json.
        try:
            app_dir = get_app_dir()
            self.cfg = load_config(app_dir)
            self.tokens = build_tokens(app_dir, self.cfg)
        except Exception:
            pass

        sch = _cfg_get_rss_schedule(self.cfg)

        enabled_var = tk.BooleanVar(value=bool(sch.get("enabled", False)))
        interval_var = tk.StringVar(value=str(sch.get("interval_minutes", 60)))
        startup_var = tk.BooleanVar(value=bool(sch.get("run_at_startup", False)))

        qh = (sch.get("quiet_hours", {}) or {})
        qh_enabled_var = tk.BooleanVar(value=bool(qh.get("enabled", False)))
        qh_start_var = tk.StringVar(value=str(qh.get("start", "22:00")))
        qh_end_var = tk.StringVar(value=str(qh.get("end", "06:00")))

        # --- Live status variables (MUST exist before any Labels bind to them) ---
        status = tk.StringVar(value="")
        next_run_var = tk.StringVar(value="")
        last_run_var = tk.StringVar(value="")
        last_detail_var = tk.StringVar(value="")
        countdown_var = tk.StringVar(value="")

        def _fmt_ts(ts: float) -> str:
            try:
                return datetime.fromtimestamp(float(ts)).strftime("%Y-%m-%d %H:%M:%S")
            except Exception:
                return ""

        def _refresh_live_labels():
            """Updates the Run Status panel once per second while the window is open."""
            # last run info
            lr = (self.state.get("rss_last_run") or {}) if isinstance(self.state.get("rss_last_run"), dict) else {}
            lr_ok = lr.get("ok")
            lr_time = _safe_str(lr.get("ended")).strip()
            feeds = lr.get("feeds_processed")
            newi = lr.get("new_inserted")
            exit_code = lr.get("exit_code")
            err = _safe_str(lr.get("error")).strip()

            if lr_ok is True:
                last_run_var.set(f"Last run: OK  •  {lr_time or 'unknown'}")
            elif lr_ok is False:
                last_run_var.set(f"Last run: FAILED  •  {lr_time or 'unknown'}")
            else:
                last_run_var.set("Last run: (none yet)")

            detail_bits = []
            if feeds is not None:
                detail_bits.append(f"Feeds: {feeds}")
            if newi is not None:
                detail_bits.append(f"New inserted: {newi}")
            if exit_code is not None:
                detail_bits.append(f"Exit: {exit_code}")
            if lr_ok is False and err:
                detail_bits.append(f"Error: {err[:140]}")
            last_detail_var.set(" • ".join(detail_bits) if detail_bits else "")

            # next run info
            sch_now = _cfg_get_rss_schedule(self.cfg)
            if not sch_now.get("enabled"):
                next_run_var.set("Next run: (scheduling disabled)")
                countdown_var.set("")
            else:
                nxt = getattr(self, "_rss_sched_next_run", None)
                if not nxt:
                    next_run_var.set("Next run: (pending)")
                    countdown_var.set("")
                else:
                    next_run_var.set(f"Next run: {_fmt_ts(nxt)}")
                    try:
                        secs = int(max(0, float(nxt) - datetime.now().timestamp()))
                        mm = secs // 60
                        ss = secs % 60
                        countdown_var.set(f"Time remaining: {mm:02d}:{ss:02d}")
                    except Exception:
                        countdown_var.set("")

            # reschedule if window still exists
            try:
                if win and win.winfo_exists():
                    win.after(1000, _refresh_live_labels)
            except Exception:
                pass

        wrap = tk.Frame(parent, bg="#111111", padx=16, pady=14)
        wrap.pack(fill="both", expand=True)

        tk.Label(wrap, text="RSS Collector Schedule", bg="#111111", fg="#ffffff",
                 font=("Segoe UI", 13, "bold")).pack(anchor="w")

        tk.Label(
            wrap,
            text="Runs RSS collection on an interval while the launcher is open.",
            bg="#111111", fg="#9aa7c0", font=("Segoe UI", 9, "normal")
        ).pack(anchor="w", pady=(2, 10))

        tk.Checkbutton(
            wrap, text="Enable scheduling", variable=enabled_var,
            bg="#111111", fg="#ffffff", selectcolor="#2a2a2a",
            activebackground="#111111", activeforeground="#ffffff"
        ).pack(anchor="w", pady=(0, 10))

        row = tk.Frame(wrap, bg="#111111")
        row.pack(fill="x", pady=(0, 8))
        tk.Label(row, text="Interval (minutes):", bg="#111111", fg="#cfcfcf",
                 font=("Segoe UI", 10, "normal")).pack(side="left")
        tk.Entry(row, textvariable=interval_var, bg="#1e1e1e", fg="#ffffff",
                 insertbackground="#ffffff", relief="flat", width=10,
                 font=("Segoe UI", 11, "normal")).pack(side="left", padx=(10, 0))

        tk.Checkbutton(
            wrap, text="Run once on startup (when enabled)", variable=startup_var,
            bg="#111111", fg="#ffffff", selectcolor="#2a2a2a",
            activebackground="#111111", activeforeground="#ffffff"
        ).pack(anchor="w", pady=(6, 12))

        tk.Label(wrap, text="Quiet hours (optional)", bg="#111111", fg="#ffffff",
                 font=("Segoe UI", 11, "bold")).pack(anchor="w", pady=(0, 6))

        tk.Checkbutton(
            wrap, text="Enable quiet hours", variable=qh_enabled_var,
            bg="#111111", fg="#ffffff", selectcolor="#2a2a2a",
            activebackground="#111111", activeforeground="#ffffff"
        ).pack(anchor="w")

        qhr = tk.Frame(wrap, bg="#111111")
        qhr.pack(fill="x", pady=(6, 10))

        tk.Label(qhr, text="Start (HH:MM):", bg="#111111", fg="#cfcfcf",
                 font=("Segoe UI", 10, "normal")).pack(side="left")
        tk.Entry(qhr, textvariable=qh_start_var, bg="#1e1e1e", fg="#ffffff",
                 insertbackground="#ffffff", relief="flat", width=8,
                 font=("Segoe UI", 11, "normal")).pack(side="left", padx=(8, 16))

        tk.Label(qhr, text="End (HH:MM):", bg="#111111", fg="#cfcfcf",
                 font=("Segoe UI", 10, "normal")).pack(side="left")
        tk.Entry(qhr, textvariable=qh_end_var, bg="#1e1e1e", fg="#ffffff",
                 insertbackground="#ffffff", relief="flat", width=8,
                 font=("Segoe UI", 11, "normal")).pack(side="left", padx=(8, 0))

        # --- Live status panel ---
        tk.Label(wrap, text="Run Status", bg="#111111", fg="#ffffff",
                 font=("Segoe UI", 11, "bold")).pack(anchor="w", pady=(10, 6))

        stbox = tk.Frame(wrap, bg="#111111")
        stbox.pack(fill="x")

        tk.Label(stbox, textvariable=last_run_var, bg="#111111", fg="#cfcfcf",
                 font=("Segoe UI", 10, "normal")).pack(anchor="w")
        tk.Label(stbox, textvariable=last_detail_var, bg="#111111", fg="#9aa7c0",
                 font=("Segoe UI", 9, "normal")).pack(anchor="w", pady=(2, 0))

        tk.Label(stbox, textvariable=next_run_var, bg="#111111", fg="#cfcfcf",
                 font=("Segoe UI", 10, "normal")).pack(anchor="w", pady=(10, 0))
        tk.Label(stbox, textvariable=countdown_var, bg="#111111", fg="#9aa7c0",
                 font=("Segoe UI", 9, "normal")).pack(anchor="w", pady=(2, 0))

        tk.Label(wrap, textvariable=status, bg="#111111", fg="#9aa7c0",
                 font=("Segoe UI", 9, "normal")).pack(anchor="w", pady=(8, 0))

        btns = tk.Frame(wrap, bg="#111111")
        btns.pack(side="bottom", fill="x", pady=(14, 0))

        def validate_hhmm(s: str) -> bool:
            try:
                s = (s or "").strip()
                hh, mm = s.split(":")
                hh = int(hh); mm = int(mm)
                return (0 <= hh <= 23) and (0 <= mm <= 59)
            except Exception:
                return False

        def on_save():
            try:
                mins = int((interval_var.get() or "0").strip())
                if mins < 5 or mins > 1440:
                    raise ValueError("Interval must be between 5 and 1440 minutes.")

                if qh_enabled_var.get():
                    if not validate_hhmm(qh_start_var.get()) or not validate_hhmm(qh_end_var.get()):
                        raise ValueError("Quiet hours start/end must be HH:MM in 24-hour format.")

                self.cfg.setdefault("rss", {})
                self.cfg["rss"].setdefault("schedule", {})
                self.cfg["rss"]["schedule"].update({
                    "enabled": bool(enabled_var.get()),
                    "mode": "interval",
                    "interval_minutes": mins,
                    "run_at_startup": bool(startup_var.get()),
                    "quiet_hours": {
                        "enabled": bool(qh_enabled_var.get()),
                        "start": (qh_start_var.get() or "").strip(),
                        "end": (qh_end_var.get() or "").strip(),
                    }
                })

                # Persist to config.json next to app
                cfg_path = get_app_dir() / "config.json"
                cfg_path.write_text(json.dumps(self.cfg, indent=2, ensure_ascii=False), encoding="utf-8")

                # Reload from disk to confirm persistence and to rehydrate UI defaults.
                try:
                    app_dir = get_app_dir()
                    self.cfg = load_config(app_dir)
                    self.tokens = build_tokens(app_dir, self.cfg)
                    sch2 = _cfg_get_rss_schedule(self.cfg)
                    enabled_var.set(bool(sch2.get("enabled", False)))
                    interval_var.set(str(sch2.get("interval_minutes", 60)))
                    startup_var.set(bool(sch2.get("run_at_startup", False)))

                    qh2 = (sch2.get("quiet_hours", {}) or {})
                    qh_enabled_var.set(bool(qh2.get("enabled", False)))
                    qh_start_var.set(str(qh2.get("start", "22:00")))
                    qh_end_var.set(str(qh2.get("end", "06:00")))
                except Exception:
                    pass

                # Apply scheduler changes immediately
                self._rss_sched_apply()
                status.set("Saved and applied.")
                _refresh_live_labels()
            except Exception as e:
                messagebox.showerror("Save Schedule", str(e))

        def on_run_now():
            self._rss_sched_kick(manual=True)
            status.set("Triggered RSS collect.")
            _refresh_live_labels()

        tk.Button(btns, text="Run Now", command=on_run_now,
                  bg="#1a1a1a", fg="#ffffff", relief="flat",
                  padx=10, pady=6, font=("Segoe UI", 10, "bold")).pack(side="left")

        tk.Button(btns, text="Save", command=on_save,
                  bg="#24304A", fg="#ffffff", relief="flat",
                  padx=10, pady=6, font=("Segoe UI", 10, "bold")).pack(side="right", padx=(8, 0))

        tk.Button(btns, text="Close", command=self._rss_review_close,
                  bg="#1a1a1a", fg="#ffffff", relief="flat",
                  padx=10, pady=6, font=("Segoe UI", 10, "bold")).pack(side="right")

        # prime the panel and start refresh loop
        _refresh_live_labels()
    def _record_recent(self, item: dict) -> None:
        if not item:
            return
        if (item.get("action") or "").lower() == "quit":
            return

        max_recent = 12

        rec = {
            "section": item.get("section", ""),
            "label": item.get("label", ""),
            "action": item.get("action", ""),
            "target": item.get("target", ""),
            "ts": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        }
        sig = item_signature(rec)

        new_list = []
        for r in self.state.get("recent", []):
            try:
                if item_signature(r) == sig:
                    continue
            except Exception:
                continue
            new_list.append(r)

        new_list.insert(0, rec)
        self.state["recent"] = new_list[:max_recent]
        save_state(self.state_path, self.state)

    def _get_recent_items_filtered_to_current_menu(self):
        current = self._flatten_menu_items()
        current_sigs = set(item_signature(it) for it in current)

        out = []
        for r in (self.state.get("recent", []) or []):
            try:
                if item_signature(r) in current_sigs:
                    out.append(r)
            except Exception:
                continue
        return out

    # ----------------------------
    # Asset + icon loading
    # ----------------------------
    def _load_image(self, path: str, required: bool = True):
        if not path:
            return None
        if not os.path.exists(path):
            if required:
                messagebox.showerror("Missing asset", f"Image not found:\n{path}")
            return None
        try:
            return Image.open(path)
        except Exception as e:
            if required:
                messagebox.showerror("Image error", f"Failed to load image:\n{path}\n\n{e}")
            return None

    def _apply_app_icon(self):
        app_cfg = self.cfg.get("app", {})
        ico = resolve_path(app_cfg.get("favicon_ico", ""), self.tokens)
        png = resolve_path(app_cfg.get("favicon_png", ""), self.tokens)

        if sys.platform.startswith("win") and ico and os.path.exists(ico):
            try:
                self.iconbitmap(ico)
                return
            except Exception:
                pass

        if png and os.path.exists(png):
            try:
                img = Image.open(png).convert("RGBA")
                self._icon_ref = ImageTk.PhotoImage(img)
                self.iconphoto(True, self._icon_ref)
            except Exception:
                pass

    # -------- icon cache + tinted icon loader --------
    def _init_icon_cache(self):
        if not hasattr(self, "_icon_cache"):
            self._icon_cache = {}  # (abs_path, size, tint_hex) -> PhotoImage

    def _hex_to_rgb(self, hx: str):
        hx = (hx or "").strip()
        if not hx.startswith("#") or len(hx) != 7:
            return None
        try:
            return (int(hx[1:3], 16), int(hx[3:5], 16), int(hx[5:7], 16))
        except Exception:
            return None

    def _load_icon_tinted(self, icon_path: str, size: int, tint_hex: str = None):
        self._init_icon_cache()

        if not icon_path:
            return None

        abs_path = resolve_path(icon_path, self.tokens)
        if not os.path.exists(abs_path):
            return None

        key = (abs_path, int(size), (tint_hex or "").lower())
        if key in self._icon_cache:
            return self._icon_cache[key]

        try:
            img = Image.open(abs_path).convert("RGBA")
            img = img.resize((int(size), int(size)), Image.Resampling.LANCZOS)

            if tint_hex:
                rgb = self._hex_to_rgb(tint_hex)
                if rgb:
                    r_t, g_t, b_t = rgb
                    px = img.load()
                    w, h = img.size
                    for y in range(h):
                        for x in range(w):
                            r, g, b, a = px[x, y]
                            if a == 0:
                                continue
                            lum = int(0.2126 * r + 0.7152 * g + 0.0722 * b)
                            nr = int(r_t * lum / 255)
                            ng = int(g_t * lum / 255)
                            nb = int(b_t * lum / 255)
                            px[x, y] = (nr, ng, nb, a)

            tk_img = ImageTk.PhotoImage(img)
            self._icon_cache[key] = tk_img
            return tk_img
        except Exception:
            return None

    # ----------------------------
    # Header building (TRUE CENTER FIX)
    # ----------------------------
    def _build_header(self):
        ui = self.cfg.get("ui", {})
        header_cfg = ui.get("header", {})
        title_cfg = ui.get("title", {})
        subtitle_cfg = ui.get("subtitle", {})

        header = tk.Frame(
            self.ui_frame,
            bg=header_cfg.get("bg", "#111111"),
            padx=int(header_cfg.get("padx", 16)),
            pady=int(header_cfg.get("pady", 12)),
        )
        header.grid(row=0, column=0, sticky="ew", pady=(18, 8), padx=18)

        # We keep 3 columns, but now we enforce symmetry:
        # minsize(left) == minsize(right) == max(logo sizes)
        left_min = int(header_cfg.get("logo_column_minsize", 180))
        right_min = int(header_cfg.get("right_logo_column_minsize", 180))
        self._header_side_minsize = max(left_min, right_min)

        header.columnconfigure(0, minsize=self._header_side_minsize, weight=0)
        header.columnconfigure(1, weight=1)
        header.columnconfigure(2, minsize=self._header_side_minsize, weight=0)

        self.logo_left_label = tk.Label(header, bg=header_cfg.get("bg", "#111111"))
        self.logo_left_label.grid(row=0, column=0, rowspan=2, sticky="w", padx=(0, 14))

        # Center text container so title/subtitle are centered as a block
        text_mid = tk.Frame(header, bg=header_cfg.get("bg", "#111111"))
        text_mid.grid(row=0, column=1, rowspan=2, sticky="nsew")
        text_mid.columnconfigure(0, weight=1)

        title = tk.Label(
            text_mid,
            text=title_cfg.get("text", "Launcher"),
            fg=title_cfg.get("fg", "#ffffff"),
            bg=header_cfg.get("bg", "#111111"),
            font=(
                title_cfg.get("font_family", "Segoe UI"),
                int(title_cfg.get("font_size", 20)),
                title_cfg.get("font_weight", "bold"),
            ),
            anchor="center",
            justify="center",
        )
        title.grid(row=0, column=0, sticky="ew")

        subtitle = tk.Label(
            text_mid,
            text=subtitle_cfg.get("text", ""),
            fg=subtitle_cfg.get("fg", "#d0d0d0"),
            bg=header_cfg.get("bg", "#111111"),
            font=(
                subtitle_cfg.get("font_family", "Segoe UI"),
                int(subtitle_cfg.get("font_size", 11)),
                subtitle_cfg.get("font_weight", "normal"),
            ),
            anchor="center",
            justify="center",
        )
        subtitle.grid(row=1, column=0, sticky="ew", pady=(4, 0))

        self.logo_right_label = tk.Label(header, bg=header_cfg.get("bg", "#111111"))
        self.logo_right_label.grid(row=0, column=2, rowspan=2, sticky="e", padx=(14, 0))

        self._set_header_logos(default_left_w=220, default_right_w=220)

    def _set_header_logos(self, default_left_w: int, default_right_w: int):
        # Keep existing sizing behavior but also keep header columns symmetric
        # so center column is truly centered.
        side_w = max(
            self._header_side_minsize if hasattr(self, "_header_side_minsize") else 180,
            int(default_left_w),
            int(default_right_w),
        )
        try:
            # make both side columns equal width
            header = self.logo_left_label.master if self.logo_left_label else None
            if header:
                header.columnconfigure(0, minsize=side_w, weight=0)
                header.columnconfigure(2, minsize=side_w, weight=0)
        except Exception:
            pass

        if self.logo_left_original and self.logo_left_label:
            try:
                w = max(120, min(320, int(default_left_w)))
                ratio = self.logo_left_original.height / self.logo_left_original.width
                img = self.logo_left_original.resize((w, int(w * ratio)), Image.Resampling.LANCZOS)
                self.logo_left_tk = ImageTk.PhotoImage(img)
                self.logo_left_label.configure(image=self.logo_left_tk)
            except Exception:
                pass

        if self.logo_right_original and self.logo_right_label:
            try:
                w = max(120, min(320, int(default_right_w)))
                ratio = self.logo_right_original.height / self.logo_right_original.width
                img = self.logo_right_original.resize((w, int(w * ratio)), Image.Resampling.LANCZOS)
                self.logo_right_tk = ImageTk.PhotoImage(img)
                self.logo_right_label.configure(image=self.logo_right_tk)
            except Exception:
                pass

    # ----------------------------
    # Cards (badge support)
    # ----------------------------
    def _card_hint(self, item: dict) -> str:
        if (item.get("action") or "").strip() == "open_rss_review":
            return "NEW: 0"
        return (item.get("tag") or item.get("section") or "").strip()[:18]

    def _card_variant(self, item: dict) -> str:
        v = (item.get("variant") or "").strip().lower()
        if v in ("primary", "danger"):
            return v
        if (item.get("action") or "").strip().lower() == "quit":
            return "danger"
        return "default"

    def _register_badge(self, item: dict, tag_label: tk.Label):
        try:
            sig = item_signature(item)
            self._badge_widgets[sig] = tag_label
            self._badge_items[sig] = dict(item)
        except Exception:
            pass

    def _make_action_card(self, parent, item: dict, ui: dict):
        card_cfg = ui.get("card", {})
        section_cfg = ui.get("section", {})

        bg = card_cfg.get("bg", section_cfg.get("panel_bg", "#111111"))
        hover_bg = card_cfg.get("hover_bg", "#1B2233")
        pressed_bg = card_cfg.get("pressed_bg", "#24304A")

        border = card_cfg.get("border", "#222A3A")
        border_hover = card_cfg.get("border_hover", "#2F3A52")

        title_fg = card_cfg.get("title_fg", "#FFFFFF")
        desc_fg = card_cfg.get("desc_fg", "#C0C0C0")

        tag_bg = card_cfg.get("tag_bg", "#0F0F0F")
        tag_fg = card_cfg.get("tag_fg", "#C0C0C0")
        tag_border = card_cfg.get("tag_border", border)

        padx = int(card_cfg.get("padx", 12))
        pady = int(card_cfg.get("pady", 10))

        title_sz = int(card_cfg.get("title_font_size", 11))
        title_wt = card_cfg.get("title_font_weight", "bold")
        desc_sz = int(card_cfg.get("desc_font_size", 10))

        accent_w = int(card_cfg.get("accent_width", 4))
        accent_default = card_cfg.get("accent_default", border)
        accent_primary = card_cfg.get("accent_primary", "#3AA6FF")
        accent_danger = card_cfg.get("accent_danger", "#E05A5A")

        shadow_enabled = bool(card_cfg.get("shadow", True))
        shadow_bg = card_cfg.get("shadow_bg", "#07080B")
        shadow_dx = int(card_cfg.get("shadow_dx", 1))
        shadow_dy = int(card_cfg.get("shadow_dy", 2))

        press_shift = int(card_cfg.get("press_shift", 1))

        chevron_enabled = bool(card_cfg.get("chevron", True))
        chevron_text = card_cfg.get("chevron_text", "›")
        chevron_fg = card_cfg.get("chevron_fg", "#7E8AA6")
        chevron_hover_fg = card_cfg.get("chevron_hover_fg", "#B7C0D6")

        icon_size = int(card_cfg.get("icon_size", 20))
        icon_padx = int(card_cfg.get("icon_padx", 10))
        icon_tint = card_cfg.get("icon_tint", "#B7C0D6") or "#B7C0D6"
        icon_hover_tint = card_cfg.get("icon_hover_tint", "#F2F5FF") or "#F2F5FF"
        icon_bg = card_cfg.get("icon_bg", None)

        variant = self._card_variant(item)
        if variant == "primary":
            accent_color = accent_primary
        elif variant == "danger":
            accent_color = accent_danger
        else:
            accent_color = accent_default

        root = tk.Frame(parent, bg=parent.cget("bg"), bd=0, highlightthickness=0)
        root.grid_columnconfigure(0, weight=1)

        shadow = None
        if shadow_enabled:
            shadow = tk.Frame(root, bg=shadow_bg, bd=0, highlightthickness=0)
            shadow.grid(row=0, column=0, sticky="ew", padx=(shadow_dx, 0), pady=(shadow_dy, 0))

        outer = tk.Frame(root, bg=border, bd=0, highlightthickness=0)
        outer.grid(row=0, column=0, sticky="ew")

        inner = tk.Frame(outer, bg=bg, bd=0, highlightthickness=0)
        inner.pack(fill="both", expand=True, padx=1, pady=1)

        inner.grid_columnconfigure(2, weight=1)

        accent = tk.Frame(inner, bg=accent_color, width=accent_w)
        accent.grid(row=0, column=0, sticky="nsw")

        icon_col = tk.Frame(inner, bg=bg)
        icon_col.grid(row=0, column=1, sticky="ns", padx=(icon_padx, 0), pady=pady)

        icon_label = None
        icon_normal = None
        icon_hover = None

        icon_path = (item.get("icon") or "").strip()
        if icon_path:
            icon_normal = self._load_icon_tinted(icon_path, icon_size, icon_tint)
            icon_hover = self._load_icon_tinted(icon_path, icon_size, icon_hover_tint) or icon_normal

            if icon_normal:
                icon_label = tk.Label(
                    icon_col,
                    image=icon_normal,
                    bg=(icon_bg if icon_bg else bg),
                    bd=0,
                    highlightthickness=0,
                )
                icon_label.pack(anchor="n", pady=(2, 0))

        content = tk.Frame(inner, bg=bg)
        content.grid(row=0, column=2, sticky="nsew")
        content.grid_columnconfigure(0, weight=1)

        text_col = tk.Frame(content, bg=bg)
        text_col.grid(row=0, column=0, sticky="nsew", padx=(padx, 8), pady=pady)

        title = tk.Label(
            text_col,
            text=item.get("label", "Action"),
            bg=bg,
            fg=title_fg,
            anchor="w",
            justify="left",
            font=("Segoe UI", title_sz, title_wt),
        )
        title.pack(fill="x")

        desc_text = (item.get("description") or "").strip()
        desc = None
        if desc_text:
            desc = tk.Label(
                text_col,
                text=desc_text,
                bg=bg,
                fg=desc_fg,
                anchor="w",
                justify="left",
                wraplength=520,
                font=("Segoe UI", desc_sz, "normal"),
            )
            desc.pack(fill="x", pady=(4, 0))

        right_col = tk.Frame(content, bg=bg)
        right_col.grid(row=0, column=1, sticky="ne", padx=(0, padx), pady=pady)

        tag_outer = tk.Frame(right_col, bg=tag_border, bd=0, highlightthickness=0)
        tag_inner = tk.Label(
            tag_outer,
            text=self._card_hint(item),
            bg=tag_bg,
            fg=tag_fg,
            padx=10,
            pady=4,
            font=("Segoe UI", 9, "bold"),
        )
        tag_inner.pack(padx=1, pady=1)
        tag_outer.pack(anchor="e")

        self._register_badge(item, tag_inner)

        chev = None
        if chevron_enabled:
            chev = tk.Label(
                right_col,
                text=chevron_text,
                bg=bg,
                fg=chevron_fg,
                font=("Segoe UI", 14, "bold"),
                padx=2,
                pady=0,
            )
            chev.pack(anchor="e", pady=(6, 0))

        state = {"pressed": False}

        def set_surface_colors(surface_bg: str, border_bg: str, tag_border_bg: str, chevron_color: str):
            outer.configure(bg=border_bg)
            inner.configure(bg=surface_bg)

            icon_col.configure(bg=surface_bg)
            if icon_label is not None:
                icon_label.configure(bg=(icon_bg if icon_bg else surface_bg))

            content.configure(bg=surface_bg)
            text_col.configure(bg=surface_bg)
            right_col.configure(bg=surface_bg)

            title.configure(bg=surface_bg)
            if desc is not None:
                desc.configure(bg=surface_bg)

            tag_outer.configure(bg=tag_border_bg)
            if chev is not None:
                chev.configure(bg=surface_bg, fg=chevron_color)

        def apply_mode(mode: str):
            if mode == "normal":
                set_surface_colors(bg, border, tag_border, chevron_fg)
                if icon_label is not None and icon_normal is not None:
                    icon_label.configure(image=icon_normal)
            elif mode == "hover":
                set_surface_colors(hover_bg, border_hover, border_hover, chevron_hover_fg)
                if icon_label is not None and icon_hover is not None:
                    icon_label.configure(image=icon_hover)
            elif mode == "pressed":
                set_surface_colors(pressed_bg, border_hover, border_hover, chevron_hover_fg)
                if icon_label is not None and icon_hover is not None:
                    icon_label.configure(image=icon_hover)

        def set_pressed_offset(is_pressed: bool):
            if press_shift <= 0:
                return
            dx = press_shift if is_pressed else 0
            dy = press_shift if is_pressed else 0
            inner.pack_configure(padx=1 + dx, pady=1 + dy)

        def on_enter(_e=None):
            if not state["pressed"]:
                apply_mode("hover")

        def on_leave(_e=None):
            if not state["pressed"]:
                apply_mode("normal")

        def on_press(_e=None):
            state["pressed"] = True
            apply_mode("pressed")
            set_pressed_offset(True)

        def on_release(_e=None):
            state["pressed"] = False
            set_pressed_offset(False)
            apply_mode("hover")
            self._dispatch(item)

        def on_key_activate(_e=None):
            apply_mode("pressed")
            set_pressed_offset(True)
            self.after(90, lambda: (set_pressed_offset(False), apply_mode("normal"), self._dispatch(item)))

        widgets = [
            root, outer, inner, accent, icon_col, content, text_col, title, right_col, tag_outer, tag_inner
        ]
        if icon_label is not None:
            widgets.append(icon_label)
        if desc is not None:
            widgets.append(desc)
        if chev is not None:
            widgets.append(chev)
        if shadow is not None:
            widgets.append(shadow)

        for wdg in widgets:
            try:
                wdg.bind("<Enter>", on_enter)
                wdg.bind("<Leave>", on_leave)
                wdg.bind("<Button-1>", on_press)
                wdg.bind("<ButtonRelease-1>", on_release)
                wdg.configure(cursor="hand2")
            except Exception:
                pass

        root.configure(takefocus=True)
        root.bind("<Return>", on_key_activate)
        root.bind("<space>", on_key_activate)

        def on_focus_in(_e=None):
            outer.configure(bg=border_hover)

        def on_focus_out(_e=None):
            outer.configure(bg=border)

        root.bind("<FocusIn>", on_focus_in)
        root.bind("<FocusOut>", on_focus_out)

        apply_mode("normal")
        return root

    # ----------------------------
    # Sections (config-driven menu)
    # ----------------------------
    def _build_sections(self):
        ui = self.cfg.get("ui", {})
        container_cfg = ui.get("container", {})
        section_cfg = ui.get("section", {})

        container = tk.Frame(
            self.ui_frame,
            padx=int(container_cfg.get("padx", 18)),
            pady=int(container_cfg.get("pady", 10)),
        )
        container.grid(row=1, column=0, sticky="nsew")

        menu = self.cfg.get("menu", [])
        gap_x = int(section_cfg.get("gap_x", 12))

        card_cfg = ui.get("card", {})
        gap_y = int(card_cfg.get("gap_y", 10))

        for i, section in enumerate(menu):
            col = tk.Frame(
                container,
                bg=section_cfg.get("panel_bg", "#111111"),
                padx=int(section_cfg.get("panel_padx", 14)),
                pady=int(section_cfg.get("panel_pady", 12)),
            )
            col.grid(row=0, column=i, sticky="nsew", padx=(0 if i == 0 else gap_x, 0))
            container.columnconfigure(i, weight=1)

            label = tk.Label(
                col,
                text=section.get("section", ""),
                fg=section_cfg.get("title_fg", "#ffffff"),
                bg=section_cfg.get("panel_bg", "#111111"),
                font=(
                    section_cfg.get("title_font_family", "Segoe UI"),
                    int(section_cfg.get("title_font_size", 13)),
                    section_cfg.get("title_font_weight", "bold"),
                ),
            )
            label.pack(anchor="w", pady=(0, 10))

            for item in section.get("buttons", []):
                # Menu: New Incident / Investigation + (Resume)
                action = (item or {}).get("action", "").strip()
                if action in ("manage_incident", "open_manage_incident", "open_incident_intake", "manage_incidents"):
                    for lab, act in (
                        ("New Incident / Investigation", "incident_intake_new"),
                    ):
                        card_item = dict(item)
                        card_item["label"] = lab
                        card_item["action"] = act
                        card_item["section"] = section.get("section", "")
                        card = self._make_action_card(col, card_item, ui)
                        card.pack(fill="x", pady=(0, gap_y))
                    continue

                card_item = dict(item)
                card_item["section"] = section.get("section", "")
                card = self._make_action_card(col, card_item, ui)
                card.pack(fill="x", pady=(0, gap_y))

    # ----------------------------
    # Command Palette (Ctrl+P)
    # ----------------------------
    def _toggle_palette(self, event=None):
        if hasattr(self, "_palette") and self._palette and self._palette.winfo_exists():
            self._close_palette()
            return
        self._open_palette()

    def _flatten_menu_items(self):
        items = []
        for section in self.cfg.get("menu", []):
            sec_name = section.get("section", "")
            for btn in section.get("buttons", []):
                label = btn.get("label", "").strip()
                action = btn.get("action", "").strip()
                target = btn.get("target", "")
                if action in ("manage_incident", "open_manage_incident", "open_incident_intake", "manage_incidents"):
                    items.append({"section": sec_name, "label": "New Incident / Investigation", "action": "incident_intake_new", "target": target})
                    items.append({"section": sec_name, "label": "Incident Intake (Resume)", "action": "incident_intake_resume", "target": target})
                    continue

                if label and action:
                    items.append({"section": sec_name, "label": label, "action": action, "target": target})
        return items

    def _open_palette(self, event=None):
        if hasattr(self, "_palette") and self._palette and self._palette.winfo_exists():
            try:
                self._palette.lift()
                self._palette.focus_force()
            except Exception:
                pass
            return

        self._palette_items = self._flatten_menu_items()
        self._palette_recent = self._get_recent_items_filtered_to_current_menu()
        self._palette_filtered = list(self._palette_items)

        pal = tk.Toplevel(self)
        self._palette = pal
        pal.title("Command Palette")
        pal.transient(self)
        pal.resizable(False, False)
        pal.configure(bg="#111111")

        self.update_idletasks()
        w, h = 800, 500
        x = self.winfo_rootx() + max((self.winfo_width() - w) // 2, 0)
        y = self.winfo_rooty() + max((self.winfo_height() - h) // 2, 0)
        pal.geometry(f"{w}x{h}+{x}+{y}")

        pal.bind("<Escape>", lambda e: self._close_palette())
        pal.bind("<Control-p>", lambda e: self._toggle_palette())
        pal.protocol("WM_DELETE_WINDOW", self._close_palette)

        header = tk.Frame(pal, bg="#111111", padx=12, pady=10)
        header.pack(fill="x")

        tk.Label(
            header,
            text="Type to search • Enter to run • ↑/↓ to navigate • Tab to autocomplete • Ctrl+P toggles",
            bg="#111111",
            fg="#cfcfcf",
            font=("Segoe UI", 10, "normal"),
        ).pack(anchor="w")

        self._palette_var = tk.StringVar()
        entry = tk.Entry(
            header,
            textvariable=self._palette_var,
            font=("Segoe UI", 12, "normal"),
            relief="flat",
            bg="#1e1e1e",
            fg="#ffffff",
            insertbackground="#ffffff",
        )
        entry.pack(fill="x", pady=(8, 0))
        entry.focus_set()

        body = tk.Frame(pal, bg="#111111", padx=12, pady=10)
        body.pack(fill="both", expand=True)

        self._palette_list = tk.Listbox(
            body,
            height=20,
            activestyle="none",
            font=("Segoe UI", 11, "normal"),
            bg="#1a1a1a",
            fg="#ffffff",
            selectbackground="#2a2a2a",
            selectforeground="#ffffff",
            relief="flat",
            highlightthickness=0,
        )
        self._palette_list.pack(fill="both", expand=True)

        self._palette_list.bind("<Double-Button-1>", lambda e: self._run_palette_selection())
        pal.bind("<Return>", lambda e: self._run_palette_selection())

        entry.bind("<KeyRelease>", lambda e: self._filter_palette())
        entry.bind("<Down>", lambda e: self._palette_move_selection(1))
        entry.bind("<Up>", lambda e: self._palette_move_selection(-1))
        entry.bind("<Tab>", self._palette_autocomplete)

        self._palette_list.bind("<Down>", lambda e: self._palette_move_selection(1))
        self._palette_list.bind("<Up>", lambda e: self._palette_move_selection(-1))

        self._render_palette_results()

    def _close_palette(self):
        try:
            if hasattr(self, "_palette") and self._palette and self._palette.winfo_exists():
                self._palette.destroy()
        except Exception:
            pass
        self._palette = None

    def _filter_palette(self):
        q = (self._palette_var.get() or "").strip().lower()

        if not q:
            self._palette_filtered = list(self._palette_items)
            self._render_palette_results()
            return

        def rank(item):
            label = (item.get("label") or "").lower()
            section = (item.get("section") or "").lower()

            if label == q:
                return (0, 0)
            if label.startswith(q):
                return (1, 0)
            if q in label:
                return (2, label.find(q))
            if section.startswith(q):
                return (3, 0)
            if q in section:
                return (4, section.find(q))
            return (9, 9999)

        filtered = []
        for it in self._palette_items:
            hay = f"{it.get('label','')} {it.get('section','')}".lower()
            if q in hay:
                filtered.append(it)

        filtered.sort(key=rank)
        self._palette_filtered = filtered
        self._render_palette_results()

    def _render_palette_results(self):
        lb = self._palette_list
        lb.delete(0, tk.END)

        def hint(action: str) -> str:
            a = (action or "").lower()
            if a == "obsidian_uri":
                return "URI"
            if a == "open_path":
                return "PATH"
            if a == "run":
                return "RUN"
            if a == "url":
                return "URL"
            if a == "open_daily_note_live":
                return "DAILY"
            if a == "run_rss_collect":
                return "RSS"
            if a == "open_rss_review":
                return "REVIEW"
            if a == "quit":
                return "QUIT"
            return a.upper() or "ACTION"

        q = ""
        try:
            q = (self._palette_var.get() or "").strip().lower()
        except Exception:
            q = ""

        display_items = []
        if not q and getattr(self, "_palette_recent", None):
            lb.insert(tk.END, "— Recent —")
            display_items.append(None)

            for it in self._palette_recent[:10]:
                h = hint(it.get("action", ""))
                lb.insert(tk.END, f"[{h}] {it.get('label','')}   —   {it.get('section','')}")
                display_items.append(it)

            lb.insert(tk.END, "— All Commands —")
            display_items.append(None)

        for it in self._palette_filtered[:250]:
            h = hint(it.get("action", ""))
            lb.insert(tk.END, f"[{h}] {it['label']}   —   {it['section']}")
            display_items.append(it)

        self._palette_display_items = display_items

        if lb.size() > 0:
            idx = 0
            while idx < lb.size():
                try:
                    if self._palette_display_items[idx] is not None:
                        break
                except Exception:
                    pass
                idx += 1

            if idx >= lb.size():
                idx = 0

            lb.selection_clear(0, tk.END)
            lb.selection_set(idx)
            lb.activate(idx)

    def _palette_move_selection(self, delta: int):
        lb = self._palette_list
        n = lb.size()
        if n <= 0:
            return "break"

        sel = lb.curselection()
        idx = sel[0] if sel else 0

        next_idx = idx
        for _ in range(n):
            next_idx = max(0, min(n - 1, next_idx + delta))
            if getattr(self, "_palette_display_items", None):
                if self._palette_display_items[next_idx] is not None:
                    break
            else:
                break
            if next_idx == 0 or next_idx == n - 1:
                break

        lb.selection_clear(0, tk.END)
        lb.selection_set(next_idx)
        lb.activate(next_idx)
        lb.see(next_idx)
        return "break"

    def _palette_autocomplete(self, event=None):
        if not getattr(self, "_palette_filtered", None):
            return "break"

        top = self._palette_filtered[0]
        label = top.get("label", "")
        if label:
            self._palette_var.set(label)
            self._filter_palette()
            try:
                w = self._palette.focus_get()
                if isinstance(w, tk.Entry):
                    w.icursor(tk.END)
            except Exception:
                pass
        return "break"

    def _run_palette_selection(self):
        lb = self._palette_list
        if lb.size() == 0:
            return

        sel = lb.curselection()
        idx = sel[0] if sel else 0

        if not getattr(self, "_palette_display_items", None):
            return

        if idx < 0 or idx >= len(self._palette_display_items):
            return

        item = self._palette_display_items[idx]
        if item is None:
            return

        self._close_palette()
        self._dispatch(item)
        
    def _open_manage_keywords(self):
        """
        Opens the Manage Keywords window (edits rss.keywords in config.json).
        """
        try:
            app_dir = get_app_dir()
            ManageKeywordsWindow(self, self.cfg, app_dir)
        except Exception as e:
            messagebox.showerror("Manage Keywords", str(e))

    def _open_shift_log_rules(self):
        """
        Opens the Shift Log Rules wizard (edits shift_log_rules in config.json).
        """
        try:
            app_dir = get_app_dir()
            ShiftLogRulesWindow(self, self.cfg, app_dir)
        except Exception as e:
            messagebox.showerror("Shift Log Rules", str(e))
        
    # ------------------------------------------------------------
    # Reports Catalog
    # ------------------------------------------------------------
    def _reports_catalog_items(self) -> list[dict]:
        items = self.cfg.get("reports_catalog", [])
        if not isinstance(items, list):
            return []

        out = []
        for it in items:
            if not isinstance(it, dict):
                continue
            title = _safe_str(it.get("title") or it.get("label") or "").strip()
            summary = _safe_str(it.get("summary") or it.get("description") or "").strip()
            last_updated = _safe_str(it.get("last_updated") or "").strip()
            action = _safe_str(it.get("action") or "").strip()
            target = it.get("target", "")
            tag = _safe_str(it.get("tag") or "").strip()

            if not title:
                continue

            out.append(
                {
                    "title": title,
                    "summary": summary,
                    "last_updated": last_updated,
                    "action": action,
                    "target": target,
                    "tag": tag,
                }
            )
        return out

    def _reports_truncate(self, text: str, max_len: int = 140) -> str:
        text = _safe_str(text).strip()
        if len(text) <= max_len:
            return text
        return text[: max_len - 3].rstrip() + "..."

    def _open_reports_catalog(self):
        if hasattr(self, "_reports_win") and self._reports_win and self._reports_win.winfo_exists():
            try:
                self._reports_win.lift()
                self._reports_win.focus_force()
            except Exception:
                pass
            return

        win = tk.Toplevel(self)
        self._reports_win = win
        win.title("Reports Catalog")
        win.configure(bg="#111111")
        win.geometry("1100x640")
        win.minsize(920, 520)
        win.resizable(True, True)

        try:
            self.update_idletasks()
            w, h = 1100, 640
            x = self.winfo_rootx() + max((self.winfo_width() - w) // 2, 0)
            y = self.winfo_rooty() + max((self.winfo_height() - h) // 2, 0)
            win.geometry(f"{w}x{h}+{x}+{y}")
        except Exception:
            pass

        win.bind("<Escape>", lambda e: win.destroy())
        win.protocol("WM_DELETE_WINDOW", win.destroy)

        # State
        self._reports_items_all = self._reports_catalog_items()
        self._reports_filtered = list(self._reports_items_all)
        self._reports_search_var = tk.StringVar(value="")
        self._reports_debounce_job = None
        self._reports_item_map = {}

        # Header
        top = tk.Frame(win, bg="#111111", padx=12, pady=10)
        top.pack(fill="x")

        left = tk.Frame(top, bg="#111111")
        left.pack(side="left", fill="x", expand=True)
        tk.Label(
            left,
            text="Reports Catalog",
            bg="#111111",
            fg="#ffffff",
            font=("Segoe UI", 14, "bold"),
        ).pack(anchor="w")

        self._reports_count_var = tk.StringVar(value="")
        tk.Label(
            left,
            textvariable=self._reports_count_var,
            bg="#111111",
            fg="#9aa7c0",
            font=("Segoe UI", 9, "normal"),
        ).pack(anchor="w", pady=(2, 0))

        # Search row
        filt = tk.Frame(win, bg="#111111", padx=12, pady=8)
        filt.pack(fill="x")
        tk.Label(
            filt, text="Search:", bg="#111111", fg="#cfcfcf", font=("Segoe UI", 10, "normal")
        ).pack(side="left")
        search = tk.Entry(
            filt,
            textvariable=self._reports_search_var,
            bg="#1e1e1e",
            fg="#ffffff",
            insertbackground="#ffffff",
            relief="flat",
            font=("Segoe UI", 11, "normal"),
            width=46,
        )
        search.pack(side="left", padx=(8, 12))
        search.focus_set()

        btns = tk.Frame(filt, bg="#111111")
        btns.pack(side="right")
        run_btn = tk.Button(
            btns,
            text="Run Selected",
            command=self._reports_run_selected,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        )
        run_btn.pack(side="left", padx=4)

        tk.Button(
            btns,
            text="Close",
            command=win.destroy,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(side="left", padx=4)

        # List area
        body = tk.Frame(win, bg="#111111", padx=12, pady=0)
        body.pack(fill="both", expand=True)

        cols = ("title", "summary", "updated")
        tree = ttk.Treeview(body, columns=cols, show="headings", selectmode="browse")
        self._reports_tree = tree
        tree.heading("title", text="Title")
        tree.heading("summary", text="Summary")
        tree.heading("updated", text="Last Updated")

        tree.column("title", width=260, stretch=False)
        tree.column("summary", width=620, stretch=True)
        tree.column("updated", width=140, stretch=False)

        ysb = ttk.Scrollbar(body, orient="vertical", command=tree.yview)
        tree.configure(yscrollcommand=ysb.set)
        tree.pack(side="left", fill="both", expand=True)
        ysb.pack(side="right", fill="y")

        tree.bind("<Double-Button-1>", lambda e: self._reports_run_selected())
        tree.bind("<Return>", lambda e: self._reports_run_selected())
        search.bind("<KeyRelease>", lambda e: self._reports_debounced_filter())

        self._reports_render_list()

    def _open_project_status_rollup(self):
        try:
            abs_path = _write_project_status_rollup_report(self.cfg, self.tokens)
        except Exception as e:
            messagebox.showerror("Project Status Rollup", f"Failed to create report.\n\n{e}")
            return

        try:
            _open_note_in_obsidian(self.cfg, self.tokens, abs_path)
        except Exception:
            pass

        messagebox.showinfo("Project Status Rollup", f"Report created:\n{abs_path}")

    def _reports_debounced_filter(self):
        if getattr(self, "_reports_debounce_job", None):
            try:
                self.after_cancel(self._reports_debounce_job)
            except Exception:
                pass
        self._reports_debounce_job = self.after(120, self._reports_apply_filter)

    def _reports_apply_filter(self):
        q = (self._reports_search_var.get() or "").strip().lower()
        if not q:
            self._reports_filtered = list(self._reports_items_all)
        else:
            filtered = []
            for it in self._reports_items_all:
                hay = f"{it.get('title','')} {it.get('summary','')} {it.get('tag','')}".lower()
                if q in hay:
                    filtered.append(it)
            self._reports_filtered = filtered

        self._reports_render_list()

    def _reports_render_list(self):
        tree = getattr(self, "_reports_tree", None)
        if not tree:
            return
        try:
            tree.delete(*tree.get_children())
        except Exception:
            return

        self._reports_item_map = {}
        for it in self._reports_filtered:
            title = it.get("title", "")
            summary = self._reports_truncate(it.get("summary", ""))
            updated = it.get("last_updated", "")
            item_id = tree.insert("", "end", values=(title, summary, updated))
            self._reports_item_map[item_id] = it

        shown = len(self._reports_filtered)
        total = len(self._reports_items_all)
        self._reports_count_var.set(f"Showing {shown} of {total} report(s)")

        if tree.get_children():
            first = tree.get_children()[0]
            tree.selection_set(first)
            tree.focus(first)

    def _reports_selected_item(self):
        tree = getattr(self, "_reports_tree", None)
        if not tree:
            return None
        sel = tree.selection()
        if not sel:
            return None
        return self._reports_item_map.get(sel[0])

    def _reports_run_selected(self):
        rep = self._reports_selected_item()
        if not rep:
            return
        action = _safe_str(rep.get("action") or "").strip()
        if not action:
            messagebox.showerror("Reports Catalog", "Selected report has no action.")
            return
        item = {
            "label": rep.get("title", "Report"),
            "action": action,
            "target": rep.get("target", ""),
            "section": "Reports",
        }
        self._dispatch(item)

    # ------------------------------------------------------------
    # IOC Enrichment (LLM)
    # ------------------------------------------------------------
    def _open_ioc_enrich(self):
        if hasattr(self, "_ioc_enrich_win") and self._ioc_enrich_win and self._ioc_enrich_win.winfo_exists():
            try:
                self._ioc_enrich_win.lift()
                self._ioc_enrich_win.focus_force()
            except Exception:
                pass
            return

        win = tk.Toplevel(self)
        self._ioc_enrich_win = win
        win.title("IOC Enrichment (AI)")
        win.configure(bg="#111111")
        win.geometry("900x640")
        win.minsize(820, 560)
        win.resizable(True, True)
        win.bind("<Escape>", lambda e: win.destroy())
        win.protocol("WM_DELETE_WINDOW", win.destroy)

        header = tk.Frame(win, bg="#111111", padx=12, pady=10)
        header.pack(fill="x")
        tk.Label(
            header,
            text="IOC Enrichment",
            bg="#111111",
            fg="#ffffff",
            font=("Segoe UI", 14, "bold"),
        ).pack(anchor="w")
        tk.Label(
            header,
            text="Paste IOCs or load a file, then generate an enrichment report.",
            bg="#111111",
            fg="#9aa7c0",
            font=("Segoe UI", 9, "normal"),
        ).pack(anchor="w", pady=(2, 0))

        body = tk.Frame(win, bg="#111111", padx=12, pady=8)
        body.pack(fill="both", expand=True)

        tk.Label(body, text="IOCs (one per line or comma-separated)", bg="#111111", fg="#cfcfcf").pack(anchor="w")
        ioc_text = tk.Text(
            body,
            height=10,
            bg="#1e1e1e",
            fg="#ffffff",
            insertbackground="#ffffff",
            relief="flat",
            font=("Segoe UI", 10, "normal"),
        )
        ioc_text.pack(fill="x", pady=(6, 10))

        tk.Label(body, text="Optional context", bg="#111111", fg="#cfcfcf").pack(anchor="w")
        context_text = tk.Text(
            body,
            height=6,
            bg="#1e1e1e",
            fg="#ffffff",
            insertbackground="#ffffff",
            relief="flat",
            font=("Segoe UI", 10, "normal"),
        )
        context_text.pack(fill="x", pady=(6, 10))

        controls = tk.Frame(body, bg="#111111")
        controls.pack(fill="x", pady=(4, 6))

        open_var = tk.BooleanVar(value=True)
        tk.Checkbutton(
            controls,
            text="Open report in Obsidian",
            variable=open_var,
            bg="#111111",
            fg="#cfcfcf",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#ffffff",
        ).pack(side="left")

        status_var = tk.StringVar(value="")
        tk.Label(controls, textvariable=status_var, bg="#111111", fg="#9aa7c0").pack(side="right")

        btn_row = tk.Frame(body, bg="#111111")
        btn_row.pack(fill="x", pady=(4, 0))

        def _load_file():
            path = filedialog.askopenfilename(
                parent=win,
                title="Select IOC file",
                filetypes=[
                    ("Text/CSV", "*.txt *.csv *.log *.ioc"),
                    ("All files", "*.*"),
                ],
            )
            if not path:
                return
            try:
                content = Path(path).read_text(encoding="utf-8", errors="ignore")
            except Exception as e:
                messagebox.showerror("IOC Enrichment", f"Failed to read file:\n{e}", parent=win)
                return
            ioc_text.delete("1.0", "end")
            ioc_text.insert("1.0", content.strip())

        def _run_enrichment():
            raw = ioc_text.get("1.0", "end").strip()
            iocs = _split_bulk(raw)
            if not iocs:
                messagebox.showerror("IOC Enrichment", "Please provide at least one IOC.", parent=win)
                return

            llm_cfg = _llm_config(self.cfg, self.tokens)
            if llm_cfg.get("provider") != "openai":
                messagebox.showerror("IOC Enrichment", "LLM provider is not set to OpenAI.", parent=win)
                return

            status_var.set("Running enrichment…")
            run_btn.config(state="disabled")

            context = context_text.get("1.0", "end").strip()

            def _worker():
                try:
                    prompt = (
                        "You are a cyber threat intelligence analyst. "
                        "Return JSON with fields: executive_summary, analyst_notes, "
                        "key_findings (array), metrics (array), actions (array), references (array)."
                    )
                    ioc_block = "\n".join(f"- {ioc}" for ioc in iocs)
                    user_msg = (
                        "IOC list:\n"
                        f"{ioc_block}\n\n"
                        "Context:\n"
                        f"{context or 'None'}\n\n"
                        "Provide concise enrichment focused on context, likely threats, and next steps."
                    )
                    messages = [
                        {"role": "system", "content": prompt},
                        {"role": "user", "content": user_msg},
                    ]
                    result = _openai_chat_json(
                        llm_cfg.get("api_key", ""),
                        llm_cfg.get("model", "gpt-4o"),
                        messages,
                    )
                    title = f"IOC Enrichment - {date.today().strftime('%Y-%m-%d')}"
                    report_path = _write_ioc_enrichment_report(
                        self.cfg,
                        self.tokens,
                        title,
                        result,
                        iocs,
                        context,
                    )

                    if open_var.get():
                        try:
                            _open_note_in_obsidian(self.cfg, self.tokens, report_path)
                        except Exception:
                            pass

                    def _done():
                        status_var.set("Report created.")
                        run_btn.config(state="normal")
                        messagebox.showinfo(
                            "IOC Enrichment",
                            f"Report created:\n{report_path}",
                            parent=win,
                        )

                    self.after(0, _done)
                except Exception as e:
                    def _err():
                        status_var.set("")
                        run_btn.config(state="normal")
                        messagebox.showerror("IOC Enrichment", str(e), parent=win)

                    self.after(0, _err)

            import threading
            threading.Thread(target=_worker, daemon=True).start()

        load_btn = tk.Button(
            btn_row,
            text="Load File",
            command=_load_file,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=12,
            pady=6,
        )
        load_btn.pack(side="left")

        run_btn = tk.Button(
            btn_row,
            text="Run Enrichment",
            command=_run_enrichment,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=12,
            pady=6,
        )
        run_btn.pack(side="right")

    def _open_threat_actor_profile(self):
        if hasattr(self, "_ta_profile_win") and self._ta_profile_win and self._ta_profile_win.winfo_exists():
            try:
                self._ta_profile_win.lift()
                self._ta_profile_win.focus_force()
            except Exception:
                pass
            return

        win = tk.Toplevel(self)
        self._ta_profile_win = win
        win.title("Threat Actor Profile (AI)")
        win.configure(bg="#111111")
        win.geometry("840x520")
        win.minsize(760, 480)
        win.resizable(True, True)
        win.bind("<Escape>", lambda e: win.destroy())
        win.protocol("WM_DELETE_WINDOW", win.destroy)

        header = tk.Frame(win, bg="#111111", padx=12, pady=10)
        header.pack(fill="x")
        tk.Label(
            header,
            text="Threat Actor Profile Generator",
            bg="#111111",
            fg="#ffffff",
            font=("Segoe UI", 14, "bold"),
        ).pack(anchor="w")
        tk.Label(
            header,
            text="Enter a MITRE ATT&CK Group ID (G####) or common name.",
            bg="#111111",
            fg="#9aa7c0",
            font=("Segoe UI", 9, "normal"),
        ).pack(anchor="w", pady=(2, 0))

        body = tk.Frame(win, bg="#111111", padx=12, pady=8)
        body.pack(fill="both", expand=True)

        tk.Label(body, text="MITRE ID or common name", bg="#111111", fg="#cfcfcf").pack(anchor="w")
        input_var = tk.StringVar(value="")
        input_entry = tk.Entry(
            body,
            textvariable=input_var,
            bg="#1e1e1e",
            fg="#ffffff",
            insertbackground="#ffffff",
            relief="flat",
            font=("Segoe UI", 11, "normal"),
        )
        input_entry.pack(fill="x", pady=(6, 10))
        input_entry.focus_set()

        tk.Label(body, text="Optional context (notes or constraints)", bg="#111111", fg="#cfcfcf").pack(anchor="w")
        context_text = tk.Text(
            body,
            height=6,
            bg="#1e1e1e",
            fg="#ffffff",
            insertbackground="#ffffff",
            relief="flat",
            font=("Segoe UI", 10, "normal"),
        )
        context_text.pack(fill="x", pady=(6, 10))

        controls = tk.Frame(body, bg="#111111")
        controls.pack(fill="x", pady=(4, 6))

        open_var = tk.BooleanVar(value=True)
        tk.Checkbutton(
            controls,
            text="Open note in Obsidian",
            variable=open_var,
            bg="#111111",
            fg="#cfcfcf",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#ffffff",
        ).pack(side="left")

        status_var = tk.StringVar(value="")
        tk.Label(controls, textvariable=status_var, bg="#111111", fg="#9aa7c0").pack(side="right")

        btn_row = tk.Frame(body, bg="#111111")
        btn_row.pack(fill="x", pady=(4, 0))

        def _run_profile():
            raw = (input_var.get() or "").strip()
            if not raw:
                messagebox.showerror("Threat Actor Profile", "Please enter a MITRE ID or common name.", parent=win)
                return

            llm_cfg = _llm_config(self.cfg, self.tokens)
            if llm_cfg.get("provider") != "openai":
                messagebox.showerror("Threat Actor Profile", "LLM provider is not set to OpenAI.", parent=win)
                return

            status_var.set("Generating profile…")
            run_btn.config(state="disabled")

            context = context_text.get("1.0", "end").strip()

            def _worker():
                try:
                    prompt = _load_vault_prompt(
                        self.cfg,
                        self.tokens,
                        "00_System/03_Prompts/04_CIPHER/SCOUT-TA - Threat Actor Profile Generator.md",
                    )
                    system_msg = (
                        prompt
                        + "\n\n"
                        + "If the input is a common name without a MITRE Group ID, "
                        + "infer the most likely Group ID and proceed. "
                        + "If you are unsure, still generate the full note but leave "
                        + "actor_id empty and explain uncertainty in Analyst Notes."
                    )
                    user_msg = f"[SCOUT-TA] {raw}"
                    if context:
                        user_msg += f"\n\nContext:\n{context}"
                    messages = [
                        {"role": "system", "content": system_msg},
                        {"role": "user", "content": user_msg},
                    ]
                    content = _openai_chat_text(
                        llm_cfg.get("api_key", ""),
                        llm_cfg.get("model", "gpt-4o"),
                        messages,
                    )

                    if content.strip().startswith("ERROR:"):
                        raise ValueError(content.strip())

                    if "entity_type: threat_actor" not in content:
                        inferred_id = ""
                        for src in (content, raw):
                            m = re.search(r"\bG\d{4}\b", src or "")
                            if m:
                                inferred_id = m.group(0)
                                break
                        content = _threat_actor_frontmatter_template(raw, inferred_id) + content

                    actor_id_match = re.search(r"^actor_id:\s*\"?([A-Za-z0-9_-]*)\"?\s*$", content, re.MULTILINE)
                    actor_id = (actor_id_match.group(1).strip() if actor_id_match else "")
                    valid_actor_id = bool(re.fullmatch(r"G\d{4}", actor_id))

                    if not valid_actor_id:
                        result = {"value": None}
                        done = threading.Event()

                        def _prompt():
                            val = simpledialog.askstring(
                                "Threat Actor Profile",
                                "Enter a 4-digit tracking ID (e.g., 1234):",
                                parent=win,
                            )
                            if val is None:
                                result["value"] = None
                                done.set()
                                return
                            digits = re.sub(r"\D", "", val or "")
                            if len(digits) != 4:
                                messagebox.showerror(
                                    "Threat Actor Profile",
                                    "Please enter exactly 4 digits.",
                                    parent=win,
                                )
                                self.after(0, _prompt)
                                return
                            result["value"] = digits
                            done.set()

                        self.after(0, _prompt)
                        done.wait()

                        if result["value"] is None:
                            raise ValueError("Threat Actor Profile canceled.")

                        content = _apply_yaml_updates(content, {"actor_id": f"G{result['value']}"})

                    label = raw
                    abs_path = _write_threat_actor_profile_note(self.cfg, self.tokens, label, content)

                    if open_var.get():
                        try:
                            _open_note_in_obsidian(self.cfg, self.tokens, abs_path)
                        except Exception:
                            pass

                    def _done():
                        status_var.set("Note created.")
                        run_btn.config(state="normal")
                        messagebox.showinfo(
                            "Threat Actor Profile",
                            f"Profile created:\n{abs_path}",
                            parent=win,
                        )

                    self.after(0, _done)
                except Exception as e:
                    err_msg = str(e)

                    def _err(msg=err_msg):
                        status_var.set("")
                        run_btn.config(state="normal")
                        messagebox.showerror("Threat Actor Profile", msg, parent=win)

                    self.after(0, _err)

            import threading
            threading.Thread(target=_worker, daemon=True).start()

        run_btn = tk.Button(
            btn_row,
            text="Generate Profile",
            command=_run_profile,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=12,
            pady=6,
        )
        run_btn.pack(side="right")

    # ------------------------------------------------------------
    # Obsidian Note -> PowerPoint
    # ------------------------------------------------------------
    def _open_obsidian_to_pptx(self):
        if hasattr(self, "_pptx_win") and self._pptx_win and self._pptx_win.winfo_exists():
            try:
                self._pptx_win.lift()
                self._pptx_win.focus_force()
            except Exception:
                pass
            return

        win = tk.Toplevel(self)
        self._pptx_win = win
        win.title("Obsidian to PowerPoint")
        win.configure(bg="#111111")
        win.geometry("820x420")
        win.minsize(760, 380)
        win.resizable(True, False)
        win.bind("<Escape>", lambda e: win.destroy())
        win.protocol("WM_DELETE_WINDOW", win.destroy)

        header = tk.Frame(win, bg="#111111", padx=12, pady=10)
        header.pack(fill="x")
        tk.Label(
            header,
            text="Obsidian Note → PowerPoint",
            bg="#111111",
            fg="#ffffff",
            font=("Segoe UI", 14, "bold"),
        ).pack(anchor="w")
        tk.Label(
            header,
            text="Select a markdown note, a PowerPoint template, and an output file.",
            bg="#111111",
            fg="#9aa7c0",
            font=("Segoe UI", 9, "normal"),
        ).pack(anchor="w", pady=(2, 0))

        body = tk.Frame(win, bg="#111111", padx=12, pady=8)
        body.pack(fill="both", expand=True)

        vault_root = (
            self.tokens.get("SCOUT_ROOT")
            or self.tokens.get("VAULT_ROOT")
            or self.tokens.get("CIPHER_ROOT")
        )
        initial_dir = str(Path(vault_root)) if vault_root else None

        note_var = tk.StringVar(value="")
        template_var = tk.StringVar(value="")
        output_var = tk.StringVar(value="")
        type_var = tk.StringVar(value="Auto (Detect)")
        status_var = tk.StringVar(value="")

        def _row(label, var, browse_cmd):
            row = tk.Frame(body, bg="#111111")
            row.pack(fill="x", pady=(0, 10))
            tk.Label(row, text=label, bg="#111111", fg="#cfcfcf").pack(anchor="w")
            inner = tk.Frame(row, bg="#111111")
            inner.pack(fill="x", pady=(6, 0))
            entry = tk.Entry(
                inner,
                textvariable=var,
                bg="#1e1e1e",
                fg="#ffffff",
                insertbackground="#ffffff",
                relief="flat",
                font=("Segoe UI", 10, "normal"),
            )
            entry.pack(side="left", fill="x", expand=True)
            tk.Button(
                inner,
                text="Browse",
                command=browse_cmd,
                bg="#2A2A2A",
                fg="#FFFFFF",
                relief="flat",
                padx=12,
                pady=4,
            ).pack(side="left", padx=(8, 0))
            return entry

        def _browse_note():
            path = filedialog.askopenfilename(
                parent=win,
                title="Select Obsidian Note",
                initialdir=initial_dir,
                filetypes=[("Markdown", "*.md"), ("All files", "*.*")],
            )
            if path:
                note_var.set(path)
                if not output_var.get():
                    out_path = Path(path).with_suffix(".pptx")
                    output_var.set(str(out_path))
                try:
                    detected = _pptx_detect_type(Path(path))
                    if detected:
                        type_var.set(PPTX_PRESENTATION_TYPES[detected]["label"])
                except Exception:
                    pass

        def _browse_template():
            path = filedialog.askopenfilename(
                parent=win,
                title="Select PowerPoint Template",
                initialdir=initial_dir,
                filetypes=[("PowerPoint", "*.pptx"), ("All files", "*.*")],
            )
            if path:
                template_var.set(path)

        def _browse_output():
            path = filedialog.asksaveasfilename(
                parent=win,
                title="Save PowerPoint As",
                initialdir=initial_dir,
                defaultextension=".pptx",
                filetypes=[("PowerPoint", "*.pptx")],
            )
            if path:
                output_var.set(path)

        _row("Obsidian note (.md)", note_var, _browse_note)
        _row("PowerPoint template (.pptx)", template_var, _browse_template)
        _row("Output presentation (.pptx)", output_var, _browse_output)

        type_row = tk.Frame(body, bg="#111111")
        type_row.pack(fill="x", pady=(0, 10))
        tk.Label(type_row, text="Presentation type", bg="#111111", fg="#cfcfcf").pack(anchor="w")
        type_opts = ["Auto (Detect)"] + [
            cfg["label"] for cfg in PPTX_PRESENTATION_TYPES.values()
        ]
        type_combo = ttk.Combobox(
            type_row,
            values=type_opts,
            textvariable=type_var,
            state="readonly",
            width=40,
        )
        type_combo.pack(anchor="w", pady=(6, 0))

        controls = tk.Frame(body, bg="#111111")
        controls.pack(fill="x", pady=(6, 4))

        open_var = tk.BooleanVar(value=True)
        tk.Checkbutton(
            controls,
            text="Open presentation after export",
            variable=open_var,
            bg="#111111",
            fg="#cfcfcf",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#ffffff",
        ).pack(side="left")

        tk.Label(controls, textvariable=status_var, bg="#111111", fg="#9aa7c0").pack(side="right")

        def _run_export():
            note_path = Path(note_var.get().strip())
            template_path = Path(template_var.get().strip())
            output_path = Path(output_var.get().strip())

            if not note_path.exists():
                messagebox.showerror("Obsidian to PowerPoint", "Select a valid Obsidian note.", parent=win)
                return
            if not template_path.exists():
                messagebox.showerror("Obsidian to PowerPoint", "Select a valid PowerPoint template.", parent=win)
                return
            if output_path.suffix.lower() != ".pptx":
                output_path = output_path.with_suffix(".pptx")
                output_var.set(str(output_path))

            status_var.set("Building presentation...")
            try:
                selection = (type_var.get() or "").strip()
                if selection == "Auto (Detect)" or not selection:
                    presentation_type = _pptx_detect_type(note_path) or "threat_intel"
                else:
                    presentation_type = "threat_intel"
                    for key, cfg in PPTX_PRESENTATION_TYPES.items():
                        if cfg.get("label") == selection:
                            presentation_type = key
                            break
                result_path = _pptx_build_presentation(
                    note_path,
                    template_path,
                    output_path,
                    presentation_type,
                )
                status_var.set("Presentation created.")
                if open_var.get():
                    try:
                        os.startfile(str(result_path))
                    except Exception:
                        pass
                messagebox.showinfo(
                    "Obsidian to PowerPoint",
                    f"Presentation created:\n{result_path}",
                    parent=win,
                )
            except Exception as e:
                status_var.set("")
                messagebox.showerror("Obsidian to PowerPoint", str(e), parent=win)

        btn_row = tk.Frame(body, bg="#111111")
        btn_row.pack(fill="x", pady=(6, 4))
        tk.Button(
            btn_row,
            text="Generate Presentation",
            command=_run_export,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=12,
            pady=6,
        ).pack(side="right")

    # ------------------------------------------------------------
    # RSS Review / Import (Treeview UI)
    # ------------------------------------------------------------
    def _open_rss_review(self):
        self._open_rss_window(select_tab="review")

    def _open_rss_window(self, select_tab: str = "review"):
        # Pre-compute watchlist stats so header diagnostics are available immediately
        try:
            kws = self._rss_watchlist_keywords()
            stats = self.state.get("rss_autoflag_stats", {}) if isinstance(self.state.get("rss_autoflag_stats", {}), dict) else {}
            stats["keywords_loaded"] = int(len(kws))
            self.state["rss_autoflag_stats"] = stats
            try:
                save_state(self.state_path, self.state)
            except Exception:
                pass
        except Exception:
            pass

        if hasattr(self, "_rss_review_win") and self._rss_review_win and self._rss_review_win.winfo_exists():
            try:
                self._rss_review_win.lift()
                self._rss_review_win.focus_force()
                if getattr(self, "_rss_notebook", None):
                    if select_tab == "schedule":
                        self._rss_notebook.select(self._rss_tab_schedule)
                    else:
                        self._rss_notebook.select(self._rss_tab_review)
            except Exception:
                pass
            return

        win = tk.Toplevel(self)
        self._rss_review_win = win
        self._rss_win = win
        win.title("RSS")
        win.configure(bg="#111111")
        win.geometry("1200x720")
        win.minsize(980, 600)

        try:
            self.update_idletasks()
            w, h = 1200, 720
            x = self.winfo_rootx() + max((self.winfo_width() - w) // 2, 0)
            y = self.winfo_rooty() + max((self.winfo_height() - h) // 2, 0)
            win.geometry(f"{w}x{h}+{x}+{y}")
        except Exception:
            pass

        win.protocol("WM_DELETE_WINDOW", self._rss_review_close)

        # stateful vars
        self._rss_items_all = self._rss_load_items()
        self._rss_filter_mode = tk.StringVar(value="new")   # new|flagged|imported|reported|all
        self._rss_search_var = tk.StringVar(value="")
        self._rss_source_var = tk.StringVar(value="All Sources")
        self._rss_status_var = tk.StringVar(value="")
        self._rss_debounce_job = None

        notebook = ttk.Notebook(win)
        self._rss_notebook = notebook
        self._rss_tab_review = tk.Frame(notebook, bg="#111111")
        self._rss_tab_schedule = tk.Frame(notebook, bg="#111111")
        self._rss_tab_keywords = tk.Frame(notebook, bg="#111111")
        notebook.add(self._rss_tab_review, text="Review")
        notebook.add(self._rss_tab_schedule, text="Schedule")
        notebook.add(self._rss_tab_keywords, text="Keywords")
        notebook.pack(fill="both", expand=True)

        self._build_rss_review_tab(self._rss_tab_review)
        self._build_rss_schedule_tab(self._rss_tab_schedule)
        self._build_rss_keywords_tab(self._rss_tab_keywords)

        if select_tab == "schedule":
            notebook.select(self._rss_tab_schedule)
        elif select_tab == "keywords":
            notebook.select(self._rss_tab_keywords)
        else:
            notebook.select(self._rss_tab_review)

        # Keyboard shortcuts
        win.bind("/", lambda e: (self._rss_focus_search(), "break"))
        win.bind("<Escape>", lambda e: self._rss_clear_search())
        win.bind("j", lambda e: self._rss_move_selection(1))
        win.bind("k", lambda e: self._rss_move_selection(-1))
        win.bind("o", lambda e: self._rss_open_selected_url())
        win.bind("f", lambda e: self._rss_toggle_flag_selected())
        win.bind("m", lambda e: self._rss_mark_read_selected())
        win.bind("r", lambda e: self._rss_toggle_report_selected())
        win.bind("g", lambda e: self._rss_generate_report())
        win.bind("w", lambda e: self._rss_generate_weekly_rollup())
        win.bind("a", lambda e: self._rss_mark_all_as_read())

        self._rss_apply_filters()
        self._rss_status_var.set(
            "Shortcuts: j/k navigate • f flag • m mark read • r report • g generate • w weekly • o open • / search • Esc clear"
        )

    # ------------------------------------------------------------
    # RSS Keywords / Watchlists (Tab)
    # ------------------------------------------------------------
    def _build_rss_keywords_tab(self, parent: tk.Frame):
        bg = "#111111"
        panel = "#141824"
        input_bg = "#1e1e1e"
        btn_bg = "#1a1a1a"
        btn_hover = "#2a2a2a"
        fg = "#ffffff"
        fg2 = "#cfcfcf"
        meta = "#9aa7c0"

        wrap = tk.Frame(parent, bg=bg, padx=14, pady=12)
        wrap.pack(fill="both", expand=True)

        tk.Label(wrap, text="RSS Watchlists", bg=bg, fg=fg, font=("Segoe UI", 13, "bold")).pack(anchor="w")
        tk.Label(
            wrap,
            text="Manage watchlists in Obsidian. Keywords, categories, and weights control RSS ranking.",
            bg=bg,
            fg=meta,
            font=("Segoe UI", 9, "normal"),
        ).pack(anchor="w", pady=(2, 10))

        body = tk.Frame(wrap, bg=bg)
        body.pack(fill="both", expand=True)

        right = tk.Frame(body, bg=panel, padx=10, pady=10)
        right.pack(fill="both", expand=True)

        # ----------------------------
        # Watchlist Keywords (Obsidian)
        # ----------------------------
        tk.Label(right, text="Watchlist Files", bg=panel, fg=fg, font=("Segoe UI", 12, "bold")).pack(anchor="w")

        self._rss_wl_dir_var = tk.StringVar(value="")
        tk.Label(right, textvariable=self._rss_wl_dir_var, bg=panel, fg=meta, font=("Segoe UI", 9, "normal")).pack(
            anchor="w", pady=(4, 6)
        )

        # Watchlists directory (configurable)
        dir_row = tk.Frame(right, bg=panel)
        dir_row.pack(fill="x", pady=(0, 6))

        self._rss_wl_dir_input_var = tk.StringVar(
            value=str(self.cfg.get("rss", {}).get("watchlists_dir") or "")
        )
        tk.Label(dir_row, text="Watchlists dir:", bg=panel, fg=fg2).pack(side="left")
        tk.Entry(
            dir_row,
            textvariable=self._rss_wl_dir_input_var,
            bg=input_bg,
            fg=fg,
            insertbackground=fg,
            relief="flat",
            font=("Segoe UI", 10, "normal"),
        ).pack(side="left", fill="x", expand=True, padx=(8, 6))

        def _wl_save_dir():
            rss_cfg = self.cfg.setdefault("rss", {})
            rss_cfg["watchlists_dir"] = self._rss_wl_dir_input_var.get().strip()
            try:
                app_dir = get_app_dir()
                cfg_path = _config_path_for_launcher(app_dir)
                _write_config_file_with_backup(cfg_path, self.cfg)
                _wl_refresh_files()
            except Exception as e:
                messagebox.showerror("Save Watchlists Dir", str(e), parent=parent)

        def _wl_set_default_dir():
            self._rss_wl_dir_input_var.set("{VAULT_ROOT}/10_Operations/13_Watchlists")

        def _wl_browse_dir():
            try:
                start = self._rss_wl_dir_input_var.get().strip()
                start = resolve_path(start, self.tokens) if start else ""
            except Exception:
                start = ""
            path = filedialog.askdirectory(initialdir=start or None, parent=parent)
            if path:
                self._rss_wl_dir_input_var.set(path)

        tk.Button(dir_row, text="Browse", command=_wl_browse_dir, bg=btn_bg, fg=fg, relief="flat").pack(
            side="right", padx=4
        )
        tk.Button(dir_row, text="Use Vault Default", command=_wl_set_default_dir, bg=btn_bg, fg=fg, relief="flat").pack(
            side="right", padx=4
        )
        tk.Button(dir_row, text="Save Path", command=_wl_save_dir, bg=btn_bg, fg=fg, relief="flat").pack(
            side="right", padx=4
        )

        wl_top = tk.Frame(right, bg=panel)
        wl_top.pack(fill="x", pady=(0, 6))

        def _wl_btn(text, cmd):
            b = tk.Button(wl_top, text=text, command=cmd, bg=btn_bg, fg=fg, relief="flat", padx=8, pady=4,
                          font=("Segoe UI", 9, "bold"))
            b.pack(side="right", padx=4)
            b.bind("<Enter>", lambda e: b.configure(bg=btn_hover))
            b.bind("<Leave>", lambda e: b.configure(bg=btn_bg))
            return b

        _wl_btn("Refresh", lambda: _wl_refresh_files())
        _wl_btn("Save", lambda: _wl_save_keywords())

        wl_body = tk.PanedWindow(right, orient="horizontal", bg=panel, sashwidth=6, sashrelief="flat")
        wl_body.pack(fill="both", expand=True)

        wl_files = tk.Frame(wl_body, bg=panel, padx=6, pady=6)
        wl_keys = tk.Frame(wl_body, bg=panel, padx=6, pady=6)
        wl_body.add(wl_files, minsize=180)
        wl_body.add(wl_keys, minsize=260)

        tk.Label(wl_files, text="Watchlists", bg=panel, fg=fg2, font=("Segoe UI", 10, "bold")).pack(anchor="w")
        self._rss_wl_listbox = tk.Listbox(
            wl_files, bg="#1a1a1a", fg=fg, selectbackground="#2a2a2a", selectforeground=fg,
            relief="flat", highlightthickness=0, font=("Segoe UI", 10, "normal")
        )
        self._rss_wl_listbox.pack(fill="both", expand=True, pady=(6, 0))

        tk.Label(wl_keys, text="Keywords", bg=panel, fg=fg2, font=("Segoe UI", 10, "bold")).pack(anchor="w")
        self._rss_wl_kw_list = tk.Listbox(
            wl_keys, bg="#1a1a1a", fg=fg, selectbackground="#2a2a2a", selectforeground=fg,
            relief="flat", highlightthickness=0, font=("Segoe UI", 10, "normal")
        )
        self._rss_wl_kw_list.pack(fill="both", expand=True, pady=(6, 6))

        wl_form = tk.Frame(wl_keys, bg=panel)
        wl_form.pack(fill="x")

        self._rss_wl_term_var = tk.StringVar(value="")
        self._rss_wl_cats_var = tk.StringVar(value="")
        self._rss_wl_weight_var = tk.StringVar(value="1")

        tk.Label(wl_form, text="Term", bg=panel, fg=fg2).grid(row=0, column=0, sticky="w")
        tk.Entry(
            wl_form, textvariable=self._rss_wl_term_var, bg=input_bg, fg=fg, insertbackground=fg, relief="flat",
            font=("Segoe UI", 10, "normal")
        ).grid(row=1, column=0, sticky="ew", padx=(0, 6), pady=(0, 6))

        tk.Label(wl_form, text="Categories (comma)", bg=panel, fg=fg2).grid(row=0, column=1, sticky="w")
        tk.Entry(
            wl_form, textvariable=self._rss_wl_cats_var, bg=input_bg, fg=fg, insertbackground=fg, relief="flat",
            font=("Segoe UI", 10, "normal")
        ).grid(row=1, column=1, sticky="ew", padx=(0, 6), pady=(0, 6))

        tk.Label(wl_form, text="Weight", bg=panel, fg=fg2).grid(row=0, column=2, sticky="w")
        tk.Entry(
            wl_form, textvariable=self._rss_wl_weight_var, bg=input_bg, fg=fg, insertbackground=fg, relief="flat",
            font=("Segoe UI", 10, "normal"), width=6
        ).grid(row=1, column=2, sticky="w", pady=(0, 6))

        wl_form.columnconfigure(0, weight=2)
        wl_form.columnconfigure(1, weight=2)

        wl_btns = tk.Frame(wl_keys, bg=panel)
        wl_btns.pack(fill="x", pady=(6, 0))
        tk.Button(wl_btns, text="+ Add", command=lambda: _wl_add_keyword(), bg=btn_bg, fg=fg, relief="flat").pack(
            side="left", padx=4
        )
        tk.Button(wl_btns, text="Update", command=lambda: _wl_update_keyword(), bg=btn_bg, fg=fg, relief="flat").pack(
            side="left", padx=4
        )
        tk.Button(wl_btns, text="Remove", command=lambda: _wl_remove_keywords(), bg=btn_bg, fg=fg, relief="flat").pack(
            side="left", padx=4
        )
        tk.Button(wl_btns, text="Bulk Add", command=lambda: _wl_bulk_add(), bg=btn_bg, fg=fg, relief="flat").pack(
            side="left", padx=4
        )
        tk.Button(wl_btns, text="Dedupe", command=lambda: _wl_dedupe(), bg=btn_bg, fg=fg, relief="flat").pack(
            side="left", padx=4
        )
        tk.Button(wl_btns, text="Normalize", command=lambda: _wl_normalize(), bg=btn_bg, fg=fg, relief="flat").pack(
            side="left", padx=4
        )

        # Frequency distribution
        freq = tk.Frame(right, bg=panel)
        freq.pack(fill="x", pady=(10, 0))
        tk.Label(freq, text="Keyword Frequency (RSS items)", bg=panel, fg=fg, font=("Segoe UI", 11, "bold")).pack(
            anchor="w"
        )
        freq_body = tk.Frame(freq, bg=panel)
        freq_body.pack(fill="x", pady=(6, 0))

        tk.Label(freq_body, text="Most frequent", bg=panel, fg=fg2).pack(anchor="w")
        self._rss_wl_freq_top = tk.Listbox(
            freq_body, bg="#1a1a1a", fg=fg, relief="flat", highlightthickness=0, font=("Segoe UI", 9, "normal"), height=6
        )
        self._rss_wl_freq_top.pack(fill="x", pady=(2, 6))

        tk.Label(freq_body, text="Least frequent", bg=panel, fg=fg2).pack(anchor="w")
        self._rss_wl_freq_bottom = tk.Listbox(
            freq_body, bg="#1a1a1a", fg=fg, relief="flat", highlightthickness=0, font=("Segoe UI", 9, "normal"), height=6
        )
        self._rss_wl_freq_bottom.pack(fill="x", pady=(2, 0))

        # ----------------------------
        # Helpers
        # ----------------------------
        self._rss_wl_current_path = None

        def _resolve_watchlists_dir() -> tuple[str, str]:
            rss_cfg = self.cfg.get("rss", {}) if isinstance(self.cfg.get("rss", {}), dict) else {}
            candidates = []
            wdir_cfg = rss_cfg.get("watchlists_dir") or rss_cfg.get("watchlist_dir") or ""
            if isinstance(wdir_cfg, str) and wdir_cfg.strip():
                candidates.append(wdir_cfg.strip())
            candidates.append("{VAULT_ROOT}/10_Operations/13_Watchlists")
            candidates.append("{VAULT_ROOT}/30_CIPHER/05_Watchlists")

            resolved_dir = ""
            selected_expr = ""
            for expr in candidates:
                try:
                    s = resolve_path(expr, self.tokens)
                except Exception:
                    s = expr
                s = (s or "").strip()
                if not s:
                    continue
                try:
                    p = Path(s)
                    if not p.is_absolute() and self.tokens.get("VAULT_ROOT"):
                        s = str((Path(self.tokens["VAULT_ROOT"]) / s).resolve())
                except Exception:
                    pass
                try:
                    s = str(Path(s).resolve())
                except Exception:
                    pass
                if Path(s).exists() and Path(s).is_dir():
                    resolved_dir = s
                    selected_expr = expr
                    break
                if not resolved_dir:
                    resolved_dir = s
                    selected_expr = expr
            return resolved_dir, selected_expr

        def _render_watchlist_keywords(items: list[dict]):
            self._rss_wl_items = items or []
            self._rss_wl_kw_list.delete(0, tk.END)
            for it in self._rss_wl_items:
                term = it.get("term") or ""
                cats = ", ".join(it.get("categories") or [])
                weight = it.get("weight", 1)
                label = term
                if cats:
                    label += f"  [{cats}]"
                label += f"  (w={weight})"
                self._rss_wl_kw_list.insert(tk.END, label)

        def _wl_dedupe_items(items: list[dict]) -> list[dict]:
            out = []
            seen = set()
            for it in items or []:
                term = _normalize_keyword(it.get("term") or "")
                if not term:
                    continue
                key = term.lower()
                if key in seen:
                    continue
                seen.add(key)
                out.append({"term": term, "categories": it.get("categories") or [], "weight": it.get("weight", 1)})
            return out

        def _wl_refresh_files():
            resolved_dir, expr = _resolve_watchlists_dir()
            self._rss_wl_dir_var.set(f"Dir: {resolved_dir or '(not resolved)'}")
            self._rss_wl_dir_resolved = resolved_dir
            self._rss_wl_dir_expr = expr
            self._rss_wl_listbox.delete(0, tk.END)
            self._rss_wl_current_path = None
            _render_watchlist_keywords([])
            if not resolved_dir or not Path(resolved_dir).exists():
                return
            files = sorted([p for p in Path(resolved_dir).rglob("*.md") if p.is_file()])
            for p in files:
                self._rss_wl_listbox.insert(tk.END, str(p))
            if files:
                self._rss_wl_listbox.selection_set(0)
                self._rss_wl_listbox.activate(0)
                self._rss_wl_listbox.see(0)
                _wl_select_file()

        def _wl_select_file(_evt=None):
            sel = self._rss_wl_listbox.curselection()
            if not sel:
                return
            path = self._rss_wl_listbox.get(sel[0])
            self._rss_wl_current_path = path
            try:
                raw = Path(path).read_text(encoding="utf-8", errors="ignore")
            except Exception:
                _render_watchlist_keywords([])
                return
            parsed = self._rss_parse_watchlist_text(raw)
            defaults = parsed.get("defaults") or {}
            items = []
            for it in parsed.get("items") or []:
                term = _normalize_keyword(it.get("term") or "")
                if not term:
                    continue
                cats = it.get("categories") or defaults.get("categories") or []
                if isinstance(cats, str):
                    cats = [c.strip() for c in cats.split(",") if c.strip()]
                weight = it.get("weight")
                if weight is None:
                    weight = defaults.get("weight", 1)
                items.append({"term": term, "categories": cats, "weight": weight})
            _render_watchlist_keywords(items)

        def _wl_select_keyword(_evt=None):
            sel = self._rss_wl_kw_list.curselection()
            if not sel:
                return
            idx = sel[0]
            if idx < 0 or idx >= len(self._rss_wl_items):
                return
            it = self._rss_wl_items[idx]
            self._rss_wl_term_var.set(it.get("term") or "")
            self._rss_wl_cats_var.set(", ".join(it.get("categories") or []))
            self._rss_wl_weight_var.set(str(it.get("weight", 1)))

        def _wl_add_keyword():
            if not self._rss_wl_current_path:
                return
            term = _normalize_keyword(self._rss_wl_term_var.get())
            if not term:
                return
            cats = [c.strip() for c in (self._rss_wl_cats_var.get() or "").split(",") if c.strip()]
            try:
                weight = int(float(self._rss_wl_weight_var.get() or 1))
            except Exception:
                weight = 1
            current = list(self._rss_wl_items or [])
            current.append({"term": term, "categories": cats, "weight": weight})
            current = _wl_dedupe_items(current)
            _render_watchlist_keywords(current)
            self._rss_wl_term_var.set("")

        def _wl_bulk_add():
            if not self._rss_wl_current_path:
                return
            raw = simpledialog.askstring("Bulk Add Watchlist Keywords", "Paste keywords (newline or comma separated):", parent=parent)
            if not raw:
                return
            items = _split_bulk(raw)
            if not items:
                return
            cats = [c.strip() for c in (self._rss_wl_cats_var.get() or "").split(",") if c.strip()]
            try:
                weight = int(float(self._rss_wl_weight_var.get() or 1))
            except Exception:
                weight = 1
            current = list(self._rss_wl_items or [])
            for t in items:
                term = _normalize_keyword(t)
                if term:
                    current.append({"term": term, "categories": cats, "weight": weight})
            current = _wl_dedupe_items(current)
            _render_watchlist_keywords(current)

        def _wl_update_keyword():
            sel = self._rss_wl_kw_list.curselection()
            if not sel:
                return
            idx = sel[0]
            if idx < 0 or idx >= len(self._rss_wl_items):
                return
            term = _normalize_keyword(self._rss_wl_term_var.get())
            if not term:
                return
            cats = [c.strip() for c in (self._rss_wl_cats_var.get() or "").split(",") if c.strip()]
            try:
                weight = int(float(self._rss_wl_weight_var.get() or 1))
            except Exception:
                weight = 1
            current = list(self._rss_wl_items or [])
            current[idx] = {"term": term, "categories": cats, "weight": weight}
            current = _wl_dedupe_items(current)
            _render_watchlist_keywords(current)

        def _wl_remove_keywords():
            if not self._rss_wl_current_path:
                return
            sel = list(self._rss_wl_kw_list.curselection())
            if not sel:
                return
            current = list(self._rss_wl_items or [])
            for i in sorted(sel, reverse=True):
                if 0 <= i < len(current):
                    current.pop(i)
            _render_watchlist_keywords(current)

        def _wl_dedupe():
            current = list(self._rss_wl_items or [])
            _render_watchlist_keywords(_wl_dedupe_items(current))

        def _wl_normalize():
            current = list(self._rss_wl_items or [])
            for it in current:
                it["term"] = _normalize_keyword(it.get("term") or "")
            _render_watchlist_keywords([it for it in current if it.get("term")])

        def _wl_save_keywords():
            if not self._rss_wl_current_path:
                return
            path = Path(self._rss_wl_current_path)
            try:
                raw = path.read_text(encoding="utf-8", errors="ignore")
            except Exception as e:
                messagebox.showerror("Save Watchlist", str(e), parent=parent)
                return
            parsed = self._rss_parse_watchlist_text(raw)
            fm_end = parsed.get("fm_end", -1)
            kw_start = parsed.get("kw_start", -1)
            kw_end = parsed.get("kw_end", -1)
            current = list(self._rss_wl_items or [])
            block = self._rss_format_watchlist_keywords_block(current)
            lines = raw.splitlines()
            if fm_end == -1:
                new_lines = ["---"] + block + ["---", ""] + lines
            else:
                if kw_start != -1 and kw_end != -1:
                    new_lines = lines[:kw_start] + block + lines[kw_end:]
                else:
                    new_lines = lines[:fm_end] + block + lines[fm_end:]
            try:
                path.write_text("\n".join(new_lines), encoding="utf-8")
            except Exception as e:
                messagebox.showerror("Save Watchlist", str(e), parent=parent)

        # bindings
        self._rss_wl_listbox.bind("<<ListboxSelect>>", _wl_select_file)
        self._rss_wl_kw_list.bind("<<ListboxSelect>>", _wl_select_keyword)

        # init
        _wl_refresh_files()
        self._rss_update_watchlist_freq(self._rss_items_all or [])

    def _rss_focus_search(self):
        try:
            entry = getattr(self, "_rss_search_entry", None)
            if entry:
                entry.focus_set()
        except Exception:
            pass

    def _build_rss_review_tab(self, parent: tk.Frame):
        # header / toolbar
        top = tk.Frame(parent, bg="#111111", padx=12, pady=10)
        top.pack(fill="x")

        lr = (self.state.get("rss_last_run") or {}) if isinstance(self.state.get("rss_last_run"), dict) else {}
        lr_ok = lr.get("ok")
        lr_time = _safe_str(lr.get("ended")).strip()
        lr_new = lr.get("new_inserted")

        if lr_ok is True:
            run_text = f"Last RSS run: OK  •  {lr_time or 'unknown'}"
        elif lr_ok is False:
            run_text = f"Last RSS run: FAILED  •  {lr_time or 'unknown'}"
        else:
            run_text = "Last RSS run: unknown"

        hdr_left = tk.Frame(top, bg="#111111")
        hdr_left.pack(side="left", fill="x", expand=True)

        tk.Label(
            hdr_left,
            text="RSS Review / Import  [obsidian-flag-export]",
            bg="#111111",
            fg="#ffffff",
            font=("Segoe UI", 14, "bold"),
        ).pack(anchor="w")

        self._rss_counts_lbl = tk.Label(
            hdr_left,
            text="",
            bg="#111111",
            fg="#cfcfcf",
            font=("Segoe UI", 10, "normal"),
        )
        self._rss_counts_lbl.pack(anchor="w", pady=(2, 0))

        # Auto-flag diagnostics (watchlists)
        try:
            stats = self.state.get("rss_autoflag_stats", {}) if isinstance(self.state.get("rss_autoflag_stats", {}), dict) else {}
            kwc = int(stats.get("keywords_loaded", 0) or 0)
            afc = int(stats.get("auto_flagged", 0) or 0)
            wdir_res = str(stats.get("watchlists_dir_resolved", "") or "")
            wdir_sel = str(stats.get("watchlists_dir_selected", "") or "")
            vroot = str(stats.get("vault_root_token", "") or "")
            scanned = int(stats.get("watchlist_files_scanned", 0) or 0)
            with_kw = int(stats.get("watchlist_files_with_keywords", 0) or 0)
            sample = stats.get("watchlist_files_sample", [])
            sample_txt = ", ".join(sample) if isinstance(sample, list) else ""

            diag = f"Auto-flag: {afc} matched • Watchlist keywords loaded: {kwc}"
            tk.Label(
                hdr_left,
                text=diag,
                bg="#111111",
                fg="#888888",
                font=("Segoe UI", 9, "normal"),
            ).pack(anchor="w", pady=(0, 0))

            tk.Label(
                hdr_left,
                text=f"Watchlists dir: {wdir_res or '(not resolved)'}",
                bg="#111111",
                fg="#777777",
                font=("Segoe UI", 9, "normal"),
            ).pack(anchor="w", pady=(0, 0))

            tk.Label(
                hdr_left,
                text=f"Watchlist files scanned: {scanned} • with keywords: {with_kw}" + (f" • sample: {sample_txt}" if sample_txt else ""),
                bg="#111111",
                fg="#666666",
                font=("Segoe UI", 9, "normal"),
            ).pack(anchor="w", pady=(0, 0))

            if wdir_sel:
                tk.Label(
                    hdr_left,
                    text=f"Watchlists dir expr: {wdir_sel}",
                    bg="#111111",
                    fg="#666666",
                    font=("Segoe UI", 9, "normal"),
                ).pack(anchor="w", pady=(0, 0))
            if vroot:
                tk.Label(
                    hdr_left,
                    text=f"VAULT_ROOT: {vroot}",
                    bg="#111111",
                    fg="#666666",
                    font=("Segoe UI", 9, "normal"),
                ).pack(anchor="w", pady=(0, 2))
        except Exception:
            pass

        tk.Label(
            hdr_left,
            text=run_text + (f"  •  New inserted: {lr_new}" if lr_new is not None else ""),
            bg="#111111",
            fg="#9aa7c0",
            font=("Segoe UI", 9, "normal"),
        ).pack(anchor="w", pady=(2, 0))

        mode_bar = tk.Frame(top, bg="#111111")
        mode_bar.pack(side="right")

        def mk_mode_btn(label, val):
            b = tk.Radiobutton(
                mode_bar,
                text=label,
                value=val,
                variable=self._rss_filter_mode,
                indicatoron=False,
                relief="flat",
                bg="#1a1a1a",
                fg="#ffffff",
                selectcolor="#2a2a2a",
                activebackground="#2a2a2a",
                activeforeground="#ffffff",
                font=("Segoe UI", 10, "bold"),
                padx=10,
                pady=6,
                command=self._rss_apply_filters,
            )
            b.pack(side="left", padx=4)
            return b

        mk_mode_btn("New", "new")
        mk_mode_btn("Flagged", "flagged")
        mk_mode_btn("Imported", "imported")
        mk_mode_btn("Reported", "reported")
        mk_mode_btn("All", "all")

        # Search + source filter row
        filt = tk.Frame(parent, bg="#111111", padx=12, pady=8)
        filt.pack(fill="x")

        tk.Label(filt, text="Search:", bg="#111111", fg="#cfcfcf", font=("Segoe UI", 10, "normal")).pack(side="left")
        search = tk.Entry(
            filt,
            textvariable=self._rss_search_var,
            bg="#1e1e1e",
            fg="#ffffff",
            insertbackground="#ffffff",
            relief="flat",
            font=("Segoe UI", 11, "normal"),
            width=40,
        )
        search.pack(side="left", padx=(8, 12))
        search.bind("<KeyRelease>", lambda e: self._rss_debounced_filter())
        self._rss_search_entry = search

        sources = sorted({(_safe_str(i.get("source")).strip() or "Unknown") for i in self._rss_items_all})
        sources = ["All Sources"] + sources

        tk.Label(filt, text="Source:", bg="#111111", fg="#cfcfcf", font=("Segoe UI", 10, "normal")).pack(side="left", padx=(10, 6))
        src_menu = tk.OptionMenu(filt, self._rss_source_var, *sources, command=lambda _v: self._rss_apply_filters())
        src_menu.configure(bg="#1a1a1a", fg="#ffffff", activebackground="#2a2a2a", activeforeground="#ffffff", relief="flat")
        src_menu["menu"].configure(bg="#1a1a1a", fg="#ffffff")
        src_menu.pack(side="left")

        btns = tk.Frame(filt, bg="#111111")
        btns.pack(side="right")

        tk.Button(
            btns,
            text="Refresh",
            command=self._rss_refresh,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(side="left", padx=4)

        tk.Button(
            btns,
            text="Mark All Read",
            command=self._rss_mark_all_as_read,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(side="left", padx=4)

        tk.Button(
            btns,
            text="Export Selection (JSON)",
            command=self._rss_export_selection_json,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(side="left", padx=4)

        # main split pane
        body = tk.PanedWindow(parent, orient="horizontal", bg="#111111", sashwidth=6, sashrelief="flat")
        body.pack(fill="both", expand=True, padx=12, pady=(0, 12))

        left = tk.Frame(body, bg="#111111")
        right = tk.Frame(body, bg="#111111")
        body.add(left, minsize=420)
        body.add(right, minsize=420)

        # Treeview list (left)
        cols = ("status", "title", "source", "published")
        self._rss_tree = ttk.Treeview(left, columns=cols, show="headings", selectmode="extended")
        self._rss_tree.heading("status", text="Status")
        self._rss_tree.heading("title", text="Title")
        self._rss_tree.heading("source", text="Source")
        self._rss_tree.heading("published", text="Published")

        self._rss_tree.column("status", width=110, stretch=False)
        self._rss_tree.column("title", width=520, stretch=True)
        self._rss_tree.column("source", width=160, stretch=False)
        self._rss_tree.column("published", width=150, stretch=False)

        ysb = ttk.Scrollbar(left, orient="vertical", command=self._rss_tree.yview)
        self._rss_tree.configure(yscrollcommand=ysb.set)
        self._rss_tree.pack(side="left", fill="both", expand=True)
        ysb.pack(side="right", fill="y")

        self._rss_tree.bind("<<TreeviewSelect>>", lambda e: self._rss_render_preview())

        # Preview (right)
        pr_top = tk.Frame(right, bg="#111111")
        pr_top.pack(fill="x")

        self._rss_prev_title = tk.Label(
            pr_top, text="", bg="#111111", fg="#ffffff", font=("Segoe UI", 13, "bold"), wraplength=520, justify="left"
        )
        self._rss_prev_title.pack(anchor="w")

        self._rss_prev_meta = tk.Label(
            pr_top, text="", bg="#111111", fg="#9aa7c0", font=("Segoe UI", 9, "normal")
        )
        self._rss_prev_meta.pack(anchor="w", pady=(2, 0))

        # Action bar
        act = tk.Frame(right, bg="#111111", pady=8)
        act.pack(fill="x")

        def mk_act_btn(text, cmd):
            b = tk.Button(
                act,
                text=text,
                command=cmd,
                bg="#1a1a1a",
                fg="#ffffff",
                relief="flat",
                padx=10,
                pady=6,
                font=("Segoe UI", 10, "bold"),
            )
            b.pack(side="left", padx=4)
            return b

        mk_act_btn("Open URL (o)", self._rss_open_selected_url)
        self._rss_flag_btn = mk_act_btn("Flag (f)", self._rss_toggle_flag_selected)
        mk_act_btn("Mark Read (m)", self._rss_mark_read_selected)
        self._rss_report_btn = mk_act_btn("Report (r)", self._rss_toggle_report_selected)
        mk_act_btn("Copy Markdown", self._rss_copy_markdown_selected)

        # Secondary action row (ensures visibility on narrow windows)
        act2 = tk.Frame(right, bg="#111111", pady=2)
        act2.pack(fill="x")
        tk.Button(
            act2,
            text="Generate Report (g)",
            command=self._rss_generate_report,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(side="left", padx=4)
        tk.Button(
            act2,
            text="Weekly Rollup (w)",
            command=self._rss_generate_weekly_rollup,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(side="left", padx=4)

        # Notes + tags
        nt = tk.Frame(right, bg="#111111")
        nt.pack(fill="x", pady=(4, 8))

        tk.Label(nt, text="Tags (comma):", bg="#111111", fg="#cfcfcf", font=("Segoe UI", 10, "normal")).pack(anchor="w")
        self._rss_tags_var = tk.StringVar(value="")
        tags_entry = tk.Entry(
            nt,
            textvariable=self._rss_tags_var,
            bg="#1e1e1e",
            fg="#ffffff",
            insertbackground="#ffffff",
            relief="flat",
            font=("Segoe UI", 11, "normal"),
        )
        tags_entry.pack(fill="x", pady=(4, 8))

        tk.Label(nt, text="Analyst Notes:", bg="#111111", fg="#cfcfcf", font=("Segoe UI", 10, "normal")).pack(anchor="w")
        self._rss_notes = tk.Text(nt, height=6, bg="#1e1e1e", fg="#ffffff", insertbackground="#ffffff", relief="flat")
        self._rss_notes.pack(fill="x", pady=(4, 0))

        tk.Button(
            nt,
            text="Save Notes/Tags",
            command=self._rss_save_notes_tags,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(anchor="e", pady=(8, 0))

        # Preview body
        pv = tk.Frame(right, bg="#111111")
        pv.pack(fill="both", expand=True)

        self._rss_prev_text = tk.Text(pv, bg="#111111", fg="#e6e6e6", insertbackground="#ffffff", relief="flat", wrap="word")
        pv_ysb = ttk.Scrollbar(pv, orient="vertical", command=self._rss_prev_text.yview)
        self._rss_prev_text.configure(yscrollcommand=pv_ysb.set)
        self._rss_prev_text.pack(side="left", fill="both", expand=True)
        pv_ysb.pack(side="right", fill="y")
        self._rss_prev_text.configure(state="disabled")

        # Bottom status
        bottom = tk.Frame(parent, bg="#111111", padx=12, pady=6)
        bottom.pack(fill="x")
        tk.Label(bottom, textvariable=self._rss_status_var, bg="#111111", fg="#cfcfcf", font=("Segoe UI", 9, "normal")).pack(anchor="w")

        # Pre-compute watchlist stats so header diagnostics are available immediately
        try:
            kws = self._rss_watchlist_keywords()
            stats = self.state.get("rss_autoflag_stats", {}) if isinstance(self.state.get("rss_autoflag_stats", {}), dict) else {}
            stats["keywords_loaded"] = int(len(kws))
            self.state["rss_autoflag_stats"] = stats
            try:
                save_state(self.state_path, self.state)
            except Exception:
                pass
        except Exception:
            pass

        if hasattr(self, "_rss_review_win") and self._rss_review_win and self._rss_review_win.winfo_exists():
            try:
                self._rss_review_win.lift()
                self._rss_review_win.focus_force()
            except Exception:
                pass
            return

        win = tk.Toplevel(self)
        self._rss_review_win = win
        win.title("RSS Review / Import")
        win.configure(bg="#111111")
        win.geometry("1200x720")
        win.minsize(980, 600)
        
                # ----------------------------
        # Run Status variables (MUST be defined before any widgets reference them)
        # ----------------------------
        self._rss_sched_vars = getattr(self, "_rss_sched_vars", {})

        def _get_var(key: str, default: str) -> tk.StringVar:
            v = self._rss_sched_vars.get(key)
            if isinstance(v, tk.StringVar):
                return v
            v = tk.StringVar(value=default)
            self._rss_sched_vars[key] = v
            return v

        last_run_var = _get_var("last_run", "Last run: (none yet)")
        last_detail_var = _get_var("last_detail", "")
        next_run_var = _get_var("next_run", "Next run: (pending)")
        countdown_var = _get_var("countdown", "Time remaining: --:--")

        try:
            self.update_idletasks()
            w, h = 1200, 720
            x = self.winfo_rootx() + max((self.winfo_width() - w) // 2, 0)
            y = self.winfo_rooty() + max((self.winfo_height() - h) // 2, 0)
            win.geometry(f"{w}x{h}+{x}+{y}")
        except Exception:
            pass

        win.protocol("WM_DELETE_WINDOW", self._rss_review_close)

        # stateful vars
        self._rss_items_all = self._rss_load_items()
        self._rss_filter_mode = tk.StringVar(value="new")   # new|flagged|imported|reported|all
        self._rss_search_var = tk.StringVar(value="")
        self._rss_source_var = tk.StringVar(value="All Sources")
        self._rss_status_var = tk.StringVar(value="")
        self._rss_debounce_job = None

        # header / toolbar
        top = tk.Frame(win, bg="#111111", padx=12, pady=10)
        top.pack(fill="x")

        lr = (self.state.get("rss_last_run") or {}) if isinstance(self.state.get("rss_last_run"), dict) else {}
        lr_ok = lr.get("ok")
        lr_time = _safe_str(lr.get("ended")).strip()
        lr_new = lr.get("new_inserted")

        if lr_ok is True:
            run_text = f"Last RSS run: OK  •  {lr_time or 'unknown'}"
        elif lr_ok is False:
            run_text = f"Last RSS run: FAILED  •  {lr_time or 'unknown'}"
        else:
            run_text = "Last RSS run: unknown"

        hdr_left = tk.Frame(top, bg="#111111")
        hdr_left.pack(side="left", fill="x", expand=True)

        tk.Label(
            hdr_left,
            text="RSS Review / Import  [obsidian-flag-export]",
            bg="#111111",
            fg="#ffffff",
            font=("Segoe UI", 14, "bold"),
        ).pack(anchor="w")

        self._rss_counts_lbl = tk.Label(
            hdr_left,
            text="",
            bg="#111111",
            fg="#cfcfcf",
            font=("Segoe UI", 10, "normal"),
        )
        self._rss_counts_lbl.pack(anchor="w", pady=(2, 0))

        # Auto-flag diagnostics (watchlists)
        try:
            stats = self.state.get("rss_autoflag_stats", {}) if isinstance(self.state.get("rss_autoflag_stats", {}), dict) else {}
            kwc = int(stats.get("keywords_loaded", 0) or 0)
            afc = int(stats.get("auto_flagged", 0) or 0)
            wdir_res = str(stats.get("watchlists_dir_resolved", "") or "")
            wdir_sel = str(stats.get("watchlists_dir_selected", "") or "")
            vroot = str(stats.get("vault_root_token", "") or "")
            scanned = int(stats.get("watchlist_files_scanned", 0) or 0)
            with_kw = int(stats.get("watchlist_files_with_keywords", 0) or 0)
            sample = stats.get("watchlist_files_sample", [])
            sample_txt = ", ".join(sample) if isinstance(sample, list) else ""

            diag = f"Auto-flag: {afc} matched • Watchlist keywords loaded: {kwc}"
            tk.Label(
                hdr_left,
                text=diag,
                bg="#111111",
                fg="#888888",
                font=("Segoe UI", 9, "normal"),
            ).pack(anchor="w", pady=(0, 0))

            tk.Label(
                hdr_left,
                text=f"Watchlists dir: {wdir_res or '(not resolved)'}",
                bg="#111111",
                fg="#777777",
                font=("Segoe UI", 9, "normal"),
            ).pack(anchor="w", pady=(0, 0))

            tk.Label(
                hdr_left,
                text=f"Watchlist files scanned: {scanned} • with keywords: {with_kw}" + (f" • sample: {sample_txt}" if sample_txt else ""),
                bg="#111111",
                fg="#666666",
                font=("Segoe UI", 9, "normal"),
            ).pack(anchor="w", pady=(0, 0))

            if wdir_sel:
                tk.Label(
                    hdr_left,
                    text=f"Watchlists dir expr: {wdir_sel}",
                    bg="#111111",
                    fg="#666666",
                    font=("Segoe UI", 9, "normal"),
                ).pack(anchor="w", pady=(0, 0))
            if vroot:
                tk.Label(
                    hdr_left,
                    text=f"VAULT_ROOT: {vroot}",
                    bg="#111111",
                    fg="#666666",
                    font=("Segoe UI", 9, "normal"),
                ).pack(anchor="w", pady=(0, 2))
        except Exception:
            pass


        tk.Label(
            hdr_left,
            text=run_text + (f"  •  New inserted: {lr_new}" if lr_new is not None else ""),
            bg="#111111",
            fg="#9aa7c0",
            font=("Segoe UI", 9, "normal"),
        ).pack(anchor="w", pady=(2, 0))

        mode_bar = tk.Frame(top, bg="#111111")
        mode_bar.pack(side="right")

        def mk_mode_btn(label, val):
            b = tk.Radiobutton(
                mode_bar,
                text=label,
                value=val,
                variable=self._rss_filter_mode,
                indicatoron=False,
                relief="flat",
                bg="#1a1a1a",
                fg="#ffffff",
                selectcolor="#2a2a2a",
                activebackground="#2a2a2a",
                activeforeground="#ffffff",
                font=("Segoe UI", 10, "bold"),
                padx=10,
                pady=6,
                command=self._rss_apply_filters,
            )
            b.pack(side="left", padx=4)
            return b

        mk_mode_btn("New", "new")
        mk_mode_btn("Flagged", "flagged")
        mk_mode_btn("Imported", "imported")
        mk_mode_btn("Reported", "reported")
        mk_mode_btn("All", "all")

        # Search + source filter row
        filt = tk.Frame(win, bg="#111111", padx=12, pady=8)
        filt.pack(fill="x")

        tk.Label(filt, text="Search:", bg="#111111", fg="#cfcfcf", font=("Segoe UI", 10, "normal")).pack(side="left")
        search = tk.Entry(
            filt,
            textvariable=self._rss_search_var,
            bg="#1e1e1e",
            fg="#ffffff",
            insertbackground="#ffffff",
            relief="flat",
            font=("Segoe UI", 11, "normal"),
            width=40,
        )
        search.pack(side="left", padx=(8, 12))
        search.bind("<KeyRelease>", lambda e: self._rss_debounced_filter())

        sources = sorted({(_safe_str(i.get("source")).strip() or "Unknown") for i in self._rss_items_all})
        sources = ["All Sources"] + sources

        tk.Label(filt, text="Source:", bg="#111111", fg="#cfcfcf", font=("Segoe UI", 10, "normal")).pack(side="left", padx=(10, 6))
        src_menu = tk.OptionMenu(filt, self._rss_source_var, *sources, command=lambda _v: self._rss_apply_filters())
        src_menu.configure(bg="#1a1a1a", fg="#ffffff", activebackground="#2a2a2a", activeforeground="#ffffff", relief="flat")
        src_menu["menu"].configure(bg="#1a1a1a", fg="#ffffff")
        src_menu.pack(side="left")

        btns = tk.Frame(filt, bg="#111111")
        btns.pack(side="right")

        tk.Button(
            btns,
            text="Refresh",
            command=self._rss_refresh,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(side="left", padx=4)

        tk.Button(
            btns,
            text="Mark All Read",
            command=self._rss_mark_all_as_read,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(side="left", padx=4)

        tk.Button(
            btns,
            text="Export Selection (JSON)",
            command=self._rss_export_selection_json,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(side="left", padx=4)

        # main split pane
        body = tk.PanedWindow(win, orient="horizontal", bg="#111111", sashwidth=6, sashrelief="flat")
        body.pack(fill="both", expand=True, padx=12, pady=(0, 12))

        left = tk.Frame(body, bg="#111111")
        right = tk.Frame(body, bg="#111111")
        body.add(left, minsize=420)
        body.add(right, minsize=420)

        # Treeview list (left)
        cols = ("status", "title", "source", "published")
        self._rss_tree = ttk.Treeview(left, columns=cols, show="headings", selectmode="extended")
        self._rss_tree.heading("status", text="Status")
        self._rss_tree.heading("title", text="Title")
        self._rss_tree.heading("source", text="Source")
        self._rss_tree.heading("published", text="Published")

        self._rss_tree.column("status", width=110, stretch=False)
        self._rss_tree.column("title", width=520, stretch=True)
        self._rss_tree.column("source", width=160, stretch=False)
        self._rss_tree.column("published", width=150, stretch=False)

        ysb = ttk.Scrollbar(left, orient="vertical", command=self._rss_tree.yview)
        self._rss_tree.configure(yscrollcommand=ysb.set)
        self._rss_tree.pack(side="left", fill="both", expand=True)
        ysb.pack(side="right", fill="y")

        self._rss_tree.bind("<<TreeviewSelect>>", lambda e: self._rss_render_preview())

        # Preview (right)
        pr_top = tk.Frame(right, bg="#111111")
        pr_top.pack(fill="x")

        self._rss_prev_title = tk.Label(
            pr_top, text="", bg="#111111", fg="#ffffff", font=("Segoe UI", 13, "bold"), wraplength=520, justify="left"
        )
        self._rss_prev_title.pack(anchor="w")

        self._rss_prev_meta = tk.Label(
            pr_top, text="", bg="#111111", fg="#9aa7c0", font=("Segoe UI", 9, "normal")
        )
        self._rss_prev_meta.pack(anchor="w", pady=(2, 0))

        # Action bar
        act = tk.Frame(right, bg="#111111", pady=8)
        act.pack(fill="x")

        def mk_act_btn(text, cmd):
            b = tk.Button(
                act,
                text=text,
                command=cmd,
                bg="#1a1a1a",
                fg="#ffffff",
                relief="flat",
                padx=10,
                pady=6,
                font=("Segoe UI", 10, "bold"),
            )
            b.pack(side="left", padx=4)
            return b

        mk_act_btn("Open URL (o)", self._rss_open_selected_url)
        self._rss_flag_btn = mk_act_btn("Flag (f)", self._rss_toggle_flag_selected)
        mk_act_btn("Mark Read (m)", self._rss_mark_read_selected)
        self._rss_report_btn = mk_act_btn("Report (r)", self._rss_toggle_report_selected)
        mk_act_btn("Copy Markdown", self._rss_copy_markdown_selected)

        # Secondary action row (ensures visibility on narrow windows)
        act2 = tk.Frame(right, bg="#111111", pady=2)
        act2.pack(fill="x")
        tk.Button(
            act2,
            text="Generate Report (g)",
            command=self._rss_generate_report,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(side="left", padx=4)
        tk.Button(
            act2,
            text="Weekly Rollup (w)",
            command=self._rss_generate_weekly_rollup,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(side="left", padx=4)

        # Notes + tags
        nt = tk.Frame(right, bg="#111111")
        nt.pack(fill="x", pady=(4, 8))

        tk.Label(nt, text="Tags (comma):", bg="#111111", fg="#cfcfcf", font=("Segoe UI", 10, "normal")).pack(anchor="w")
        self._rss_tags_var = tk.StringVar(value="")
        tags_entry = tk.Entry(
            nt,
            textvariable=self._rss_tags_var,
            bg="#1e1e1e",
            fg="#ffffff",
            insertbackground="#ffffff",
            relief="flat",
            font=("Segoe UI", 11, "normal"),
        )
        tags_entry.pack(fill="x", pady=(4, 8))

        tk.Label(nt, text="Analyst Notes:", bg="#111111", fg="#cfcfcf", font=("Segoe UI", 10, "normal")).pack(anchor="w")
        self._rss_notes = tk.Text(nt, height=6, bg="#1e1e1e", fg="#ffffff", insertbackground="#ffffff", relief="flat")
        self._rss_notes.pack(fill="x", pady=(4, 0))

        tk.Button(
            nt,
            text="Save Notes/Tags",
            command=self._rss_save_notes_tags,
            bg="#1a1a1a",
            fg="#ffffff",
            relief="flat",
            padx=10,
            pady=6,
            font=("Segoe UI", 10, "bold"),
        ).pack(anchor="e", pady=(8, 0))

        # Preview body
        pv = tk.Frame(right, bg="#111111")
        pv.pack(fill="both", expand=True)

        self._rss_prev_text = tk.Text(pv, bg="#111111", fg="#e6e6e6", insertbackground="#ffffff", relief="flat", wrap="word")
        pv_ysb = ttk.Scrollbar(pv, orient="vertical", command=self._rss_prev_text.yview)
        self._rss_prev_text.configure(yscrollcommand=pv_ysb.set)
        self._rss_prev_text.pack(side="left", fill="both", expand=True)
        pv_ysb.pack(side="right", fill="y")
        self._rss_prev_text.configure(state="disabled")

        # Bottom status
        bottom = tk.Frame(win, bg="#111111", padx=12, pady=6)
        bottom.pack(fill="x")
        tk.Label(bottom, textvariable=self._rss_status_var, bg="#111111", fg="#cfcfcf", font=("Segoe UI", 9, "normal")).pack(anchor="w")

        # Keyboard shortcuts
        win.bind("/", lambda e: (search.focus_set(), "break"))
        win.bind("<Escape>", lambda e: self._rss_clear_search())
        win.bind("j", lambda e: self._rss_move_selection(1))
        win.bind("k", lambda e: self._rss_move_selection(-1))
        win.bind("o", lambda e: self._rss_open_selected_url())
        win.bind("f", lambda e: self._rss_toggle_flag_selected())
        win.bind("m", lambda e: self._rss_mark_read_selected())
        win.bind("r", lambda e: self._rss_toggle_report_selected())
        win.bind("g", lambda e: self._rss_generate_report())
        win.bind("w", lambda e: self._rss_generate_weekly_rollup())
        win.bind("a", lambda e: self._rss_mark_all_as_read())

        self._rss_apply_filters()
        self._rss_status_var.set(
            "Shortcuts: j/k navigate • f flag • m mark read • r report • g generate • w weekly • o open • / search • Esc clear"
        )

    def _rss_review_close(self):
        try:
            win = getattr(self, "_rss_win", None) or getattr(self, "_rss_review_win", None)
            if win and win.winfo_exists():
                win.destroy()
        except Exception:
            pass
        self._rss_review_win = None
        self._rss_win = None
        self._rss_notebook = None
        self._rss_tab_review = None
        self._rss_tab_schedule = None
        self._rss_tab_keywords = None

    def _rss_refresh(self):
        self._rss_items_cache = []  # force reload
        self._rss_items_all = self._rss_load_items()
        self._rss_apply_filters()
        self._refresh_rss_cache_and_badges()

    def _rss_debounced_filter(self):
        if self._rss_debounce_job is not None:
            try:
                self.after_cancel(self._rss_debounce_job)
            except Exception:
                pass
        self._rss_debounce_job = self.after(160, self._rss_apply_filters)

    def _rss_clear_search(self):
        try:
            self._rss_search_var.set("")
            self._rss_apply_filters()
        except Exception:
            pass

    def _rss_apply_filters(self):
        self._rss_debounce_job = None

        mode = (self._rss_filter_mode.get() or "new").strip().lower()
        q = (_safe_str(self._rss_search_var.get()).strip().lower())
        src = (_safe_str(self._rss_source_var.get()).strip())

        items = self._rss_items_all or []
        watch_items = []
        try:
            watch_items = self._rss_watchlist_items()
        except Exception:
            watch_items = []
        out = []

        for it in items:
            fp = _rss_fp(it)
            st = self._rss_state_get(fp)

            if mode == "new" and st.get("read"):
                continue
            if mode == "flagged" and not st.get("flagged"):
                continue
            if mode == "imported" and not st.get("imported"):
                continue
            if mode == "reported" and not _safe_str(st.get("reported_on")).strip():
                continue

            it_src = (_safe_str(it.get("source")).strip() or "Unknown")
            if src != "All Sources" and it_src != src:
                continue

            if q:
                hay = " ".join([
                    _safe_str(it.get("title")),
                    _safe_str(it.get("source")),
                    _safe_str(it.get("summary")),
                    _safe_str(it.get("content")),
                    _safe_str(it.get("url")),
                ]).lower()
                if q not in hay:
                    continue

            score, hits, cats = self._rss_compute_watchlist_hits(it, watch_items)
            it["_watchlist_score"] = score
            it["_watchlist_hits"] = hits
            it["_watchlist_categories"] = cats
            out.append(it)

        out.sort(
            key=lambda x: (
                int(x.get("_watchlist_score") or 0),
                _safe_str(x.get("published")).strip(),
            ),
            reverse=True,
        )
        self._rss_items_filtered = out
        self._rss_update_watchlist_freq(self._rss_items_all or [])

        c = self._rss_counts(items)
        try:
            self._rss_counts_lbl.configure(
                text=(
                    f"New: {c['new']}  •  Flagged: {c['flagged']}  •  Imported: {c['imported']}  •  "
                    f"Reported: {c['reported']}  •  Total: {c['total']}"
                )
            )
        except Exception:
            pass

        tv = self._rss_tree
        tv.delete(*tv.get_children())

        for idx, it in enumerate(out):
            fp = _rss_fp(it)
            st = self._rss_state_get(fp)

            flags = []
            if not st.get("read"):
                flags.append("N")
            if st.get("flagged"):
                flags.append("F")
            if st.get("queued_for_report") or _safe_str(st.get("reported_on")).strip():
                flags.append("R")
            status = "".join([f"[{x}]" for x in flags]) if flags else ""

            title = _safe_str(it.get("title")).strip()
            source = _safe_str(it.get("source")).strip() or "Unknown"
            published = _fmt_when(it.get("published"))

            tv.insert("", "end", iid=str(idx), values=(status, title, source, published))

        kids = tv.get_children()
        if kids:
            tv.selection_set(kids[0])
            tv.focus(kids[0])
            tv.see(kids[0])
            self._rss_render_preview()
        else:
            self._rss_prev_title.configure(text="")
            self._rss_prev_meta.configure(text="")
            self._rss_tags_var.set("")
            self._rss_notes.delete("1.0", "end")
            self._rss_prev_text.configure(state="normal")
            self._rss_prev_text.delete("1.0", "end")
            self._rss_prev_text.insert("1.0", "No items match the current filters.")
            self._rss_prev_text.configure(state="disabled")
            try:
                if getattr(self, "_rss_flag_btn", None):
                    self._rss_flag_btn.configure(text="Flag (f)")
                if getattr(self, "_rss_report_btn", None):
                    self._rss_report_btn.configure(text="Report (r)")
            except Exception:
                pass

    def _rss_get_selected_items(self):
        tv = self._rss_tree
        sel = list(tv.selection() or [])
        out = []
        for iid in sel:
            try:
                i = int(iid)
                if 0 <= i < len(self._rss_items_filtered):
                    out.append(self._rss_items_filtered[i])
            except Exception:
                continue
        return out

    def _rss_get_primary_item(self):
        items = self._rss_get_selected_items()
        return items[0] if items else None

    def _rss_render_preview(self):
        it = self._rss_get_primary_item()
        if not it:
            return

        fp = _rss_fp(it)
        st = self._rss_state_get(fp)

        title = _safe_str(it.get("title")).strip() or "Untitled"
        url = _safe_str(it.get("url")).strip()
        source = _safe_str(it.get("source")).strip() or "Unknown"
        published = _fmt_when(it.get("published"))

        self._rss_prev_title.configure(text=title)
        meta = f"{source}  •  {published}" + (f"  •  {url}" if url else "")
        wl_score = int(it.get("_watchlist_score") or 0)
        wl_hits = it.get("_watchlist_hits") or []
        wl_cats = it.get("_watchlist_categories") or []
        if wl_score > 0:
            hit_txt = ", ".join(wl_hits[:6]) + ("…" if len(wl_hits) > 6 else "")
            cat_txt = ", ".join(wl_cats[:6]) + ("…" if len(wl_cats) > 6 else "")
            meta += f"  •  WL score: {wl_score}"
            if hit_txt:
                meta += f"  •  Hits: {hit_txt}"
            if cat_txt:
                meta += f"  •  Categories: {cat_txt}"
        self._rss_prev_meta.configure(text=meta)

        tags = st.get("tags") if isinstance(st.get("tags"), list) else []
        self._rss_tags_var.set(", ".join([_safe_str(t).strip() for t in tags if _safe_str(t).strip()]))

        self._rss_notes.delete("1.0", "end")
        self._rss_notes.insert("1.0", _safe_str(st.get("notes")))

        summary = _safe_str(it.get("summary")).strip()
        content = _safe_str(it.get("content")).strip()
        if not content:
            content = summary
        if not content:
            content = "(No content/summary available for this item.)"

        self._rss_prev_text.configure(state="normal")
        self._rss_prev_text.delete("1.0", "end")
        self._rss_prev_text.insert("1.0", content)
        self._rss_prev_text.configure(state="disabled")

        # Toggle action button labels based on status
        try:
            if getattr(self, "_rss_flag_btn", None):
                self._rss_flag_btn.configure(text=("Unflag (f)" if st.get("flagged") else "Flag (f)"))
            if getattr(self, "_rss_report_btn", None):
                has_report = bool(st.get("queued_for_report") or _safe_str(st.get("reported_on")).strip())
                self._rss_report_btn.configure(text=("Unreport (r)" if has_report else "Report (r)"))
        except Exception:
            pass

    def _rss_save_notes_tags(self):
        it = self._rss_get_primary_item()
        if not it:
            return

        fp = _rss_fp(it)
        tags_raw = _safe_str(self._rss_tags_var.get()).strip()
        tags = [t.strip() for t in tags_raw.split(",") if t.strip()] if tags_raw else []

        notes = _safe_str(self._rss_notes.get("1.0", "end")).rstrip("\n")
        self._rss_state_set(fp, tags=tags, notes=notes)
        self._rss_status_var.set("Saved notes/tags.")

    def _rss_open_selected_url(self):
        it = self._rss_get_primary_item()
        if not it:
            return
        url = _safe_str(it.get("url")).strip()
        if not url:
            messagebox.showerror("Open URL", "No URL available for this item.")
            return
        open_url(url)

    def _rss_toggle_flag_selected(self):
        items = self._rss_get_selected_items()
        if not items:
            return
        primary_fp = _rss_fp(items[0])
        to_flag = not bool(self._rss_state_get(primary_fp).get("flagged"))

        for it in items:
            fp = _rss_fp(it)
            self._rss_state_set(fp, flagged=to_flag)

        self._rss_apply_filters()
        self._rss_status_var.set(f"{'Flagged' if to_flag else 'Unflagged'} {len(items)} item(s).")

    def _rss_mark_read_selected(self):
        items = self._rss_get_selected_items()
        if not items:
            return
        for it in items:
            fp = _rss_fp(it)
            self._rss_state_set(fp, read=True)
        self._rss_apply_filters()
        self._rss_status_var.set(f"Marked read: {len(items)} item(s).")
        self._refresh_rss_cache_and_badges()

    def _rss_mark_all_as_read(self):
        """
        Marks every unread RSS item as read (across all sources).
        This updates state in bulk for performance, refreshes the list, and updates the launcher RSS badge.
        """
        items = self._rss_items_all or []
        if not items:
            try:
                self._rss_status_var.set("No RSS items loaded.")
            except Exception:
                pass
            return

        bucket = self._rss_state_bucket()
        changed = 0

        for it in items:
            fp = _rss_fp(it)
            st = bucket.get(fp)
            if not isinstance(st, dict):
                st = {
                    "read": False,
                    "flagged": False,
                    "imported": False,
                    "tags": [],
                    "notes": "",
                    "last_seen": "",
                }
                bucket[fp] = st

            if not st.get("read"):
                st["read"] = True
                changed += 1

        if changed:
            save_state(self.state_path, self.state)

        self._rss_apply_filters()
        self._refresh_rss_cache_and_badges()

        try:
            self._rss_status_var.set(f"Marked all as read: {changed} item(s).")
        except Exception:
            pass

    def _rss_mark_imported_selected(self):
        items = self._rss_get_selected_items()
        if not items:
            return
        for it in items:
            fp = _rss_fp(it)
            self._rss_state_set(fp, imported=True, read=True)
        self._rss_apply_filters()
        self._rss_status_var.set(f"Marked imported: {len(items)} item(s).")
        self._refresh_rss_cache_and_badges()

    def _rss_report_queue(self) -> list[dict]:
        if "rss_report_queue" not in self.state or not isinstance(self.state.get("rss_report_queue"), list):
            self.state["rss_report_queue"] = []
        return self.state["rss_report_queue"]

    def _rss_unreport_selected(self):
        items = self._rss_get_selected_items()
        if not items:
            return

        queue = self._rss_report_queue()
        removed = 0
        for it in items:
            fp = _rss_fp(it)
            st = self._rss_state_get(fp)
            if not (st.get("queued_for_report") or _safe_str(st.get("reported_on")).strip()):
                continue

            url = _safe_str(it.get("url")).strip()
            queue = [q for q in queue if _safe_str(q.get("fp")).strip() != fp and _safe_str(q.get("url")).strip() != url]
            self._rss_state_set(
                fp,
                queued_for_report=False,
                queued_on="",
                reported_on="",
                reported_url="",
                printed=False,
            )
            removed += 1

        self.state["rss_report_queue"] = queue
        save_state(self.state_path, self.state)
        self._rss_apply_filters()
        self._refresh_rss_cache_and_badges()
        self._rss_status_var.set(f"Unreported: {removed} item(s).")

    def _rss_toggle_report_selected(self):
        it = self._rss_get_primary_item()
        if not it:
            return
        fp = _rss_fp(it)
        st = self._rss_state_get(fp)
        if st.get("queued_for_report") or _safe_str(st.get("reported_on")).strip():
            return self._rss_unreport_selected()
        return self._rss_report_selected()

    def _rss_report_selected(self):
        items = self._rss_get_selected_items()
        if not items:
            return

        queued = 0
        skipped_existing = 0
        errors = 0
        today = date.today().strftime("%Y-%m-%d")
        queue = self._rss_report_queue()
        queued_fps = {q.get("fp") for q in queue if isinstance(q, dict)}
        queued_urls = {q.get("url") for q in queue if isinstance(q, dict)}

        for it in items:
            fp = _rss_fp(it)

            url = _safe_str(it.get("url")).strip()
            if not url:
                errors += 1
                messagebox.showerror("Report", "Selected item has no URL to report.")
                continue
            st = self._rss_state_get(fp)
            if not st.get("flagged"):
                errors += 1
                messagebox.showerror("Report", "You must flag the item before reporting it.")
                continue
            if fp in queued_fps or url in queued_urls:
                skipped_existing += 1
                continue

            title = _safe_str(it.get("title")).strip() or "Untitled"
            source = _safe_str(it.get("source")).strip()
            published = _safe_str(it.get("published")).strip()

            queue.append(
                {
                    "fp": fp,
                    "title": title,
                    "url": url,
                    "source": source,
                    "published": published,
                    "queued_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                }
            )
            self._rss_state_set(
                fp,
                queued_for_report=True,
                queued_on=today,
                reported_on=today,
                reported_url=url,
                printed=False,
            )
            queued += 1

        self._rss_apply_filters()
        self._refresh_rss_cache_and_badges()

        save_state(self.state_path, self.state)

        msg = f"Queued for report: {queued} item(s)."
        if skipped_existing:
            msg += f" Skipped existing: {skipped_existing}."
        if errors:
            msg += f" Errors: {errors}."
        self._rss_status_var.set(msg)

    def _rss_resolve_vault_root(self) -> str:
        paths = self.cfg.get("paths", {}) or {}
        candidates = [
            paths.get("SCOUT_ROOT", ""),
            paths.get("VAULT_ROOT", ""),
            self.cfg.get("obsidian_vault_path", ""),
        ]
        for c in candidates:
            if not isinstance(c, str) or not c.strip():
                continue
            p = resolve_path(c, self.tokens)
            if p and os.path.isdir(p):
                return p
        return ""

    def _rss_write_report_note(self, entries: list[dict]) -> str:
        def fetch_clean_markdown(url: str) -> str:
            return self._rss_fetch_clean_markdown(url)

        scout_root = self._rss_resolve_vault_root()
        if not scout_root:
            raise ValueError("Obsidian vault root is not set. Configure paths.SCOUT_ROOT or paths.VAULT_ROOT.")

        report_rel = "90_Views/02_Reports"
        folder_abs = os.path.join(scout_root, *report_rel.split("/"))
        os.makedirs(folder_abs, exist_ok=True)

        today = date.today().strftime("%Y-%m-%d")
        fname = f"{today} - RSS Report.md"
        abs_path = os.path.join(folder_abs, fname)

        if os.path.exists(abs_path):
            stem = f"{today} - RSS Report"
            i = 2
            while True:
                cand = os.path.join(folder_abs, f"{stem} ({i}).md")
                if not os.path.exists(cand):
                    abs_path = cand
                    break
                i += 1

        frontmatter = [
            "---",
            'entity_type: "report"',
            'report_type: "rss"',
            f'report_date: "{today}"',
            f"report_count: {len(entries)}",
            'report_status: "final"',
            "tags: []",
            "---",
            "",
        ]

        body = [
            f"# RSS Report — {today}",
            "",
            "## Reported Articles",
            "",
        ]
        for e in entries:
            title = _safe_str(e.get("title")).strip() or "Untitled"
            url = _safe_str(e.get("url")).strip()
            source = _safe_str(e.get("source")).strip()
            published = _safe_str(e.get("published")).strip()
            try:
                content = fetch_clean_markdown(url) if url else ""
            except Exception as e:
                content = f"_Failed to fetch content: {e}_"

            body.append(f"### {title}")
            if url:
                body.append(f"- URL: {url}")
            if source:
                body.append(f"- Source: {source}")
            if published:
                body.append(f"- Published: {published}")
            body.append("")
            body.append(content or "_No content extracted._")
            body.append("")

        with open(abs_path, "w", encoding="utf-8") as f:
            f.write("\n".join(frontmatter + body + [""]))

        return abs_path

    def _rss_fetch_clean_markdown(self, url: str) -> str:
        if not requests or not Document or not BeautifulSoup:
            raise ValueError(
                "Missing dependencies for clean extraction. Install: pip install requests readability-lxml beautifulsoup4"
            )
        timeout = int((self.cfg.get("rss", {}) or {}).get("http_timeout", 20))
        headers = {
            "User-Agent": "SCOUT-RSS/1.0 (+local)",
            "Accept": "text/html,application/xhtml+xml",
        }
        r = requests.get(url, headers=headers, timeout=timeout)
        r.raise_for_status()
        doc = Document(r.text)
        html = doc.summary(html_partial=True)
        soup = BeautifulSoup(html, "html.parser")
        parts = []
        for node in soup.find_all(["p", "li"]):
            t = node.get_text(" ", strip=True)
            if t:
                parts.append(t)
        if not parts:
            text = soup.get_text("\n", strip=True)
        else:
            text = "\n\n".join(parts)
        max_chars = int((self.cfg.get("rss_review", {}) or {}).get("report_max_chars", 4000))
        if len(text) > max_chars:
            text = text[:max_chars].rstrip() + "…"
        return text

    def _rss_write_weekly_rollup_note(self, entries: list[dict], start_date: date, end_date: date) -> str:
        scout_root = self._rss_resolve_vault_root()
        if not scout_root:
            raise ValueError("Obsidian vault root is not set. Configure paths.SCOUT_ROOT or paths.VAULT_ROOT.")

        report_rel = "90_Views/02_Reports"
        folder_abs = os.path.join(scout_root, *report_rel.split("/"))
        os.makedirs(folder_abs, exist_ok=True)

        today = date.today().strftime("%Y-%m-%d")
        fname = f"{today} - Weekly RSS Rollup.md"
        abs_path = os.path.join(folder_abs, fname)

        if os.path.exists(abs_path):
            stem = f"{today} - Weekly RSS Rollup"
            i = 2
            while True:
                cand = os.path.join(folder_abs, f"{stem} ({i}).md")
                if not os.path.exists(cand):
                    abs_path = cand
                    break
                i += 1

        frontmatter = [
            "---",
            'entity_type: "report"',
            'report_type: "rss_weekly"',
            f'report_date: "{today}"',
            f'report_period_start: "{start_date.strftime("%Y-%m-%d")}"',
            f'report_period_end: "{end_date.strftime("%Y-%m-%d")}"',
            f"report_count: {len(entries)}",
            'report_status: "final"',
            "tags: []",
            "---",
            "",
        ]

        body = [
            f"# Weekly RSS Rollup — {start_date.strftime('%Y-%m-%d')} to {end_date.strftime('%Y-%m-%d')}",
            "",
            "## Reported Articles",
            "",
        ]

        def _sort_key(e: dict) -> str:
            published = _safe_str(e.get("published")).strip()
            if published:
                return published
            queued_at = _safe_str(e.get("queued_at")).strip()
            return queued_at

        for e in sorted(entries, key=_sort_key, reverse=True):
            title = _safe_str(e.get("title")).strip() or "Untitled"
            url = _safe_str(e.get("url")).strip()
            source = _safe_str(e.get("source")).strip() or "Unknown"
            published = _safe_str(e.get("published")).strip()
            pub_date = published[:10] if published else "Unknown"
            body.append(f"- {pub_date} | {title} | {source} | {url}")

        with open(abs_path, "w", encoding="utf-8") as f:
            f.write("\n".join(frontmatter + body + [""]))

        return abs_path

    def _rss_generate_report(self):
        queue = list(self._rss_report_queue())
        selected_items = self._rss_get_selected_items()
        if len(selected_items) <= 1:
            selected_items = []
        active = []
        seen_fps = set()
        seen_urls = set()
        selected_fps = set()
        errors = 0

        if selected_items:
            queue_by_fp = {str(q.get("fp")).strip(): q for q in queue if isinstance(q, dict) and q.get("fp")}
            queue_by_url = {str(q.get("url")).strip(): q for q in queue if isinstance(q, dict) and q.get("url")}

            for it in selected_items:
                fp = _rss_fp(it)
                st = self._rss_state_get(fp)

                url = _safe_str(it.get("url")).strip()
                if not url:
                    errors += 1
                    messagebox.showerror("Report", "Selected item has no URL to report.")
                    continue

                if fp and fp in seen_fps:
                    continue
                if url and url in seen_urls:
                    continue

                entry = queue_by_fp.get(fp) or (queue_by_url.get(url) if url else None)
                if not entry:
                    entry = {
                        "fp": fp,
                        "title": _safe_str(it.get("title")).strip() or "Untitled",
                        "url": url,
                        "source": _safe_str(it.get("source")).strip(),
                        "published": _safe_str(it.get("published")).strip(),
                        "queued_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    }
                active.append(entry)
                if fp:
                    seen_fps.add(fp)
                    selected_fps.add(fp)
                if url:
                    seen_urls.add(url)

            if not active:
                msg = "No selected items available for report."
                if errors:
                    msg += f" Skipped: errors={errors}."
                self._rss_status_var.set(msg)
                return
        else:
            if not queue:
                self._rss_status_var.set("Report queue is empty.")
                return

            # Drop any stale entries that are no longer queued
            stale = 0
            for e in queue:
                fp = _safe_str(e.get("fp")).strip()
                url = _safe_str(e.get("url")).strip()
                st = self._rss_state_get(fp) if fp else {}
                if not st.get("queued_for_report"):
                    stale += 1
                    continue
                if st.get("printed"):
                    stale += 1
                    continue
                if fp and fp in seen_fps:
                    stale += 1
                    continue
                if url and url in seen_urls:
                    stale += 1
                    continue
                if fp:
                    seen_fps.add(fp)
                if url:
                    seen_urls.add(url)
                active.append(e)

            if stale:
                self.state["rss_report_queue"] = active
                save_state(self.state_path, self.state)

            if not active:
                self._rss_status_var.set("Report queue is empty (stale entries cleared).")
                return

        open_after = bool((self.cfg.get("rss_review", {}) or {}).get("open_after_import", False))
        today = date.today().strftime("%Y-%m-%d")
        try:
            abs_path = self._rss_write_report_note(active)
        except Exception as e:
            messagebox.showerror("Generate Report", f"Failed to create report note.\n\n{e}")
            return

        for e in active:
            fp = _safe_str(e.get("fp")).strip()
            if not fp:
                continue
            self._rss_state_set(
                fp,
                read=True,
                queued_for_report=False,
                reported_on=today,
                reported_url=_safe_str(e.get("url")).strip(),
                printed=True,
            )

        if selected_items:
            if selected_fps:
                self.state["rss_report_queue"] = [
                    e for e in queue if _safe_str(e.get("fp")).strip() not in selected_fps
                ]
            else:
                self.state["rss_report_queue"] = queue
        else:
            self.state["rss_report_queue"] = []
        save_state(self.state_path, self.state)

        if open_after:
            try:
                _open_note_in_obsidian(self.cfg, self.tokens, abs_path)
            except Exception:
                pass

        self._rss_apply_filters()
        self._refresh_rss_cache_and_badges()
        if selected_items:
            msg = f"Generated report with {len(active)} item(s)."
            if errors:
                msg += f" Skipped: errors={errors}."
            self._rss_status_var.set(msg)
        else:
            self._rss_status_var.set(f"Generated report with {len(active)} item(s). Queue cleared.")
        try:
            messagebox.showinfo("Generate Report", f"Report created:\n{abs_path}")
        except Exception:
            pass

    def _rss_generate_weekly_rollup(self):
        queue = list(self._rss_report_queue())
        if not queue:
            self._rss_status_var.set("Report queue is empty.")
            return

        today = date.today()
        start_date = today - timedelta(days=6)
        end_date = today

        recent = []
        for e in queue:
            queued_at = _safe_str(e.get("queued_at")).strip()
            if not queued_at:
                continue
            fp = _safe_str(e.get("fp")).strip()
            if fp:
                st = self._rss_state_get(fp)
                if st.get("printed"):
                    continue
            try:
                qd = datetime.strptime(queued_at[:10], "%Y-%m-%d").date()
            except Exception:
                continue
            if start_date <= qd <= end_date:
                recent.append(e)

        if not recent:
            self._rss_status_var.set("No queued items in the last 7 days.")
            return

        open_after = bool((self.cfg.get("rss_review", {}) or {}).get("open_after_import", False))
        today_str = today.strftime("%Y-%m-%d")
        try:
            abs_path = self._rss_write_weekly_rollup_note(recent, start_date=start_date, end_date=end_date)
        except Exception as e:
            messagebox.showerror("Weekly Rollup", f"Failed to create weekly rollup.\n\n{e}")
            return

        recent_fps = {(_safe_str(e.get("fp")).strip()) for e in recent}
        for e in recent:
            fp = _safe_str(e.get("fp")).strip()
            if not fp:
                continue
            self._rss_state_set(
                fp,
                read=True,
                queued_for_report=False,
                reported_on=today_str,
                reported_url=_safe_str(e.get("url")).strip(),
                printed=True,
            )

        self.state["rss_report_queue"] = [
            e for e in queue if _safe_str(e.get("fp")).strip() not in recent_fps
        ]
        save_state(self.state_path, self.state)

        if open_after:
            try:
                _open_note_in_obsidian(self.cfg, self.tokens, abs_path)
            except Exception:
                pass

        self._rss_apply_filters()
        self._refresh_rss_cache_and_badges()
        self._rss_status_var.set(f"Weekly rollup created with {len(recent)} item(s). Queue updated.")
        try:
            messagebox.showinfo("Weekly Rollup", f"Weekly rollup created:\n{abs_path}")
        except Exception:
            pass

    def _rss_copy_markdown_selected(self):
        it = self._rss_get_primary_item()
        if not it:
            return

        fp = _rss_fp(it)
        st = self._rss_state_get(fp)
        tags = st.get("tags") if isinstance(st.get("tags"), list) else []
        notes = _safe_str(st.get("notes"))

        md = _mk_markdown_from_item(it, tags=tags, notes=notes)
        try:
            self.clipboard_clear()
            self.clipboard_append(md)
            self._rss_status_var.set("Copied markdown for selected item to clipboard.")
        except Exception as e:
            messagebox.showerror("Clipboard", f"Failed to copy markdown.\n\n{e}")

    def _rss_export_selection_json(self):
        items = self._rss_get_selected_items()
        if not items:
            messagebox.showinfo("Export", "Select one or more items first.")
            return

        export = []
        for it in items:
            fp = _rss_fp(it)
            st = self._rss_state_get(fp)
            export.append({"item": it, "state": st})

        try:
            base = self.state_path.parent
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            out_path = base / f"rss_export_{ts}.json"
            out_path.write_text(json.dumps(export, indent=2), encoding="utf-8")
            messagebox.showinfo("Export", f"Exported {len(items)} item(s) to:\n{out_path}")
        except Exception as e:
            messagebox.showerror("Export", f"Failed to export.\n\n{e}")

    def _rss_move_selection(self, delta: int):
        tv = self._rss_tree
        kids = list(tv.get_children())
        if not kids:
            return "break"

        sel = list(tv.selection() or [])
        if not sel:
            idx = 0
        else:
            try:
                idx = kids.index(sel[0])
            except Exception:
                idx = 0

        nxt = max(0, min(len(kids) - 1, idx + delta))
        tv.selection_set(kids[nxt])
        tv.focus(kids[nxt])
        tv.see(kids[nxt])
        self._rss_render_preview()
        return "break"

    # ----------------------------
    # Actions
    # ----------------------------
    def _open_obsidian_vault(self) -> None:
        paths_cfg = self.cfg.get("paths", {}) or {}
        tokens = self.tokens
        vault_name = expand_value(paths_cfg.get("OBSIDIAN_VAULT", ""), tokens).strip()
        vault_root = expand_value(
            paths_cfg.get("SCOUT_ROOT", "") or paths_cfg.get("VAULT_ROOT", ""),
            tokens,
        ).strip()
        obsidian_exe = expand_value(paths_cfg.get("OBSIDIAN_EXE", ""), tokens).strip()

        if vault_name:
            uri = f"obsidian://open?vault={quote(vault_name, safe='')}"
            open_obsidian_uri(uri)
            return

        if obsidian_exe and vault_root:
            run_command(f"\"{obsidian_exe}\" \"{vault_root}\"")
            return

        messagebox.showerror(
            "Obsidian Vault",
            "Obsidian vault is not configured. Set paths.OBSIDIAN_VAULT or paths.SCOUT_ROOT in config.json.",
            parent=self,
        )

    def _dispatch(self, item: dict):
        action = (item.get("action", "") or "").strip()
        target = expand_value(item.get("target", ""), self.tokens)

        # Goal Intake
        if action in ("new_goal", "goal_new", "open_new_goal", "goal_intake_new") or (
            str(item.get("label", "")).strip().lower() in ("new goal", "goal")
        ):
            try:
                self._open_goal_intake()
            except Exception as e:
                messagebox.showerror("Goal Intake", f"Failed to open Goal Intake UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)
            return

        # FAQ Intake
        if action in ("new_faq", "faq_new", "open_new_faq", "faq_intake_new") or (
            str(item.get("label","")).strip().lower() in ("new faq", "new frequently asked question")
        ):
            try:
                self._open_faq_intake()
            except Exception as e:
                messagebox.showerror("FAQ Intake", f"Failed to open FAQ Intake UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)
            return

        # How-To Intake
        if action in ("new_howto", "howto_new", "open_new_howto") or (
            str(item.get("label", "")).strip().lower() in ("new how-to", "new howto", "new how to", "how-to", "how to")
        ):
            try:
                self._open_howto_intake()
            except Exception as e:
                messagebox.showerror("How-To Intake", f"Failed to open How-To Intake UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)
            return

        # Threat Actor Intake
        if action in ("new_threat_actor", "threat_actor_new", "open_new_threat_actor") or (
            str(item.get("label", "")).strip().lower() in ("new threat actor", "threat actor")
        ):
            try:
                self._open_threat_actor_intake()
            except Exception as e:
                messagebox.showerror("Threat Actor Intake", f"Failed to open Threat Actor UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)
            return

        # Campaign Intake
        if action in ("new_campaign", "campaign_new", "open_new_campaign") or (
            str(item.get("label", "")).strip().lower() in ("new hacking campaign", "new campaign", "campaign")
        ):
            try:
                self._open_campaign_intake()
            except Exception as e:
                messagebox.showerror("Campaign Intake", f"Failed to open Campaign UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)
            return

        # MITRE TTP Intake
        if action in ("new_mitre_ttp", "mitre_ttp_new", "open_new_mitre_ttp") or (
            str(item.get("label", "")).strip().lower() in ("new mitre att&ck ttp", "new mitre ttp", "mitre ttp")
        ):
            try:
                self._open_mitre_ttp_intake()
            except Exception as e:
                messagebox.showerror("MITRE TTP Intake", f"Failed to open MITRE TTP UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)
            return

        # Malware Intake
        if action in ("new_malware", "malware_new", "open_new_malware") or (
            str(item.get("label", "")).strip().lower() in ("new malware family / strain", "new malware", "malware")
        ):
            try:
                self._open_malware_intake()
            except Exception as e:
                messagebox.showerror("Malware Intake", f"Failed to open Malware UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)
            return

        # IOC Intake
        if action in ("new_ioc", "ioc_new", "open_new_ioc") or (
            str(item.get("label", "")).strip().lower() in ("new indicator of compromise (ioc)", "new ioc", "ioc")
        ):
            try:
                self._open_ioc_intake()
            except Exception as e:
                messagebox.showerror("IOC Intake", f"Failed to open IOC UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)
            return

        if action == "open_path":
            open_path(resolve_path(target, self.tokens))
            self._record_recent(item)

        elif action == "run":
            run_command(expand_value(item.get("target", ""), self.tokens))
            self._record_recent(item)

        elif action == "open_rss_schedule":
            self._open_rss_schedule()
            self._record_recent(item)

        elif action == "url":
            open_url(target)
            self._record_recent(item)

        elif action == "obsidian_uri":
            open_obsidian_uri(target)
            self._record_recent(item)

        elif action == "open_obsidian_vault":
            self._open_obsidian_vault()
            self._record_recent(item)

        elif action == "obsidian_to_pptx":
            try:
                self._open_obsidian_to_pptx()
            except Exception as e:
                messagebox.showerror(
                    "Obsidian to PowerPoint",
                    f"Failed to open PowerPoint builder.\n\nAction: {action}\nError: {e}",
                    parent=self,
                )
            self._record_recent(item)

        elif action == "ioc_enrich":
            self._open_ioc_enrich()
            self._record_recent(item)

        elif action == "threat_actor_profile":
            self._open_threat_actor_profile()
            self._record_recent(item)

        elif action in ("open_daily_note_live", "daily_note_intake", "daily_note"):
            try:
                self._open_daily_note_intake()
            except Exception as e:
                messagebox.showerror("Daily Note Intake", f"Failed to open Daily Note Intake UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)

        elif action == "open_manage_keywords":
            if hasattr(self, "_open_manage_keywords"):
                self._open_manage_keywords()
                self._record_recent(item)
            else:
                messagebox.showerror(
                    "Manage Keywords",
                    "Launcher is missing method _open_manage_keywords().\n"
                    "Ensure it is defined inside class LauncherApp (not at top-level)."
                )

        elif action == "open_shift_log_rules":
            if hasattr(self, "_open_shift_log_rules"):
                self._open_shift_log_rules()
                self._record_recent(item)
            else:
                messagebox.showerror(
                    "Shift Log Rules",
                    "Launcher is missing method _open_shift_log_rules().\n"
                    "Ensure it is defined inside class LauncherApp (not at top-level)."
                )


        elif action == "open_reports_catalog":
            self._open_reports_catalog()
            self._record_recent(item)

        elif action == "project_status_rollup":
            self._open_project_status_rollup()
            self._record_recent(item)

        elif action == "open_rss_review":
            self._open_rss_review()
            self._record_recent(item)

        elif action == "run_rss_collect":
            python_cmd = (target or "").strip() or "python"
            app_dir = get_app_dir()
            result = run_rss_collect_and_capture(app_dir, python_cmd=python_cmd)

            self.state["rss_last_run"] = {
                "ok": bool(result.get("ok")),
                "ended": result.get("ended"),
                "feeds_processed": result.get("feeds_processed"),
                "new_inserted": result.get("new_inserted"),
                "exit_code": result.get("exit_code"),
                "error": result.get("error", ""),
                "stdout": result.get("stdout", ""),
                "stderr": result.get("stderr", ""),
            }
            save_state(self.state_path, self.state)

            ok = bool(result.get("ok"))
            feeds = result.get("feeds_processed")
            newi = result.get("new_inserted")

            msg = f"Command:\n{result.get('cmd','')}\n\nEnded: {result.get('ended','')}\n"
            if feeds is not None:
                msg += f"Feeds processed: {feeds}\n"
            if newi is not None:
                msg += f"New items inserted: {newi}\n"

            if ok:
                messagebox.showinfo("RSS Collect", msg)
            else:
                msg += f"\nError:\n{result.get('error','')}\n"
                messagebox.showerror("RSS Collect Failed", msg)

            self.after(30, self._refresh_rss_cache_and_badges)
            self._record_recent(item)

        
        elif action in (
            "manage_incident",
            "open_manage_incident",
            "open_incident_intake",
            "manage_incidents",
            "incident_intake_new",
            "incident_intake_resume",
        ):
            # Incident Intake entry points
            if not hasattr(self, "_open_incident_intake"):
                messagebox.showerror(
                    "Incident Intake",
                    "Launcher is missing method _open_incident_intake().",
                )
                return

            if action == "incident_intake_new":
                self._open_incident_intake(resume=False)
            else:
                # Default to resume to avoid data loss for legacy actions
                self._open_incident_intake(resume=True)

            self._record_recent(item)

        elif action in ("new_project", "project_new", "open_new_project", "project_intake_new") or (
            str(item.get("label", "")).strip().lower() in ("new project", "project intake", "project")
        ):
            try:
                self._open_project_intake(resume=False)
            except Exception as e:
                messagebox.showerror("Project Intake", f"Failed to open Project Intake UI.\n\nAction: {action}\nError: {e}")
                return
            self._record_recent(item)

        elif action in ("project_intake_resume", "open_project_resume", "project_resume"):
            try:
                self._open_project_intake(resume=True)
            except Exception as e:
                messagebox.showerror("Project Resume", f"Failed to open Project Intake UI (resume).\n\nAction: {action}\nError: {e}")
                return
            self._record_recent(item)

        elif action in ("new_meeting", "open_new_meeting", "meeting_intake_new", "meeting_new", "new_meeting_note", "create_meeting", "meeting"):
            try:
                self._open_meeting_intake(resume=False)
            except Exception as e:
                messagebox.showerror("New Meeting", f"Failed to open Meeting Intake UI.\n\nAction: {action}\nError: {e}")
                return
            self._record_recent(item)

        elif action in ("meeting_intake_resume", "open_meeting_resume", "meeting_resume"):
            try:
                self._open_meeting_intake(resume=True)
            except Exception as e:
                messagebox.showerror("Meeting Resume", f"Failed to open Meeting Intake UI (resume).\n\nAction: {action}\nError: {e}")
                return
            self._record_recent(item)

        elif action in ("new_playbook", "playbook_new", "open_new_playbook", "playbook_intake_new") or (str(item.get("label","")).strip().lower() == "new playbook"):
            try:
                self._open_playbook_intake(resume=False)
            except Exception as e:
                messagebox.showerror("Playbook Intake", f"Failed to open Playbook Intake UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)

        elif action in ("playbook_intake_resume", "open_playbook_resume", "playbook_resume"):
            try:
                self._open_playbook_intake(resume=True)
            except Exception as e:
                messagebox.showerror("Playbook Intake", f"Failed to open Playbook Intake UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)

        elif action in ("new_procedure", "procedure_new", "open_new_procedure", "procedure_intake_new") or (str(item.get("label","")).strip().lower() == "new procedure"):
            try:
                self._open_procedure_intake(resume=False)
            except Exception as e:
                messagebox.showerror("Procedure Intake", f"Failed to open Procedure Intake UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)

        elif action in ("procedure_intake_resume", "open_procedure_resume", "procedure_resume"):
            try:
                self._open_procedure_intake(resume=True)
            except Exception as e:
                messagebox.showerror("Procedure Intake", f"Failed to open Procedure Intake UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)
        elif action in ("new_itid", "itid_new", "open_new_itid", "itid_intake_new") or (str(item.get("label","")).strip().lower() == "new itid"):
            try:
                self._open_itid_intake(resume=False)
            except Exception as e:
                messagebox.showerror("ITID Intake", f"Failed to open ITID Intake UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)

        elif action in ("itid_intake_resume", "open_itid_resume", "itid_resume"):
            try:
                self._open_itid_intake(resume=True)
            except Exception as e:
                messagebox.showerror("ITID Intake", f"Failed to open ITID Intake UI.\n\nAction: {action}\nError: {e}", parent=self)
            self._record_recent(item)


        elif action in ("new_sla", "sla_new", "new_service_level_agreement", "sla_intake_new", "open_new_sla") or (
            str(item.get("label", "")).strip().lower() == "new service level agreement"
        ):
            self._open_sla_intake(resume=False)


        elif action == "quit":
            self.destroy()


        else:
            # Label-based safety net: if config label is "New Meeting" but action differs, still open meeting intake.
            lab = (item.get("label", "") or "").strip().lower()
            if lab == "new meeting":
                try:
                    self._open_meeting_intake(resume=False)
                    self._record_recent(item)
                    return
                except Exception as e:
                    messagebox.showerror("New Meeting", f"Failed to open Meeting Intake UI.\n\nAction: {action}\nError: {e}")
                    return

            messagebox.showerror("Unsupported action", f"Unknown action: {action}")


    # ----------------------------
    # Centering logic
    # ----------------------------
    def _recenter_ui(self, w: int, h: int) -> None:
        self.ui_frame.update_idletasks()

        ui_w = self.ui_frame.winfo_width()
        ui_h = self.ui_frame.winfo_height()

        if ui_w <= 1 or ui_h <= 1:
            ui_w = self.ui_frame.winfo_reqwidth()
            ui_h = self.ui_frame.winfo_reqheight()

        avail_w = max(w - 2 * self.bg_margin_x, 1)
        avail_h = max(h - 2 * self.bg_margin_y, 1)

        x = max((avail_w - ui_w) // 2 + self.bg_margin_x, 0)
        y = max((avail_h - ui_h) // 2 + self.bg_margin_y + self.center_offset_y, 0)

        self.canvas.coords(self.ui_window_id, x, y)

    # ----------------------------
    # Initial render + resizing (DEBOUNCED)
    # ----------------------------
    def _initial_render(self):
        self.update_idletasks()
        self._on_resize(None)

    def _on_resize(self, event):
        w = self.winfo_width()
        h = self.winfo_height()

        if getattr(self, "print_resize", False):
            print("RESIZE:", w, h, "BG_LOADED:", bool(self.bg_original))

        self._pending_resize = (w, h)

        if self._resize_job is not None:
            try:
                self.after_cancel(self._resize_job)
            except Exception:
                pass

        self._resize_job = self.after(60, self._apply_resize)

    def _apply_resize(self):
        self._resize_job = None
        if not self._pending_resize:
            return

        w, h = self._pending_resize
        self._pending_resize = None

        if (w, h) != self._bg_last_size:
            self._bg_last_size = (w, h)

            if self.bg_original:
                now_ts = datetime.now().timestamp()
                resample = Image.Resampling.BILINEAR if now_ts < self._startup_fast_resample_until else Image.Resampling.LANCZOS
                bg = self.bg_original.resize((w, h), resample)
                self.bg_tk = ImageTk.PhotoImage(bg)

                if self.bg_id is None:
                    self.bg_id = self.canvas.create_image(0, 0, anchor="nw", image=self.bg_tk)
                else:
                    self.canvas.itemconfig(self.bg_id, image=self.bg_tk)

                self.canvas.tag_lower(self.bg_id)

        left_w = max(140, min(260, int(w * 0.16)))
        right_w = max(140, min(260, int(w * 0.16)))
        self._set_header_logos(default_left_w=left_w, default_right_w=right_w)

        self._recenter_ui(w, h)
        

    # ----------------------------
    # Incident Intake (preliminary data capture)
    # ----------------------------
    def _open_incident_intake(self, resume: bool = True):
        """Opens the Incident Intake UI and persists the resulting draft."""
        try:
            app_dir = get_app_dir()
            if not resume:
                delete_incident_draft(Path(app_dir))


            # Best-effort vault root (filesystem root of the Obsidian vault).
            # We support multiple config conventions to reduce brittleness:
            #   - cfg["paths"]["SCOUT_ROOT"] (preferred in this launcher)
            #   - cfg["paths"]["VAULT_ROOT"]
            #   - cfg["obsidian_vault_path"] (top-level convenience key)
            vault_root = None
            vault_root_err = None

            def _try_dir(candidate: str):
                nonlocal vault_root, vault_root_err
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)
                else:
                    vault_root_err = f"Vault root is not a directory: {cand}"

            try:
                paths = self.cfg.get("paths", {}) or {}
                _try_dir(paths.get("SCOUT_ROOT", ""))
                if vault_root is None:
                    _try_dir(paths.get("VAULT_ROOT", ""))
                if vault_root is None:
                    _try_dir(self.cfg.get("obsidian_vault_path", ""))
            except Exception as e:
                vault_root_err = str(e)
                vault_root = None

            def _after_submit(payload: dict):
                # Optionally open the newly created note (if any)
                note_abs = (payload or {}).get("obsidian_note_path")
                open_note = bool((payload or {}).get("open_note", True))
                if note_abs and open_note:
                    try:
                        _open_note_in_obsidian(self.cfg, self.tokens, str(note_abs))
                    except Exception as e:
                        try:
                            messagebox.showwarning(
                                "Incident Intake",
                                f"Incident saved, but the note could not be opened in Obsidian:\n{e}",
                                parent=self,
                            )
                        except Exception:
                            pass

                # Lightweight confirmation for operator
                try:
                    inc_id = (payload or {}).get("incident_id", "") or ""
                    title = (payload or {}).get("title", "") or ""
                    self.state["incident_last_intake"] = {
                        "incident_id": inc_id,
                        "title": title,
                        "ts": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        "payload": payload or {},
                    }
                    save_state(self.state_path, self.state)
                except Exception:
                    pass

            IncidentIntakeWindow(
                self,
                app_dir=Path(app_dir),
                vault_root=vault_root,
                on_submit=_after_submit,
                title="Incident Intake",
            )

        except Exception as e:
            messagebox.showerror("Incident Intake", str(e))


    def _open_meeting_intake(self, resume: bool = False):
        app_dir = get_app_dir()

        # Best-effort vault root (filesystem root of the Obsidian vault).
        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
            if vault_root is None:
                _try_dir(self.cfg.get("obsidian_vault_path", ""))
        except Exception:
            vault_root = None

        def _after_meeting(meeting: dict):
            # Store state; no success dialog on creation.
            self.current_meeting = meeting
            try:
                title = (meeting or {}).get("title", "") or ""
                start_time = (meeting or {}).get("start_time", "") or ""
                self.state["meeting_last_intake"] = {
                    "title": title,
                    "start_time": start_time,
                    "ts": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "payload": meeting or {},
                }
                save_state(self.state_path, self.state)
            except Exception:
                pass


        MeetingIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_meeting,
            resume=resume,
        )


    def _open_project_intake(self, resume: bool = False):
        app_dir = get_app_dir()

        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            candidates = [
                paths.get("SCOUT_ROOT", ""),
                paths.get("VAULT_ROOT", ""),
                self.cfg.get("obsidian_vault_path", ""),
            ]

            resolved = []
            for candidate in candidates:
                cand = (candidate or "").strip()
                if not cand:
                    continue
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    resolved.append(Path(cand))

            # Prefer the vault that actually contains the Project template
            for root in resolved:
                tpl = root / "00_System" / "00_Templates" / "02_Operations" / "Project Template.md"
                if tpl.exists():
                    vault_root = root
                    break

            if vault_root is None and resolved:
                vault_root = resolved[0]
        except Exception:
            vault_root = None

        def _after_project(project: dict):
            self.current_project = project
            try:
                name = (project or {}).get("project_name", "") or ""
                self.state["project_last_intake"] = {
                    "project_name": name,
                    "ts": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "payload": project or {},
                }
                save_state(self.state_path, self.state)
            except Exception:
                pass

        ProjectIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_project,
            resume=resume,
        )


    def _open_playbook_intake(self, resume: bool = False):
        app_dir = get_app_dir()

        # Best-effort vault root (filesystem root of the Obsidian vault).
        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
            if vault_root is None:
                _try_dir(self.cfg.get("obsidian_vault_path", ""))
        except Exception:
            vault_root = None

        def _after_playbook(pb: dict):
            self.current_playbook = pb
            try:
                title = (pb or {}).get("title", "") or ""
                self.state["playbook_last_intake"] = {
                    "title": title,
                    "ts": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "payload": pb or {},
                }
                save_state(self.state_path, self.state)
            except Exception:
                pass

        PlaybookIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_playbook,
            resume=resume,
        )


    def _open_procedure_intake(self, resume: bool = False):
        app_dir = get_app_dir()

        # Best-effort vault root (filesystem root of the Obsidian vault).
        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
            if vault_root is None:
                _try_dir(self.cfg.get("obsidian_vault_path", ""))
        except Exception:
            vault_root = None

        def _after_procedure(proc: dict):
            self.current_procedure = proc
            try:
                title = (proc or {}).get("title", "") or ""
                self.state["procedure_last_intake"] = {
                    "title": title,
                    "ts": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "payload": proc or {},
                }
                save_state(self.state_path, self.state)
            except Exception:
                pass

        ProcedureIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_procedure,
            resume=resume,
        )


    def _open_sla_intake(self, resume: bool = False):
        app_dir = get_app_dir()

        # Best-effort vault root (filesystem root of the Obsidian vault).
        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
            if vault_root is None:
                _try_dir(self.cfg.get("obsidian_vault_path", ""))
        except Exception:
            vault_root = None

        def _after_sla(sla: dict):
            self.current_sla = sla
            try:
                sla_id = (sla or {}).get("sla_id", "") or ""
                sla_title = (sla or {}).get("sla_title", "") or ""
                self.state["sla_last_intake"] = {
                    "sla_id": sla_id,
                    "sla_title": sla_title,
                    "ts": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "payload": sla or {},
                }
                save_state(self.state_path, self.state)
            except Exception:
                pass

        SLAIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_sla,
            resume=resume,
        )

    def _open_itid_intake(self, resume: bool = False):
        app_dir = get_app_dir()

        # Best-effort vault root (filesystem root of the Obsidian vault).
        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
            if vault_root is None:
                _try_dir(self.cfg.get("obsidian_vault_path", ""))
        except Exception:
            vault_root = None

        def _after_itid(itid: dict):
            self.current_itid = itid
            try:
                itid_id = (itid or {}).get("itid_id", "") or ""
                itid_name = (itid or {}).get("itid_name", "") or ""
                self.state["itid_last_intake"] = {
                    "itid_id": itid_id,
                    "itid_name": itid_name,
                    "ts": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "payload": itid or {},
                }
                save_state(self.state_path, self.state)
            except Exception:
                pass

        ITIDIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_itid,
            resume=resume,
        )



    def _open_daily_note_intake(self):
        if TkinterDnD and hasattr(TkinterDnD, "Toplevel"):
            win = TkinterDnD.Toplevel(self)
        else:
            win = tk.Toplevel(self)
        win.title("Daily Note Intake")
        win.configure(bg="#111111")
        win.geometry("760x720")
        win.minsize(720, 620)

        # Best-effort vault root for relative image paths.
        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
            if vault_root is None:
                _try_dir(self.cfg.get("obsidian_vault_path", ""))
        except Exception:
            vault_root = None

        container = tk.Frame(win, bg="#111111")
        container.pack(fill="both", expand=True)

        canvas = tk.Canvas(container, bg="#111111", highlightthickness=0)
        vsb = tk.Scrollbar(container, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vsb.set)

        vsb.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        inner = tk.Frame(canvas, bg="#111111")
        inner_id = canvas.create_window((0, 0), window=inner, anchor="nw")

        def _on_inner_config(event=None):
            try:
                canvas.configure(scrollregion=canvas.bbox("all"))
            except Exception:
                pass

        def _on_canvas_config(event):
            try:
                canvas.itemconfigure(inner_id, width=event.width)
            except Exception:
                pass

        inner.bind("<Configure>", _on_inner_config)
        canvas.bind("<Configure>", _on_canvas_config)

        hdr = tk.Frame(inner, bg="#111111")
        hdr.pack(fill="x", padx=14, pady=(14, 6))
        tk.Label(hdr, text="Daily Note Intake", fg="#FFFFFF", bg="#111111", font=("Segoe UI", 16, "bold")).pack(anchor="w")
        tk.Label(
            hdr,
            text="Capture shift metadata and ratings before creating/opening today’s journal entry.",
            fg="#AAAAAA",
            bg="#111111",
            font=("Segoe UI", 10),
        ).pack(anchor="w", pady=(4, 0))

        form = tk.Frame(inner, bg="#111111")
        form.pack(fill="both", expand=True, padx=14, pady=10)

        def add_labeled_entry(parent, label, default="", *, width=24, readonly=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, width=width, anchor="w", fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            var = tk.StringVar(value=default)
            ent = tk.Entry(row, textvariable=var, bg="#1A1A1A", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat")
            if readonly:
                ent.configure(state="readonly")
            ent.pack(side="left", fill="x", expand=True)
            return var

        def add_labeled_combo(parent, label, values, default="", *, width=24):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, width=width, anchor="w", fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            var = tk.StringVar(value=default)
            cb = ttk.Combobox(row, textvariable=var, values=list(values), state="readonly")
            cb.pack(side="left", fill="x", expand=True)
            return var

        def add_labeled_slider(parent, label, default=5, *, width=24):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=10)
            tk.Label(row, text=label, width=width, anchor="w", fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            var = tk.IntVar(value=int(default))
            scale = tk.Scale(
                row,
                from_=1,
                to=10,
                orient="horizontal",
                variable=var,
                showvalue=True,
                resolution=1,
                bg="#111111",
                fg="#FFFFFF",
                highlightthickness=0,
                troughcolor="#2A2A2A",
                relief="flat",
            )
            scale.pack(side="left", fill="x", expand=True)
            return var

        now_dt = datetime.now()
        today = now_dt.strftime("%Y-%m-%d")
        now_ts = now_dt.strftime("%Y-%m-%d %H:%M:%S")

        analyst_default = (self.cfg.get("analyst_profile", {}) or {}).get("name", "")
        analyst_var = add_labeled_entry(form, "Analyst", analyst_default)
        shift_type_var = add_labeled_combo(form, "Shift Type", ["Daytime", "Swing", "Night", "On-Call"], default="Daytime")
        shift_start_var = add_labeled_entry(form, "Shift Start", now_ts)

        energy_var = add_labeled_slider(form, "Energy Level", default=6)
        cynicism_var = add_labeled_slider(form, "Cynicism Level", default=4)
        productivity_var = add_labeled_slider(form, "Productivity Level", default=6)

        music_var = add_labeled_entry(form, "Music Mood", "")
        tags_var = add_labeled_entry(form, "Tags", "#log/daily")

        picture_row = tk.Frame(form, bg="#111111")
        picture_row.pack(fill="x", pady=6)
        tk.Label(
            picture_row,
            text="Picture of the Day",
            width=24,
            anchor="w",
            fg="#DDDDDD",
            bg="#111111",
            font=("Segoe UI", 10),
        ).pack(side="left", anchor="n")
        picture_var = tk.StringVar(value="")
        # Drop zone frame - prominent drag-and-drop area
        drop_zone = tk.Frame(
            picture_row,
            bg="#1A1A1A",
            highlightbackground="#444444",
            highlightthickness=2,
            highlightcolor="#666666",
            cursor="hand2",
        )
        drop_zone.pack(side="left", fill="both", expand=True, padx=(0, 0))
        drop_zone_inner = tk.Frame(drop_zone, bg="#1A1A1A")
        drop_zone_inner.pack(fill="both", expand=True, padx=12, pady=10)
        drop_zone_label = tk.Label(
            drop_zone_inner,
            text="Drop image here (file or from browser) or click Browse",
            fg="#888888",
            bg="#1A1A1A",
            font=("Segoe UI", 10),
            cursor="hand2",
        )
        drop_zone_label.pack(anchor="w")
        picture_entry = tk.Entry(
            drop_zone_inner,
            textvariable=picture_var,
            bg="#1A1A1A",
            fg="#FFFFFF",
            insertbackground="#FFFFFF",
            relief="flat",
        )
        picture_entry.pack(fill="x", pady=(4, 0))

        def _pick_picture_from_path(file_path: str):
            try:
                if vault_root is None:
                    messagebox.showerror(
                        "Picture of the Day",
                        "Vault root is not configured. Set SCOUT_ROOT/VAULT_ROOT in config.json.",
                        parent=win,
                    )
                    return
                if not file_path:
                    return
                src_path = Path(os.path.normpath(file_path))
                if not src_path.exists() or not src_path.is_file():
                    messagebox.showerror(
                        "Picture of the Day",
                        "Select a valid image file.",
                        parent=win,
                    )
                    return
                if src_path.suffix.lower() not in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}:
                    messagebox.showerror(
                        "Picture of the Day",
                        "Unsupported file type. Use an image file.",
                        parent=win,
                    )
                    return
                attach_dir = Path(vault_root) / "99_Attachments" / "Picture_of_the_Day"
                try:
                    attach_dir.mkdir(parents=True, exist_ok=True)
                except Exception as e:
                    messagebox.showerror(
                        "Picture of the Day",
                        f"Could not create attachment folder:\n{attach_dir}\n\nError: {e}",
                        parent=win,
                    )
                    return

                date_prefix = datetime.now().strftime("%Y-%m-%d")
                base = f"{date_prefix} - {src_path.stem}"
                ext = src_path.suffix
                dest_path = attach_dir / f"{base}{ext}"
                counter = 1
                while dest_path.exists():
                    dest_path = attach_dir / f"{base}_{counter}{ext}"
                    counter += 1
                try:
                    shutil.copy2(src_path, dest_path)
                except Exception as e:
                    messagebox.showerror(
                        "Picture of the Day",
                        f"Failed to copy image into the vault.\n\nError: {e}",
                        parent=win,
                    )
                    return

                rel = dest_path.relative_to(Path(vault_root)).as_posix()
                picture_var.set(f"[[{rel}]]")
            except Exception:
                pass

        def _pick_picture_from_url(url: str):
            try:
                if vault_root is None:
                    messagebox.showerror(
                        "Picture of the Day",
                        "Vault root is not configured. Set SCOUT_ROOT/VAULT_ROOT in config.json.",
                        parent=win,
                    )
                    return
                url = (url or "").strip()
                if not (url.startswith("http://") or url.startswith("https://")):
                    return
                parsed = urlparse(url)
                filename = unquote(os.path.basename(parsed.path)) or f"image_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
                if "." not in filename:
                    filename += ".png"
                # Sanitize filename for Windows/Obsidian (remove invalid chars)
                filename = re.sub(r'[<>:"/\\|?*]', "_", filename)
                filename = filename.strip(". ") or f"image_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"

                attach_dir = Path(vault_root) / "99_Attachments" / "Picture_of_the_Day"
                attach_dir.mkdir(parents=True, exist_ok=True)
                dest_path = attach_dir / filename
                counter = 1
                while dest_path.exists():
                    stem = dest_path.stem
                    suffix = dest_path.suffix or ".png"
                    dest_path = attach_dir / f"{stem}_{counter}{suffix}"
                    counter += 1

                # Use browser-like headers to avoid 403 (many sites block non-browser requests)
                headers = {
                    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                    "Accept": "image/avif,image/webp,image/apng,image/svg+xml,image/*,*/*;q=0.8",
                    "Accept-Language": "en-US,en;q=0.9",
                    "Referer": f"{parsed.scheme}://{parsed.netloc}/",
                }
                data = None
                content_type = ""
                if requests:
                    resp = requests.get(url, headers=headers, timeout=15)
                    resp.raise_for_status()
                    data = resp.content
                    content_type = (resp.headers.get("Content-Type") or "").lower()
                else:
                    import urllib.request

                    req = urllib.request.Request(url, headers=headers)
                    with urllib.request.urlopen(req, timeout=15) as resp:
                        data = resp.read()
                        content_type = (resp.headers.get("Content-Type") or "").lower()

                if not data:
                    raise ValueError("No image data downloaded.")

                if "image/" not in content_type:
                    # best-effort: still save if extension looks like image
                    ext = Path(dest_path).suffix.lower()
                    if ext not in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}:
                        raise ValueError("URL does not appear to be an image.")

                dest_path.write_bytes(data)
                rel = dest_path.relative_to(Path(vault_root)).as_posix()
                picture_var.set(f"[[{rel}]]")
            except Exception as e:
                messagebox.showerror("Picture of the Day", f"Failed to fetch image from URL.\n\n{e}", parent=win)

        def _paste_picture_from_clipboard():
            try:
                if vault_root is None:
                    messagebox.showerror(
                        "Picture of the Day",
                        "Vault root is not configured. Set SCOUT_ROOT/VAULT_ROOT in config.json.",
                        parent=win,
                    )
                    return
                try:
                    from PIL import ImageGrab  # type: ignore
                except Exception:
                    messagebox.showerror("Picture of the Day", "Clipboard image support not available.", parent=win)
                    return
                grabbed = ImageGrab.grabclipboard()
                if grabbed is None:
                    messagebox.showerror("Picture of the Day", "Clipboard does not contain an image.", parent=win)
                    return
                attach_dir = Path(vault_root) / "99_Attachments" / "Picture_of_the_Day"
                attach_dir.mkdir(parents=True, exist_ok=True)
                date_prefix = datetime.now().strftime("%Y-%m-%d")
                dest_path = attach_dir / f"{date_prefix} - clipboard.png"
                counter = 1
                while dest_path.exists():
                    dest_path = attach_dir / f"{date_prefix} - clipboard_{counter}.png"
                    counter += 1
                if hasattr(grabbed, "save"):
                    grabbed.save(dest_path)
                elif isinstance(grabbed, list) and grabbed:
                    _pick_picture_from_path(grabbed[0])
                    return
                else:
                    messagebox.showerror("Picture of the Day", "Clipboard does not contain an image.", parent=win)
                    return
                rel = dest_path.relative_to(Path(vault_root)).as_posix()
                picture_var.set(f"[[{rel}]]")
            except Exception as e:
                messagebox.showerror("Picture of the Day", f"Failed to paste image.\n\n{e}", parent=win)

        def _pick_picture():
            try:
                initial_dir = str(vault_root) if vault_root else None
                file_path = filedialog.askopenfilename(
                    parent=win,
                    title="Select Picture of the Day",
                    initialdir=initial_dir,
                    filetypes=[
                        ("Image files", "*.png;*.jpg;*.jpeg;*.gif;*.bmp;*.webp"),
                        ("All files", "*.*"),
                    ],
                )
                if not file_path:
                    return
                _pick_picture_from_path(file_path)
            except Exception:
                pass

        def _extract_url_from_drop(text: str) -> str | None:
            """Extract image URL from browser drop data (may include 'URL\\n' prefix or title)."""
            if not text or not isinstance(text, str):
                return None
            text = text.strip()
            for line in text.splitlines():
                line = line.strip()
                if line.upper().startswith("URL"):
                    line = line[3:].strip()
                if line.startswith("http://") or line.startswith("https://"):
                    return line
            if text.startswith("http://") or text.startswith("https://"):
                return text.splitlines()[0].strip()
            return None

        def _drop_picture(event):
            action = getattr(event, "action", "copy")
            try:
                raw = event.data
                if not raw:
                    return action
                # Handle URL/text from browser (DND_TEXT) - e.g. dragging image from web page
                url = _extract_url_from_drop(raw)
                if url:
                    _pick_picture_from_url(url)
                    return action
                # Handle file paths (DND_FILES)
                paths = []
                if hasattr(win, "tk"):
                    try:
                        paths = list(win.tk.splitlist(raw))
                    except Exception:
                        paths = [raw] if isinstance(raw, str) else [raw]
                else:
                    paths = [raw] if isinstance(raw, str) else [raw]
                if not paths:
                    return action
                first = paths[0]
                if isinstance(first, str) and (first.startswith("http://") or first.startswith("https://")):
                    _pick_picture_from_url(first)
                else:
                    _pick_picture_from_path(first)
            except Exception:
                pass
            return action

        def _update_drop_zone_label(*_):
            val = (picture_var.get() or "").strip()
            if val:
                # Show shortened path for display
                display = val.replace("[[", "").replace("]]", "")
                if len(display) > 50:
                    display = "..." + display[-47:]
                drop_zone_label.config(text=display, fg="#AAAAAA")
            else:
                drop_zone_label.config(text="Drop image here (file or from browser) or click Browse", fg="#888888")

        def _drop_enter(event):
            drop_zone.config(highlightbackground="#6A9BD8", highlightcolor="#6A9BD8")
            drop_zone_label.config(fg="#AAAAAA")
            return event.action if hasattr(event, "action") else "copy"

        def _drop_leave(event):
            drop_zone.config(highlightbackground="#444444", highlightcolor="#666666")
            _update_drop_zone_label()
            return event.action if hasattr(event, "action") else "copy"

        picture_var.trace_add("write", _update_drop_zone_label)

        btn_row = tk.Frame(picture_row, bg="#111111")
        btn_row.pack(side="left", padx=(8, 0), anchor="n")
        tk.Button(
            btn_row,
            text="Browse",
            command=_pick_picture,
            bg="#2A2A2A",
            fg="#FFFFFF",
            relief="flat",
            padx=12,
            pady=2,
        ).pack(side="left", padx=(0, 4))
        tk.Button(
            btn_row,
            text="Paste",
            command=_paste_picture_from_clipboard,
            bg="#2A2A2A",
            fg="#FFFFFF",
            relief="flat",
            padx=12,
            pady=2,
        ).pack(side="left")

        # Register drop zone as drag-and-drop target (DND_FILES + DND_TEXT for browser drops)
        if TkinterDnD and DND_FILES and hasattr(drop_zone, "drop_target_register"):
            try:
                types = [DND_FILES]
                if DND_TEXT:
                    types.append(DND_TEXT)
                drop_zone.drop_target_register(*types)
                drop_zone.dnd_bind("<<Drop>>", _drop_picture)
                drop_zone.dnd_bind("<<DropEnter>>", _drop_enter)
                drop_zone.dnd_bind("<<DropLeave>>", _drop_leave)
            except Exception:
                pass
        # Also allow drop on entry for users who drag onto the path field
        if TkinterDnD and DND_FILES and hasattr(picture_entry, "drop_target_register"):
            try:
                types = [DND_FILES]
                if DND_TEXT:
                    types.append(DND_TEXT)
                picture_entry.drop_target_register(*types)
                picture_entry.dnd_bind("<<Drop>>", _drop_picture)
            except Exception:
                pass

        def _on_drop_zone_click(event=None):
            if event and event.widget == picture_entry:
                return
            _pick_picture()

        drop_zone_label.bind("<Button-1>", _on_drop_zone_click)
        drop_zone.bind("<Button-1>", _on_drop_zone_click)

        def _on_paste_event(_event=None):
            _paste_picture_from_clipboard()
            return "break"

        try:
            picture_entry.bind("<Control-v>", _on_paste_event)
            picture_entry.bind("<Control-V>", _on_paste_event)
        except Exception:
            pass

        btns = tk.Frame(inner, bg="#111111")
        btns.pack(fill="x", padx=14, pady=(0, 14))

        def _cancel():
            try:
                win.destroy()
            except Exception:
                pass

        def _save_and_open():
            journal_date = today
            analyst = (analyst_var.get() or "").strip()
            shift_type = (shift_type_var.get() or "").strip() or "Daytime"
            shift_start = (shift_start_var.get() or "").strip()
            music = (music_var.get() or "").strip()
            tags = (tags_var.get() or "").strip()
            picture_of_the_day = (picture_var.get() or "").strip()

            yaml_updates = {
                "journal_date": journal_date,
                "analyst": analyst,
                "shift_type": shift_type,
                "shift_start": shift_start,
                "energy_rating": int(energy_var.get()),
                "cynicism_rating": int(cynicism_var.get()),
                "productivity_rating": int(productivity_var.get()),
                "music_selection": music,
                "picture_of_the_day": picture_of_the_day,
            }
            intake_md = ""

            try:
                win.destroy()
            except Exception:
                pass

            open_daily_note_live_with_intake(
                self.cfg,
                self.tokens,
                intake_block=intake_md,
                tags_csv=tags,
                yaml_updates=yaml_updates,
            )

        tk.Button(btns, text="Cancel", command=_cancel, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8).pack(side="right")
        tk.Button(btns, text="Save & Open Daily Note", command=_save_and_open, bg="#3A6EA5", fg="#FFFFFF", relief="flat", padx=14, pady=8).pack(side="right", padx=(0, 10))

        def _on_mousewheel(event):
            try:
                canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")
            except Exception:
                pass
        canvas.bind_all("<MouseWheel>", _on_mousewheel)

        try:
            win.transient(self)
            win.grab_set()
            win.focus_force()
        except Exception:
            pass




    def _open_faq_intake(self):
        app_dir = get_app_dir()

        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
        except Exception:
            vault_root = None

        if vault_root is None:
            messagebox.showerror("FAQ Intake", "Could not resolve Obsidian vault root (SCOUT_ROOT/VAULT_ROOT).", parent=self)
            return

        def _after_faq(faq: dict):
            try:
                note_path = create_obsidian_faq_note(vault_root, faq)
                faq["obsidian_note_path"] = str(note_path)
                _log_shift_entry_from_launcher(self, vault_root, "faq", faq, note_path)
            except Exception as e:
                messagebox.showerror("FAQ Intake", f"Failed to create FAQ note.\n\nError: {e}", parent=self)
                return

            try:
                self.state["last_faq"] = faq
                save_state(self.state_path, self.state)
            except Exception:
                pass

            try:
                if faq.get("open_after_save"):
                    open_obsidian_file(vault_root, Path(faq["obsidian_note_path"]))
            except Exception:
                pass

        FAQIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_faq,
        )

    def _open_goal_intake(self):
        app_dir = get_app_dir()

        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
        except Exception:
            vault_root = None

        if vault_root is None:
            messagebox.showerror("Goal Intake", "Could not resolve Obsidian vault root (SCOUT_ROOT/VAULT_ROOT).", parent=self)
            return

        def _after_goal(goal: dict):
            try:
                note_path = create_obsidian_goal_note(vault_root, goal)
                goal["obsidian_note_path"] = str(note_path)
                _log_shift_entry_from_launcher(self, vault_root, "goal", goal, note_path)
            except Exception as e:
                messagebox.showerror("Goal Intake", f"Failed to create Goal note.\n\nError: {e}", parent=self)
                return

            try:
                self.state["last_goal"] = goal
                save_state(self.state_path, self.state)
            except Exception:
                pass

            try:
                if goal.get("open_after_save"):
                    open_obsidian_file(vault_root, Path(goal["obsidian_note_path"]))
            except Exception:
                pass

        GoalIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_goal,
        )

    def _open_howto_intake(self):
        app_dir = get_app_dir()

        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                try:
                    p = Path(cand).expanduser()
                    if p.exists() and p.is_dir():
                        vault_root = p
                except Exception:
                    pass

            _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
        except Exception:
            vault_root = None

        if vault_root is None:
            messagebox.showerror("How-To Intake", "Could not resolve Obsidian vault root (SCOUT_ROOT/VAULT_ROOT).", parent=self)
            return

        def _after_howto(howto: dict):
            try:
                note_path = create_obsidian_howto_note(vault_root, howto)
                howto["obsidian_note_path"] = str(note_path)
                _log_shift_entry_from_launcher(self, vault_root, "how_to", howto, note_path)
            except Exception as e:
                messagebox.showerror("How-To Intake", f"Failed to create How-To note.\n\nError: {e}", parent=self)
                return

            try:
                self.state["last_howto"] = howto
                save_state(self.state_path, self.state)
            except Exception:
                pass

            try:
                if howto.get("open_after_save"):
                    open_obsidian_file(vault_root, Path(howto["obsidian_note_path"]))
            except Exception:
                pass

        HowToIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            cfg=self.cfg,
            on_submit=_after_howto,
        )

    def _open_threat_actor_intake(self):
        app_dir = get_app_dir()

        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("CIPHER_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
            if vault_root is None:
                _try_dir(self.cfg.get("obsidian_vault_path", ""))
        except Exception:
            vault_root = None

        if vault_root is None:
            messagebox.showerror("Threat Actor Intake", "Could not resolve CIPHER vault root.", parent=self)
            return

        def _after_actor(actor: dict):
            try:
                note_path = create_obsidian_threat_actor_note(vault_root, actor)
                actor["obsidian_note_path"] = str(note_path)
                _log_shift_entry_from_launcher(self, vault_root, "threat_actor", actor, note_path)
            except Exception as e:
                messagebox.showerror("Threat Actor Intake", f"Failed to create Threat Actor note.\n\nError: {e}", parent=self)
                return

            try:
                if actor.get("open_after_save"):
                    open_obsidian_file(vault_root, Path(actor["obsidian_note_path"]))
            except Exception:
                pass

        ThreatActorIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_actor,
        )

    def _open_campaign_intake(self):
        app_dir = get_app_dir()

        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("CIPHER_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
            if vault_root is None:
                _try_dir(self.cfg.get("obsidian_vault_path", ""))
        except Exception:
            vault_root = None

        if vault_root is None:
            messagebox.showerror("Campaign Intake", "Could not resolve CIPHER vault root.", parent=self)
            return

        def _after_campaign(campaign: dict):
            try:
                note_path = create_obsidian_campaign_note(vault_root, campaign)
                campaign["obsidian_note_path"] = str(note_path)
                _log_shift_entry_from_launcher(self, vault_root, "campaign", campaign, note_path)
            except Exception as e:
                messagebox.showerror("Campaign Intake", f"Failed to create Campaign note.\n\nError: {e}", parent=self)
                return

            try:
                if campaign.get("open_after_save"):
                    open_obsidian_file(vault_root, Path(campaign["obsidian_note_path"]))
            except Exception:
                pass

        CampaignIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_campaign,
        )

    def _open_mitre_ttp_intake(self):
        app_dir = get_app_dir()

        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("CIPHER_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
            if vault_root is None:
                _try_dir(self.cfg.get("obsidian_vault_path", ""))
        except Exception:
            vault_root = None

        if vault_root is None:
            messagebox.showerror("MITRE TTP Intake", "Could not resolve CIPHER vault root.", parent=self)
            return

        def _after_ttp(ttp: dict):
            try:
                note_path = create_obsidian_mitre_ttp_note(vault_root, ttp)
                ttp["obsidian_note_path"] = str(note_path)
                _log_shift_entry_from_launcher(self, vault_root, "mitre_ttp", ttp, note_path)
            except Exception as e:
                messagebox.showerror("MITRE TTP Intake", f"Failed to create MITRE TTP note.\n\nError: {e}", parent=self)
                return

            try:
                if ttp.get("open_after_save"):
                    open_obsidian_file(vault_root, Path(ttp["obsidian_note_path"]))
            except Exception:
                pass

        MITRETTPIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_ttp,
        )

    def _open_malware_intake(self):
        app_dir = get_app_dir()

        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("CIPHER_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
            if vault_root is None:
                _try_dir(self.cfg.get("obsidian_vault_path", ""))
        except Exception:
            vault_root = None

        if vault_root is None:
            messagebox.showerror("Malware Intake", "Could not resolve CIPHER vault root.", parent=self)
            return

        def _after_malware(malware: dict):
            try:
                note_path = create_obsidian_malware_note(vault_root, malware)
                malware["obsidian_note_path"] = str(note_path)
                _log_shift_entry_from_launcher(self, vault_root, "malware", malware, note_path)
            except Exception as e:
                messagebox.showerror("Malware Intake", f"Failed to create Malware note.\n\nError: {e}", parent=self)
                return

            try:
                if malware.get("open_after_save"):
                    open_obsidian_file(vault_root, Path(malware["obsidian_note_path"]))
            except Exception:
                pass

        MalwareIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_malware,
        )

    def _open_ioc_intake(self):
        app_dir = get_app_dir()

        vault_root = None
        try:
            paths = self.cfg.get("paths", {}) or {}

            def _try_dir(candidate: str):
                nonlocal vault_root
                cand = (candidate or "").strip()
                if not cand:
                    return
                try:
                    cand = resolve_path(cand, self.tokens)
                except Exception:
                    cand = expand_value(cand, self.tokens)
                cand = os.path.normpath(cand)
                if os.path.isdir(cand):
                    vault_root = Path(cand)

            _try_dir(paths.get("CIPHER_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("SCOUT_ROOT", ""))
            if vault_root is None:
                _try_dir(paths.get("VAULT_ROOT", ""))
            if vault_root is None:
                _try_dir(self.cfg.get("obsidian_vault_path", ""))
        except Exception:
            vault_root = None

        if vault_root is None:
            messagebox.showerror("IOC Intake", "Could not resolve CIPHER vault root.", parent=self)
            return

        def _after_ioc(ioc: dict):
            try:
                note_path = create_obsidian_ioc_note(vault_root, ioc)
                ioc["obsidian_note_path"] = str(note_path)
                _log_shift_entry_from_launcher(self, vault_root, "ioc", ioc, note_path)
            except Exception as e:
                messagebox.showerror("IOC Intake", f"Failed to create IOC note.\n\nError: {e}", parent=self)
                return

            try:
                if ioc.get("open_after_save"):
                    open_obsidian_file(vault_root, Path(ioc["obsidian_note_path"]))
            except Exception:
                pass

        IOCIntakeWindow(
            self,
            app_dir=app_dir,
            vault_root=vault_root,
            on_submit=_after_ioc,
        )




    # -----------------------------
# Incident draft persistence
# -----------------------------
def _incident_store_path(app_dir: Path) -> Path:
    """
    Stores incident drafts in the app directory so it persists across launches.
    """
    return app_dir / "incident_draft.json"

def open_daily_note_live_with_intake(cfg: dict, tokens: dict, intake_block: str = "", tags_csv: str = "", yaml_updates: dict | None = None) -> None:
    """
    Creates/opens today's daily note (live mode) and optionally injects an intake section.
    Intake is inserted after YAML frontmatter (if present) or near the top otherwise.
    """
    uri = None
    created = False
    try:
        vault_name, scout_root, file_within_vault, abs_note_path, template_abs, create_if_missing, apply_when = _daily_note_compute(cfg, tokens)

        created = _ensure_daily_note_exists(abs_note_path, template_abs, apply_when, create_if_missing)

        uri = (
            "obsidian://advanced-uri?"
            f"vault={quote(vault_name, safe='')}"
            f"&filepath={quote(file_within_vault if file_within_vault.lower().endswith('.md') else file_within_vault + '.md', safe='')}"
            f"&viewmode=live"
        )

        # Optionally merge tags / YAML updates into frontmatter
        if tags_csv or yaml_updates:
            try:
                content = Path(abs_note_path).read_text(encoding="utf-8", errors="ignore")
            except Exception:
                content = ""

            def _ensure_frontmatter(text: str) -> str:
                s = text.lstrip("\ufeff")
                if s.startswith("---"):
                    return text
                # create minimal frontmatter
                return "---\n---\n\n" + text

            # Ensure YAML exists if we need to update keys
            if yaml_updates:
                content = _ensure_frontmatter(content)
            # Apply YAML field updates from the intake form
            try:
                if yaml_updates:

                    # Daily note: set created/updated dates (YYYY-MM-DD)
                    try:
                        _today = datetime.now().strftime("%Y-%m-%d")
                        yaml_updates["updated"] = _today
                        if not (yaml_updates.get("created") or "").strip():
                            yaml_updates["created"] = _today
                    except Exception:
                        pass

                    content = _apply_yaml_updates(content, yaml_updates)

                    # Replace the picture placeholder with an inline embed (vault-relative path)
                    try:
                        pic_val = (yaml_updates.get("picture_of_the_day") or "").strip()
                        if pic_val:
                            if pic_val.startswith("![[") and pic_val.endswith("]]"):
                                embed = pic_val
                            elif pic_val.startswith("[[") and pic_val.endswith("]]"):
                                embed = "!" + pic_val
                            else:
                                embed = f"![[{pic_val}]]"

                            for token in (
                                "![]({{picture_of_the_day}})",
                                "![]({{ picture_of_the_day }})",
                                "!{{picture_of_the_day}}",
                                "!{{ picture_of_the_day }}",
                                "{{picture_of_the_day}}",
                                "{{ picture_of_the_day }}",
                            ):
                                content = content.replace(token, embed)
                    except Exception:
                        pass

                    # Ensure the H1 title includes the current date
                    try:
                        title_date = (yaml_updates.get("journal_date") or "").strip()
                        if not title_date:
                            title_date = datetime.now().strftime("%Y-%m-%d")
                        lines = content.splitlines()
                        for i, line in enumerate(lines):
                            if line.startswith("# Daily Journal"):
                                if "{{" in line and "}}" in line:
                                    line = re.sub(r"\{\{\s*created\s*\}\}", title_date, line)
                                    line = re.sub(r"\{\{\s*journal_date\s*\}\}", title_date, line)
                                if "—" in line:
                                    head, _ = line.split("—", 1)
                                    line = f"{head.strip()} — {title_date}"
                                else:
                                    line = f"# Daily Journal — {title_date}"
                                lines[i] = line
                                break
                        content = "\n".join(lines)
                    except Exception:
                        pass
            except Exception:
                pass


            tags = []
            if isinstance(tags_csv, str) and tags_csv.strip():
                # accept comma-separated or space-separated hashtags
                raw = tags_csv.replace(";", ",")
                parts = []
                for chunk in raw.split(","):
                    chunk = chunk.strip()
                    if not chunk:
                        continue
                    parts.extend(chunk.split())
                for t in parts:
                    t = t.strip()
                    if not t:
                        continue
                    if not t.startswith("#"):
                        t = "#" + t
                    tags.append(t)

            def _merge_tags_into_frontmatter(text: str, new_tags: list[str]) -> str:
                if not new_tags:
                    return text
                s = text.lstrip("\ufeff")
                if not s.startswith("---"):
                    return text
                lines = s.splitlines()
                # find closing ---
                end_idx = None
                for i in range(1, len(lines)):
                    if lines[i].strip() == "---":
                        end_idx = i
                        break
                if end_idx is None:
                    return text
                fm = lines[: end_idx + 1]
                body = "\n".join(lines[end_idx + 1 :]).lstrip("\n")

                # locate tags: key
                tag_line = None
                for i in range(1, end_idx):
                    if lines[i].strip().lower().startswith("tags:"):
                        tag_line = i
                        break

                # collect existing tags (simple YAML list)
                existing = []
                if tag_line is not None:
                    # handle tags: [a, b]
                    rest = lines[tag_line].split(":", 1)[1].strip()
                    if rest.startswith("[") and rest.endswith("]"):
                        inner = rest[1:-1].strip()
                        if inner:
                            for tok in inner.split(","):
                                t = tok.strip().strip('"').strip("'")
                                if t:
                                    existing.append(t)
                        # rewrite as block list below
                        fm[tag_line] = "tags:"
                        insert_at = tag_line + 1
                    else:
                        fm[tag_line] = "tags:"
                        insert_at = tag_line + 1
                        # gather subsequent - items
                        j = tag_line + 1
                        while j < end_idx and lines[j].lstrip().startswith("-"):
                            t = lines[j].lstrip()[1:].strip().strip('"').strip("'")
                            if t:
                                existing.append(t)
                            j += 1
                        # remove existing list items
                        del fm[insert_at:j]
                else:
                    # insert tags: before closing ---
                    insert_at = end_idx
                    fm.insert(insert_at, "tags:")
                    end_idx += 1  # shifted

                merged = []
                seen = set()
                for t in existing + new_tags:
                    if not t:
                        continue
                    if t not in seen:
                        seen.add(t)
                        merged.append(t)
                for k, t in enumerate(merged):
                    fm.insert(insert_at + k, f"  - \"{t}\"")

                return "\n".join(fm) + "\n\n" + body

            if tags:
                content2 = _merge_tags_into_frontmatter(content, tags)
            else:
                content2 = content
            try:
                Path(abs_note_path).write_text(content2, encoding="utf-8")
            except Exception:
                pass
    except Exception as e:
        messagebox.showerror("Daily Notes Error", str(e))
    finally:
        if uri:
            if created:
                opened = open_note_in_obsidian(vault_name, Path(scout_root), Path(abs_note_path))
                if not opened:
                    open_obsidian_uri(uri)
            else:
                open_obsidian_uri(uri)



def load_incident_draft(app_dir: Path) -> dict:
    p = _incident_store_path(app_dir)
    if not p.exists():
        return {}
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return {}

def save_incident_draft(app_dir: Path, data: dict) -> None:
    p = _incident_store_path(app_dir)
    p.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")

def delete_incident_draft(app_dir: Path) -> None:
    p = _incident_store_path(app_dir)
    try:
        if p.exists():
            p.unlink()
    except Exception:
        pass


# -----------------------------
# Optional: Obsidian stub creation
# -----------------------------
def _safe_slug(s: str) -> str:
    # Conservative filename sanitizer
    keep = []
    for ch in (s or "").strip():
        if ch.isalnum() or ch in (" ", "-", "_"):
            keep.append(ch)
    out = "".join(keep).strip().replace("  ", " ")
    return out[:120] if out else "Untitled"


def _normalize_tags(tags_csv: str) -> list:
    """Normalize user tag input into a YAML-friendly list of tags."""
    raw = (tags_csv or "").strip()
    if not raw:
        return []
    # Split on commas and whitespace, preserve hashtags
    parts = []
    for chunk in raw.replace("\n", " ").split(","):
        chunk = chunk.strip()
        if not chunk:
            continue
        parts.extend([p.strip() for p in chunk.split() if p.strip()])

    out = []
    seen = set()
    for t in parts:
        # allow users to type without leading '#'
        if not t.startswith("#"):
            t = "#" + t
        # collapse duplicates (case-insensitive)
        key = t.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(t)
    return out


def _read_text_safe(p: Path) -> str:
    try:
        return p.read_text(encoding="utf-8")
    except Exception:
        # Some templates may be authored with UTF-8 BOM or other encodings
        try:
            return p.read_text(encoding="utf-8-sig")
        except Exception:
            return p.read_text(errors="replace")


def _extract_note_title(md_text: str, fallback: str) -> str:
    """
    Best-effort title extraction for Obsidian markdown notes:
      1) YAML frontmatter key `title:` (flat, single-line)
      2) First H1 heading: '# ...'
      3) Filename fallback
    """
    # 1) YAML frontmatter title
    if md_text.lstrip().startswith("---"):
        # frontmatter block is between first two '---' lines
        lines = md_text.splitlines()
        if len(lines) > 2:
            # find end delimiter
            end_idx = None
            for i in range(1, len(lines)):
                if lines[i].strip() == "---":
                    end_idx = i
                    break
            if end_idx:
                for ln in lines[1:end_idx]:
                    m = re.match(r'^\s*title\s*:\s*(.+?)\s*$', ln)
                    if m:
                        val = m.group(1).strip()
                        # strip surrounding quotes if present
                        if (val.startswith('"') and val.endswith('"')) or (val.startswith("'") and val.endswith("'")):
                            val = val[1:-1].strip()
                        if val:
                            return val

    # 2) First H1
    for ln in md_text.splitlines():
        s = ln.strip()
        if s.startswith("# "):
            t = s[2:].strip()
            if t:
                return t

    return fallback


def list_itid_definition_titles(vault_root: Path) -> list[str]:
    """
    Returns a sorted list of ITID definition entries under:
      40_ITIDs/01_Definitions

    Output labels are designed to be both human-friendly and complete:
    - Each note contributes its filename stem (often the ITID id, e.g., ITID-11112).
    - If an extracted title differs from the stem, the label becomes: "<STEM> — <Title>".
      This prevents "missing" ITIDs when the note's H1/frontmatter title is not the ID.
    - De-duplicates files by real path to avoid double-counting via symlinks/junctions.
    - If multiple notes share the same stem (rare), disambiguates with relative path.
    """
    base = vault_root / "40_ITIDs" / "01_Definitions"
    if not base.exists():
        return []

    seen_files: set[str] = set()
    items: list[tuple[str, str, str]] = []  # (stem, title, relative_path)

    for p in sorted(base.rglob("*")):
        if not p.is_file():
            continue

        suf = p.suffix.lower()
        if suf not in (".md", ".markdown"):
            continue

        # De-dupe by real path (handles symlinks/junctions)
        try:
            real = str(p.resolve())
        except Exception:
            real = str(p.absolute())
        real_key = real.lower()
        if real_key in seen_files:
            continue
        seen_files.add(real_key)

        stem = p.stem.strip()
        rel = str(p.relative_to(base)).replace("\\", "/").strip()

        try:
            md = _read_text_safe(p)
        except Exception:
            md = ""

        title = _extract_note_title(md, stem).strip()

        if stem and rel:
            items.append((stem, title, rel))

    if not items:
        return []

    # Group by stem (so every file's identifier is preserved)
    groups: dict[str, list[tuple[str, str, str]]] = {}
    for stem, title, rel in items:
        key = stem.lower()
        groups.setdefault(key, []).append((stem, title, rel))

    out: list[str] = []
    seen_out: set[str] = set()

    for key, group in groups.items():
        # If a stem appears multiple times, disambiguate by relative path
        group_sorted = sorted(group, key=lambda x: x[2].lower())
        for stem, title, rel in group_sorted:
            label = stem
            if title and title.strip() and title.strip().lower() != stem.lower():
                label = f"{stem} — {title.strip()}"
            if len(group_sorted) > 1:
                label = f"{label} ({rel})"

            lk = label.lower()
            if lk in seen_out:
                continue
            seen_out.add(lk)
            out.append(label)

    return sorted(out, key=lambda s: s.lower())



def _resolve_investigation_template(vault_root: Path) -> Path | None:
    """
    Resolves the user's Incident template path inside the vault.
    Primary: 00_System/00_Templates/02_Operations/Incident.md Template.md
    Fallback: first matching file in the same folder containing 'Incident' and 'Template'.
    """
    folder = vault_root / "00_System" / "00_Templates" / "02_Operations"
    primary = folder / "Incident.md Template.md"
    if primary.exists():
        return primary

    # Fallback search (case-insensitive) in the expected folder
    if folder.exists():
        for p in sorted(folder.glob("*.md")):
            name = p.name.lower()
            if "incident" in name and "template" in name:
                return p

    return None


def _incident_chronos_marker(created: str) -> str:
    """Build Chronos marker line for incident created date. Format: = [YYYY-MM-DDThh:mm:ss] Incident Created"""
    created = (created or "").strip()
    if not created:
        created = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    # Chronos expects YYYY-MM-DDThh:mm:ss (replace space with T)
    chronos_date = created.replace(" ", "T", 1)
    return f"= [{chronos_date}] Incident Created"


def _apply_template_tokens(template_text: str, intake: dict) -> str:
    """
    Replaces common placeholder token forms in templates.
    Supported patterns include:
      {{incident_id}}, {{INCIDENT_ID}}, ${incident_id}, [[incident_id]]
    """
    # Normalize values
    created = intake.get("created", "").strip()
    tokens = {
        "incident_id": (intake.get("incident_id", "").strip() or "INC-UNASSIGNED"),
        "title": (intake.get("title", "").strip() or "Untitled Incident"),
        "status": intake.get("status", "").strip(),
        "severity": intake.get("severity", "").strip(),
        "tlp": intake.get("tlp", "").strip(),
        "discovered": intake.get("discovered", "").strip(),
        "created": created,
        "source": intake.get("source", "").strip(),
        "reporter": intake.get("reporter", "").strip(),
        "tags": intake.get("tags", "").strip(),
        "affected_assets": intake.get("affected_assets", "").strip(),
        "summary": intake.get("summary", "").strip(),
        "notes": intake.get("notes", "").strip(),
        "incident_chronos_marker": _incident_chronos_marker(created),
    }

    out = template_text
    for k, v in tokens.items():
        # Common template conventions
        out = out.replace(f"{{{{{k}}}}}", v)
        out = out.replace(f"{{{{{k.upper()}}}}}", v)
        out = out.replace(f"${{{k}}}", v)
        out = out.replace(f"[[{k}]]", v)
        out = out.replace(f"[[{k.upper()}]]", v)

    return out


INCIDENT_ID_RE = re.compile(r"INC-(\d{4})-(\d{1,6})", re.IGNORECASE)
PROJECT_ID_RE = re.compile(r"PRJ-(\d{4})-(\d{1,4})", re.IGNORECASE)
GOAL_ID_RE = re.compile(r"GOAL-(\d{4})-(\d{1,4})", re.IGNORECASE)


def _normalize_incident_id(raw: str, *, default_year: int | None = None) -> str:
    val = (raw or "").strip().upper()
    if not val:
        return ""

    m = re.search(r"INC-(\d{4})-(\d{1,6})", val, flags=re.IGNORECASE)
    if m:
        year = int(m.group(1))
        seq = int(m.group(2))
        return f"INC-{year}-{seq:06d}"

    m = re.search(r"(\d{4})-(\d{1,6})", val)
    if m:
        year = int(m.group(1))
        seq = int(m.group(2))
        return f"INC-{year}-{seq:06d}"

    m = re.search(r"(\d{1,6})", val)
    if m:
        seq = int(m.group(1))
        year = default_year or datetime.now().year
        return f"INC-{year}-{seq:06d}"

    return val


def _find_highest_incident_id(incidents_root: Path) -> tuple[int, int] | None:
    if incidents_root is None or not incidents_root.exists():
        return None

    best: tuple[int, int] | None = None
    try:
        for p in incidents_root.rglob("*"):
            name = p.name
            for m in INCIDENT_ID_RE.finditer(name):
                year = int(m.group(1))
                seq = int(m.group(2))
                if best is None or (year, seq) > best:
                    best = (year, seq)
    except Exception:
        return None

    return best


def _next_incident_id_from_vault(vault_root: Path | None, *, default_year: int | None = None) -> str:
    year = default_year or datetime.now().year
    if vault_root is None:
        return f"INC-{year}-000001"

    incidents_root = vault_root / "10_Operations" / "04_Incidents"
    best = _find_highest_incident_id(incidents_root)
    if best is None:
        return f"INC-{year}-000001"

    best_year, best_seq = best
    return f"INC-{best_year}-{best_seq + 1:06d}"


def _normalize_project_id(raw: str, *, default_year: int | None = None) -> str:
    val = (raw or "").strip().upper()
    if not val:
        return ""

    m = re.search(r"PRJ-(\d{4})-(\d{1,4})", val, flags=re.IGNORECASE)
    if m:
        year = int(m.group(1))
        seq = int(m.group(2))
        return f"PRJ-{year}-{seq:04d}"

    m = re.search(r"(\d{4})-(\d{1,4})", val)
    if m:
        year = int(m.group(1))
        seq = int(m.group(2))
        return f"PRJ-{year}-{seq:04d}"

    m = re.search(r"(\d{1,4})", val)
    if m:
        seq = int(m.group(1))
        year = default_year or datetime.now().year
        return f"PRJ-{year}-{seq:04d}"

    return val


def _find_highest_project_id(projects_root: Path) -> tuple[int, int] | None:
    if projects_root is None or not projects_root.exists():
        return None

    best: tuple[int, int] | None = None
    try:
        for p in projects_root.rglob("*"):
            name = p.name
            for m in PROJECT_ID_RE.finditer(name):
                year = int(m.group(1))
                seq = int(m.group(2))
                if best is None or (year, seq) > best:
                    best = (year, seq)
    except Exception:
        return None

    return best


def _next_project_id_from_vault(vault_root: Path | None, *, default_year: int | None = None) -> str:
    year = default_year or datetime.now().year
    if vault_root is None:
        return f"PRJ-{year}-0001"

    projects_root = vault_root / "10_Operations" / "14_Projects"
    best = _find_highest_project_id(projects_root)
    if best is None:
        return f"PRJ-{year}-0001"

    best_year, best_seq = best
    return f"PRJ-{best_year}-{best_seq + 1:04d}"


def _normalize_goal_id(raw: str, *, default_year: int | None = None) -> str:
    val = (raw or "").strip().upper()
    if not val:
        return ""

    m = re.search(r"GOAL-(\d{4})-(\d{1,4})", val, flags=re.IGNORECASE)
    if m:
        year = int(m.group(1))
        seq = int(m.group(2))
        return f"GOAL-{year}-{seq:04d}"

    m = re.search(r"(\d{4})-(\d{1,4})", val)
    if m:
        year = int(m.group(1))
        seq = int(m.group(2))
        return f"GOAL-{year}-{seq:04d}"

    m = re.search(r"(\d{1,4})", val)
    if m:
        seq = int(m.group(1))
        year = default_year or datetime.now().year
        return f"GOAL-{year}-{seq:04d}"

    return val


def _find_highest_goal_id(goals_root: Path) -> tuple[int, int] | None:
    if goals_root is None or not goals_root.exists():
        return None

    best: tuple[int, int] | None = None
    try:
        for p in goals_root.rglob("*"):
            name = p.name
            for m in GOAL_ID_RE.finditer(name):
                year = int(m.group(1))
                seq = int(m.group(2))
                if best is None or (year, seq) > best:
                    best = (year, seq)
    except Exception:
        return None

    return best


def _next_goal_id_from_vault(vault_root: Path | None, *, default_year: int | None = None) -> str:
    year = default_year or datetime.now().year
    if vault_root is None:
        return f"GOAL-{year}-0001"

    goals_root = vault_root / "10_Operations" / "15_Goals"
    best = _find_highest_goal_id(goals_root)
    if best is None:
        return f"GOAL-{year}-0001"

    best_year, best_seq = best
    return f"GOAL-{best_year}-{best_seq + 1:04d}"


def _build_incident_frontmatter(intake: dict) -> str:
    inc_id = intake.get("incident_id", "").strip() or "INC-UNASSIGNED"
    title = intake.get("title", "").strip() or "Untitled Incident"
    created = intake.get("created", "").strip() or datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    discovered = intake.get("discovered", "").strip() or created

    # Flat YAML (Obsidian-friendly)
    return f"""---
entity_type: incident
incident_id: "{inc_id}"
title: "{title}"
incident_type: "{intake.get('itid_definition', '').strip()}"
status: "{intake.get('status', '').strip()}"
severity: "{intake.get('severity', '').strip()}"
detected_time: "{discovered}"
created: "{created}"
created_time: "{created}"
tlp_classification: "{intake.get('tlp', '').strip()}"
alert_sources: "{intake.get('source', '').strip()}"
primary_user: "{intake.get('reporter', '').strip()}"
tags: "{intake.get('tags', '').strip()}"
affected_assets: "{intake.get('affected_assets', '').strip()}"
impact_summary: "{intake.get('summary', '').strip()}"
updated: "{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
---
"""

def _merge_template_frontmatter(template_text: str, intake: dict) -> str:
    """
    If template_text begins with YAML frontmatter, merge/overwrite selected keys from intake.
        - title is populated from intake title
        Keeps the rest of the template frontmatter intact.
    """
    s = template_text.lstrip("\ufeff")  # handle BOM
    # Require frontmatter at very start (after optional whitespace/newlines)
    # We'll operate on the original string positions by finding the first '---' after leading whitespace.
    lead_len = len(template_text) - len(template_text.lstrip())
    if not template_text.lstrip().startswith("---"):
        return template_text

    # Find frontmatter delimiters line-wise
    lines = template_text.splitlines(True)  # keepends
    # Identify start: first non-empty/whitespace line must be ---
    i0 = None
    for i, ln in enumerate(lines):
        if ln.strip() == "":
            continue
        if ln.strip() == "---":
            i0 = i
        break
    if i0 is None:
        return template_text

    # Find end delimiter
    i1 = None
    for j in range(i0 + 1, len(lines)):
        if lines[j].strip() == "---":
            i1 = j
            break
    if i1 is None:
        return template_text

    fm_lines = lines[i0 + 1 : i1]
    fm_map = {}
    key_line_index = {}

    key_re = re.compile(r"^([A-Za-z0-9_-]+)\s*:\s*(.*)$")
    for idx, ln in enumerate(fm_lines):
        raw = ln.rstrip("\n")
        if not raw.strip() or raw.lstrip().startswith("#"):
            continue
        m = key_re.match(raw)
        if m:
            k = m.group(1)
            fm_map[k] = m.group(2)
            key_line_index[k] = idx

    # Values to apply (quoted as strings for safety)
    incident_id_val = intake.get("incident_id", "").strip() or "INC-UNASSIGNED"
    title_val = intake.get("title", "").strip() or "Untitled Incident"

    updates = {
        # New field mappings (Python Intake -> YAML)
        "incident_id": f"\"{incident_id_val}\"",
        "title": f"\"{title_val}\"",
        "incident_type": f"\"{intake.get('itid_definition', '').strip()}\"",
        "status": f"\"{intake.get('status', '').strip()}\"",
        "severity": f"\"{intake.get('severity', '').strip()}\"",
        "detected_time": f"\"{(intake.get('discovered', '').strip() or (intake.get('created', '').strip() or datetime.now().strftime('%Y-%m-%d %H:%M:%S')))}\"",
        "created": f"\"{(intake.get('created', '').strip() or datetime.now().strftime('%Y-%m-%d %H:%M:%S'))}\"",
        "created_time": f"\"{(intake.get('created', '').strip() or datetime.now().strftime('%Y-%m-%d %H:%M:%S'))}\"",
        "tlp_classification": f"\"{intake.get('tlp', '').strip()}\"",
        "alert_sources": f"\"{intake.get('source', '').strip()}\"",
        "primary_user": f"\"{intake.get('reporter', '').strip()}\"",
        "tags": f"\"{intake.get('tags', '').strip()}\"",
        "affected_assets": f"\"{intake.get('affected_assets', '').strip()}\"",
        "impact_summary": f"\"{intake.get('summary', '').strip()}\"",
        "updated": f"\"{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\"",
    }

    def set_key(k: str, v: str):
        if k in key_line_index:
            fm_lines[key_line_index[k]] = f"{k}: {v}\n"
        else:
            # Insert near the top (after entity_type if it exists), else at top
            insert_at = 0
            if "entity_type" in key_line_index:
                insert_at = key_line_index["entity_type"] + 1
            fm_lines.insert(insert_at, f"{k}: {v}\n")

    for k, v in updates.items():
        set_key(k, v)

    # Rebuild full document
    new_lines = []
    new_lines.extend(lines[: i0 + 1])
    new_lines.extend(fm_lines)
    new_lines.append("---\n")
    new_lines.extend(lines[i1 + 1 :])

    return "".join(new_lines)



def create_obsidian_incident_stub(vault_root: Path, intake: dict) -> Path:
    """
    Creates a preliminary incident note in the incidents folder (10_Operations/04_Incidents) using the Incident template.
    If the template is missing frontmatter, frontmatter is prepended.
    If the template includes placeholders, they are replaced from the intake payload.
    """
    base_folder = vault_root / "10_Operations" / "04_Incidents"
    base_folder.mkdir(parents=True, exist_ok=True)

    inc_id = intake.get("incident_id", "").strip() or "INC-UNASSIGNED"
    title = intake.get("title", "").strip() or "Untitled Incident"
    slug = _safe_slug(f"{inc_id} - {title}")
    incident_dir = base_folder / slug
    incident_dir.mkdir(parents=True, exist_ok=True)
    (incident_dir / "Evidence").mkdir(parents=True, exist_ok=True)
    note_path = incident_dir / f"{slug}.md"

    tpl_path = _resolve_investigation_template(vault_root)

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)
        tpl_text = _apply_template_tokens(tpl_text, intake)

        # If template does not start with YAML frontmatter, prepend it.
        if not tpl_text.lstrip().startswith("---"):
            md = _build_incident_frontmatter(intake) + "\n" + tpl_text.lstrip()
        else:
            # Merge intake values into the template YAML frontmatter (flat keys).
            md = _merge_template_frontmatter(tpl_text, intake)
    else:
        # Fallback to a minimal stub if the template cannot be found.
        md = _build_incident_frontmatter(intake) + f"""
## Summary
{intake.get('summary', '').strip()}

## Initial Notes
{intake.get('notes', '').strip()}

## Next Actions
- [ ] Triage
- [ ] Scope
- [ ] Containment
- [ ] Eradication
- [ ] Recovery
- [ ] Lessons Learned
"""

    note_path.write_text(md, encoding="utf-8")
    return note_path


def _append_daily_shift_notes_incident(vault_root: Path, intake: dict) -> None:
    """
    Append a Shift Notes event block to the daily journal note for the incident's created date.
    Uses the SHIFT_NOTES_LOG marker when available; falls back to Shift Notes header.
    """
    try:
        created_raw = (intake or {}).get("created", "").strip()
        if created_raw:
            parts = created_raw.split()
            date_str = parts[0]
            time_str = parts[1][:5] if len(parts) > 1 else datetime.now().strftime("%H:%M")
        else:
            date_str = datetime.now().strftime("%Y-%m-%d")
            time_str = datetime.now().strftime("%H:%M")

        year = date_str.split("-")[0]
        month = date_str.split("-")[1]
        daily_path = (
            Path(vault_root)
            / "10_Operations"
            / "09_Journals"
            / "01_Daily"
            / year
            / month
            / f"{date_str}.md"
        )
        if not daily_path.exists():
            return

        incident_id = (intake or {}).get("incident_id", "").strip()
        title = (intake or {}).get("title", "").strip() or "Untitled Incident"
        header_title = " - ".join([p for p in [incident_id, title] if p])
        header_line = f"### {time_str} — {header_title}"

        block = "\n".join(
            [
                header_line,
                "",
                "Context:",
                "",
                "- What triggered this work?",
                "- Why did you look into it?",
                "",
                "Actions:",
                "",
                "- What did you do?",
                "- Queries executed?",
                "- Tools used?",
                "",
                "Findings:",
                "",
                "- What did you learn?",
                "- Any anomalies or notable patterns?",
                "",
                "Links:",
                "- [Alert link]",
                "- [Incident link]",
                "- [Investigation link]",
                "- [Observation link]",
                "",
            ]
        )

        marker = "<!-- SHIFT_NOTES_LOG -->"
        shift_heading = "## 2. Shift Notes (Chronological Log)"
        content = daily_path.read_text(encoding="utf-8")
        updated = content

        if marker in content:
            marker_idx = content.index(marker)
            section_start = marker_idx + len(marker)
            next_heading_idx = content.find("\n## ", section_start)
            if next_heading_idx == -1:
                next_heading_idx = len(content)
            section_body = content[section_start:next_heading_idx].rstrip("\n")
            insert = f"{section_body}\n{block}\n" if section_body.strip() else f"\n{block}\n"
            updated = content[:section_start] + insert + content[next_heading_idx:]
        elif shift_heading in content:
            idx = content.index(shift_heading)
            next_heading_idx = content.find("\n## ", idx + len(shift_heading))
            if next_heading_idx == -1:
                next_heading_idx = len(content)
            section_body = content[idx + len(shift_heading):next_heading_idx].rstrip("\n")
            insert = f"{section_body}\n{block}\n" if section_body.strip() else f"\n{block}\n"
            updated = content[: idx + len(shift_heading)] + insert + content[next_heading_idx:]
        else:
            updated = content + f"\n\n{block}\n"

        if updated != content:
            daily_path.write_text(updated, encoding="utf-8")
    except Exception:
        return


def _append_daily_shift_notes_meeting(vault_root: Path, meeting: dict) -> None:
    """
    Append a Shift Notes event line to the daily journal note for the meeting date.
    Uses the SHIFT_NOTES_LOG marker when available; falls back to Shift Notes header.
    Format: Timestamp: Meeting Title - Meeting Start Date - Meeting End Date
    """
    try:
        start_raw = (meeting or {}).get("start_time", "").strip()
        end_raw = (meeting or {}).get("end_time", "").strip()
        created_time = datetime.now().strftime("%H:%M")
        if start_raw:
            parts = start_raw.split()
            date_str = parts[0]
            start_datetime = f"{parts[0]} {parts[1][:5]}" if len(parts) > 1 else f"{parts[0]} 00:00"
        else:
            date_str = datetime.now().strftime("%Y-%m-%d")
            start_datetime = f"{date_str} 00:00"
        if end_raw:
            end_parts = end_raw.split()
            if len(end_parts) >= 2:
                end_datetime = f"{end_parts[0]} {end_parts[1][:5]}"
            elif end_parts:
                end_datetime = f"{date_str} {end_parts[0][:5]}"
            else:
                end_datetime = start_datetime
        else:
            end_datetime = start_datetime

        year = date_str.split("-")[0]
        month = date_str.split("-")[1]
        daily_path = (
            Path(vault_root)
            / "10_Operations"
            / "09_Journals"
            / "01_Daily"
            / year
            / month
            / f"{date_str}.md"
        )
        if not daily_path.exists():
            return

        title = (meeting or {}).get("title", "").strip() or "Untitled Meeting"
        header_line = f"### {created_time} - {title} - {start_datetime} - {end_datetime}"

        marker = "<!-- SHIFT_NOTES_LOG -->"
        shift_heading = "## 2. Shift Notes (Chronological Log)"
        content = daily_path.read_text(encoding="utf-8")
        updated = content

        if marker in content:
            marker_idx = content.index(marker)
            section_start = marker_idx + len(marker)
            next_heading_idx = content.find("\n## ", section_start)
            if next_heading_idx == -1:
                next_heading_idx = len(content)
            section_body = content[section_start:next_heading_idx].rstrip("\n")
            insert = f"{section_body}\n{header_line}\n" if section_body.strip() else f"\n{header_line}\n"
            updated = content[:section_start] + insert + content[next_heading_idx:]
        elif shift_heading in content:
            idx = content.index(shift_heading)
            next_heading_idx = content.find("\n## ", idx + len(shift_heading))
            if next_heading_idx == -1:
                next_heading_idx = len(content)
            section_body = content[idx + len(shift_heading):next_heading_idx].rstrip("\n")
            insert = f"{section_body}\n{header_line}\n" if section_body.strip() else f"\n{header_line}\n"
            updated = content[: idx + len(shift_heading)] + insert + content[next_heading_idx:]
        else:
            updated = content + f"\n\n{header_line}\n"

        if updated != content:
            daily_path.write_text(updated, encoding="utf-8")
    except Exception:
        return


# -----------------------------
# Incident Intake UI
# -----------------------------
class IncidentIntakeWindow(tk.Toplevel):
    """Incident Intake UI.

    Key UX goals:
      - Scrollable content area so fields never fall off-screen
      - Fixed bottom action bar so buttons are always visible
      - Draft persistence (Save Draft / Clear Draft)
      - Optional Obsidian stub creation + open-note toggle
    """

    def __init__(self, master, *, app_dir: Path, vault_root: Path | None, on_submit, title="Incident Intake"):
        super().__init__(master)
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("New Incident Intake")
        self.configure(bg="#111111")
        self.geometry("980x720")
        self.minsize(860, 640)
        self.resizable(True, True)

        self._build_styles()
        self._build_widgets()

        # Load persisted draft (if any)
        draft = load_incident_draft(self.app_dir)
        self._apply_draft(draft)

        # Default timestamps if empty
        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not self.vars["created"].get().strip():
            self.vars["created"].set(now)
        if not self.vars["discovered"].get().strip():
            self.vars["discovered"].set(now)

        self.transient(master)
        self.grab_set()
        self.focus_force()

    def _build_styles(self):
        style = ttk.Style(self)
        try:
            style.theme_use("clam")
        except Exception:
            pass

    # -----------------------------
    # Scroll helpers
    # -----------------------------
    def _on_canvas_configure(self, _event=None):
        try:
            self._canvas.configure(scrollregion=self._canvas.bbox("all"))
        except Exception:
            pass

    def _on_frame_configure(self, _event=None):
        self._on_canvas_configure()

    def _on_mousewheel(self, event):
        # Windows/macOS wheel delta normalization
        try:
            if event.delta:
                self._canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")
        except Exception:
            pass

    def _bind_mousewheel(self, widget):
        # Bind wheel events to all descendants
        try:
            widget.bind_all("<MouseWheel>", self._on_mousewheel)
        except Exception:
            pass

    # -----------------------------
    # UI build
    # -----------------------------
    def _build_widgets(self):
        cfg = getattr(self.master, "cfg", {}) or {}
        analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")

        self.vars = {
            "incident_id": tk.StringVar(),
            "title": tk.StringVar(),
            "status": tk.StringVar(value="new"),
            "severity": tk.StringVar(value="low"),
            "tlp": tk.StringVar(value="TLP:GREEN"),
            "discovered": tk.StringVar(),
            "created": tk.StringVar(),
            "source": tk.StringVar(value="ticket"),
            "reporter": tk.StringVar(value=analyst_default),
            "tags": tk.StringVar(),
            "affected_assets": tk.StringVar(),
            "itid_primary": tk.StringVar(),
        }
        self.sequential_mode_var = tk.BooleanVar(value=False)

        # Root layout: scrollable content + fixed bottom bar
        root = tk.Frame(self, bg="#111111")
        root.pack(fill="both", expand=True)

        self._canvas = tk.Canvas(root, bg="#111111", highlightthickness=0)
        vsb = tk.Scrollbar(root, orient="vertical", command=self._canvas.yview)
        self._canvas.configure(yscrollcommand=vsb.set)

        vsb.pack(side="right", fill="y")
        self._canvas.pack(side="left", fill="both", expand=True)

        self._content = tk.Frame(self._canvas, bg="#111111")
        self._content_id = self._canvas.create_window((0, 0), window=self._content, anchor="nw")

        # Ensure content frame tracks canvas width
        def _sync_width(event):
            try:
                self._canvas.itemconfigure(self._content_id, width=event.width)
            except Exception:
                pass

        self._canvas.bind("<Configure>", _sync_width)
        self._content.bind("<Configure>", self._on_frame_configure)
        self._bind_mousewheel(self)

        # ----- Header -----
        top = tk.Frame(self._content, bg="#111111")
        top.pack(fill="x", padx=16, pady=(14, 6))

        tk.Label(
            top,
            text="New Incident Intake",
            font=("Segoe UI", 16, "bold"),
            fg="#FFFFFF",
            bg="#111111",
        ).pack(anchor="w")

        tk.Label(
            top,
            text="Capture preliminary details before entering the incident workflow.",
            font=("Segoe UI", 10),
            fg="#BBBBBB",
            bg="#111111",
        ).pack(anchor="w", pady=(2, 0))

        # ----- Two-column form -----
        form = tk.Frame(self._content, bg="#111111")
        form.pack(fill="x", padx=16, pady=(10, 6))

        form.grid_columnconfigure(0, weight=1)
        form.grid_columnconfigure(1, weight=1)

        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.grid(row=0, column=0, sticky="nsew", padx=(0, 10))
        right.grid(row=0, column=1, sticky="nsew", padx=(10, 0))

        def add_labeled_entry(parent, label, var, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            lbl = tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10))
            lbl.pack(side="left", anchor="w")
            if required:
                tk.Label(row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left", anchor="w")
            ent = tk.Entry(
                row,
                textvariable=var,
                bg="#1B1B1B",
                fg="#FFFFFF",
                insertbackground="#FFFFFF",
                relief="flat",
            )
            ent.pack(fill="x", ipady=6)
            return ent

        def add_labeled_combo(parent, label, var, values):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=2)
            return cb

        # Left: identity + triage
        incident_id_entry = add_labeled_entry(left, "Incident ID (optional)", self.vars["incident_id"])
        incident_id_entry.bind("<FocusOut>", self._on_incident_id_focus_out)

        seq_row = tk.Frame(left, bg="#111111")
        seq_row.pack(fill="x", pady=(2, 6))
        tk.Checkbutton(
            seq_row,
            text="Sequential mode",
            variable=self.sequential_mode_var,
            command=self._on_sequential_toggle,
            bg="#111111",
            fg="#DDDDDD",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#DDDDDD",
        ).pack(anchor="w")
        add_labeled_entry(left, "Title", self.vars["title"], required=True)
        add_labeled_combo(left, "Severity", self.vars["severity"], ["low", "medium", "high", "critical"])
        add_labeled_combo(left, "TLP", self.vars["tlp"], ["TLP:CLEAR", "TLP:GREEN", "TLP:AMBER", "TLP:AMBER+STRICT", "TLP:RED"])
        add_labeled_entry(left, "Discovered (YYYY-MM-DD HH:MM:SS)", self.vars["discovered"])
        add_labeled_entry(left, "Created (YYYY-MM-DD HH:MM:SS)", self.vars["created"])

        # Right: context
        add_labeled_combo(right, "Source", self.vars["source"], ["system", "ticket", "url"])
        add_labeled_entry(right, "Reporter / Owner", self.vars["reporter"])
        add_labeled_entry(right, "Tags (comma-separated)", self.vars["tags"])
        add_labeled_entry(right, "Affected Assets (brief)", self.vars["affected_assets"])

        # ITID Definition dropdown (loaded from 40_ITIDs/01_Definitions)
        itid_values = []
        if self.vault_root is not None:
            try:
                itid_values = list_itid_definition_titles(self.vault_root)
            except Exception:
                itid_values = []
        add_labeled_combo(right, "ITID Definition", self.vars["itid_primary"], itid_values or ["(No ITID definitions found)"])

        # ----- Narrative fields -----
        narrative = tk.Frame(self._content, bg="#111111")
        narrative.pack(fill="x", padx=16, pady=(6, 10))

        sum_lbl_row = tk.Frame(narrative, bg="#111111")
        sum_lbl_row.pack(anchor="w")
        tk.Label(sum_lbl_row, text="Summary", fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
        tk.Label(sum_lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
        self.summary_txt = tk.Text(
            narrative,
            height=6,
            bg="#1B1B1B",
            fg="#FFFFFF",
            insertbackground="#FFFFFF",
            relief="flat",
            wrap="word",
        )
        self.summary_txt.pack(fill="x", pady=(4, 10))

        # Spacer so last text widget doesn't hide behind fixed bar when scrolled to bottom
        tk.Frame(self._content, bg="#111111", height=16).pack(fill="x")

        # ----- Fixed bottom action bar -----
        self._btns = tk.Frame(self, bg="#111111")
        self._btns.pack(fill="x", padx=16, pady=(6, 14))

        # Obsidian options (hidden, both always enabled)
        self.create_stub_var = tk.BooleanVar(value=True)
        self.open_note_var = tk.BooleanVar(value=True)

        # Right-aligned actions (always visible)
        tk.Button(
            self._btns,
            text="Cancel",
            command=self.destroy,
            bg="#2A2A2A",
            fg="#FFFFFF",
            relief="flat",
            padx=14,
            pady=8,
        ).pack(side="right", padx=(8, 0))

        tk.Button(
            self._btns,
            text="Save & Continue",
            command=self._save_continue,
            bg="#3A3A3A",
            fg="#FFFFFF",
            relief="flat",
            padx=14,
            pady=8,
        ).pack(side="right", padx=(8, 0))

        # Initialize scroll region
        self._on_canvas_configure()

    # -----------------------------
    # Data handling
    # -----------------------------
    def _apply_draft(self, draft: dict):
        for k, v in (draft or {}).items():
            if k in self.vars and isinstance(v, str):
                self.vars[k].set(v)
        if "summary" in (draft or {}):
            self.summary_txt.delete("1.0", "end")
            self.summary_txt.insert("1.0", draft.get("summary", ""))
        # notes field removed from form
        if "sequential_mode" in (draft or {}):
            try:
                self.sequential_mode_var.set(bool(draft.get("sequential_mode")))
            except Exception:
                pass
            if self.sequential_mode_var.get():
                self._apply_sequential_suggestion()

    def _collect(self) -> dict:
        data = {k: v.get().strip() for k, v in self.vars.items()}
        data["status"] = "new"
        data["sequential_mode"] = bool(getattr(self, "sequential_mode_var", tk.BooleanVar(value=False)).get())
        if getattr(self, "sequential_mode_var", None) is not None and self.sequential_mode_var.get():
            year = datetime.now().year
            inc_val = data.get("incident_id", "")
            if inc_val:
                data["incident_id"] = _normalize_incident_id(inc_val, default_year=year)
            else:
                data["incident_id"] = _next_incident_id_from_vault(self.vault_root, default_year=year)
        data["summary"] = self.summary_txt.get("1.0", "end").strip()
        data["notes"] = ""
        return data

    def _apply_sequential_suggestion(self) -> None:
        year = datetime.now().year
        current = self.vars["incident_id"].get().strip()
        if current:
            normalized = _normalize_incident_id(current, default_year=year)
            if normalized != current:
                self.vars["incident_id"].set(normalized)
            return
        suggested = _next_incident_id_from_vault(self.vault_root, default_year=year)
        self.vars["incident_id"].set(suggested)

    def _on_sequential_toggle(self) -> None:
        if getattr(self, "sequential_mode_var", None) is None:
            return
        if self.sequential_mode_var.get():
            self._apply_sequential_suggestion()

    def _on_incident_id_focus_out(self, _event=None) -> None:
        if getattr(self, "sequential_mode_var", None) is None:
            return
        if not self.sequential_mode_var.get():
            return
        year = datetime.now().year
        current = self.vars["incident_id"].get().strip()
        if current:
            self.vars["incident_id"].set(_normalize_incident_id(current, default_year=year))

    def _validate(self, data: dict) -> tuple[bool, str]:
        if not data.get("title"):
            return False, "Title is required."
        if not data.get("summary"):
            return False, "Summary is required."
        return True, ""

    def _save_draft(self):
        data = self._collect()
        save_incident_draft(self.app_dir, data)
        messagebox.showinfo("Incident Intake", "Draft saved.", parent=self)

    def _clear_draft(self):
        delete_incident_draft(self.app_dir)
        # Clear UI fields
        for v in self.vars.values():
            v.set("")
        self.vars["status"].set("new")
        self.vars["severity"].set("medium")
        self.vars["tlp"].set("TLP:AMBER")
        if getattr(self, "sequential_mode_var", None) is not None:
            self.sequential_mode_var.set(False)
        self.summary_txt.delete("1.0", "end")
        messagebox.showinfo("Incident Intake", "Draft cleared.", parent=self)

    def _save_continue(self):
        data = self._collect()
        ok, msg = self._validate(data)
        if not ok:
            messagebox.showwarning("Incident Intake", msg, parent=self)
            return

        # 1) Persist draft
        draft_path = _incident_store_path(self.app_dir)
        try:
            save_incident_draft(self.app_dir, data)
        except Exception as e:
            messagebox.showerror(
                "Incident Intake",
                f"Failed to save draft to disk:\n{draft_path}\n\n{e}",
                parent=self,
            )
            return

        # 2) Optional Obsidian stub creation
        note_path = None
        wants_stub = bool(getattr(self, "create_stub_var", tk.BooleanVar(value=False)).get())
        if wants_stub and self.vault_root is not None:
            try:
                note_path = create_obsidian_incident_stub(self.vault_root, data)
            except Exception as e:
                messagebox.showwarning(
                    "Incident Intake",
                    f"Draft saved, but failed to create Obsidian staging note:\n{e}",
                    parent=self,
                )

        # 3) Append to Daily Note Shift Log (always when vault_root available)
        if self.vault_root is not None:
            try:
                _log_shift_entry_from_launcher(self.master, self.vault_root, "incident", data, note_path)
            except Exception:
                pass

        # 4) Hand off payload
        payload = dict(data)
        payload["draft_path"] = str(draft_path)
        payload["open_note"] = bool(getattr(self, "open_note_var", tk.BooleanVar(value=True)).get())
        if note_path is not None:
            payload["obsidian_note_path"] = str(note_path)

            # Open in Obsidian (avoid OS default markdown handler)
            if self.open_note_var.get():
                try:
                    vault_name = (getattr(getattr(self, "master", None), "cfg", {}) or {}).get("paths", {}).get("OBSIDIAN_VAULT", "")
                    open_note_in_obsidian(vault_name, self.vault_root, note_path)
                except Exception:
                    pass

        try:
            self.on_submit(payload)
        finally:
            self.destroy()
# ------------------------------------------------------------
# Entry point
# ------------------------------------------------------------
# -----------------------------
# Meeting draft persistence
# -----------------------------
def _meeting_store_path(app_dir: Path) -> Path:
    return app_dir / "meeting_draft.json"


def load_meeting_draft(app_dir: Path) -> dict:
    p = _meeting_store_path(app_dir)
    if not p.exists():
        return {}
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return {}


def save_meeting_draft(app_dir: Path, data: dict) -> None:
    p = _meeting_store_path(app_dir)
    p.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def clear_meeting_draft(app_dir: Path) -> None:
    p = _meeting_store_path(app_dir)
    try:
        if p.exists():
            p.unlink()
    except Exception:
        pass


def _project_store_path(app_dir: Path) -> Path:
    return app_dir / "project_draft.json"


def load_project_draft(app_dir: Path) -> dict:
    p = _project_store_path(app_dir)
    if not p.exists():
        return {}
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return {}


def save_project_draft(app_dir: Path, data: dict) -> None:
    p = _project_store_path(app_dir)
    p.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def clear_project_draft(app_dir: Path) -> None:
    p = _project_store_path(app_dir)
    try:
        if p.exists():
            p.unlink()
    except Exception:
        pass


def _resolve_meeting_template(vault_root: Path) -> Path | None:
    """
    Primary: 00_System/00_Templates/01_Entities/Meeting Template.md
    Fallback: first matching file in the same folder containing 'meeting' and 'template'.
    """
    folder = vault_root / "00_System" / "00_Templates" / "01_Entities"
    primary = folder / "Meeting Template.md"
    if primary.exists():
        return primary

    if folder.exists():
        for p in sorted(folder.glob("*.md")):
            name = p.name.lower()
            if "meeting" in name and "template" in name:
                return p

    return None


def _resolve_project_template(vault_root: Path) -> Path | None:
    """
    Primary: 00_System/00_Templates/02_Operations/Project Template.md
    Fallback: first matching file in the same folder containing 'project' and 'template'.
    """
    folder = vault_root / "00_System" / "00_Templates" / "02_Operations"
    primary = folder / "Project Template.md"
    if primary.exists():
        return primary

    if folder.exists():
        for p in sorted(folder.glob("*.md")):
            name = p.name.lower()
            if "project" in name and "template" in name:
                return p

    return None




def _merge_meeting_frontmatter(template_text: str, meeting: dict) -> str:
    """
    Merge meeting fields into YAML frontmatter if present.
    Maps:
      title -> title
      start_time -> start_time
      end_time -> end_time
      location -> location
      attendees -> attendees
      tags -> tags
      updated -> updated (now)
      created -> created (from start_time if provided else now)
    """
    if not template_text.lstrip().startswith("---"):
        return template_text

    lines = template_text.splitlines(True)
    # find start delimiter after leading blanks
    i0 = None
    for i, ln in enumerate(lines):
        if ln.strip() == "":
            continue
        if ln.strip() == "---":
            i0 = i
        break
    if i0 is None:
        return template_text

    i1 = None
    for j in range(i0 + 1, len(lines)):
        if lines[j].strip() == "---":
            i1 = j
            break
    if i1 is None:
        return template_text

    fm_lines = lines[i0 + 1 : i1]
    key_line_index = {}
    key_re = re.compile(r"^([A-Za-z0-9_-]+)\s*:\s*(.*)$")
    for idx, ln in enumerate(fm_lines):
        raw = ln.rstrip("\n")
        if not raw.strip() or raw.lstrip().startswith("#"):
            continue
        m = key_re.match(raw)
        if m:
            key_line_index[m.group(1)] = idx

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    title = (meeting.get("title") or "").strip()
    start_time = (meeting.get("start_time") or "").strip()
    end_time = (meeting.get("end_time") or "").strip()
    location = (meeting.get("location") or "").strip()
    attendees = (meeting.get("attendees") or "").strip()
    tags = (meeting.get("tags") or "").strip()

    created = start_time or now

    updates = {
        "title": f"\"{title}\"",
        "start_time": f"\"{start_time}\"",
        "end_time": f"\"{end_time}\"",
        "location": f"\"{location}\"",
        "attendees": f"\"{attendees}\"",
        "tags": f"\"{tags}\"",
        "created": f"\"{created}\"",
        "updated": f"\"{now}\"",
    }

    def set_key(k: str, v: str):
        if k in key_line_index:
            fm_lines[key_line_index[k]] = f"{k}: {v}\n"
        else:
            fm_lines.append(f"{k}: {v}\n")

    for k, v in updates.items():
        set_key(k, v)

    new_lines = []
    new_lines.extend(lines[: i0 + 1])
    new_lines.extend(fm_lines)
    new_lines.append("---\n")
    new_lines.extend(lines[i1 + 1 :])
    return "".join(new_lines)


def _merge_project_frontmatter(template_text: str, project: dict) -> str:
    """
    Merge project fields into YAML frontmatter if present.
    Maps:
      project_name -> project_name/title
      project_id -> project_id
      project_owner -> project_owner/owner
      project_status -> project_status/status
    """
    s = template_text.lstrip("\ufeff")
    if not s.lstrip().startswith("---"):
        return template_text

    lines = template_text.splitlines(True)
    i0 = None
    for i, ln in enumerate(lines):
        if ln.strip() == "":
            continue
        if ln.strip() == "---":
            i0 = i
        break
    if i0 is None:
        return template_text

    i1 = None
    for j in range(i0 + 1, len(lines)):
        if lines[j].strip() == "---":
            i1 = j
            break
    if i1 is None:
        return template_text

    fm_lines = lines[i0 + 1 : i1]
    key_line_index = {}
    key_re = re.compile(r"^([A-Za-z0-9_-]+)\s*:\s*(.*)$")
    for idx, ln in enumerate(fm_lines):
        raw = ln.rstrip("\n")
        if not raw.strip() or raw.lstrip().startswith("#"):
            continue
        m = key_re.match(raw)
        if m:
            key_line_index[m.group(1)] = idx

    project_id = _normalize_project_id((project.get("project_id") or "").strip())
    project_name = (project.get("project_name") or project.get("title") or "").strip()
    owner = (project.get("project_owner") or project.get("owner") or "").strip()
    status = (project.get("project_status") or project.get("status") or "").strip()
    priority = (project.get("project_priority") or project.get("priority") or "").strip()
    start_date = (project.get("start_date") or "").strip()
    end_date = (project.get("end_date") or "").strip()
    tags = (project.get("tags") or "").strip()
    summary = (project.get("summary") or "").strip()
    links = (project.get("links") or "").strip()
    risks = (project.get("risks") or "").strip()

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    created = (project.get("created") or "").strip() or now
    project["created"] = created
    project["project_id"] = project_id

    updates = {
        "entity_type": "\"project\"",
        "project_id": f"\"{project_id}\"",
        "project_name": f"\"{project_name}\"",
        "project_owner": f"\"{owner}\"",
        "project_status": f"\"{status}\"",
        "project_priority": f"\"{priority}\"",
        "start_date": f"\"{start_date}\"",
        "end_date": f"\"{end_date}\"",
        "tags": f"\"{tags}\"",
        "summary": f"\"{summary}\"",
        "links": f"\"{links}\"",
        "risks": f"\"{risks}\"",
        "created": f"\"{created}\"",
        "updated": f"\"{now}\"",
        "title": f"\"{project_name}\"",
        "owner": f"\"{owner}\"",
        "status": f"\"{status}\"",
        "priority": f"\"{priority}\"",
    }

    def set_key(k: str, v: str):
        if k in key_line_index:
            fm_lines[key_line_index[k]] = f"{k}: {v}\n"
        else:
            fm_lines.append(f"{k}: {v}\n")

    for k, v in updates.items():
        set_key(k, v)

    new_lines = []
    new_lines.extend(lines[: i0 + 1])
    new_lines.extend(fm_lines)
    new_lines.append("---\n")
    new_lines.extend(lines[i1 + 1 :])
    return "".join(new_lines)


def _resolve_goal_template(vault_root: Path) -> Path | None:
    """Resolve Goal Template: 02_Operations first, then 01_Entities."""
    primary = vault_root / "00_System" / "00_Templates" / "02_Operations" / "Goal Template.md"
    if primary.exists():
        return primary
    fallback = vault_root / "00_System" / "00_Templates" / "01_Entities" / "Goal Template.md"
    if fallback.exists():
        return fallback
    folder = vault_root / "00_System" / "00_Templates" / "02_Operations"
    if folder.exists():
        for p in sorted(folder.glob("*.md")):
            if "goal" in p.name.lower() and "template" in p.name.lower():
                return p
    folder = vault_root / "00_System" / "00_Templates" / "01_Entities"
    if folder.exists():
        for p in sorted(folder.glob("*.md")):
            if "goal" in p.name.lower() and "template" in p.name.lower():
                return p
    return None


def create_obsidian_goal_note(vault_root: Path, goal: dict) -> Path:
    """Create a Goal note under 10_Operations/15_Goals using the Goal template when available."""
    base_folder = vault_root / "10_Operations" / "15_Goals"
    base_folder.mkdir(parents=True, exist_ok=True)

    raw_id = (goal.get("goal_id") or goal.get("goal_ID") or "").strip()
    goal_id = _normalize_goal_id(raw_id)
    if not goal_id:
        goal_id = _next_goal_id_from_vault(vault_root)

    title = (goal.get("title") or "").strip() or "Untitled Goal"
    status = (goal.get("status") or "").strip() or "planned"
    owner = (goal.get("owner") or "").strip()
    area = (goal.get("area") or "").strip()
    start_date = (goal.get("start_date") or "").strip()
    end_date = (goal.get("end_date") or "").strip()
    cadence = (goal.get("cadence") or "").strip()
    success_metric = (goal.get("success_metric") or "").strip()

    def _list_from_value(value) -> list[str]:
        if isinstance(value, (list, tuple)):
            return [str(v).strip() for v in value if str(v).strip()]
        return _split_bulk(value or "")

    key_results = _list_from_value(goal.get("key_results"))
    milestones = _list_from_value(goal.get("milestones"))
    related_projects = _list_from_value(goal.get("related_projects"))
    related_tasks = _list_from_value(goal.get("related_tasks"))
    tags = _list_from_value(goal.get("tags"))

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    created = (goal.get("created") or "").strip() or now
    if len(created) == 10:
        created = created + " 00:00:00"
    updated = now

    goal["goal_id"] = goal_id
    goal["created"] = created
    goal["updated"] = updated

    name_bits = [goal_id, title] if goal_id else [title]
    slug = _safe_slug(" - ".join(name_bits)) or "Goal"
    note_path = base_folder / f"{slug}.md"

    tpl_path = _resolve_goal_template(vault_root)

    token_map = {
        "goal_id": goal_id,
        "title": title,
        "status": status,
        "owner": owner,
        "area": area,
        "start_date": start_date,
        "end_date": end_date,
        "target_date": end_date,
        "cadence": cadence,
        "review_cadence": cadence or "quarterly",
        "specific_outcome": title,
        "success_metric": success_metric,
        "created": created,
        "updated": updated,
    }

    updates = {
        "entity_type": "goal",
        "goal_id": goal_id,
        "title": title,
        "status": status,
        "owner": owner,
        "area": area,
        "start_date": start_date,
        "end_date": end_date,
        "cadence": cadence,
        "success_metric": success_metric,
        "key_results": key_results,
        "milestones": milestones,
        "related_projects": related_projects,
        "related_tasks": related_tasks,
        "tags": tags,
        "created": created,
        "updated": updated,
    }

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)
        for k, v in token_map.items():
            tpl_text = tpl_text.replace(f"{{{{{k}}}}}", v)
            tpl_text = tpl_text.replace(f"{{{{{k.upper()}}}}}", v)
            tpl_text = tpl_text.replace(f"${{{k}}}", v)
        md = _apply_yaml_updates(tpl_text, updates)
    else:
        md = f"""---
entity_type: goal
goal_id: "{goal_id}"
title: "{title}"
status: "{status}"
owner: "{owner}"
area: "{area}"
start_date: "{start_date}"
end_date: "{end_date}"
cadence: "{cadence}"
success_metric: "{success_metric}"
key_results: []
milestones: []
related_projects: []
related_tasks: []
tags: []
created: "{created}"
updated: "{updated}"
---

# Goal — {title}
"""
        md = _apply_yaml_updates(md, updates)

    note_path.write_text(md, encoding="utf-8")
    return note_path


def create_obsidian_meeting_note(vault_root: Path, meeting: dict) -> Path:
    base_folder = vault_root / "10_Operations" / "10_Meetings"
    base_folder.mkdir(parents=True, exist_ok=True)

    title = (meeting.get("title") or "").strip() or "Untitled Meeting"
    start_time = (meeting.get("start_time") or "").strip() or datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    end_time = (meeting.get("end_time") or "").strip()
    location = (meeting.get("location") or "").strip()
    attendees = (meeting.get("attendees") or "").strip()
    tags = (meeting.get("tags") or "").strip()
    agenda = (meeting.get("agenda") or "").strip()
    notes = (meeting.get("notes") or "").strip()

    date_prefix = start_time[:10] if len(start_time) >= 10 else datetime.now().strftime("%Y-%m-%d")
    slug = _safe_slug(f"{date_prefix} - {title}")
    note_path = base_folder / f"{slug}.md"

    tpl_path = _resolve_meeting_template(vault_root)
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)

        # token replacement
        token_map = {
            "title": title, "start_time": start_time, "end_time": end_time,
            "location": location, "attendees": attendees, "tags": tags,
            "agenda": agenda, "notes": notes, "created": now, "updated": now,
        }
        for k, v in token_map.items():
            tpl_text = tpl_text.replace(f"{{{{{k}}}}}", v)
            tpl_text = tpl_text.replace(f"{{{{{k.upper()}}}}}", v)
            tpl_text = tpl_text.replace(f"${{{k}}}", v)

        # Explicit token replacements requested by user
        tpl_text = tpl_text.replace("{{Agenda}}", agenda)
        tpl_text = tpl_text.replace("{{Notes}}", notes)

        if tpl_text.lstrip().startswith("---"):
            md = _merge_meeting_frontmatter(tpl_text, meeting)
        else:
            md = f"""---
entity_type: meeting
title: "{title}"
start_time: "{start_time}"
end_time: "{end_time}"
location: "{location}"
attendees: "{attendees}"
tags: "{tags}"
created: "{now}"
updated: "{now}"
---

""" + tpl_text.lstrip()
    else:
        md = f"""---
entity_type: meeting
title: "{title}"
start_time: "{start_time}"
end_time: "{end_time}"
location: "{location}"
attendees: "{attendees}"
tags: "{tags}"
created: "{now}"
updated: "{now}"
---

## Agenda
{agenda}

## Notes
{notes}
"""
    note_path.write_text(md, encoding="utf-8")
    return note_path


def create_obsidian_project_note(vault_root: Path, project: dict) -> Path:
    """
    Creates a project note under 10_Operations/14_Projects (created if missing),
    based on the Project Template when available.
    """
    base_folder = vault_root / "10_Operations" / "14_Projects"
    base_folder.mkdir(parents=True, exist_ok=True)

    project_id = _normalize_project_id((project.get("project_id") or "").strip())
    project_name = (project.get("project_name") or project.get("title") or "").strip() or "Untitled Project"
    owner = (project.get("project_owner") or project.get("owner") or "").strip()
    status = (project.get("project_status") or project.get("status") or "").strip()
    priority = (project.get("project_priority") or project.get("priority") or "").strip()
    start_date = (project.get("start_date") or "").strip()
    end_date = (project.get("end_date") or "").strip()
    tags = (project.get("tags") or "").strip()
    summary = (project.get("summary") or "").strip()
    links = (project.get("links") or "").strip()
    risks = (project.get("risks") or "").strip()

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    created = (project.get("created") or "").strip() or now
    project["created"] = created
    project["project_id"] = project_id

    name_bits = []
    if project_id:
        name_bits.append(project_id)
    name_bits.append(project_name)
    slug = _safe_slug(" - ".join(name_bits))
    note_path = base_folder / f"{slug}.md"

    tpl_path = _resolve_project_template(vault_root)

    token_map = {
        "project_id": project_id,
        "project_name": project_name,
        "project_owner": owner,
        "project_status": status,
        "project_priority": priority,
        "start_date": start_date,
        "end_date": end_date,
        "tags": tags,
        "summary": summary,
        "created": created,
        "updated": now,
        "title": project_name,
        "owner": owner,
        "status": status,
        "priority": priority,
    }
    timeline_lines = ["```chronos"]
    if start_date:
        timeline_lines.append(f"- [{start_date}] Project Start")
    if end_date:
        timeline_lines.append(f"- [{end_date}] Target Completion")
    timeline_lines.append("```")
    token_map["project_timeline"] = "\n".join(timeline_lines)

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)

        for k, v in token_map.items():
            tpl_text = tpl_text.replace(f"{{{{{k}}}}}", v)
            tpl_text = tpl_text.replace(f"{{{{{k.upper()}}}}}", v)
            tpl_text = tpl_text.replace(f"${{{k}}}", v)

        if tpl_text.lstrip().startswith("---"):
            md = _merge_project_frontmatter(tpl_text, project)
        else:
            fm = f"""---
entity_type: project
project_id: "{project_id}"
project_name: "{project_name}"
project_owner: "{owner}"
project_status: "{status}"
project_priority: "{priority}"
start_date: "{start_date}"
end_date: "{end_date}"
tags: "{tags}"
summary: "{summary}"
created: "{created}"
updated: "{now}"
---

"""
            md = fm + tpl_text.lstrip()
    else:
        md = f"""---
entity_type: project
project_id: "{project_id}"
project_name: "{project_name}"
project_owner: "{owner}"
project_status: "{status}"
project_priority: "{priority}"
start_date: "{start_date}"
end_date: "{end_date}"
tags: "{tags}"
summary: "{summary}"
created: "{created}"
updated: "{now}"
---

## Summary
{summary}

## Project Timeline (Chronos Timeline)
{token_map["project_timeline"]}
"""

    note_path.write_text(md, encoding="utf-8")
    return note_path


def _parse_datetime_for_picker(value: str) -> tuple[date, int, int, int]:
    """Parse 'YYYY-MM-DD HH:MM:SS' or 'YYYY-MM-DD' into (date, hour, minute, second)."""
    value = (value or "").strip()
    now = datetime.now()
    default_date = now.date()
    default_h, default_m, default_s = now.hour, now.minute, 0  # seconds default to 00
    if not value:
        return default_date, default_h, default_m, default_s
    parts = value.split()
    try:
        d = datetime.strptime(parts[0], "%Y-%m-%d").date()
    except (ValueError, IndexError):
        d = default_date
    h, m, s = default_h, default_m, 0  # seconds default to 00
    if len(parts) >= 2:
        tparts = parts[1].split(":")
        if len(tparts) >= 1 and tparts[0].isdigit():
            h = min(23, max(0, int(tparts[0])))
        if len(tparts) >= 2 and tparts[1].isdigit():
            m = min(59, max(0, int(tparts[1])))
        if len(tparts) >= 3 and tparts[2].isdigit():
            s = min(59, max(0, int(tparts[2])))
    return d, h, m, s


def _format_datetime_from_picker(d: date, h: int, m: int, s: int) -> str:
    return f"{d.isoformat()} {h:02d}:{m:02d}:{s:02d}"


_MEETING_MINUTES = (0, 10, 15, 20, 30, 40, 45, 50)


def _round_to_meeting_minute(m: int) -> int:
    """Round minute to nearest of 00, 10, 15, 20, 30, 40, 45, 50."""
    return min(_MEETING_MINUTES, key=lambda x: abs(x - min(59, max(0, m))))


class MeetingIntakeWindow(tk.Toplevel):
    def __init__(self, master, *, app_dir: Path, vault_root: Path | None, on_submit, resume: bool = False):
        super().__init__(master)
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("Meeting Intake — v20")
        self.configure(bg="#111111")
        self.geometry("860x620")
        self.minsize(760, 560)
        self.resizable(True, True)

        self._build_widgets()

        if resume:
            self._apply_draft(load_meeting_draft(self.app_dir))

        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not self.vars["start_time"].get().strip():
            self.vars["start_time"].set(now)

        self.transient(master)
        self.grab_set()
        self.focus_force()

    def _build_widgets(self):
        cfg = getattr(self.master, "cfg", {}) or {}
        analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")

        container = tk.Frame(self, bg="#111111")
        container.pack(fill="both", expand=True)

        canvas = tk.Canvas(container, bg="#111111", highlightthickness=0)
        vbar = tk.Scrollbar(container, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vbar.set)
        vbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        content = tk.Frame(canvas, bg="#111111")
        win_id = canvas.create_window((0, 0), window=content, anchor="nw")

        def _sync(_=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfigure(win_id, width=canvas.winfo_width())

        content.bind("<Configure>", _sync)
        canvas.bind("<Configure>", _sync)
        canvas.bind_all("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1*(e.delta/120)), "units"))

        hdr = tk.Frame(content, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 10))
        tk.Label(hdr, text="Meeting Intake", font=("Segoe UI", 16, "bold"), fg="#FFFFFF", bg="#111111").pack(anchor="w")
        tk.Label(hdr, text="Capture meeting details and generate an Obsidian note.", font=("Segoe UI", 10),
                 fg="#BBBBBB", bg="#111111").pack(anchor="w")

        form = tk.Frame(content, bg="#111111")
        form.pack(fill="both", expand=True, padx=16)
        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        right.pack(side="right", fill="both", expand=True, padx=(10, 0))

        self.vars = {
            "title": tk.StringVar(),
            "start_time": tk.StringVar(),
            "end_time": tk.StringVar(),
            "location": tk.StringVar(),
            "attendees": tk.StringVar(),
            "tags": tk.StringVar(),
        }
        self._dt_pickers: dict[str, dict] = {}

        def add_entry(parent, label, var):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF",
                           insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_labeled_combo(parent, label, var, values):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        def add_datetime_picker(parent, label: str, key: str, initial_value: str):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            picker_row = tk.Frame(row, bg="#111111")
            picker_row.pack(fill="x")
            d, h, m, s = _parse_datetime_for_picker(initial_value)
            ent = None
            if DateEntry is not None:
                try:
                    style = ttk.Style(self)
                    style.configure(
                        "Meeting.DateEntry",
                        fieldbackground="#1B1B1B",
                        background="#2A2A2A",
                        foreground="#FFFFFF",
                        arrowcolor="#FFFFFF",
                    )
                    date_entry = DateEntry(
                        picker_row, date_pattern="y-mm-dd", year=d.year, month=d.month, day=d.day,
                        style="Meeting.DateEntry",
                    )
                except Exception:
                    date_entry = DateEntry(
                        picker_row, date_pattern="y-mm-dd", year=d.year, month=d.month, day=d.day,
                    )
                date_entry.pack(side="left", padx=(0, 8))
            else:
                date_entry = None
                ent = tk.Entry(picker_row, bg="#1B1B1B", fg="#FFFFFF", relief="flat", width=12)
                ent.insert(0, d.isoformat())
                ent.pack(side="left", padx=(0, 8))
            m = _round_to_meeting_minute(m)
            hour_var = tk.StringVar(value=str(h))
            min_var = tk.StringVar(value=f"{m:02d}")
            sp = tk.Spinbox(
                picker_row, from_=0, to=23, increment=1, textvariable=hour_var, width=3,
                bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat",
                buttonbackground="#1B1B1B",
            )
            sp.pack(side="left", padx=(0, 4))
            min_values = tuple(f"{x:02d}" for x in _MEETING_MINUTES)
            min_sp = tk.Spinbox(
                picker_row, values=min_values, textvariable=min_var, width=3,
                bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat",
                buttonbackground="#1B1B1B",
            )
            min_sp.pack(side="left", padx=(0, 4))
            self._dt_pickers[key] = {
                "date_entry": date_entry,
                "fallback_entry": ent if DateEntry is None else None,
                "hour_var": hour_var,
                "min_var": min_var,
            }

        add_entry(left, "Title", self.vars["title"])
        add_datetime_picker(left, "Meeting Start Date", "start_time", self.vars["start_time"].get())
        add_datetime_picker(left, "Meeting End Date", "end_time", self.vars["end_time"].get())
        add_entry(right, "Location", self.vars["location"])
        add_entry(right, "Attendees (comma-separated)", self.vars["attendees"])
        add_entry(right, "Tags (comma-separated)", self.vars["tags"])

        tk.Label(content, text="Agenda", fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w", padx=16, pady=(12, 0))
        self.agenda_txt = tk.Text(content, height=6, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat", wrap="word")
        self.agenda_txt.pack(fill="x", padx=16, pady=(4, 10))

        tk.Label(content, text="Notes", fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w", padx=16, pady=(0, 0))
        self.notes_txt = tk.Text(content, height=8, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat", wrap="word")
        self.notes_txt.pack(fill="x", padx=16, pady=(4, 10))

        bar = tk.Frame(self, bg="#111111")
        bar.pack(side="bottom", fill="x", padx=16, pady=(6, 14))

        self.open_note_var = tk.BooleanVar(value=True)
        if self.vault_root is not None:
            tk.Checkbutton(
                bar, text="Open created note automatically",
                variable=self.open_note_var,
                bg="#111111", fg="#DDDDDD", selectcolor="#111111",
                activebackground="#111111", activeforeground="#DDDDDD",
            ).pack(side="left")

        def btn(t, cmd):
            return tk.Button(bar, text=t, command=cmd, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right", padx=(8, 0))
        btn("Save & Continue", self._save_continue).pack(side="right")
        btn("Clear Draft", self._clear_draft).pack(side="right", padx=(8, 0))
        btn("Save Draft", self._save_draft).pack(side="right", padx=(8, 0))

    def _apply_draft(self, draft: dict):
        if not draft:
            return
        for k, v in draft.items():
            if k in self.vars and isinstance(v, str):
                self.vars[k].set(v)
        for key in ("start_time", "end_time"):
            if key in draft and key in self._dt_pickers:
                self._set_picker_datetime(key, str(draft[key]))
        if "agenda" in draft:
            self.agenda_txt.delete("1.0", "end")
            self.agenda_txt.insert("1.0", draft.get("agenda", ""))
        if "notes" in draft:
            self.notes_txt.delete("1.0", "end")
            self.notes_txt.insert("1.0", draft.get("notes", ""))

    def _get_picker_datetime(self, key: str) -> str:
        p = self._dt_pickers.get(key)
        if not p:
            return self.vars[key].get().strip()
        try:
            if p["date_entry"] is not None:
                d = p["date_entry"].get_date()
            elif p.get("fallback_entry"):
                d = datetime.strptime(p["fallback_entry"].get().strip(), "%Y-%m-%d").date()
            else:
                return self.vars[key].get().strip()
            h = min(23, max(0, int(p["hour_var"].get() or 0)))
            m_raw = min(59, max(0, int(p["min_var"].get() or 0)))
            m = _round_to_meeting_minute(m_raw)
            s = 0  # seconds always 00 for meeting times
            return _format_datetime_from_picker(d, h, m, s)
        except (ValueError, TypeError):
            return self.vars[key].get().strip()

    def _set_picker_datetime(self, key: str, value: str) -> None:
        p = self._dt_pickers.get(key)
        if not p:
            self.vars[key].set(value)
            return
        d, h, m, s = _parse_datetime_for_picker(value)
        m = _round_to_meeting_minute(m)
        if p["date_entry"] is not None:
            p["date_entry"].set_date(d)
        elif p.get("fallback_entry"):
            p["fallback_entry"].delete(0, "end")
            p["fallback_entry"].insert(0, d.isoformat())
        p["hour_var"].set(str(h))
        p["min_var"].set(f"{m:02d}")

    def _collect(self) -> dict:
        data = {k: v.get().strip() for k, v in self.vars.items()}
        for key in ("start_time", "end_time"):
            if key in self._dt_pickers:
                data[key] = self._get_picker_datetime(key)
        data["agenda"] = self.agenda_txt.get("1.0", "end").strip()
        data["notes"] = self.notes_txt.get("1.0", "end").strip()
        return data

    def _save_draft(self):
        save_meeting_draft(self.app_dir, self._collect())
        messagebox.showinfo("Meeting Intake", "Draft saved.", parent=self)

    def _clear_draft(self):
        clear_meeting_draft(self.app_dir)
        for v in self.vars.values():
            v.set("")
        self.agenda_txt.delete("1.0", "end")
        self.notes_txt.delete("1.0", "end")
        messagebox.showinfo("Meeting Intake", "Draft cleared.", parent=self)

    def _save_continue(self):
        data = self._collect()
        if not data.get("title"):
            messagebox.showwarning("Meeting Intake", "Title is required.", parent=self)
            return

        start_raw = (data.get("start_time") or "").strip()
        end_raw = (data.get("end_time") or "").strip()
        if start_raw and end_raw:
            start_dt = end_dt = None
            for fmt in ("%Y-%m-%d %H:%M:%S", "%Y-%m-%d %H:%M", "%Y-%m-%d"):
                try:
                    start_dt = datetime.strptime(start_raw, fmt)
                    end_dt = datetime.strptime(end_raw, fmt)
                    break
                except ValueError:
                    continue
            if start_dt is not None and end_dt is not None and end_dt <= start_dt:
                messagebox.showwarning(
                    "Meeting Intake",
                    "End date/time must be later than start date/time.",
                    parent=self,
                )
                return

        save_meeting_draft(self.app_dir, data)

        note_path = None
        if self.vault_root is not None:
            try:
                note_path = create_obsidian_meeting_note(self.vault_root, data)
            except Exception as e:
                messagebox.showwarning("Meeting Intake", f"Saved draft, but failed to create meeting note:\n{e}", parent=self)
            try:
                _log_shift_entry_from_launcher(self.master, self.vault_root, "meeting", data, note_path)
            except Exception as e:
                messagebox.showwarning(
                    "Shift Log",
                    f"Meeting note created, but failed to add entry to daily shift log:\n{e}",
                    parent=self,
                )

        if note_path is not None and self.open_note_var.get():
            try:
                open_note_in_obsidian(getattr(getattr(self, "master", None), "cfg", {}).get("paths", {}).get("OBSIDIAN_VAULT", ""), self.vault_root, note_path) if self.vault_root is not None else False
            except Exception:
                pass

        try:
            self.on_submit({**data, **({"obsidian_note_path": str(note_path)} if note_path else {})})
        finally:
            self.destroy()



# -----------------------------
# Playbook draft persistence
# -----------------------------
def _playbook_store_path(app_dir: Path) -> Path:
    return app_dir / "playbook_draft.json"


def load_playbook_draft(app_dir: Path) -> dict:
    p = _playbook_store_path(app_dir)
    if not p.exists():
        return {}
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return {}


def save_playbook_draft(app_dir: Path, data: dict) -> None:
    p = _playbook_store_path(app_dir)
    p.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def clear_playbook_draft(app_dir: Path) -> None:
    p = _playbook_store_path(app_dir)
    try:
        if p.exists():
            p.unlink()
    except Exception:
        pass


def _resolve_playbook_template(vault_root: Path) -> Path | None:
    """
    Primary (recommended): 00_System/00_Templates/01_Entities/Playbook Template.md
    Fallback: first matching file in the same folder containing 'playbook' and 'template'.
    """
    folder = vault_root / "00_System" / "00_Templates" / "01_Entities"
    primary = folder / "Playbook Template.md"
    if primary.exists():
        return primary

    if folder.exists():
        for p in sorted(folder.glob("*.md")):
            name = p.name.lower()
            if "playbook" in name and "template" in name:
                return p

    # Secondary fallback: Operations templates
    folder2 = vault_root / "00_System" / "00_Templates" / "02_Operations"
    if folder2.exists():
        for p in sorted(folder2.glob("*.md")):
            name = p.name.lower()
            if "playbook" in name and "template" in name:
                return p

    return None


def _format_playbook_id(value: str) -> str:
    raw = (value or "").strip()
    if not raw:
        return ""
    if raw.upper().startswith("PB-"):
        return f"PB-{raw[3:]}"
    return f"PB-{raw}"


def _merge_playbook_frontmatter(template_text: str, pb: dict) -> str:
    """
    Merge playbook fields into YAML frontmatter if present.
    Fields:
      playbook_id, title, status, owner_primary, tags, itid_primary, created, updated
    """
    if not template_text.lstrip().startswith("---"):
        return template_text

    lines = template_text.splitlines(True)
    # locate fm bounds
    i0 = None
    for i, ln in enumerate(lines):
        if ln.strip() == "":
            continue
        if ln.strip() == "---":
            i0 = i
        break
    if i0 is None:
        return template_text

    i1 = None
    for j in range(i0 + 1, len(lines)):
        if lines[j].strip() == "---":
            i1 = j
            break
    if i1 is None:
        return template_text

    fm_lines = lines[i0 + 1 : i1]
    key_line_index = {}
    key_re = re.compile(r"^([A-Za-z0-9_-]+)\s*:\s*(.*)$")
    for idx, ln in enumerate(fm_lines):
        raw = ln.rstrip("\n")
        if not raw.strip() or raw.lstrip().startswith("#"):
            continue
        m = key_re.match(raw)
        if m:
            key_line_index[m.group(1)] = idx

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    playbook_id = _format_playbook_id((pb.get("playbook_id") or "").strip())
    title = (pb.get("title") or "").strip()
    status = (pb.get("status") or "").strip()
    owner_primary = (pb.get("owner_primary") or "").strip()
    tags = (pb.get("tags") or "").strip()
    itid_primary = (pb.get("itid_primary") or "").strip()
    created = (pb.get("created") or "").strip() or now

    updates = {
        "entity_type": "\"playbook\"",
        "playbook_id": f"\"{playbook_id}\"",
        "title": f"\"{title}\"",
        "status": f"\"{status}\"",
        "owner_primary": f"\"{owner_primary}\"",
        "tags": f"\"{tags}\"",
        "created": f"\"{created}\"",
        "updated": f"\"{now}\"",
    }

    def set_key(k: str, v: str):
        if k in key_line_index:
            fm_lines[key_line_index[k]] = f"{k}: {v}\n"
        else:
            fm_lines.append(f"{k}: {v}\n")

    for k, v in updates.items():
        set_key(k, v)

    new_lines = []
    new_lines.extend(lines[: i0 + 1])
    new_lines.extend(fm_lines)
    new_lines.append("---\n")
    new_lines.extend(lines[i1 + 1 :])
    return "".join(new_lines)



def create_obsidian_playbook_note(vault_root: Path, pb: dict) -> Path:
    """
    Creates a playbook note under 10_Operations/05_Playbooks (created if missing).
    Uses Playbook Template when available; otherwise writes a minimal stub.

    Field mappings:
      Playbook ID -> playbook_id
      Title -> title
      Status -> status
      Owner -> owner_primary
      Tags -> tags
      Purpose -> {{Purpose}}
    """
    base_folder = vault_root / "10_Operations" / "05_Playbooks"
    base_folder.mkdir(parents=True, exist_ok=True)

    playbook_id = _format_playbook_id((pb.get("playbook_id") or "").strip())
    title = (pb.get("title") or "").strip() or "Untitled Playbook"
    status = (pb.get("status") or "").strip() or "draft"
    owner_primary = (pb.get("owner_primary") or pb.get("owner") or "").strip()
    tags = (pb.get("tags") or "").strip()
    itid_primary = (pb.get("itid_primary") or "").strip()
    purpose = (pb.get("purpose") or "").strip()

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    created = (pb.get("created") or "").strip() or now
    pb["created"] = created
    pb["playbook_id"] = playbook_id

    name_bits = []
    if playbook_id:
        name_bits.append(playbook_id)
    name_bits.append(title)
    slug = _safe_slug(" - ".join(name_bits))
    note_path = base_folder / f"{slug}.md"

    tpl_path = _resolve_playbook_template(vault_root)

    token_map = {
        "playbook_id": playbook_id,
        "title": title,
        "status": status,
        "owner_primary": owner_primary,
        "tags": tags,
        "itid_primary": itid_primary,
        "purpose": purpose,
        "created": created,
        "updated": now,
        "title": title,
        "owner": owner_primary,
    }

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)

        for k, v in token_map.items():
            tpl_text = tpl_text.replace(f"{{{{{k}}}}}", v)
            tpl_text = tpl_text.replace(f"{{{{{k.upper()}}}}}", v)
            tpl_text = tpl_text.replace(f"${{{k}}}", v)

        tpl_text = tpl_text.replace("{{Purpose}}", purpose)
        tpl_text = tpl_text.replace("{{PURPOSE}}", purpose)

        if tpl_text.lstrip().startswith("---"):
            md = _merge_playbook_frontmatter(tpl_text, token_map)
        else:
            fm = f"""---
entity_type: playbook
playbook_id: \"{playbook_id}\"
title: \"{title}\"
status: \"{status}\"
owner_primary: \"{owner_primary}\"
tags: \"{tags}\"
created: \"{created}\"
updated: \"{now}\"
---

"""
            md = fm + tpl_text.lstrip()
    else:
        md = f"""---
entity_type: playbook
playbook_id: \"{playbook_id}\"
title: \"{title}\"
status: \"{status}\"
owner_primary: \"{owner_primary}\"
tags: \"{tags}\"
created: \"{created}\"
updated: \"{now}\"
---

## Purpose

{purpose}
"""

    note_path.write_text(md, encoding="utf-8")
    return note_path


class ProjectIntakeWindow(tk.Toplevel):
    def __init__(self, master, *, app_dir: Path, vault_root: Path | None, on_submit, resume: bool = False):
        super().__init__(master)
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("Project Intake")
        self.configure(bg="#111111")
        self.geometry("900x680")
        self.minsize(800, 600)
        self.resizable(True, True)

        self._build_widgets()

        if resume:
            self._apply_draft(load_project_draft(self.app_dir))

        now = datetime.now().strftime("%Y-%m-%d")
        if not self.vars["start_date"].get().strip():
            self.vars["start_date"].set(now)

        self.transient(master)
        self.grab_set()
        self.focus_force()

    def _build_widgets(self):
        container = tk.Frame(self, bg="#111111")
        container.pack(fill="both", expand=True)

        content = tk.Frame(container, bg="#111111")
        content.pack(fill="both", expand=True)

        form = tk.Frame(content, bg="#111111")
        form.pack(fill="x", padx=16, pady=12)
        form.grid_columnconfigure(0, weight=1)
        form.grid_columnconfigure(1, weight=1)

        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.grid(row=0, column=0, sticky="nsew", padx=(0, 10))
        right.grid(row=0, column=1, sticky="nsew", padx=(10, 0))

        cfg = getattr(self.master, "cfg", {}) or {}
        analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")

        self.vars = {
            "project_id": tk.StringVar(),
            "project_name": tk.StringVar(),
            "project_owner": tk.StringVar(value=analyst_default),
            "project_status": tk.StringVar(value="planned"),
            "project_priority": tk.StringVar(value="medium"),
            "start_date": tk.StringVar(),
            "end_date": tk.StringVar(),
            "tags": tk.StringVar(),
        }
        self.summary_var = tk.StringVar()
        self.sequential_mode_var = tk.BooleanVar(value=False)

        def add_entry(parent, label, var, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_combo(parent, label, var, values, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        project_id_entry = add_entry(left, "Project ID", self.vars["project_id"])
        project_id_entry.bind("<FocusOut>", self._on_project_id_focus_out)

        seq_row = tk.Frame(left, bg="#111111")
        seq_row.pack(fill="x", pady=(2, 6))
        tk.Checkbutton(
            seq_row,
            text="Sequential mode",
            variable=self.sequential_mode_var,
            command=self._on_sequential_toggle,
            bg="#111111",
            fg="#DDDDDD",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#DDDDDD",
        ).pack(anchor="w")

        add_entry(left, "Project Name", self.vars["project_name"], required=True)
        add_entry(left, "Owner/Lead", self.vars["project_owner"], required=True)
        add_combo(
            left,
            "Status",
            self.vars["project_status"],
            ["ideation", "planned", "active", "blocked", "done", "archived"],
            required=True,
        )

        add_combo(
            right,
            "Priority",
            self.vars["project_priority"],
            ["low", "medium", "high", "critical"],
            required=True,
        )
        add_entry(right, "Start Date (YYYY-MM-DD)", self.vars["start_date"], required=True)
        add_entry(right, "Target Completion Date (YYYY-MM-DD)", self.vars["end_date"], required=True)
        add_entry(right, "Tags (comma-separated)", self.vars["tags"])

        sum_lbl = tk.Frame(content, bg="#111111")
        sum_lbl.pack(anchor="w", padx=16, pady=(8, 0))
        tk.Label(sum_lbl, text="Summary / Objective", fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
        tk.Label(sum_lbl, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
        self.summary_entry = tk.Entry(
            content,
            textvariable=self.summary_var,
            bg="#1B1B1B",
            fg="#FFFFFF",
            insertbackground="#FFFFFF",
            relief="flat",
        )
        self.summary_entry.pack(fill="x", padx=16, pady=(4, 8), ipady=6)

        bar = tk.Frame(self, bg="#111111")
        bar.pack(side="bottom", fill="x", padx=16, pady=(6, 14))

        self.open_note_var = tk.BooleanVar(value=True)

        def btn(t, cmd):
            return tk.Button(bar, text=t, command=cmd, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right", padx=(8, 0))
        btn("Save & Continue", self._save_continue).pack(side="right")

    def _apply_draft(self, draft: dict):
        if not draft:
            return
        for k, v in draft.items():
            if k in self.vars and isinstance(v, str):
                self.vars[k].set(v)
        if "summary" in draft:
            self.summary_var.set(draft.get("summary", ""))
        if "sequential_mode" in draft:
            try:
                self.sequential_mode_var.set(bool(draft.get("sequential_mode")))
            except Exception:
                pass
            if self.sequential_mode_var.get():
                self._apply_sequential_suggestion()

    def _collect(self) -> dict:
        data = {k: v.get().strip() for k, v in self.vars.items()}
        data["summary"] = self.summary_var.get().strip()
        data["sequential_mode"] = bool(self.sequential_mode_var.get())
        if self.sequential_mode_var.get():
            year = datetime.now().year
            current = data.get("project_id", "")
            if current:
                data["project_id"] = _normalize_project_id(current, default_year=year)
            else:
                data["project_id"] = _next_project_id_from_vault(self.vault_root, default_year=year)
        return data

    def _apply_sequential_suggestion(self) -> None:
        year = datetime.now().year
        current = self.vars["project_id"].get().strip()
        if current:
            normalized = _normalize_project_id(current, default_year=year)
            if normalized != current:
                self.vars["project_id"].set(normalized)
            return
        suggested = _next_project_id_from_vault(self.vault_root, default_year=year)
        self.vars["project_id"].set(suggested)

    def _on_sequential_toggle(self) -> None:
        if self.sequential_mode_var.get():
            self._apply_sequential_suggestion()

    def _on_project_id_focus_out(self, _event=None) -> None:
        if not self.sequential_mode_var.get():
            return
        year = datetime.now().year
        current = self.vars["project_id"].get().strip()
        if current:
            self.vars["project_id"].set(_normalize_project_id(current, default_year=year))

    def _save_draft(self):
        save_project_draft(self.app_dir, self._collect())
        messagebox.showinfo("Project Intake", "Draft saved.", parent=self)

    def _clear_draft(self):
        clear_project_draft(self.app_dir)
        for v in self.vars.values():
            v.set("")
        self.vars["project_status"].set("planned")
        self.vars["project_priority"].set("medium")
        self.sequential_mode_var.set(False)
        self.summary_var.set("")
        messagebox.showinfo("Project Intake", "Draft cleared.", parent=self)

    def _save_continue(self):
        data = self._collect()
        missing = []
        if not data.get("project_name"):
            missing.append("Project Name")
        if not data.get("project_owner"):
            missing.append("Owner/Lead")
        if not data.get("project_status"):
            missing.append("Status")
        if not data.get("project_priority"):
            missing.append("Priority")
        if not data.get("start_date") or not data.get("end_date"):
            missing.append("Start and End Dates")
        if not data.get("summary"):
            missing.append("Summary")
        if missing:
            messagebox.showwarning(
                "Project Intake",
                "Missing required fields: " + ", ".join(missing),
                parent=self,
            )
            return

        save_project_draft(self.app_dir, data)

        note_path = None
        if self.vault_root is not None:
            try:
                note_path = create_obsidian_project_note(self.vault_root, data)
            except Exception as e:
                messagebox.showerror("Project Intake", f"Failed to create note:\n\n{e}", parent=self)
                try:
                    _log_shift_entry_from_launcher(self.master, self.vault_root, "project", data, None)
                except Exception:
                    pass
                return

        if note_path:
            data["obsidian_note_path"] = str(note_path)

        try:
            self.on_submit(data)
        except Exception:
            pass

        if note_path and self.open_note_var.get():
            try:
                vault_name = (
                    getattr(getattr(self, "master", None), "cfg", {}).get("paths", {}).get("OBSIDIAN_VAULT", "") or ""
                )
                open_note_in_obsidian(vault_name, self.vault_root, note_path)
            except Exception:
                pass

        if self.vault_root is not None:
            try:
                _log_shift_entry_from_launcher(self.master, self.vault_root, "project", data, note_path)
            except Exception:
                pass

        messagebox.showinfo("Project Intake", "Saved.", parent=self)
        self.destroy()


class PlaybookIntakeWindow(tk.Toplevel):
    """
    Playbook Intake UI (similar to Incident/Meeting Intake).
    """

    def __init__(self, master, *, app_dir: Path, vault_root: Path | None, on_submit, resume: bool = False):
        super().__init__(master)
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("Playbook Intake")
        self.configure(bg="#111111")
        self.geometry("900x640")
        self.minsize(780, 560)
        self.resizable(True, True)

        self._build_widgets()

        if resume:
            self._apply_draft(load_playbook_draft(self.app_dir))

        self.transient(master)
        self.grab_set()
        self.focus_force()

    def _build_widgets(self):
        cfg = getattr(self.master, "cfg", {}) or {}
        analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")

        container = tk.Frame(self, bg="#111111")
        container.pack(fill="both", expand=True)

        canvas = tk.Canvas(container, bg="#111111", highlightthickness=0)
        vbar = tk.Scrollbar(container, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vbar.set)
        vbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        content = tk.Frame(canvas, bg="#111111")
        win_id = canvas.create_window((0, 0), window=content, anchor="nw")

        def _sync(_=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfigure(win_id, width=canvas.winfo_width())

        content.bind("<Configure>", _sync)
        canvas.bind("<Configure>", _sync)
        canvas.bind_all("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1*(e.delta/120)), "units"))

        hdr = tk.Frame(content, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 10))
        tk.Label(hdr, text="Playbook Intake", font=("Segoe UI", 16, "bold"), fg="#FFFFFF", bg="#111111").pack(anchor="w")
        tk.Label(
            hdr,
            text="Capture playbook metadata and generate an Obsidian playbook note.",
            font=("Segoe UI", 10),
            fg="#BBBBBB",
            bg="#111111",
        ).pack(anchor="w")

        form = tk.Frame(content, bg="#111111")
        form.pack(fill="both", expand=True, padx=16)
        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        right.pack(side="right", fill="both", expand=True, padx=(10, 0))

        self.vars = {
            "playbook_id": tk.StringVar(),
            "title": tk.StringVar(),
            "status": tk.StringVar(value="draft"),
            "owner_primary": tk.StringVar(value=analyst_default),
            "tags": tk.StringVar(),
            "itid_primary": tk.StringVar(),
        }

        def add_entry(parent, label, var):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_combo(parent, label, var, values):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        add_entry(left, "Playbook ID (numerals only)", self.vars["playbook_id"])
        add_entry(left, "Title", self.vars["title"])
        add_combo(left, "Status", self.vars["status"], ["draft", "active", "deprecated"])

        add_entry(right, "Owner", self.vars["owner_primary"])
        add_entry(right, "Tags (comma-separated)", self.vars["tags"])

        # Supporting ITID dropdown (loaded from 40_ITIDs/01_Definitions)
        itid_values = []
        if self.vault_root is not None:
            try:
                itid_values = list_itid_definition_titles(self.vault_root)
            except Exception:
                itid_values = []
        add_combo(right, "Supporting ITID", self.vars["itid_primary"], itid_values)

        tk.Label(content, text="Purpose", fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w", padx=16, pady=(12, 0))
        self.purpose_txt = tk.Text(content, height=10, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat", wrap="word")
        self.purpose_txt.pack(fill="x", padx=16, pady=(4, 10))

        bar = tk.Frame(self, bg="#111111")
        bar.pack(side="bottom", fill="x", padx=16, pady=(6, 14))

        self.open_note_var = tk.BooleanVar(value=True)
        if self.vault_root is not None:
            tk.Checkbutton(
                bar,
                text="Open created note automatically",
                variable=self.open_note_var,
                bg="#111111",
                fg="#DDDDDD",
                selectcolor="#111111",
                activebackground="#111111",
                activeforeground="#DDDDDD",
            ).pack(side="left")

        def btn(t, cmd):
            return tk.Button(bar, text=t, command=cmd, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right", padx=(8, 0))
        btn("Save & Continue", self._save_continue).pack(side="right")
        btn("Clear Draft", self._clear_draft).pack(side="right", padx=(8, 0))
        btn("Save Draft", self._save_draft).pack(side="right", padx=(8, 0))

    def _apply_draft(self, draft: dict):
        if not draft:
            return
        for k, v in draft.items():
            if k in self.vars and isinstance(v, str):
                self.vars[k].set(v)
        if "purpose" in draft:
            self.purpose_txt.delete("1.0", "end")
            self.purpose_txt.insert("1.0", draft.get("purpose", ""))

    def _collect(self) -> dict:
        data = {k: v.get().strip() for k, v in self.vars.items()}
        data["purpose"] = self.purpose_txt.get("1.0", "end").strip()
        return data

    def _save_draft(self):
        save_playbook_draft(self.app_dir, self._collect())
        messagebox.showinfo("Playbook Intake", "Draft saved.", parent=self)

    def _clear_draft(self):
        clear_playbook_draft(self.app_dir)
        for v in self.vars.values():
            v.set("")
        self.purpose_txt.delete("1.0", "end")
        messagebox.showinfo("Playbook Intake", "Draft cleared.", parent=self)

    def _save_continue(self):
        data = self._collect()
        if not data.get("title"):
            messagebox.showwarning("Playbook Intake", "Title is required.", parent=self)
            return

        save_playbook_draft(self.app_dir, data)

        note_path = None
        if self.vault_root is not None:
            try:
                note_path = create_obsidian_playbook_note(self.vault_root, data)
            except Exception as e:
                messagebox.showwarning("Playbook Intake", f"Saved draft, but failed to create playbook note:\n{e}", parent=self)
            try:
                _log_shift_entry_from_launcher(self.master, self.vault_root, "playbook", data, note_path)
            except Exception:
                pass

        if note_path is not None and self.open_note_var.get():
            try:
                vault_name = (getattr(getattr(self, "master", None), "cfg", {}) or {}).get("paths", {}).get("OBSIDIAN_VAULT", "")
                open_note_in_obsidian(vault_name, self.vault_root, note_path)
            except Exception:
                pass

        try:
            self.on_submit({**data, **({"obsidian_note_path": str(note_path)} if note_path else {})})
        finally:
            self.destroy()



# -----------------------------
# Procedure (SOP) draft persistence
# -----------------------------
def _procedure_store_path(app_dir: Path) -> Path:
    return app_dir / "procedure_draft.json"


def load_procedure_draft(app_dir: Path) -> dict:
    p = _procedure_store_path(app_dir)
    if not p.exists():
        return {}
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return {}


def save_procedure_draft(app_dir: Path, data: dict) -> None:
    p = _procedure_store_path(app_dir)
    p.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def clear_procedure_draft(app_dir: Path) -> None:
    p = _procedure_store_path(app_dir)
    try:
        if p.exists():
            p.unlink()
    except Exception:
        pass


def _resolve_procedure_template(vault_root: Path) -> Path | None:
    """
    Primary: 00_System/00_Templates/02_Operations/SOP Template.md
    """
    primary = vault_root / "00_System" / "00_Templates" / "02_Operations" / "SOP Template.md"
    if primary.exists():
        return primary

    # Fallback: first matching file containing 'sop' and 'template' in 02_Operations
    folder = vault_root / "00_System" / "00_Templates" / "02_Operations"
    if folder.exists():
        for p in sorted(folder.glob("*.md")):
            name = p.name.lower()
            if "sop" in name and "template" in name:
                return p

    return None


def _merge_procedure_frontmatter(template_text: str, proc: dict) -> str:
    """
    Merge procedure fields into YAML frontmatter if present.
    Fields:
      procedure_id, title, status, owner, risk_level, tags, tlp_classification, created, updated
    """
    if not template_text.lstrip().startswith("---"):
        return template_text

    lines = template_text.splitlines(True)

    # locate frontmatter bounds
    i0 = None
    for i, ln in enumerate(lines):
        if ln.strip() == "":
            continue
        if ln.strip() == "---":
            i0 = i
        break
    if i0 is None:
        return template_text

    i1 = None
    for j in range(i0 + 1, len(lines)):
        if lines[j].strip() == "---":
            i1 = j
            break
    if i1 is None:
        return template_text

    fm_lines = lines[i0 + 1 : i1]
    key_line_index = {}
    key_re = re.compile(r"^([A-Za-z0-9_-]+)\s*:\s*(.*)$")
    for idx, ln in enumerate(fm_lines):
        raw = ln.rstrip("\n")
        if not raw.strip() or raw.lstrip().startswith("#"):
            continue
        m = key_re.match(raw)
        if m:
            key_line_index[m.group(1)] = idx

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    procedure_id = (proc.get("procedure_id") or "").strip()
    title = (proc.get("title") or "").strip()
    status = (proc.get("status") or "").strip()
    owner = (proc.get("owner") or "").strip()
    risk_level = (proc.get("risk_level") or "").strip()
    tags = (proc.get("tags") or "").strip()
    tlp = (proc.get("tlp_classification") or "").strip()
    created = (proc.get("created") or "").strip() or now

    updates = {
        "entity_type": "\"procedure\"",
        "procedure_id": f"\"{procedure_id}\"",
        "title": f"\"{title}\"",
        "status": f"\"{status}\"",
        "owner": f"\"{owner}\"",
        "risk_level": f"\"{risk_level}\"",
        "tags": f"\"{tags}\"",
        "tlp_classification": f"\"{tlp}\"",
        "created": f"\"{created}\"",
        "updated": f"\"{now}\"",
    }

    def set_key(k: str, v: str):
        if k in key_line_index:
            fm_lines[key_line_index[k]] = f"{k}: {v}\n"
        else:
            fm_lines.append(f"{k}: {v}\n")

    for k, v in updates.items():
        set_key(k, v)

    new_lines = []
    new_lines.extend(lines[: i0 + 1])
    new_lines.extend(fm_lines)
    new_lines.append("---\n")
    new_lines.extend(lines[i1 + 1 :])
    return "".join(new_lines)


def create_obsidian_procedure_note(vault_root: Path, proc: dict) -> Path:
    """
    Creates a procedure/SOP note under 10_Operations/07_Procedures (created if missing),
    based on the SOP Template.
    """
    base_folder = vault_root / "10_Operations" / "07_Procedures"
    base_folder.mkdir(parents=True, exist_ok=True)

    procedure_id = (proc.get("procedure_id") or "").strip()
    title = (proc.get("title") or "").strip() or "Untitled Procedure"
    status = (proc.get("status") or "").strip() or "draft"
    owner = (proc.get("owner") or "").strip()
    risk_level = (proc.get("risk_level") or "").strip()
    tags = (proc.get("tags") or "").strip()
    tlp = (proc.get("tlp_classification") or "").strip()

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    created = (proc.get("created") or "").strip() or now
    proc["created"] = created

    year = datetime.now().strftime("%Y")
    name_bits = []
    if procedure_id:
        name_bits.append(f"PRC-{year}-{procedure_id}")
    name_bits.append(title)
    slug = _safe_slug(" - ".join(name_bits))
    note_path = base_folder / f"{slug}.md"

    tpl_path = _resolve_procedure_template(vault_root)

    token_map = {
        "procedure_id": procedure_id,
        "title": title,
        "status": status,
        "owner": owner,
        "risk_level": risk_level,
        "tags": tags,
        "tlp_classification": tlp,
        "created": created,
        "updated": now,
    }

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)

        # Replace common token styles ({{key}}, {{KEY}}, ${key})
        for k, v in token_map.items():
            tpl_text = tpl_text.replace(f"{{{{{k}}}}}", v)
            tpl_text = tpl_text.replace(f"{{{{{k.upper()}}}}}", v)
            tpl_text = tpl_text.replace(f"${{{k}}}", v)

        # Title variants
        tpl_text = tpl_text.replace("{{Procedure ID}}", procedure_id)
        tpl_text = tpl_text.replace("{{Procedure Title}}", title)

        if tpl_text.lstrip().startswith("---"):
            md = _merge_procedure_frontmatter(tpl_text, token_map)
        else:
            fm = f"""---
entity_type: procedure
procedure_id: "{procedure_id}"
title: "{title}"
status: "{status}"
owner: "{owner}"
risk_level: "{risk_level}"
tags: "{tags}"
tlp_classification: "{tlp}"
created: "{created}"
updated: "{now}"
---

"""
            md = fm + tpl_text.lstrip()
    else:
        md = f"""---
entity_type: procedure
procedure_id: "{procedure_id}"
title: "{title}"
status: "{status}"
owner: "{owner}"
risk_level: "{risk_level}"
tags: "{tags}"
tlp_classification: "{tlp}"
created: "{created}"
updated: "{now}"
---

"""

    note_path.write_text(md, encoding="utf-8")
    return note_path


class ProcedureIntakeWindow(tk.Toplevel):
    def __init__(self, master, *, app_dir: Path, vault_root: Path | None, on_submit, resume: bool = False):
        super().__init__(master)
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("Procedure Intake")
        self.configure(bg="#111111")
        self.geometry("900x680")
        self.minsize(780, 560)
        self.resizable(True, True)

        self._build_widgets()

        if resume:
            self._apply_draft(load_procedure_draft(self.app_dir))

        self.transient(master)
        self.grab_set()
        self.focus_force()

    def _build_widgets(self):
        cfg = getattr(self.master, "cfg", {}) or {}
        analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")

        container = tk.Frame(self, bg="#111111")
        container.pack(fill="both", expand=True)

        canvas = tk.Canvas(container, bg="#111111", highlightthickness=0)
        vbar = tk.Scrollbar(container, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vbar.set)
        vbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        content = tk.Frame(canvas, bg="#111111")
        win_id = canvas.create_window((0, 0), window=content, anchor="nw")

        def _sync(_=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfigure(win_id, width=canvas.winfo_width())

        content.bind("<Configure>", _sync)
        canvas.bind("<Configure>", _sync)
        canvas.bind_all("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1*(e.delta/120)), "units"))

        hdr = tk.Frame(content, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 10))
        tk.Label(hdr, text="Procedure Intake", font=("Segoe UI", 16, "bold"), fg="#FFFFFF", bg="#111111").pack(anchor="w")
        tk.Label(hdr, text="Capture procedure/SOP metadata and generate an Obsidian note.", font=("Segoe UI", 10),
                 fg="#BBBBBB", bg="#111111").pack(anchor="w")

        form = tk.Frame(content, bg="#111111")
        form.pack(fill="both", expand=True, padx=16)
        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        right.pack(side="right", fill="both", expand=True, padx=(10, 0))

        self.vars = {
            "procedure_id": tk.StringVar(),
            "title": tk.StringVar(),
            "status": tk.StringVar(value="draft"),
            "owner": tk.StringVar(value=analyst_default),
            "risk_level": tk.StringVar(value=""),
            "tags": tk.StringVar(),
            "tlp_classification": tk.StringVar(value="TLP:GREEN"),
        }

        def add_entry(parent, label, var):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF",
                           insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_labeled_combo(parent, label, var, values):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        add_entry(left, "Procedure ID (#### format, numerals only)", self.vars["procedure_id"])
        add_entry(left, "Procedure Title", self.vars["title"])

        add_labeled_combo(left, "Status", self.vars["status"], ["draft", "active", "deprecated"])
        add_entry(right, "Owner", self.vars["owner"])
        add_labeled_combo(right, "Risk Level", self.vars["risk_level"], ["", "low", "medium", "high", "critical"])
        add_entry(right, "Tags (comma-separated)", self.vars["tags"])
        add_labeled_combo(right, "TLP Classification", self.vars["tlp_classification"], ["TLP:WHITE", "TLP:GREEN", "TLP:AMBER", "TLP:RED"])

        bar = tk.Frame(self, bg="#111111")
        bar.pack(side="bottom", fill="x", padx=16, pady=(6, 14))

        self.open_note_var = tk.BooleanVar(value=True)
        if self.vault_root is not None:
            tk.Checkbutton(
                bar, text="Open created note automatically",
                variable=self.open_note_var,
                bg="#111111", fg="#DDDDDD", selectcolor="#111111",
                activebackground="#111111", activeforeground="#DDDDDD",
            ).pack(side="left")

        def btn(t, cmd):
            return tk.Button(bar, text=t, command=cmd, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right", padx=(8, 0))
        btn("Save & Continue", self._save_continue).pack(side="right")
        btn("Clear Draft", self._clear_draft).pack(side="right", padx=(8, 0))
        btn("Save Draft", self._save_draft).pack(side="right", padx=(8, 0))

    def _apply_draft(self, draft: dict):
        if not draft:
            return
        for k, v in draft.items():
            if k in self.vars and isinstance(v, str):
                self.vars[k].set(v)

    def _collect(self) -> dict:
        return {k: v.get().strip() for k, v in self.vars.items()}

    def _save_draft(self):
        save_procedure_draft(self.app_dir, self._collect())
        messagebox.showinfo("Procedure Intake", "Draft saved.", parent=self)

    def _clear_draft(self):
        clear_procedure_draft(self.app_dir)
        for v in self.vars.values():
            v.set("")
        self.vars["status"].set("draft")
        self.vars["tlp_classification"].set("TLP:GREEN")
        messagebox.showinfo("Procedure Intake", "Draft cleared.", parent=self)

    def _save_continue(self):
        data = self._collect()
        if not data.get("title"):
            messagebox.showwarning("Procedure Intake", "Procedure Title is required.", parent=self)
            return
        proc_id = (data.get("procedure_id") or "").strip()
        if proc_id and not re.match(r"^\d{4}$", proc_id):
            messagebox.showwarning(
                "Procedure Intake",
                "Procedure ID must be 4 numerals only (#### format).",
                parent=self,
            )
            return

        save_procedure_draft(self.app_dir, data)

        note_path = None
        if self.vault_root is not None:
            try:
                note_path = create_obsidian_procedure_note(self.vault_root, data)
            except Exception as e:
                messagebox.showwarning("Procedure Intake", f"Saved draft, but failed to create procedure note:\n{e}", parent=self)
            try:
                _log_shift_entry_from_launcher(self.master, self.vault_root, "procedure", data, note_path)
            except Exception:
                pass

        if note_path is not None and self.open_note_var.get():
            try:
                vault_name = (getattr(getattr(self, "master", None), "cfg", {}) or {}).get("paths", {}).get("OBSIDIAN_VAULT", "")
                open_note_in_obsidian(vault_name, self.vault_root, note_path)
            except Exception:
                pass

        try:
            self.on_submit({**data, **({"obsidian_note_path": str(note_path)} if note_path else {})})
        finally:
            self.destroy()



# -----------------------------
# ITID draft persistence
# -----------------------------
def _itid_store_path(app_dir: Path) -> Path:
    return app_dir / "itid_draft.json"


def load_itid_draft(app_dir: Path) -> dict:
    p = _itid_store_path(app_dir)
    if not p.exists():
        return {}
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return {}


def save_itid_draft(app_dir: Path, data: dict) -> None:
    p = _itid_store_path(app_dir)
    p.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def clear_itid_draft(app_dir: Path) -> None:
    p = _itid_store_path(app_dir)
    try:
        if p.exists():
            p.unlink()
    except Exception:
        pass


def _resolve_itid_template(vault_root: Path) -> Path | None:
    """
    Primary: 00_System/00_Templates/04_ITID/ITID_template.md
    """
    p = vault_root / "00_System" / "00_Templates" / "04_ITID" / "ITID_template.md"
    return p if p.exists() else None


def _merge_itid_frontmatter(template_text: str, itid: dict) -> str:
    """
    Merge ITID fields into YAML frontmatter if present.
    """
    if not template_text.lstrip().startswith("---"):
        return template_text

    lines = template_text.splitlines(True)

    i0 = None
    for i, ln in enumerate(lines):
        if ln.strip() == "":
            continue
        if ln.strip() == "---":
            i0 = i
        break
    if i0 is None:
        return template_text

    i1 = None
    for j in range(i0 + 1, len(lines)):
        if lines[j].strip() == "---":
            i1 = j
            break
    if i1 is None:
        return template_text

    fm_lines = lines[i0 + 1 : i1]
    key_line_index = {}
    key_re = re.compile(r"^([A-Za-z0-9_-]+)\s*:\s*(.*)$")
    for idx, ln in enumerate(fm_lines):
        raw = ln.rstrip("\n")
        if not raw.strip() or raw.lstrip().startswith("#"):
            continue
        m = key_re.match(raw)
        if m:
            key_line_index[m.group(1)] = idx

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    created = (itid.get("created") or "").strip() or now

    if len(created) == 10:
        created = created + " 00:00:00"

    if len(created) == 10:
        created = created + " 00:00:00"
    updates = {
        "entity_type": "\"itid\"",
        "itid_id": f"\"{(itid.get('itid_id') or '').strip()}\"",
        "itid_name": f"\"{(itid.get('itid_name') or '').strip()}\"",
        "definition": f"\"{(itid.get('definition') or '').strip()}\"",
        "severity_guidance": f"\"{(itid.get('severity_guidance') or '').strip()}\"",
        "owner": f"\"{(itid.get('owner') or '').strip()}\"",
        "tags": f"\"{(itid.get('tags') or '').strip()}\"",
        "created": f"\"{created}\"",
        "updated": f"\"{now}\"",
    }

    def set_key(k: str, v: str):
        if k in key_line_index:
            fm_lines[key_line_index[k]] = f"{k}: {v}\n"
        else:
            fm_lines.append(f"{k}: {v}\n")

    for k, v in updates.items():
        set_key(k, v)

    new_lines = []
    new_lines.extend(lines[: i0 + 1])
    new_lines.extend(fm_lines)
    new_lines.append("---\n")
    new_lines.extend(lines[i1 + 1 :])
    return "".join(new_lines)


ITID_ID_RE = re.compile(r"^(?:ITID-)?(\d{5})$", re.IGNORECASE)


def _format_itid_id(value: str) -> str:
    raw = (value or "").strip()
    if not raw:
        return ""
    m = ITID_ID_RE.match(raw)
    if m:
        return f"ITID-{m.group(1)}"
    return raw


def create_obsidian_itid_note(vault_root: Path, itid: dict) -> Path:
    """
    Creates an ITID note in 40_ITIDs/01_Definitions using ITID_template.md.
    """
    base_folder = vault_root / "40_ITIDs" / "01_Definitions"
    base_folder.mkdir(parents=True, exist_ok=True)

    itid_id_raw = (itid.get("itid_id") or "").strip()
    itid_id = _format_itid_id(itid_id_raw)
    itid_name = (itid.get("itid_name") or "").strip() or "Untitled ITID"

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    updated = now
    created = (itid.get("created") or "").strip() or datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    # Always stamp updated at save time
    itid["updated"] = now
    itid["created"] = created

    name_bits = []
    if itid_id:
        name_bits.append(itid_id)
    name_bits.append(itid_name)
    slug = _safe_slug(" - ".join(name_bits))
    note_path = base_folder / f"{slug}.md"

    tpl_path = _resolve_itid_template(vault_root)

    token_map = {
        "itid_id": itid_id,
        "itid_name": itid_name,
        "definition": (itid.get("definition") or "").strip(),
        "severity_guidance": (itid.get("severity_guidance") or "").strip(),
        "owner": (itid.get("owner") or "").strip(),
        "tags": (itid.get("tags") or "").strip(),
        "created": created,
        "updated": now,
        # convenience aliases
        "title": itid_name,
    }

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)

        for k, v in token_map.items():
            tpl_text = tpl_text.replace(f"{{{{{k}}}}}", v)
            tpl_text = tpl_text.replace(f"{{{{{k.upper()}}}}}", v)
            tpl_text = tpl_text.replace(f"${{{k}}}", v)

        if tpl_text.lstrip().startswith("---"):
            md = _merge_itid_frontmatter(tpl_text, token_map)
        else:
            fm = f"""---
entity_type: itid
itid_id: "{itid_id}"
itid_name: "{itid_name}"
definition: "{token_map['definition']}"
severity_guidance: "{token_map['severity_guidance']}"
owner: "{token_map['owner']}"
tags: "{token_map['tags']}"
created: "{created}"
updated: "{now}"
---

"""
            md = fm + tpl_text.lstrip()
    else:
        md = f"""---
entity_type: itid
itid_id: "{itid_id}"
itid_name: "{itid_name}"
definition: "{token_map['definition']}"
severity_guidance: "{token_map['severity_guidance']}"
owner: "{token_map['owner']}"
tags: "{token_map['tags']}"
created: "{created}"
updated: "{now}"
---

# {itid_name}

## Definition
{token_map['definition']}

## Default Severity
{token_map['severity_guidance']}
"""

    note_path.write_text(md, encoding="utf-8")
    return note_path


class ITIDIntakeWindow(tk.Toplevel):
    """
    ITID Intake UI (consistent with other intake forms).
    """

    def __init__(self, master, *, app_dir: Path, vault_root: Path | None, on_submit, resume: bool = False):
        super().__init__(master)
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("ITID Intake")
        self.configure(bg="#111111")
        self.geometry("900x720")
        self.minsize(780, 600)
        self.resizable(True, True)

        self._build_widgets()

        if resume:
            self._apply_draft(load_itid_draft(self.app_dir))
        else:
            # defaults
            self.vars["created"].set(datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
            self.vars["updated"].set(datetime.now().strftime("%Y-%m-%d %H:%M:%S"))

        self.transient(master)
        self.grab_set()
        self.focus_force()

    def _build_widgets(self):
        cfg = getattr(self.master, "cfg", {}) or {}
        analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")

        container = tk.Frame(self, bg="#111111")
        container.pack(fill="both", expand=True)

        canvas = tk.Canvas(container, bg="#111111", highlightthickness=0)
        vbar = tk.Scrollbar(container, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vbar.set)
        vbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        content = tk.Frame(canvas, bg="#111111")
        win_id = canvas.create_window((0, 0), window=content, anchor="nw")

        def _sync(_=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfigure(win_id, width=canvas.winfo_width())

        content.bind("<Configure>", _sync)
        canvas.bind("<Configure>", _sync)
        canvas.bind_all("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1*(e.delta/120)), "units"))

        hdr = tk.Frame(content, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 10))
        tk.Label(hdr, text="ITID Intake", font=("Segoe UI", 16, "bold"), fg="#FFFFFF", bg="#111111").pack(anchor="w")
        tk.Label(
            hdr,
            text="Create a new ITID note from the ITID template.",
            font=("Segoe UI", 10),
            fg="#BBBBBB",
            bg="#111111",
        ).pack(anchor="w")

        form = tk.Frame(content, bg="#111111")
        form.pack(fill="both", expand=True, padx=16)
        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        right.pack(side="right", fill="both", expand=True, padx=(10, 0))

        self.vars = {
            "itid_id": tk.StringVar(),
            "itid_name": tk.StringVar(),
            "severity_guidance": tk.StringVar(value="Low"),
            "owner": tk.StringVar(value=analyst_default),
            "tags": tk.StringVar(),
            "created": tk.StringVar(),
            "updated": tk.StringVar(),
        }

        def add_entry(parent, label, var, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF",
                           insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_combo(parent, label, var, values, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        add_entry(left, "ITID ID (5 digits, numerals only)", self.vars["itid_id"], required=True)
        add_entry(left, "ITID Title", self.vars["itid_name"], required=True)
        add_combo(left, "Default Severity", self.vars["severity_guidance"], ["Low", "Medium", "High", "Critical"])
        add_entry(right, "Owner", self.vars["owner"])
        add_entry(right, "Tags", self.vars["tags"])
        add_entry(right, "Date Created", self.vars["created"])
        add_entry(right, "Date Updated", self.vars["updated"])

        def_lbl = tk.Frame(content, bg="#111111")
        def_lbl.pack(anchor="w", padx=16, pady=(12, 0))
        tk.Label(def_lbl, text="Brief Description", fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
        tk.Label(def_lbl, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
        self.definition_txt = tk.Text(content, height=8, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat", wrap="word")
        self.definition_txt.pack(fill="x", padx=16, pady=(4, 10))

        bar = tk.Frame(self, bg="#111111")
        bar.pack(side="bottom", fill="x", padx=16, pady=(6, 14))

        self.open_note_var = tk.BooleanVar(value=True)  # Hidden; always True

        def btn(t, cmd):
            return tk.Button(bar, text=t, command=cmd, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right", padx=(8, 0))
        btn("Save & Continue", self._save_continue).pack(side="right")

    def _apply_draft(self, draft: dict):
        if not draft:
            return
        for k, v in draft.items():
            if k in self.vars and isinstance(v, str):
                self.vars[k].set(v)
        if "definition" in draft:
            self.definition_txt.delete("1.0", "end")
            self.definition_txt.insert("1.0", draft.get("definition", ""))

    def _collect(self) -> dict:
        data = {k: v.get().strip() for k, v in self.vars.items()}
        data["definition"] = self.definition_txt.get("1.0", "end").strip()
        return data

    def _save_draft(self):
        save_itid_draft(self.app_dir, self._collect())
        messagebox.showinfo("ITID Intake", "Draft saved.", parent=self)

    def _clear_draft(self):
        clear_itid_draft(self.app_dir)
        for v in self.vars.values():
            v.set("")
        self.definition_txt.delete("1.0", "end")
        messagebox.showinfo("ITID Intake", "Draft cleared.", parent=self)

    def _save_continue(self):
        data = self._collect()
        if not data.get("itid_id") and not data.get("itid_name"):
            messagebox.showwarning("ITID Intake", "Provide at least ITID ID or ITID Title.", parent=self)
            return
        itid_id_raw = (data.get("itid_id") or "").strip()
        if itid_id_raw and not ITID_ID_RE.match(itid_id_raw):
            messagebox.showwarning(
                "ITID Intake",
                "ITID ID must be exactly 5 digits (numerals only).",
                parent=self,
            )
            return
        if not (data.get("definition") or "").strip():
            messagebox.showwarning("ITID Intake", "Brief Description is required.", parent=self)
            return

        # stamp updated now regardless
        data["updated"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not data.get("created"):
            data["created"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        save_itid_draft(self.app_dir, data)

        note_path = None
        if self.vault_root is not None:
            try:
                note_path = create_obsidian_itid_note(self.vault_root, data)
            except Exception as e:
                messagebox.showwarning("ITID Intake", f"Saved draft, but failed to create ITID note:\n{e}", parent=self)
            try:
                _log_shift_entry_from_launcher(self.master, self.vault_root, "itid", data, note_path)
            except Exception:
                pass

        if note_path is not None and self.open_note_var.get():
            try:
                vault_name = (getattr(getattr(self, "master", None), "cfg", {}) or {}).get("paths", {}).get("OBSIDIAN_VAULT", "")
                open_note_in_obsidian(vault_name, self.vault_root, note_path)
            except Exception:
                pass

        try:
            self.on_submit({**data, **({"obsidian_note_path": str(note_path)} if note_path else {})})
        finally:
            self.destroy()



# -----------------------------
# SLA draft persistence
# -----------------------------
def _sla_store_path(app_dir: Path) -> Path:
    return app_dir / "sla_draft.json"



def _format_faq_id(value: str) -> str:
    raw = (value or "").strip()
    if not raw:
        return ""
    if raw.upper().startswith("FAQ-"):
        return f"FAQ-{raw[4:]}"
    return f"FAQ-{raw}"


def _merge_faq_frontmatter(template_text: str, faq: dict) -> str:
    """Merge FAQ fields into YAML frontmatter if present (flat keys only)."""
    if not template_text.lstrip().startswith("---"):
        return template_text

    lines = template_text.splitlines(True)
    # locate frontmatter bounds
    i0 = None
    for i, ln in enumerate(lines):
        if ln.strip() == "":
            continue
        if ln.strip() == "---":
            i0 = i
        break
    if i0 is None:
        return template_text

    i1 = None
    for j in range(i0 + 1, len(lines)):
        if lines[j].strip() == "---":
            i1 = j
            break
    if i1 is None:
        return template_text

    fm_lines = lines[i0 + 1 : i1]
    key_line_index = {}
    key_re = re.compile(r"^([A-Za-z0-9_-]+)\s*:\s*(.*)$")
    for idx, ln in enumerate(fm_lines):
        raw = ln.rstrip("\n")
        if not raw.strip() or raw.lstrip().startswith("#"):
            continue
        m = key_re.match(raw)
        if m:
            key_line_index[m.group(1)] = idx

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    faq_id = _format_faq_id((faq.get("faq_ID") or "").strip())
    faq_title = (faq.get("faq_title") or "").strip()
    status = (faq.get("status") or "").strip() or "Draft"
    category = (faq.get("category") or "").strip()
    owner = (faq.get("owner") or "").strip()
    tags = (faq.get("tags") or "").strip()
    created = (faq.get("created") or "").strip() or now
    if len(created) == 10:
        created = created + " 00:00:00"

    updates = {
        "entity_type": '"faq"',
        "faq_id": f'"{faq_id}"',
        "faq_title": f'"{faq_title}"',
        "status": f'"{status}"',
        "category": f'"{category}"',
        "owner": f'"{owner}"',
        "tags": f'"{tags}"',
        "created": f'"{created}"',
        "updated": f'"{now}"',
    }

    def set_key(k: str, v: str):
        if k in key_line_index:
            fm_lines[key_line_index[k]] = f"{k}: {v}\n"
        else:
            fm_lines.append(f"{k}: {v}\n")

    for k, v in updates.items():
        set_key(k, v)

    new_lines = []
    new_lines.extend(lines[: i0 + 1])
    new_lines.extend(fm_lines)
    new_lines.append("---\n")
    new_lines.extend(lines[i1 + 1 :])
    return "".join(new_lines)


def create_obsidian_faq_note(vault_root: Path, faq: dict) -> Path:
    """Creates a FAQ note under 60_Knowledge/FAQs using the FAQ template when available."""
    base_folder = vault_root / "60_Knowledge" / "FAQs"
    base_folder.mkdir(parents=True, exist_ok=True)

    faq_id = _format_faq_id((faq.get("faq_ID") or "").strip())
    faq_title = (faq.get("faq_title") or "").strip() or "Untitled FAQ"
    status = (faq.get("status") or "").strip() or "Draft"
    category = (faq.get("category") or "").strip()
    owner = (faq.get("owner") or "").strip()
    tags = (faq.get("tags") or "").strip()

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    created = (faq.get("created") or "").strip() or now
    if len(created) == 10:
        created = created + " 00:00:00"
    updated = now
    faq["updated"] = updated
    faq["created"] = created
    faq["faq_id"] = faq_id

    name_bits = []
    if faq_id:
        name_bits.append(faq_id)
    name_bits.append(faq_title)
    slug = _safe_slug(" - ".join(name_bits))
    note_path = base_folder / f"{slug}.md"

    tpl_path = vault_root / "00_System" / "00_Templates" / "06_Knowledge_Base" / "FAQ.md"

    token_map = {
        "faq_ID": faq_id,
        "faq_title": faq_title,
        "status": status,
        "category": category,
        "owner": owner,
        "tags": tags,
        "created": created,
        "updated": updated,
        # convenience aliases
        "title": faq_title,
        "faq_id": faq_id,
    }

    if tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)

        for k, v in token_map.items():
            tpl_text = tpl_text.replace(f"{{{{{k}}}}}", v)
            tpl_text = tpl_text.replace(f"{{{{{k.upper()}}}}}", v)
            tpl_text = tpl_text.replace(f"${{{k}}}", v)

        if tpl_text.lstrip().startswith("---"):
            md = _merge_faq_frontmatter(tpl_text, token_map)
        else:
            fm = f"""---
entity_type: faq
faq_id: "{faq_id}"
faq_title: "{faq_title}"
status: "{status}"
category: "{category}"
owner: "{owner}"
tags: "{tags}"
created: "{created}"
updated: "{updated}"
---

"""
            md = fm + tpl_text.lstrip()
    else:
        md = f"""---
entity_type: faq
faq_id: "{faq_id}"
faq_title: "{faq_title}"
status: "{status}"
category: "{category}"
owner: "{owner}"
tags: "{tags}"
created: "{created}"
updated: "{updated}"
---

# {faq_title}

## Question
{faq_title}

## Answer
-
"""

    note_path.write_text(md, encoding="utf-8")
    return note_path

HOWTO_ID_RE = re.compile(r"^(?:HOW-)?(\d{5})$", re.IGNORECASE)


def _format_howto_id(value: str) -> str:
    raw = (value or "").strip()
    if not raw:
        return ""
    if raw.upper().startswith("HOW-"):
        return f"HOW-{raw[4:]}"
    return f"HOW-{raw}"


def create_obsidian_howto_note(vault_root: Path, howto: dict) -> Path:
    """Create a How-To note under 60_Knowledge/How-Tos using KB_How-To template when available."""
    base_folder = vault_root / "60_Knowledge" / "How-Tos"
    base_folder.mkdir(parents=True, exist_ok=True)

    how_to_id = _format_howto_id((howto.get("how_to_id") or "").strip())
    title = (howto.get("title") or "").strip() or "Untitled How-To"
    owner = (howto.get("owner") or "").strip()
    status = (howto.get("status") or "").strip() or "Draft"
    tlp = (howto.get("tlp_classification") or "").strip() or "TLP:CLEAR"

    tags_csv = (howto.get("tags") or "").strip()
    tags_list = _normalize_tags(tags_csv)

    created = (howto.get("created") or "").strip() or datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    updated = (howto.get("updated") or "").strip() or created
    howto["created"] = created
    howto["updated"] = updated
    howto["how_to_id"] = how_to_id

    name_bits = []
    if how_to_id:
        name_bits.append(how_to_id)
    name_bits.append(title)
    slug = _safe_slug(" - ".join(name_bits))
    note_path = base_folder / f"{slug}.md"

    tpl_path = vault_root / "00_System" / "00_Templates" / "06_Knowledge_Base" / "KB_How-To.md"

    if tpl_path.exists():
        tpl_text = tpl_path.read_text(encoding="utf-8", errors="ignore")
        yaml_updates = {
            "entity_type": "how_to",
            "how_to_id": how_to_id,
            "title": title,
            "owner": owner,
            "status": status,
            "tlp_classification": tlp,
            "tags": tags_list,
            "created": created,
            "updated": updated,
        }
        md = _apply_yaml_updates(tpl_text, yaml_updates)
        for k, v in [("how_to_id", how_to_id), ("title", title)]:
            md = md.replace(f"{{{{{k}}}}}", v)
            md = md.replace(f"{{{{{k.upper()}}}}}", v)
    else:
        tags_block = "\n".join([f"  - {t}" for t in tags_list]) if tags_list else ""
        md = f"""---
entity_type: how_to
how_to_id: {how_to_id}
title: {title}
owner: {owner}
status: {status}
tlp_classification: {tlp}
tags:
{tags_block}
created: {created}
updated: {updated}
---

# {title}

## Purpose
-

## Procedure
-

## Validation
-

## References
-
"""

    note_path.write_text(md, encoding="utf-8")
    return note_path


MITRE_ACTOR_ID_RE = re.compile(r"^\d{4}$")
INTERNAL_ACTOR_ID_RE = re.compile(r"^(?:TA-)?\d{4}$", re.IGNORECASE)


def _format_threat_actor_id(value: str, *, internal_tracked: bool) -> str:
    """Format actor ID: internal tracked -> TA-####, otherwise -> G####."""
    raw = (value or "").strip()
    if not raw:
        return ""
    digits = raw[4:] if raw.upper().startswith("TA-") else (raw[1:] if raw.upper().startswith("G") and len(raw) > 1 and raw[1:].isdigit() else raw)
    if not (digits.isdigit() and len(digits) == 4):
        digits = raw
    if internal_tracked:
        return f"TA-{digits}"
    return f"G{digits}"


def _resolve_threat_actor_template(vault_root: Path) -> Path | None:
    """
    Primary: 00_System/00_Templates/03_CIPHER/Threat Actor Template.md
    Fallback: first matching file in the same folder containing 'threat actor' and 'template'.
    """
    folder = vault_root / "00_System" / "00_Templates" / "03_CIPHER"
    primary = folder / "Threat Actor Template.md"
    if primary.exists():
        return primary
    if folder.exists():
        for p in sorted(folder.glob("*.md")):
            name = p.name.lower()
            if "threat actor" in name and "template" in name:
                return p
    return None


def create_obsidian_threat_actor_note(vault_root: Path, actor: dict) -> Path:
    """
    Creates a Threat Actor note under 30_CIPHER/03_Threat_Actors using the Threat Actor template when available.
    """
    base_folder = vault_root / "30_CIPHER" / "03_Threat_Actors"
    base_folder.mkdir(parents=True, exist_ok=True)

    actor_id = (actor.get("actor_id") or "").strip()
    actor_name = (actor.get("actor_name") or "").strip() or "Untitled Threat Actor"
    common_name = (actor.get("common_name") or "").strip()
    actor_type = (actor.get("actor_type") or "").strip()
    first_seen = (actor.get("first_seen") or "").strip()
    last_seen = (actor.get("last_seen") or "").strip()
    status = (actor.get("status") or "").strip()
    attribution_confidence = (actor.get("attribution_confidence") or "").strip()

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    created = (actor.get("created") or "").strip() or now
    updated = now

    actor["created"] = created
    actor["updated"] = updated

    name_bits = []
    if actor_id:
        name_bits.append(actor_id)
    name_bits.append(actor_name)
    slug = _safe_slug(" - ".join(name_bits))
    note_path = base_folder / f"{slug}.md"

    tpl_path = _resolve_threat_actor_template(vault_root)
    updates = {
        "entity_type": "threat_actor",
        "actor_id": actor_id,
        "actor_name": actor_name,
        "common_name": common_name,
        "actor_type": actor_type,
        "first_seen": first_seen,
        "last_seen": last_seen,
        "status": status,
        "attribution_confidence": attribution_confidence,
        "created": created,
        "updated": updated,
    }

    def _fmt_list(val, bullets=False):
        if isinstance(val, (list, tuple)):
            items = [str(x).strip() for x in val if str(x).strip()]
            if bullets:
                return "\n".join(f"- {x}" for x in items) if items else "- (none recorded)"
            return ", ".join(items) if items else "(none recorded)"
        return str(val).strip() if val else "(none recorded)"

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)
        md = _apply_yaml_updates(tpl_text, updates)
        token_map = {
            "actor_id": actor_id or "(not specified)",
            "actor_name": actor_name or "Threat Actor",
            "common_name": common_name,
            "actor_type": actor_type,
            "first_seen": first_seen,
            "last_seen": last_seen,
            "first_identified": first_seen or (actor.get("first_identified") or "").strip() or "(not specified)",
            "active_period": last_seen or (actor.get("active_period") or "").strip() or "(not specified)",
            "status": status,
            "attribution_confidence": attribution_confidence or "(not specified)",
            "nation_state": (actor.get("nation_state") or "").strip() or "(not specified)",
            "sponsor_type": _fmt_list(actor.get("sponsor_type")),
            "motivation": _fmt_list(actor.get("motivation")),
            "aliases": _fmt_list(actor.get("aliases")),
            "target_sectors": _fmt_list(actor.get("target_sectors"), bullets=True),
            "target_regions": _fmt_list(actor.get("target_regions"), bullets=True),
            "target_technologies": _fmt_list(actor.get("target_technologies"), bullets=True),
            "ttp_profile": _fmt_list(actor.get("ttp_profile"), bullets=True),
            "malware_used": _fmt_list(actor.get("malware_used"), bullets=True),
            "tools_used": _fmt_list(actor.get("tools_used"), bullets=True),
            "infrastructure_profile": _fmt_list(actor.get("infrastructure_profile"), bullets=True),
            "associated_campaigns": _fmt_list(actor.get("associated_campaigns"), bullets=True),
            "related_incidents": _fmt_list(actor.get("related_incidents"), bullets=True),
            "intel_sources": _fmt_list(actor.get("intel_sources"), bullets=True),
            "risk_level": (actor.get("risk_level") or "").strip() or "(not specified)",
            "threat_score": str(actor.get("threat_score", 1)),
            "tlp_classification": (actor.get("tlp_classification") or "").strip() or "(not specified)",
        }
        for k, v in token_map.items():
            md = md.replace(f"{{{{{k}}}}}", v)
            md = md.replace(f"{{{{{k.upper()}}}}}", v)
    else:
        md = f"""---
entity_type: threat_actor
actor_id: "{actor_id}"
actor_name: "{actor_name}"
common_name: "{common_name}"
actor_type: "{actor_type}"
first_seen: "{first_seen}"
last_seen: "{last_seen}"
status: "{status}"
attribution_confidence: "{attribution_confidence}"
created: "{created}"
updated: "{updated}"
---

# {actor_name}
"""

    note_path.write_text(md, encoding="utf-8")
    return note_path


CAMPAIGN_ID_RE = re.compile(r"^\d{4}$")


def _format_campaign_id(value: str, *, internal_tracked: bool) -> str:
    raw = (value or "").strip()
    if not raw:
        return ""
    if not CAMPAIGN_ID_RE.match(raw):
        return raw
    return f"CAMP-{raw}" if internal_tracked else f"C{raw}"


def _resolve_campaign_template(vault_root: Path) -> Path | None:
    """
    Primary: 00_System/00_Templates/03_CIPHER/Campaign Template.md
    Fallback: first matching file in the same folder containing 'campaign' and 'template'.
    """
    folder = vault_root / "00_System" / "00_Templates" / "03_CIPHER"
    primary = folder / "Campaign Template.md"
    if primary.exists():
        return primary
    if folder.exists():
        for p in sorted(folder.glob("*.md")):
            name = p.name.lower()
            if "campaign" in name and "template" in name:
                return p
    return None


def create_obsidian_campaign_note(vault_root: Path, campaign: dict) -> Path:
    """
    Creates a Campaign note under 30_CIPHER/04_Campaigns using the Campaign template when available.
    """
    base_folder = vault_root / "30_CIPHER" / "04_Campaigns"
    base_folder.mkdir(parents=True, exist_ok=True)

    campaign_id = _format_campaign_id(
        (campaign.get("campaign_id") or "").strip(),
        internal_tracked=bool(campaign.get("internal_tracked", False)),
    )
    campaign_name = (campaign.get("campaign_name") or "").strip() or "Untitled Campaign"
    tlp = (campaign.get("tlp_classification") or "").strip()
    attribution_confidence = (campaign.get("attribution_confidence") or "").strip()
    admirality_source_reliability = (campaign.get("admirality_source_reliability") or "").strip()
    admirality_information_credibility = (campaign.get("admirality_information_credibility") or "").strip()
    first_observed = (campaign.get("first_observed") or "").strip()
    last_observed = (campaign.get("last_observed") or "").strip()
    risk_level = (campaign.get("risk_level") or "").strip()

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    created = (campaign.get("created") or "").strip() or now
    updated = now

    campaign["created"] = created
    campaign["updated"] = updated

    name_bits = []
    if campaign_id:
        name_bits.append(campaign_id)
    name_bits.append(campaign_name)
    slug = _safe_slug(" - ".join(name_bits))
    note_path = base_folder / f"{slug}.md"

    tpl_path = _resolve_campaign_template(vault_root)
    updates = {
        "entity_type": "campaign",
        "campaign_id": campaign_id,
        "campaign_name": campaign_name,
        "tlp_classification": tlp,
        "attribution_confidence": attribution_confidence,
        "admirality_source_reliability": admirality_source_reliability,
        "admirality_information_credibility": admirality_information_credibility,
        "first_observed": first_observed,
        "last_observed": last_observed,
        "risk_level": risk_level,
        "created": created,
        "updated": updated,
    }

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)
        md = _apply_yaml_updates(tpl_text, updates)
    else:
        md = f"""---
entity_type: campaign
campaign_id: "{campaign_id}"
campaign_name: "{campaign_name}"
tlp_classification: "{tlp}"
attribution_confidence: "{attribution_confidence}"
admirality_source_reliability: "{admirality_source_reliability}"
admirality_information_credibility: "{admirality_information_credibility}"
first_observed: "{first_observed}"
last_observed: "{last_observed}"
risk_level: "{risk_level}"
created: "{created}"
updated: "{updated}"
---

# {campaign_name}
"""

    note_path.write_text(md, encoding="utf-8")
    return note_path

def load_sla_draft(app_dir: Path) -> dict:
    p = _sla_store_path(app_dir)
    if not p.exists():
        return {}
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return {}


def save_sla_draft(app_dir: Path, data: dict) -> None:
    p = _sla_store_path(app_dir)
    p.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def clear_sla_draft(app_dir: Path) -> None:
    p = _sla_store_path(app_dir)
    try:
        if p.exists():
            p.unlink()
    except Exception:
        pass


def _resolve_sla_template(vault_root: Path) -> Path | None:
    # Exact requested template path
    primary = vault_root / "00_System" / "00_Templates" / "02_Operations" / "Service Level Agreement Template.md"
    if primary.exists():
        return primary

    # Fallback: search operations templates
    folder = vault_root / "00_System" / "00_Templates" / "02_Operations"
    if folder.exists():
        for p in sorted(folder.glob("*.md")):
            n = p.name.lower()
            if "service level agreement" in n or ("sla" in n and "template" in n):
                return p
    return None


SLA_ID_RE = re.compile(r"^(?:SLA-)?(\d{5})$", re.IGNORECASE)


def _format_sla_id(value: str) -> str:
    raw = (value or "").strip()
    if not raw:
        return ""
    if raw.upper().startswith("SLA-"):
        return f"SLA-{raw[4:]}"
    return f"SLA-{raw}"


def _merge_sla_frontmatter(template_text: str, sla: dict) -> str:
    """
    Merge SLA fields into YAML frontmatter if present.
    Keys:
      sla_id, sla_title, sla_status, sla_type, sla_owner_primary, sla_owning_team,
      sla_provider_contact, sla_customer_contact, tags, tlp_classification, created, updated
    """
    if not template_text.lstrip().startswith("---"):
        return template_text

    lines = template_text.splitlines(True)
    # locate fm bounds
    i0 = None
    for i, ln in enumerate(lines):
        if ln.strip() == "":
            continue
        if ln.strip() == "---":
            i0 = i
        break
    if i0 is None:
        return template_text

    i1 = None
    for j in range(i0 + 1, len(lines)):
        if lines[j].strip() == "---":
            i1 = j
            break
    if i1 is None:
        return template_text

    fm_lines = lines[i0 + 1 : i1]
    key_line_index = {}
    key_re = re.compile(r"^([A-Za-z0-9_-]+)\s*:\s*(.*)$")
    for idx, ln in enumerate(fm_lines):
        raw = ln.rstrip("\n")
        if not raw.strip() or raw.lstrip().startswith("#"):
            continue
        m = key_re.match(raw)
        if m:
            key_line_index[m.group(1)] = idx

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    def _get(k: str) -> str:
        return (sla.get(k) or "").strip()

    created = _get("created") or now
    if len(created) == 10:
        created = created + " 00:00:00"

    updated = now  # always stamp updated at save time

    sla_id = _format_sla_id(_get("sla_id"))
    eff_date = _get("sla_effective_date")
    exp_date = _get("sla_expiration_date")
    if eff_date and " " in eff_date:
        eff_date = eff_date.split()[0]
    if exp_date and " " in exp_date:
        exp_date = exp_date.split()[0]
    updates = {
        "entity_type": "\"sla\"",
        "sla_id": f"\"{sla_id}\"",
        "sla_title": f"\"{_get('sla_title')}\"",
        "sla_status": f"\"{_get('sla_status')}\"",
        "sla_type": f"\"{_get('sla_type')}\"",
        "sla_owner_primary": f"\"{_get('sla_owner_primary')}\"",
        "sla_owning_team": f"\"{_get('sla_owning_team')}\"",
        "sla_provider_contact": f"\"{_get('sla_provider_contact')}\"",
        "sla_customer_contact": f"\"{_get('sla_customer_contact')}\"",
        "sla_effective_date": f"\"{eff_date}\"",
        "sla_expiration_date": f"\"{exp_date}\"",
        "sla_service_description": f"\"{_get('sla_service_description').replace(chr(10), ' ').replace('\"', chr(39))}\"",
        "tags": f"\"{_get('tags')}\"",
        "tlp_classification": f"\"{_get('tlp_classification')}\"",
        "created": f"\"{created}\"",
        "updated": f"\"{updated}\"",
    }

    def set_key(k: str, v: str):
        if k in key_line_index:
            fm_lines[key_line_index[k]] = f"{k}: {v}\n"
        else:
            fm_lines.append(f"{k}: {v}\n")

    for k, v in updates.items():
        set_key(k, v)

    new_lines = []
    new_lines.extend(lines[: i0 + 1])
    new_lines.extend(fm_lines)
    new_lines.append("---\n")
    new_lines.extend(lines[i1 + 1 :])
    return "".join(new_lines)


def create_obsidian_sla_note(vault_root: Path, sla: dict) -> Path:
    """
    Creates an SLA note under 10_Operations/11_Service_Level_Agreements.
    Uses Service Level Agreement Template when available; otherwise writes a minimal stub.
    """
    base_folder = vault_root / "10_Operations" / "11_Service_Level_Agreements"
    base_folder.mkdir(parents=True, exist_ok=True)

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    def _get(k: str) -> str:
        return (sla.get(k) or "").strip()

    created = _get("created") or now
    if len(created) == 10:
        created = created + " 00:00:00"

    updated = now
    sla["created"] = created
    sla["updated"] = updated

    sla_id = _format_sla_id(_get("sla_id"))
    sla["sla_id"] = sla_id
    sla_title = _get("sla_title") or "Untitled SLA"

    name_bits = []
    if sla_id:
        name_bits.append(sla_id)
    name_bits.append(sla_title)
    slug = _safe_slug(" - ".join(name_bits))
    note_path = base_folder / f"{slug}.md"

    tpl_path = _resolve_sla_template(vault_root)

    token_map = {
        "sla_id": sla_id,
        "sla_title": sla_title,
        "sla_status": _get("sla_status"),
        "sla_type": _get("sla_type"),
        "sla_owner_primary": _get("sla_owner_primary"),
        "sla_owning_team": _get("sla_owning_team"),
        "sla_provider_contact": _get("sla_provider_contact"),
        "sla_customer_contact": _get("sla_customer_contact"),
        "sla_effective_date": _get("sla_effective_date").split()[0] if _get("sla_effective_date") else "",
        "sla_expiration_date": _get("sla_expiration_date").split()[0] if _get("sla_expiration_date") else "",
        "sla_service_description": _get("sla_service_description"),
        "tags": _get("tags"),
        "tlp_classification": _get("tlp_classification"),
        "created": created,
        "updated": updated,
        # common aliases
        "title": sla_title,
        "status": _get("sla_status"),
    }

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)

        for k, v in token_map.items():
            tpl_text = tpl_text.replace(f"{{{{{k}}}}}", v)
            tpl_text = tpl_text.replace(f"{{{{{k.upper()}}}}}", v)
            tpl_text = tpl_text.replace(f"${{{k}}}", v)

        if tpl_text.lstrip().startswith("---"):
            md = _merge_sla_frontmatter(tpl_text, token_map)
        else:
            fm = f"""---
entity_type: sla
sla_id: "{sla_id}"
sla_title: "{sla_title}"
sla_status: "{token_map['sla_status']}"
sla_type: "{token_map['sla_type']}"
sla_owner_primary: "{token_map['sla_owner_primary']}"
sla_owning_team: "{token_map['sla_owning_team']}"
sla_provider_contact: "{token_map['sla_provider_contact']}"
sla_customer_contact: "{token_map['sla_customer_contact']}"
tags: "{token_map['tags']}"
tlp_classification: "{token_map['tlp_classification']}"
created: "{created}"
updated: "{updated}"
---

"""
            md = fm + tpl_text.lstrip()
    else:
        md = f"""---
entity_type: sla
sla_id: "{sla_id}"
sla_title: "{sla_title}"
sla_status: "{token_map['sla_status']}"
sla_type: "{token_map['sla_type']}"
sla_owner_primary: "{token_map['sla_owner_primary']}"
sla_owning_team: "{token_map['sla_owning_team']}"
sla_provider_contact: "{token_map['sla_provider_contact']}"
sla_customer_contact: "{token_map['sla_customer_contact']}"
tags: "{token_map['tags']}"
tlp_classification: "{token_map['tlp_classification']}"
created: "{created}"
updated: "{updated}"
---

# {sla_title}
"""
    note_path.write_text(md, encoding="utf-8")
    return note_path


def _resolve_mitre_ttp_template(vault_root: Path) -> Path | None:
    """
    Primary: 00_System/00_Templates/01_Entities/MITRE TTP Template.md
    Fallback: first matching file in the same folder containing 'mitre' and 'ttp' or 'technique'.
    """
    folder = vault_root / "00_System" / "00_Templates" / "01_Entities"
    primary = folder / "MITRE TTP Template.md"
    if primary.exists():
        return primary
    if folder.exists():
        for p in sorted(folder.glob("*.md")):
            name = p.name.lower()
            if ("mitre" in name and "ttp" in name) or ("technique" in name and "template" in name):
                return p
    return None


TTP_ID_RE = re.compile(r"^\d{4}$")
SUBTTP_ID_RE = re.compile(r"^\d{4}\.\d{3}$")


def _format_ttp_id(value: str) -> str:
    raw = (value or "").strip()
    if not raw:
        return ""
    if raw.upper().startswith("T"):
        raw = raw[1:]
    if TTP_ID_RE.match(raw):
        return f"T{raw}"
    return raw


def _format_subttp_id(value: str) -> str:
    raw = (value or "").strip()
    if not raw:
        return ""
    if raw.upper().startswith("T"):
        raw = raw[1:]
    if SUBTTP_ID_RE.match(raw):
        return f"T{raw}"
    return raw


def create_obsidian_mitre_ttp_note(vault_root: Path, ttp: dict) -> Path:
    """
    Creates a MITRE ATT&CK TTP note under 20_Entities/07_TTPs using the MITRE TTP template when available.
    """
    base_folder = vault_root / "20_Entities" / "07_TTPs"
    base_folder.mkdir(parents=True, exist_ok=True)

    technique_id = _format_ttp_id((ttp.get("technique_id") or "").strip())
    subtechnique_id = _format_subttp_id((ttp.get("subtechnique_id") or "").strip())
    technique_name = (ttp.get("technique_name") or "").strip() or "Untitled TTP"
    mitre_version = (ttp.get("mitre_version") or "").strip()
    detection_priority = (ttp.get("detection_priority") or "").strip()
    threat_score = (ttp.get("threat_score") or "").strip()

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    created = (ttp.get("created") or "").strip() or now
    updated = now

    ttp["created"] = created
    ttp["updated"] = updated
    ttp["technique_id"] = technique_id
    ttp["subtechnique_id"] = subtechnique_id

    primary_id = subtechnique_id or technique_id
    name_bits = [primary_id] if primary_id else []
    name_bits.append(technique_name)
    slug = _safe_slug(" - ".join(name_bits))
    note_path = base_folder / f"{slug}.md"

    tpl_path = _resolve_mitre_ttp_template(vault_root)
    updates = {
        "entity_type": "mitre_ttp",
        "technique_id": technique_id,
        "subtechnique_id": subtechnique_id,
        "technique_name": technique_name,
        "mitre_version": mitre_version,
        "detection_priority": detection_priority,
        "threat_score": threat_score,
        "created": created,
        "updated": updated,
    }

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)
        md = _apply_yaml_updates(tpl_text, updates)
    else:
        md = f"""---
entity_type: mitre_ttp
technique_id: "{technique_id}"
subtechnique_id: "{subtechnique_id}"
technique_name: "{technique_name}"
mitre_version: "{mitre_version}"
detection_priority: "{detection_priority}"
threat_score: "{threat_score}"
created: "{created}"
updated: "{updated}"
---

# {technique_name}
"""

    note_path.write_text(md, encoding="utf-8")
    return note_path


def _resolve_malware_template(vault_root: Path) -> Path | None:
    """
    Primary: 00_System/00_Templates/03_CIPHER/Malware Template.md
    Fallback: first matching file in the same folder containing 'malware' and 'template'.
    """
    folder = vault_root / "00_System" / "00_Templates" / "03_CIPHER"
    primary = folder / "Malware Template.md"
    if primary.exists():
        return primary
    if folder.exists():
        for p in sorted(folder.glob("*.md")):
            name = p.name.lower()
            if "malware" in name and "template" in name:
                return p
    return None


MALWARE_ID_RE = re.compile(r"^\d{4}$")


def _format_malware_id(value: str, *, internal_tracked: bool) -> str:
    raw = (value or "").strip()
    if not raw:
        return ""
    if not MALWARE_ID_RE.match(raw):
        return raw
    return f"MAL-{raw}" if internal_tracked else f"S{raw}"


def create_obsidian_malware_note(vault_root: Path, malware: dict) -> Path:
    """
    Creates a Malware note under 30_CIPHER/05_Malware using the Malware template when available.
    """
    base_folder = vault_root / "30_CIPHER" / "05_Malware"
    base_folder.mkdir(parents=True, exist_ok=True)

    malware_name = (malware.get("malware_name") or "").strip() or "Untitled Malware"
    malware_id = _format_malware_id(
        (malware.get("malware_id") or "").strip(),
        internal_tracked=bool(malware.get("internal_tracked", False)),
    )
    malware_type = (malware.get("malware_type") or "").strip()
    category = (malware.get("category") or "").strip()
    first_seen = (malware.get("first_seen") or "").strip()
    active_status = (malware.get("active_status") or "").strip()
    attribution_confidence = (malware.get("attribution_confidence") or "").strip()
    risk_level = (malware.get("risk_level") or "").strip()
    threat_score = (malware.get("threat_score") or "").strip()
    tlp = (malware.get("tlp_classification") or "").strip()

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    created = (malware.get("created") or "").strip() or now
    updated = now

    malware["created"] = created
    malware["updated"] = updated
    malware["malware_id"] = malware_id

    name_bits = []
    if malware_id:
        name_bits.append(malware_id)
    name_bits.append(malware_name)
    slug = _safe_slug(" - ".join(name_bits))
    note_path = base_folder / f"{slug}.md"

    tpl_path = _resolve_malware_template(vault_root)
    updates = {
        "entity_type": "malware",
        "malware_name": malware_name,
        "malware_id": malware_id,
        "malware_type": malware_type,
        "category": category,
        "first_seen": first_seen,
        "active_status": active_status,
        "attribution_confidence": attribution_confidence,
        "risk_level": risk_level,
        "threat_score": threat_score,
        "tlp_classification": tlp,
        "created": created,
        "updated": updated,
    }

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)
        md = _apply_yaml_updates(tpl_text, updates)
    else:
        md = f"""---
entity_type: malware
malware_name: "{malware_name}"
malware_id: "{malware_id}"
malware_type: "{malware_type}"
category: "{category}"
first_seen: "{first_seen}"
active_status: "{active_status}"
attribution_confidence: "{attribution_confidence}"
risk_level: "{risk_level}"
threat_score: "{threat_score}"
tlp_classification: "{tlp}"
created: "{created}"
updated: "{updated}"
---

# {malware_name}
"""

    note_path.write_text(md, encoding="utf-8")
    return note_path


IOC_TEMPLATE_MAP = {
    "domain": "IOC - Domain Template.md",
    "email_address": "IOC - Email Address Template.md",
    "file_path": "IOC - File Path Template.md",
    "hash": "IOC - Hash Template.md",
    "ip": "IOC - IP Template.md",
    "registry_key": "IOC - Registry Key Template.md",
    "url": "IOC - URL Template.md",
}


def _resolve_ioc_template(vault_root: Path, ioc_type: str) -> Path | None:
    """
    Resolve IOC template based on type within 00_System/00_Templates/01_Entities.
    """
    folder = vault_root / "00_System" / "00_Templates" / "01_Entities"
    if not folder.exists():
        return None
    filename = IOC_TEMPLATE_MAP.get((ioc_type or "").strip().lower(), "")
    if filename:
        candidate = folder / filename
        if candidate.exists():
            return candidate
    return None


def create_obsidian_ioc_note(vault_root: Path, ioc: dict) -> Path:
    """
    Creates an IOC note under 20_Entities/06_IOCs using an IOC template when available.
    """
    base_folder = vault_root / "20_Entities" / "06_IOCs"
    base_folder.mkdir(parents=True, exist_ok=True)

    ioc_type = (ioc.get("ioc_type") or "").strip()
    value = (ioc.get("value") or "").strip()
    classification = (ioc.get("classification") or "").strip()
    confidence = (ioc.get("confidence") or "").strip()
    last_seen = (ioc.get("last_seen") or "").strip()
    notes = (ioc.get("notes") or "").strip()

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    created = (ioc.get("created") or "").strip() or now
    updated = now

    ioc["created"] = created
    ioc["updated"] = updated

    name_bits = []
    if ioc_type:
        name_bits.append(ioc_type)
    if value:
        name_bits.append(value)
    slug = _safe_slug(" - ".join(name_bits)) or "IOC"
    note_path = base_folder / f"{slug}.md"

    tpl_path = _resolve_ioc_template(vault_root, ioc_type)
    updates = {
        "entity_type": "ioc",
        "ioc_type": ioc_type,
        "value": value,
        "classification": classification,
        "confidence": confidence,
        "last_seen": last_seen,
        "notes": notes,
        "created": created,
        "updated": updated,
    }

    if tpl_path and tpl_path.exists():
        tpl_text = _read_text_safe(tpl_path)
        md = _apply_yaml_updates(tpl_text, updates)
    else:
        md = f"""---
entity_type: ioc
ioc_type: "{ioc_type}"
value: "{value}"
classification: "{classification}"
confidence: "{confidence}"
last_seen: "{last_seen}"
created: "{created}"
updated: "{updated}"
---

## Notes
{notes}
"""

    note_path.write_text(md, encoding="utf-8")
    return note_path


class FAQIntakeWindow(tk.Toplevel):
    def __init__(self, master, *, app_dir: Path, vault_root: Path, on_submit):
        super().__init__(master)
        self.master = master
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("FAQ Intake")
        self.configure(bg="#111111")
        self.geometry("880x720")
        self.minsize(820, 620)

        self.draft_path = self.app_dir / "faq_draft.json"

        self.vars = {
            "faq_ID": tk.StringVar(),
            "faq_title": tk.StringVar(),
            "status": tk.StringVar(value="Draft"),
            "category": tk.StringVar(value="Operations"),
            "owner": tk.StringVar(),
            "tags": tk.StringVar(),
            "created": tk.StringVar(),
            "updated": tk.StringVar(),
            "open_after_save": tk.BooleanVar(value=True),
        }

        cfg = getattr(self.master, "cfg", {}) or {}
        analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")
        self.vars["owner"].set(analyst_default)

        self._build_widgets()
        self._load_draft()

    def _build_widgets(self):
        # Header
        hdr = tk.Frame(self, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 8))
        tk.Label(
            hdr,
            text="New FAQ",
            font=("Segoe UI", 16, "bold"),
            fg="#FFFFFF",
            bg="#111111",
            wraplength=820,
            justify="left",
        ).pack(anchor="w", fill="x")
        tk.Label(
            hdr,
            text="Create a FAQ note from the knowledge base template.",
            font=("Segoe UI", 10),
            fg="#BBBBBB",
            bg="#111111",
            wraplength=820,
            justify="left",
        ).pack(anchor="w", fill="x", pady=(4, 0))

        # Auto-flag diagnostics (watchlists)
        try:
            stats = self.state.get("rss_autoflag_stats", {}) if isinstance(self.state.get("rss_autoflag_stats", {}), dict) else {}
            kwc = int(stats.get("keywords_loaded", 0) or 0)
            afc = int(stats.get("auto_flagged", 0) or 0)
            diag = f"Auto-flag: {afc} matched • Watchlist keywords loaded: {kwc}"
            tk.Label(hdr, text=diag, fg="#888888", bg="#111111", font=("Segoe UI", 9), wraplength=820, justify="left").pack(anchor="w", fill="x", pady=(2, 0))
        except Exception:
            pass


        # Form container (no scrollbar)
        form = tk.Frame(self, bg="#111111")
        form.pack(fill="both", expand=True, padx=10, pady=(0, 10))

        # Helpers
        def add_entry(parent, label, var):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10), wraplength=380, justify="left").pack(anchor="w", fill="x")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_combo(parent, label, var, values):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10), wraplength=380, justify="left").pack(anchor="w", fill="x")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        def add_hint(parent, text_):
            tk.Label(parent, text=text_, fg="#888888", bg="#111111", font=("Segoe UI", 9), wraplength=760, justify="left").pack(anchor="w", pady=(0, 6))

        # Layout two columns
        cols = tk.Frame(form, bg="#111111")
        cols.pack(fill="x", padx=10, pady=10)

        left = tk.Frame(cols, bg="#111111")
        right = tk.Frame(cols, bg="#111111")
        cols.grid_columnconfigure(0, weight=1)
        cols.grid_columnconfigure(1, weight=0, minsize=125)
        left.grid(row=0, column=0, sticky="nsew", padx=(0, 10))
        right.grid(row=0, column=1, sticky="nsew", padx=(10, 0))
        right.configure(width=125)
        right.grid_propagate(False)

        add_entry(left, "FAQ ID", self.vars["faq_ID"])
        add_entry(left, "FAQ Question", self.vars["faq_title"])

        add_combo(right, "Status", self.vars["status"], ["Draft", "Active", "Deprecated"])
        add_combo(right, "Category", self.vars["category"], ["Operations", "Security", "Tooling", "Process", "Governance", "Other"])

        add_entry(right, "Owner", self.vars["owner"])
        add_entry(right, "Tags (comma-separated)", self.vars["tags"])

        add_entry(left, "Created Date", self.vars["created"])
        add_hint(left, "If blank, Created Date is set automatically. You can also enter YYYY-MM-DD or YYYY-MM-DD HH:MM:SS.")
        # updated is stamped automatically; show read-only field for transparency
        add_entry(left, "Updated Date", self.vars["updated"])
        self.vars["updated"].set(datetime.now().strftime("%Y-%m-%d %H:%M:%S"))

        # Options
        opts = tk.Frame(form, bg="#111111")
        opts.pack(fill="x", padx=20, pady=(8, 0))
        tk.Checkbutton(
            opts,
            text="Open note in Obsidian after save",
            variable=self.vars["open_after_save"],
            bg="#111111",
            fg="#DDDDDD",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#DDDDDD",
        ).pack(anchor="w")

        # Action bar
        bar = tk.Frame(self, bg="#0D0D0D")
        bar.pack(fill="x", side="bottom")
        inner = tk.Frame(bar, bg="#0D0D0D")
        inner.pack(fill="x", padx=12, pady=10)

        def btn(txt_, cmd, primary=False):
            bg = "#2A2A2A" if not primary else "#3A3A3A"
            return tk.Button(inner, text=txt_, command=cmd, bg=bg, fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Save Draft", self._save_draft).pack(side="left")
        btn("Clear Draft", self._clear_draft).pack(side="left", padx=(8, 0))
        btn("Cancel", self.destroy).pack(side="right")
        btn("Save & Continue", self._save_continue, primary=True).pack(side="right", padx=(0, 8))

    def _collect(self) -> dict:
        data = {}
        for k, v in self.vars.items():
            if isinstance(v, tk.BooleanVar):
                data[k] = bool(v.get())
            else:
                data[k] = (v.get() or "").strip()

        # Stamp updated now
        data["updated"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not data.get("created"):
            data["created"] = data["updated"]
        return data

    def _load_draft(self):
        try:
            if self.draft_path.exists():
                d = json.loads(self.draft_path.read_text(encoding="utf-8"))
                for k, val in d.items():
                    if k in self.vars and not isinstance(self.vars[k], tk.BooleanVar):
                        self.vars[k].set(val)
                    elif k in self.vars and isinstance(self.vars[k], tk.BooleanVar):
                        self.vars[k].set(bool(val))
        except Exception:
            pass
        try:
            cfg = getattr(self.master, "cfg", {}) or {}
            analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")
            if analyst_default:
                self.vars["owner"].set(analyst_default)
        except Exception:
            pass

    def _save_draft(self):
        try:
            data = self._collect()
            # Do not overwrite updated on draft save; keep current snapshot
            self.draft_path.write_text(json.dumps(data, indent=2), encoding="utf-8")
            messagebox.showinfo("FAQ Draft", "Draft saved.", parent=self)
        except Exception as e:
            messagebox.showerror("FAQ Draft", f"Failed to save draft.\n\nError: {e}", parent=self)

    def _clear_draft(self):
        try:
            if self.draft_path.exists():
                self.draft_path.unlink()
            # Clear fields
            for k, v in self.vars.items():
                if isinstance(v, tk.BooleanVar):
                    v.set(False)
                else:
                    v.set("")
            self.vars["status"].set("Draft")
            self.vars["category"].set("Operations")
            self.vars["open_after_save"].set(True)
            self.vars["updated"].set(datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
            try:
                cfg = getattr(self.master, "cfg", {}) or {}
                analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")
                self.vars["owner"].set(analyst_default)
            except Exception:
                pass
        except Exception:
            pass

    def _save_continue(self):
        data = self._collect()
        if not data.get("faq_title"):
            messagebox.showerror("FAQ Intake", "FAQ Question is required.", parent=self)
            return

        try:
            if callable(self.on_submit):
                self.on_submit(data)
        finally:
            # Clear draft after successful submit best-effort
            try:
                if self.draft_path.exists():
                    self.draft_path.unlink()
            except Exception:
                pass
            self.destroy()


class GoalIntakeWindow(tk.Toplevel):
    def __init__(self, master, *, app_dir: Path, vault_root: Path, on_submit):
        super().__init__(master)
        self.master = master
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("Goal Intake")
        self.configure(bg="#111111")
        self.geometry("900x760")
        self.minsize(820, 640)

        self.draft_path = self.app_dir / "goal_draft.json"

        self.vars = {
            "goal_id": tk.StringVar(),
            "title": tk.StringVar(),
            "status": tk.StringVar(value="planned"),
            "owner": tk.StringVar(),
            "area": tk.StringVar(),
            "start_date": tk.StringVar(),
            "end_date": tk.StringVar(),
            "cadence": tk.StringVar(value="quarterly"),
            "related_projects": tk.StringVar(),
            "related_tasks": tk.StringVar(),
            "tags": tk.StringVar(),
            "open_after_save": tk.BooleanVar(value=True),
        }

        cfg = getattr(self.master, "cfg", {}) or {}
        analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")
        self.vars["owner"].set(analyst_default)

        self._build_widgets()
        self._load_draft()

    def _build_widgets(self):
        hdr = tk.Frame(self, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 8))
        tk.Label(
            hdr,
            text="New Goal",
            font=("Segoe UI", 16, "bold"),
            fg="#FFFFFF",
            bg="#111111",
            wraplength=820,
            justify="left",
        ).pack(anchor="w", fill="x")
        tk.Label(
            hdr,
            text="Create a goal note and track progress over time.",
            font=("Segoe UI", 10),
            fg="#BBBBBB",
            bg="#111111",
            wraplength=820,
            justify="left",
        ).pack(anchor="w", fill="x", pady=(4, 0))

        form = tk.Frame(self, bg="#111111")
        form.pack(fill="both", expand=True, padx=10, pady=(0, 10))

        def add_entry(parent, label, var, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10), wraplength=380, justify="left").pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_combo(parent, label, var, values, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10), wraplength=380, justify="left").pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        def add_hint(parent, text_):
            tk.Label(parent, text=text_, fg="#888888", bg="#111111", font=("Segoe UI", 9), wraplength=760, justify="left").pack(anchor="w", pady=(0, 6))

        cols = tk.Frame(form, bg="#111111")
        cols.pack(fill="x", padx=10, pady=10)

        left = tk.Frame(cols, bg="#111111")
        right = tk.Frame(cols, bg="#111111")
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        right.pack(side="left", fill="both", expand=True, padx=(10, 0))

        add_entry(left, "Goal ID (auto if blank)", self.vars["goal_id"])
        add_entry(left, "Goal Title", self.vars["title"], required=True)
        add_entry(left, "Owner", self.vars["owner"])
        add_entry(left, "Area", self.vars["area"])
        add_entry(left, "Start Date (YYYY-MM-DD)", self.vars["start_date"])
        add_entry(left, "End Date (YYYY-MM-DD)", self.vars["end_date"])

        add_combo(right, "Status", self.vars["status"], ["planned", "active", "blocked", "complete", "archived"])
        add_combo(right, "Cadence", self.vars["cadence"], ["daily", "weekly", "monthly", "quarterly", "yearly"])
        add_entry(right, "Related Projects (comma or line separated)", self.vars["related_projects"])
        add_entry(right, "Related Tasks (comma or line separated)", self.vars["related_tasks"])
        add_entry(right, "Tags (comma or line separated)", self.vars["tags"])

        # open_after_save stays True (hidden)

        bar = tk.Frame(self, bg="#0D0D0D")
        bar.pack(fill="x", side="bottom")
        inner = tk.Frame(bar, bg="#0D0D0D")
        inner.pack(fill="x", padx=12, pady=10)

        def btn(txt_, cmd, primary=False):
            bg = "#2A2A2A" if not primary else "#3A3A3A"
            return tk.Button(inner, text=txt_, command=cmd, bg=bg, fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right")
        btn("Save & Continue", self._save_continue, primary=True).pack(side="right", padx=(0, 8))

    def _collect(self) -> dict:
        data = {}
        for k, v in self.vars.items():
            if isinstance(v, tk.BooleanVar):
                data[k] = bool(v.get())
            else:
                data[k] = (v.get() or "").strip()

        data["updated"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not data.get("created"):
            data["created"] = data["updated"]
        return data

    def _load_draft(self):
        try:
            if self.draft_path.exists():
                d = json.loads(self.draft_path.read_text(encoding="utf-8"))
                for k, val in d.items():
                    if k in self.vars and not isinstance(self.vars[k], tk.BooleanVar):
                        self.vars[k].set(val)
                    elif k in self.vars and isinstance(self.vars[k], tk.BooleanVar):
                        self.vars[k].set(bool(val))
        except Exception:
            pass
        try:
            cfg = getattr(self.master, "cfg", {}) or {}
            analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")
            if analyst_default:
                self.vars["owner"].set(analyst_default)
        except Exception:
            pass
        try:
            if not self.vars["goal_id"].get().strip():
                self.vars["goal_id"].set(_next_goal_id_from_vault(self.vault_root))
        except Exception:
            pass
        try:
            if not self.vars["start_date"].get().strip():
                self.vars["start_date"].set(datetime.now().strftime("%Y-%m-%d"))
        except Exception:
            pass

    def _save_draft(self):
        try:
            data = self._collect()
            self.draft_path.write_text(json.dumps(data, indent=2), encoding="utf-8")
            messagebox.showinfo("Goal Draft", "Draft saved.", parent=self)
        except Exception as e:
            messagebox.showerror("Goal Draft", f"Failed to save draft.\n\nError: {e}", parent=self)

    def _clear_draft(self):
        try:
            if self.draft_path.exists():
                self.draft_path.unlink()
            for k, v in self.vars.items():
                if isinstance(v, tk.BooleanVar):
                    v.set(False)
                else:
                    v.set("")
            self.vars["status"].set("planned")
            self.vars["cadence"].set("quarterly")
            self.vars["open_after_save"].set(True)
            try:
                cfg = getattr(self.master, "cfg", {}) or {}
                analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")
                self.vars["owner"].set(analyst_default)
            except Exception:
                pass
        except Exception:
            pass

    def _save_continue(self):
        data = self._collect()
        if not data.get("title"):
            messagebox.showerror("Goal Intake", "Goal Title is required.", parent=self)
            return

        start_raw = (data.get("start_date") or "").strip()
        end_raw = (data.get("end_date") or "").strip()
        if start_raw and end_raw:
            try:
                start_d = datetime.strptime(start_raw, "%Y-%m-%d").date()
                end_d = datetime.strptime(end_raw, "%Y-%m-%d").date()
                if end_d <= start_d:
                    messagebox.showwarning(
                        "Goal Intake",
                        "End Date must be greater than Start Date.",
                        parent=self,
                    )
                    return
            except ValueError:
                pass

        try:
            if callable(self.on_submit):
                self.on_submit(data)
        finally:
            try:
                if self.draft_path.exists():
                    self.draft_path.unlink()
            except Exception:
                pass
            self.destroy()


class HowToIntakeWindow(tk.Toplevel):
    def __init__(self, master, *, app_dir: Path, vault_root: Path, cfg: dict, on_submit):
        super().__init__(master)
        self.master = master
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.cfg = cfg or {}
        self.on_submit = on_submit

        self.title("How-To Intake")
        self.configure(bg="#111111")
        self.geometry("900x720")
        self.minsize(780, 600)

        self.draft_path = self.app_dir / "howto_draft.json"

        self.vars = {
            "how_to_id": tk.StringVar(),
            "title": tk.StringVar(),
            "status": tk.StringVar(value="Draft"),
            "tlp_classification": tk.StringVar(value="TLP:CLEAR"),
            "tags": tk.StringVar(),
            "owner": tk.StringVar(),
            "created": tk.StringVar(),
            "updated": tk.StringVar(),
            "open_after_save": tk.BooleanVar(value=True),
        }

        self._build_widgets()
        self._load_draft()

        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        self.vars["created"].set(now)
        self.vars["updated"].set(now)

    def _build_widgets(self):
        cfg = self.cfg or {}
        analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")
        if analyst_default:
            self.vars["owner"].set(analyst_default)

        container = tk.Frame(self, bg="#111111")
        container.pack(fill="both", expand=True)

        canvas = tk.Canvas(container, bg="#111111", highlightthickness=0)
        vbar = tk.Scrollbar(container, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vbar.set)
        vbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        content = tk.Frame(canvas, bg="#111111")
        win_id = canvas.create_window((0, 0), window=content, anchor="nw")

        def _sync(_=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfigure(win_id, width=canvas.winfo_width())

        content.bind("<Configure>", _sync)
        canvas.bind("<Configure>", _sync)
        canvas.bind_all("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1 * (e.delta / 120)), "units"))

        hdr = tk.Frame(content, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 10))
        tk.Label(hdr, text="New How-To", font=("Segoe UI", 16, "bold"), fg="#FFFFFF", bg="#111111").pack(anchor="w")
        tk.Label(
            hdr,
            text="Create a How-To note in Obsidian using your template.",
            font=("Segoe UI", 10),
            fg="#BBBBBB",
            bg="#111111",
        ).pack(anchor="w", pady=(2, 0))

        form = tk.Frame(content, bg="#111111")
        form.pack(fill="both", expand=True, padx=16)

        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        right.pack(side="right", fill="both", expand=True, padx=(10, 0))

        def add_entry(parent, label, var, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=7)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10), wraplength=760, justify="left").pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=7)
            return ent

        def add_combo(parent, label, var, values, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=7)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10), wraplength=760, justify="left").pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=4)
            return cb

        add_entry(left, "How-to ID (5 digits, numerals only)", self.vars["how_to_id"], required=True)
        add_entry(left, "Title", self.vars["title"], required=True)
        add_combo(left, "Status", self.vars["status"], ["Draft", "Published", "Deprecated"])
        add_combo(left, "TLP Classification", self.vars["tlp_classification"], ["TLP:CLEAR", "TLP:GREEN", "TLP:AMBER", "TLP:AMBER+STRICT", "TLP:RED"])
        add_entry(left, "Tags", self.vars["tags"])

        add_entry(right, "Owner", self.vars["owner"])
        add_entry(right, "Created Date", self.vars["created"])
        add_entry(right, "Updated Date", self.vars["updated"])

        bar = tk.Frame(self, bg="#111111")
        bar.pack(side="bottom", fill="x", padx=16, pady=(6, 14))

        # open_after_save stays True (hidden)

        def btn(t, cmd):
            return tk.Button(bar, text=t, command=cmd, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right", padx=(8, 0))
        btn("Save & Continue", self._submit).pack(side="right")

    def _collect(self) -> dict:
        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        return {
            "how_to_id": (self.vars["how_to_id"].get() or "").strip(),
            "title": (self.vars["title"].get() or "").strip(),
            "owner": (self.vars["owner"].get() or "").strip(),
            "status": (self.vars["status"].get() or "").strip(),
            "tlp_classification": (self.vars["tlp_classification"].get() or "").strip(),
            "tags": (self.vars["tags"].get() or "").strip(),
            "created": now,
            "updated": now,
            "open_after_save": bool(self.vars["open_after_save"].get()),
        }

    def _save_draft(self):
        try:
            save_json(self.draft_path, self._collect())
        except Exception:
            pass

    def _clear_draft(self):
        try:
            if self.draft_path.exists():
                self.draft_path.unlink()
        except Exception:
            pass
        try:
            for k, v in self.vars.items():
                if isinstance(v, tk.BooleanVar):
                    v.set(False)
                else:
                    v.set("")
            self.vars["status"].set("Draft")
            self.vars["tlp_classification"].set("TLP:CLEAR")
            self.vars["open_after_save"].set(True)
            cfg = self.cfg or {}
            analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")
            if analyst_default:
                self.vars["owner"].set(analyst_default)
            now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            self.vars["created"].set(now)
            self.vars["updated"].set(now)
        except Exception:
            pass

    def _load_draft(self):
        try:
            if self.draft_path.exists():
                data = json.loads(self.draft_path.read_text(encoding="utf-8"))
                if isinstance(data, dict):
                    for k in ["status", "tlp_classification", "tags", "owner", "created", "updated"]:
                        if k in self.vars and k in data:
                            try:
                                self.vars[k].set(data.get(k, ""))
                            except Exception:
                                pass
                    try:
                        self.vars["open_after_save"].set(bool(data.get("open_after_save", True)))
                    except Exception:
                        pass
        except Exception:
            pass

    def _submit(self):
        data = self._collect()
        if not data.get("title"):
            messagebox.showerror("How-To Intake", "Title is required.", parent=self)
            return
        howto_id_raw = (data.get("how_to_id") or "").strip()
        if not howto_id_raw:
            messagebox.showwarning("How-To Intake", "How-to ID is required.", parent=self)
            return
        if not HOWTO_ID_RE.match(howto_id_raw):
            messagebox.showwarning(
                "How-To Intake",
                "How-to ID must be exactly 5 digits (numerals only).",
                parent=self,
            )
            return
        self._save_draft()
        try:
            self.on_submit(data)
        finally:
            self.destroy()


class ThreatActorIntakeWindow(tk.Toplevel):
    def __init__(self, master, *, app_dir: Path, vault_root: Path | None, on_submit):
        super().__init__(master)
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("Threat Actor Intake")
        self.configure(bg="#111111")
        self.geometry("900x720")
        self.minsize(780, 600)

        self.draft_path = self.app_dir / "threat_actor_draft.json"

        self.vars = {
            "actor_id": tk.StringVar(),
            "actor_name": tk.StringVar(),
            "common_name": tk.StringVar(),
            "actor_type": tk.StringVar(value="Unknown"),
            "first_seen": tk.StringVar(),
            "last_seen": tk.StringVar(),
            "status": tk.StringVar(value="Unknown"),
            "attribution_confidence": tk.StringVar(value="low"),
            "internal_tracked": tk.BooleanVar(value=True),
            "open_after_save": tk.BooleanVar(value=True),
        }

        self._build_widgets()
        self._load_draft()

    def _build_widgets(self):
        container = tk.Frame(self, bg="#111111")
        container.pack(fill="both", expand=True)

        canvas = tk.Canvas(container, bg="#111111", highlightthickness=0)
        vbar = tk.Scrollbar(container, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vbar.set)
        vbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        content = tk.Frame(canvas, bg="#111111")
        win_id = canvas.create_window((0, 0), window=content, anchor="nw")

        def _sync(_=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfigure(win_id, width=canvas.winfo_width())

        content.bind("<Configure>", _sync)
        canvas.bind("<Configure>", _sync)
        canvas.bind_all("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1 * (e.delta / 120)), "units"))

        hdr = tk.Frame(content, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 10))
        tk.Label(hdr, text="New Threat Actor", font=("Segoe UI", 16, "bold"), fg="#FFFFFF", bg="#111111").pack(anchor="w")
        tk.Label(
            hdr,
            text="Capture threat actor metadata and create a CIPHER note.",
            font=("Segoe UI", 10),
            fg="#BBBBBB",
            bg="#111111",
        ).pack(anchor="w")

        form = tk.Frame(content, bg="#111111")
        form.pack(fill="both", expand=True, padx=16)
        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        right.pack(side="right", fill="both", expand=True, padx=(10, 0))

        def add_entry(parent, label, var, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_combo(parent, label, var, values, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        add_entry(left, "Actor ID (4 digits, numerals only)", self.vars["actor_id"], required=True)
        tk.Checkbutton(
            left,
            text="Internal, tracked actor",
            variable=self.vars["internal_tracked"],
            bg="#111111",
            fg="#DDDDDD",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#DDDDDD",
        ).pack(anchor="w", pady=(0, 6))

        add_entry(left, "Actor Name", self.vars["actor_name"], required=True)
        add_entry(left, "Common Name", self.vars["common_name"])
        add_combo(left, "Actor Type", self.vars["actor_type"], ["Unknown", "Nation/State", "State-Sponsored", "Cybercrime", "Hacktivist"])

        add_entry(right, "First Seen (YYYY-MM-DD)", self.vars["first_seen"])
        add_entry(right, "Last Seen (YYYY-MM-DD)", self.vars["last_seen"])
        add_combo(right, "Status", self.vars["status"], ["Active", "Active (reported)", "Revoked", "Unknown"])
        add_combo(right, "Attribution Confidence", self.vars["attribution_confidence"], ["high", "medium", "low"])

        bar = tk.Frame(self, bg="#111111")
        bar.pack(side="bottom", fill="x", padx=16, pady=(6, 14))

        # open_after_save stays True (hidden)

        def btn(t, cmd):
            return tk.Button(bar, text=t, command=cmd, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right", padx=(8, 0))
        btn("Save & Continue", self._save_continue).pack(side="right")

    def _collect(self) -> dict:
        return {
            "actor_id": (self.vars["actor_id"].get() or "").strip(),
            "actor_name": (self.vars["actor_name"].get() or "").strip(),
            "common_name": (self.vars["common_name"].get() or "").strip(),
            "actor_type": (self.vars["actor_type"].get() or "").strip(),
            "first_seen": (self.vars["first_seen"].get() or "").strip(),
            "last_seen": (self.vars["last_seen"].get() or "").strip(),
            "status": (self.vars["status"].get() or "").strip(),
            "attribution_confidence": (self.vars["attribution_confidence"].get() or "").strip(),
            "internal_untracked": not bool(self.vars["internal_tracked"].get()),
            "open_after_save": bool(self.vars["open_after_save"].get()),
        }

    def _validate(self, data: dict) -> tuple[bool, str]:
        actor_id = data.get("actor_id", "")
        internal_tracked = not bool(data.get("internal_untracked"))
        if actor_id:
            if internal_tracked:
                if not INTERNAL_ACTOR_ID_RE.match(actor_id):
                    return False, "Internal actor IDs must be TA-#### or ####."
            else:
                if not MITRE_ACTOR_ID_RE.match(actor_id):
                    return False, "Actor IDs must be 4 digits."
        return True, ""

    def _save_draft(self):
        try:
            self.draft_path.write_text(json.dumps(self._collect(), indent=2), encoding="utf-8")
            messagebox.showinfo("Threat Actor Intake", "Draft saved.", parent=self)
        except Exception as e:
            messagebox.showerror("Threat Actor Intake", f"Failed to save draft.\n\nError: {e}", parent=self)

    def _clear_draft(self):
        try:
            if self.draft_path.exists():
                self.draft_path.unlink()
        except Exception:
            pass
        try:
            for v in self.vars.values():
                if isinstance(v, tk.BooleanVar):
                    v.set(False)
                else:
                    v.set("")
            self.vars["internal_tracked"].set(True)
            self.vars["open_after_save"].set(True)
        except Exception:
            pass

    def _load_draft(self):
        try:
            if self.draft_path.exists():
                data = json.loads(self.draft_path.read_text(encoding="utf-8"))
                if isinstance(data, dict):
                    for k, val in data.items():
                        if k == "internal_untracked":
                            self.vars["internal_tracked"].set(not bool(val))
                        elif k in self.vars:
                            self.vars[k].set(val)
        except Exception:
            pass

    def _save_continue(self):
        data = self._collect()
        ok, msg = self._validate(data)
        if not ok:
            messagebox.showwarning("Threat Actor Intake", msg, parent=self)
            return

        data["actor_id"] = _format_threat_actor_id(
            data.get("actor_id", ""),
            internal_tracked=not bool(data.get("internal_untracked")),
        )
        data["updated"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not data.get("created"):
            data["created"] = data["updated"]

        try:
            if callable(self.on_submit):
                self.on_submit(data)
        finally:
            self.destroy()


class CampaignIntakeWindow(tk.Toplevel):
    def __init__(self, master, *, app_dir: Path, vault_root: Path | None, on_submit):
        super().__init__(master)
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("Campaign Intake")
        self.configure(bg="#111111")
        self.geometry("900x720")
        self.minsize(780, 600)

        self.draft_path = self.app_dir / "campaign_draft.json"

        self.vars = {
            "campaign_id": tk.StringVar(),
            "campaign_name": tk.StringVar(),
            "tlp_classification": tk.StringVar(value="TLP:CLEAR"),
            "attribution_confidence": tk.StringVar(value="low"),
            "admirality_source_reliability": tk.StringVar(),
            "admirality_information_credibility": tk.StringVar(),
            "first_observed": tk.StringVar(),
            "last_observed": tk.StringVar(),
            "risk_level": tk.StringVar(value="low"),
            "internal_tracked": tk.BooleanVar(value=False),
            "open_after_save": tk.BooleanVar(value=True),
        }

        self._build_widgets()
        self._load_draft()

    def _build_widgets(self):
        container = tk.Frame(self, bg="#111111")
        container.pack(fill="both", expand=True)

        canvas = tk.Canvas(container, bg="#111111", highlightthickness=0)
        vbar = tk.Scrollbar(container, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vbar.set)
        vbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        content = tk.Frame(canvas, bg="#111111")
        win_id = canvas.create_window((0, 0), window=content, anchor="nw")

        def _sync(_=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfigure(win_id, width=canvas.winfo_width())

        content.bind("<Configure>", _sync)
        canvas.bind("<Configure>", _sync)
        canvas.bind_all("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1 * (e.delta / 120)), "units"))

        hdr = tk.Frame(content, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 10))
        tk.Label(hdr, text="New Hacking Campaign", font=("Segoe UI", 16, "bold"), fg="#FFFFFF", bg="#111111").pack(anchor="w")
        tk.Label(
            hdr,
            text="Capture campaign metadata and create a CIPHER note.",
            font=("Segoe UI", 10),
            fg="#BBBBBB",
            bg="#111111",
        ).pack(anchor="w")

        form = tk.Frame(content, bg="#111111")
        form.pack(fill="both", expand=True, padx=16)
        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        right.pack(side="right", fill="both", expand=True, padx=(10, 0))

        def add_entry(parent, label, var):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_combo(parent, label, var, values):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        add_entry(left, "Campaign ID (numerals only ####)", self.vars["campaign_id"])
        tk.Checkbutton(
            left,
            text="Internal, tracked campaign",
            variable=self.vars["internal_tracked"],
            bg="#111111",
            fg="#DDDDDD",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#DDDDDD",
        ).pack(anchor="w", pady=(0, 6))
        add_entry(left, "Campaign Name", self.vars["campaign_name"])
        add_combo(left, "TLP Classification", self.vars["tlp_classification"], ["TLP:CLEAR", "TLP:GREEN", "TLP:AMBER", "TLP:AMBER+STRICT", "TLP:RED"])
        add_combo(left, "Attribution Confidence", self.vars["attribution_confidence"], ["high", "medium", "low"])
        add_combo(left, "Admirality Source Reliability", self.vars["admirality_source_reliability"], ["A", "B", "C", "D", "E", "F"])
        add_combo(left, "Admirality Information Credibility", self.vars["admirality_information_credibility"], ["1", "2", "3", "4", "5", "6"])

        add_entry(right, "First Observed (YYYY-MM-DD)", self.vars["first_observed"])
        add_entry(right, "Last Observed (YYYY-MM-DD)", self.vars["last_observed"])
        add_combo(right, "Risk Level", self.vars["risk_level"], ["low", "medium", "high", "critical"])

        bar = tk.Frame(self, bg="#111111")
        bar.pack(side="bottom", fill="x", padx=16, pady=(6, 14))

        tk.Checkbutton(
            bar,
            text="Open created note automatically",
            variable=self.vars["open_after_save"],
            bg="#111111",
            fg="#DDDDDD",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#DDDDDD",
        ).pack(side="left")

        def btn(t, cmd):
            return tk.Button(bar, text=t, command=cmd, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right", padx=(8, 0))
        btn("Save & Continue", self._save_continue).pack(side="right")
        btn("Clear Draft", self._clear_draft).pack(side="right", padx=(8, 0))
        btn("Save Draft", self._save_draft).pack(side="right", padx=(8, 0))

    def _collect(self) -> dict:
        return {
            "campaign_id": (self.vars["campaign_id"].get() or "").strip(),
            "campaign_name": (self.vars["campaign_name"].get() or "").strip(),
            "tlp_classification": (self.vars["tlp_classification"].get() or "").strip(),
            "attribution_confidence": (self.vars["attribution_confidence"].get() or "").strip(),
            "admirality_source_reliability": (self.vars["admirality_source_reliability"].get() or "").strip(),
            "admirality_information_credibility": (self.vars["admirality_information_credibility"].get() or "").strip(),
            "first_observed": (self.vars["first_observed"].get() or "").strip(),
            "last_observed": (self.vars["last_observed"].get() or "").strip(),
            "risk_level": (self.vars["risk_level"].get() or "").strip(),
            "internal_tracked": bool(self.vars["internal_tracked"].get()),
            "open_after_save": bool(self.vars["open_after_save"].get()),
        }

    def _save_draft(self):
        try:
            self.draft_path.write_text(json.dumps(self._collect(), indent=2), encoding="utf-8")
            messagebox.showinfo("Campaign Intake", "Draft saved.", parent=self)
        except Exception as e:
            messagebox.showerror("Campaign Intake", f"Failed to save draft.\n\nError: {e}", parent=self)

    def _clear_draft(self):
        try:
            if self.draft_path.exists():
                self.draft_path.unlink()
        except Exception:
            pass
        try:
            for v in self.vars.values():
                if isinstance(v, tk.BooleanVar):
                    v.set(False)
                else:
                    v.set("")
            self.vars["tlp_classification"].set("TLP:CLEAR")
            self.vars["attribution_confidence"].set("low")
            self.vars["admirality_source_reliability"].set("")
            self.vars["admirality_information_credibility"].set("")
            self.vars["risk_level"].set("low")
            self.vars["internal_tracked"].set(False)
            self.vars["open_after_save"].set(True)
        except Exception:
            pass

    def _load_draft(self):
        try:
            if self.draft_path.exists():
                data = json.loads(self.draft_path.read_text(encoding="utf-8"))
                if isinstance(data, dict):
                    for k, val in data.items():
                        if k in self.vars:
                            self.vars[k].set(val)
        except Exception:
            pass

    def _save_continue(self):
        data = self._collect()
        if data.get("campaign_id") and not CAMPAIGN_ID_RE.match(data.get("campaign_id", "")):
            messagebox.showwarning("Campaign Intake", "Campaign ID must be 4 digits.", parent=self)
            return

        data["campaign_id"] = _format_campaign_id(
            data.get("campaign_id", ""),
            internal_tracked=bool(data.get("internal_tracked")),
        )
        data["updated"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not data.get("created"):
            data["created"] = data["updated"]

        try:
            if callable(self.on_submit):
                self.on_submit(data)
        finally:
            self.destroy()


class MITRETTPIntakeWindow(tk.Toplevel):
    def __init__(self, master, *, app_dir: Path, vault_root: Path | None, on_submit):
        super().__init__(master)
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("MITRE ATT&CK TTP Intake")
        self.configure(bg="#111111")
        self.geometry("900x720")
        self.minsize(780, 600)

        self.draft_path = self.app_dir / "mitre_ttp_draft.json"

        self.vars = {
            "technique_id": tk.StringVar(),
            "subtechnique_id": tk.StringVar(),
            "technique_name": tk.StringVar(),
            "mitre_version": tk.StringVar(),
            "detection_priority": tk.StringVar(value="low"),
            "threat_score": tk.StringVar(value="5"),
            "open_after_save": tk.BooleanVar(value=True),
        }

        self._build_widgets()
        self._load_draft()

    def _build_widgets(self):
        container = tk.Frame(self, bg="#111111")
        container.pack(fill="both", expand=True)

        canvas = tk.Canvas(container, bg="#111111", highlightthickness=0)
        vbar = tk.Scrollbar(container, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vbar.set)
        vbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        content = tk.Frame(canvas, bg="#111111")
        win_id = canvas.create_window((0, 0), window=content, anchor="nw")

        def _sync(_=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfigure(win_id, width=canvas.winfo_width())

        content.bind("<Configure>", _sync)
        canvas.bind("<Configure>", _sync)
        canvas.bind_all("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1 * (e.delta / 120)), "units"))

        hdr = tk.Frame(content, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 10))
        tk.Label(hdr, text="New MITRE ATT&CK TTP", font=("Segoe UI", 16, "bold"), fg="#FFFFFF", bg="#111111").pack(anchor="w")
        tk.Label(
            hdr,
            text="Capture MITRE TTP metadata and create an entity note.",
            font=("Segoe UI", 10),
            fg="#BBBBBB",
            bg="#111111",
        ).pack(anchor="w")

        form = tk.Frame(content, bg="#111111")
        form.pack(fill="both", expand=True, padx=16)
        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        right.pack(side="right", fill="both", expand=True, padx=(10, 0))

        def add_entry(parent, label, var):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_combo(parent, label, var, values):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        add_entry(left, "Technique ID (numerals only ####)", self.vars["technique_id"])
        add_entry(left, "Subtechnique ID (numerals only ####.###)", self.vars["subtechnique_id"])
        add_entry(left, "Technique Name", self.vars["technique_name"])

        add_entry(right, "MITRE Version", self.vars["mitre_version"])
        add_combo(right, "Detection Priority", self.vars["detection_priority"], ["low", "medium", "high", "critical"])
        add_combo(right, "Threat Score", self.vars["threat_score"], [str(i) for i in range(1, 11)])

        bar = tk.Frame(self, bg="#111111")
        bar.pack(side="bottom", fill="x", padx=16, pady=(6, 14))

        tk.Checkbutton(
            bar,
            text="Open created note automatically",
            variable=self.vars["open_after_save"],
            bg="#111111",
            fg="#DDDDDD",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#DDDDDD",
        ).pack(side="left")

        def btn(t, cmd):
            return tk.Button(bar, text=t, command=cmd, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right", padx=(8, 0))
        btn("Save & Continue", self._save_continue).pack(side="right")
        btn("Clear Draft", self._clear_draft).pack(side="right", padx=(8, 0))
        btn("Save Draft", self._save_draft).pack(side="right", padx=(8, 0))

    def _collect(self) -> dict:
        return {
            "technique_id": (self.vars["technique_id"].get() or "").strip(),
            "subtechnique_id": (self.vars["subtechnique_id"].get() or "").strip(),
            "technique_name": (self.vars["technique_name"].get() or "").strip(),
            "mitre_version": (self.vars["mitre_version"].get() or "").strip(),
            "detection_priority": (self.vars["detection_priority"].get() or "").strip(),
            "threat_score": (self.vars["threat_score"].get() or "").strip(),
            "open_after_save": bool(self.vars["open_after_save"].get()),
        }

    def _save_draft(self):
        try:
            self.draft_path.write_text(json.dumps(self._collect(), indent=2), encoding="utf-8")
            messagebox.showinfo("MITRE TTP Intake", "Draft saved.", parent=self)
        except Exception as e:
            messagebox.showerror("MITRE TTP Intake", f"Failed to save draft.\n\nError: {e}", parent=self)

    def _clear_draft(self):
        try:
            if self.draft_path.exists():
                self.draft_path.unlink()
        except Exception:
            pass
        try:
            for v in self.vars.values():
                if isinstance(v, tk.BooleanVar):
                    v.set(False)
                else:
                    v.set("")
            self.vars["open_after_save"].set(True)
        except Exception:
            pass

    def _load_draft(self):
        try:
            if self.draft_path.exists():
                data = json.loads(self.draft_path.read_text(encoding="utf-8"))
                if isinstance(data, dict):
                    for k, val in data.items():
                        if k in self.vars:
                            self.vars[k].set(val)
        except Exception:
            pass

    def _save_continue(self):
        data = self._collect()
        raw_tid = (data.get("technique_id") or "").strip()
        raw_stid = (data.get("subtechnique_id") or "").strip()
        if raw_tid:
            tid_check = raw_tid[1:] if raw_tid.upper().startswith("T") else raw_tid
            if not TTP_ID_RE.match(tid_check):
                messagebox.showwarning("MITRE TTP Intake", "Technique ID must be 4 digits.", parent=self)
                return
        if raw_stid:
            stid_check = raw_stid[1:] if raw_stid.upper().startswith("T") else raw_stid
            if not SUBTTP_ID_RE.match(stid_check):
                messagebox.showwarning("MITRE TTP Intake", "Subtechnique ID must be ####.###.", parent=self)
                return

        data["technique_id"] = _format_ttp_id(raw_tid)
        data["subtechnique_id"] = _format_subttp_id(raw_stid)
        data["updated"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not data.get("created"):
            data["created"] = data["updated"]

        try:
            if callable(self.on_submit):
                self.on_submit(data)
        finally:
            self.destroy()


class MalwareIntakeWindow(tk.Toplevel):
    def __init__(self, master, *, app_dir: Path, vault_root: Path | None, on_submit):
        super().__init__(master)
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("Malware Intake")
        self.configure(bg="#111111")
        self.geometry("900x720")
        self.minsize(780, 600)

        self.draft_path = self.app_dir / "malware_draft.json"

        self.vars = {
            "malware_name": tk.StringVar(),
            "malware_id": tk.StringVar(),
            "malware_type": tk.StringVar(value="infostealer"),
            "category": tk.StringVar(value="commodity"),
            "first_seen": tk.StringVar(),
            "active_status": tk.StringVar(),
            "attribution_confidence": tk.StringVar(value="low"),
            "risk_level": tk.StringVar(value="low"),
            "threat_score": tk.StringVar(value="5"),
            "tlp_classification": tk.StringVar(value="TLP:CLEAR"),
            "internal_tracked": tk.BooleanVar(value=False),
            "open_after_save": tk.BooleanVar(value=True),
        }

        self._build_widgets()
        self._load_draft()

    def _build_widgets(self):
        container = tk.Frame(self, bg="#111111")
        container.pack(fill="both", expand=True)

        canvas = tk.Canvas(container, bg="#111111", highlightthickness=0)
        vbar = tk.Scrollbar(container, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vbar.set)
        vbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        content = tk.Frame(canvas, bg="#111111")
        win_id = canvas.create_window((0, 0), window=content, anchor="nw")

        def _sync(_=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfigure(win_id, width=canvas.winfo_width())

        content.bind("<Configure>", _sync)
        canvas.bind("<Configure>", _sync)
        canvas.bind_all("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1 * (e.delta / 120)), "units"))

        hdr = tk.Frame(content, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 10))
        tk.Label(hdr, text="New Malware Family / Strain", font=("Segoe UI", 16, "bold"), fg="#FFFFFF", bg="#111111").pack(anchor="w")
        tk.Label(
            hdr,
            text="Capture malware metadata and create a CIPHER note.",
            font=("Segoe UI", 10),
            fg="#BBBBBB",
            bg="#111111",
        ).pack(anchor="w")

        form = tk.Frame(content, bg="#111111")
        form.pack(fill="both", expand=True, padx=16)
        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        right.pack(side="right", fill="both", expand=True, padx=(10, 0))

        def add_entry(parent, label, var):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_combo(parent, label, var, values):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        add_entry(left, "Malware ID (numerals only ####)", self.vars["malware_id"])
        tk.Checkbutton(
            left,
            text="Internal, tracked malware strain",
            variable=self.vars["internal_tracked"],
            bg="#111111",
            fg="#DDDDDD",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#DDDDDD",
        ).pack(anchor="w", pady=(0, 6))
        add_entry(left, "Malware Name", self.vars["malware_name"])
        add_combo(
            left,
            "Malware Type",
            self.vars["malware_type"],
            ["ransomware", "loader", "backdoor", "trojan", "infostealer", "worm", "botnet", "RAT", "rootkit", "wiper", "dropper", "exploit_kit", "other"],
        )
        add_combo(
            left,
            "Category",
            self.vars["category"],
            ["commodity", "custom", "APT-exclusive", "crimeware", "open-source", "red-team tool"],
        )

        add_entry(right, "First Seen (YYYY-MM-DD)", self.vars["first_seen"])
        add_combo(right, "Active Status", self.vars["active_status"], ["active", "dormant", "unknown"])
        add_combo(right, "Attribution Confidence", self.vars["attribution_confidence"], ["high", "medium", "low"])
        add_combo(right, "Risk Level", self.vars["risk_level"], ["low", "medium", "high", "critical"])
        add_combo(right, "Threat Score", self.vars["threat_score"], [str(i) for i in range(1, 11)])
        add_combo(right, "TLP Classification", self.vars["tlp_classification"], ["TLP:CLEAR", "TLP:GREEN", "TLP:AMBER", "TLP:AMBER+STRICT", "TLP:RED"])

        bar = tk.Frame(self, bg="#111111")
        bar.pack(side="bottom", fill="x", padx=16, pady=(6, 14))

        tk.Checkbutton(
            bar,
            text="Open created note automatically",
            variable=self.vars["open_after_save"],
            bg="#111111",
            fg="#DDDDDD",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#DDDDDD",
        ).pack(side="left")

        def btn(t, cmd):
            return tk.Button(bar, text=t, command=cmd, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right", padx=(8, 0))
        btn("Save & Continue", self._save_continue).pack(side="right")
        btn("Clear Draft", self._clear_draft).pack(side="right", padx=(8, 0))
        btn("Save Draft", self._save_draft).pack(side="right", padx=(8, 0))

    def _collect(self) -> dict:
        return {
            "malware_name": (self.vars["malware_name"].get() or "").strip(),
            "malware_id": (self.vars["malware_id"].get() or "").strip(),
            "malware_type": (self.vars["malware_type"].get() or "").strip(),
            "category": (self.vars["category"].get() or "").strip(),
            "first_seen": (self.vars["first_seen"].get() or "").strip(),
            "active_status": (self.vars["active_status"].get() or "").strip(),
            "attribution_confidence": (self.vars["attribution_confidence"].get() or "").strip(),
            "risk_level": (self.vars["risk_level"].get() or "").strip(),
            "threat_score": (self.vars["threat_score"].get() or "").strip(),
            "tlp_classification": (self.vars["tlp_classification"].get() or "").strip(),
            "internal_tracked": bool(self.vars["internal_tracked"].get()),
            "open_after_save": bool(self.vars["open_after_save"].get()),
        }

    def _save_draft(self):
        try:
            self.draft_path.write_text(json.dumps(self._collect(), indent=2), encoding="utf-8")
            messagebox.showinfo("Malware Intake", "Draft saved.", parent=self)
        except Exception as e:
            messagebox.showerror("Malware Intake", f"Failed to save draft.\n\nError: {e}", parent=self)

    def _clear_draft(self):
        try:
            if self.draft_path.exists():
                self.draft_path.unlink()
        except Exception:
            pass
        try:
            for v in self.vars.values():
                if isinstance(v, tk.BooleanVar):
                    v.set(False)
                else:
                    v.set("")
            self.vars["malware_type"].set("infostealer")
            self.vars["category"].set("commodity")
            self.vars["attribution_confidence"].set("low")
            self.vars["risk_level"].set("low")
            self.vars["threat_score"].set("5")
            self.vars["tlp_classification"].set("TLP:CLEAR")
            self.vars["internal_tracked"].set(False)
            self.vars["open_after_save"].set(True)
        except Exception:
            pass

    def _load_draft(self):
        try:
            if self.draft_path.exists():
                data = json.loads(self.draft_path.read_text(encoding="utf-8"))
                if isinstance(data, dict):
                    for k, val in data.items():
                        if k in self.vars:
                            self.vars[k].set(val)
        except Exception:
            pass

    def _save_continue(self):
        data = self._collect()
        if data.get("malware_id") and not MALWARE_ID_RE.match(data.get("malware_id", "")):
            messagebox.showwarning("Malware Intake", "Malware ID must be 4 digits.", parent=self)
            return

        data["malware_id"] = _format_malware_id(
            data.get("malware_id", ""),
            internal_tracked=bool(data.get("internal_tracked")),
        )
        data["updated"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not data.get("created"):
            data["created"] = data["updated"]

        try:
            if callable(self.on_submit):
                self.on_submit(data)
        finally:
            self.destroy()


class IOCIntakeWindow(tk.Toplevel):
    def __init__(self, master, *, app_dir: Path, vault_root: Path | None, on_submit):
        super().__init__(master)
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("IOC Intake")
        self.configure(bg="#111111")
        self.geometry("900x720")
        self.minsize(780, 600)

        self.draft_path = self.app_dir / "ioc_draft.json"

        self.vars = {
            "ioc_type": tk.StringVar(value="domain"),
            "value": tk.StringVar(),
            "classification": tk.StringVar(),
            "confidence": tk.StringVar(),
            "last_seen": tk.StringVar(),
            "open_after_save": tk.BooleanVar(value=True),
        }

        self._build_widgets()
        self._load_draft()

    def _build_widgets(self):
        container = tk.Frame(self, bg="#111111")
        container.pack(fill="both", expand=True)

        canvas = tk.Canvas(container, bg="#111111", highlightthickness=0)
        vbar = tk.Scrollbar(container, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vbar.set)
        vbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        content = tk.Frame(canvas, bg="#111111")
        win_id = canvas.create_window((0, 0), window=content, anchor="nw")

        def _sync(_=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfigure(win_id, width=canvas.winfo_width())

        content.bind("<Configure>", _sync)
        canvas.bind("<Configure>", _sync)
        canvas.bind_all("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1 * (e.delta / 120)), "units"))

        hdr = tk.Frame(content, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 10))
        tk.Label(hdr, text="New Indicator of Compromise (IOC)", font=("Segoe UI", 16, "bold"), fg="#FFFFFF", bg="#111111").pack(anchor="w")
        tk.Label(
            hdr,
            text="Capture IOC metadata and create an entity note.",
            font=("Segoe UI", 10),
            fg="#BBBBBB",
            bg="#111111",
        ).pack(anchor="w")

        form = tk.Frame(content, bg="#111111")
        form.pack(fill="both", expand=True, padx=16)
        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        right.pack(side="right", fill="both", expand=True, padx=(10, 0))

        def add_entry(parent, label, var):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_combo(parent, label, var, values):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        add_combo(left, "IOC Type", self.vars["ioc_type"], ["domain", "email_address", "file_path", "hash", "ip", "registry_key", "url"])
        add_combo(left, "Classification", self.vars["classification"], ["malicious", "suspicious", "benign", "unknown"])
        add_combo(left, "Confidence", self.vars["confidence"], ["low", "medium", "high"])

        add_entry(right, "Value", self.vars["value"])
        add_entry(right, "Last Seen (YYYY-MM-DD)", self.vars["last_seen"])
        tk.Label(right, text="Notes", fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w", pady=(6, 0))
        self.notes_txt = tk.Text(
            right,
            height=6,
            bg="#1B1B1B",
            fg="#FFFFFF",
            insertbackground="#FFFFFF",
            relief="flat",
            wrap="word",
        )
        self.notes_txt.pack(fill="x", pady=(4, 0))

        bar = tk.Frame(self, bg="#111111")
        bar.pack(side="bottom", fill="x", padx=16, pady=(6, 14))

        tk.Checkbutton(
            bar,
            text="Open created note automatically",
            variable=self.vars["open_after_save"],
            bg="#111111",
            fg="#DDDDDD",
            selectcolor="#111111",
            activebackground="#111111",
            activeforeground="#DDDDDD",
        ).pack(side="left")

        def btn(t, cmd):
            return tk.Button(bar, text=t, command=cmd, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right", padx=(8, 0))
        btn("Save & Continue", self._save_continue).pack(side="right")
        btn("Clear Draft", self._clear_draft).pack(side="right", padx=(8, 0))
        btn("Save Draft", self._save_draft).pack(side="right", padx=(8, 0))

    def _collect(self) -> dict:
        return {
            "ioc_type": (self.vars["ioc_type"].get() or "").strip(),
            "value": (self.vars["value"].get() or "").strip(),
            "classification": (self.vars["classification"].get() or "").strip(),
            "confidence": (self.vars["confidence"].get() or "").strip(),
            "last_seen": (self.vars["last_seen"].get() or "").strip(),
            "notes": self.notes_txt.get("1.0", "end").strip(),
            "open_after_save": bool(self.vars["open_after_save"].get()),
        }

    def _save_draft(self):
        try:
            data = self._collect()
            self.draft_path.write_text(json.dumps(data, indent=2), encoding="utf-8")
            messagebox.showinfo("IOC Intake", "Draft saved.", parent=self)
        except Exception as e:
            messagebox.showerror("IOC Intake", f"Failed to save draft.\n\nError: {e}", parent=self)

    def _clear_draft(self):
        try:
            if self.draft_path.exists():
                self.draft_path.unlink()
        except Exception:
            pass
        try:
            for v in self.vars.values():
                if isinstance(v, tk.BooleanVar):
                    v.set(False)
                else:
                    v.set("")
            self.vars["ioc_type"].set("domain")
            self.vars["open_after_save"].set(True)
            self.notes_txt.delete("1.0", "end")
        except Exception:
            pass

    def _load_draft(self):
        try:
            if self.draft_path.exists():
                data = json.loads(self.draft_path.read_text(encoding="utf-8"))
                if isinstance(data, dict):
                    for k, val in data.items():
                        if k in self.vars:
                            self.vars[k].set(val)
                    if "notes" in data:
                        self.notes_txt.delete("1.0", "end")
                        self.notes_txt.insert("1.0", data.get("notes", ""))
        except Exception:
            pass

    def _save_continue(self):
        data = self._collect()
        data["updated"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not data.get("created"):
            data["created"] = data["updated"]

        try:
            if callable(self.on_submit):
                self.on_submit(data)
        finally:
            self.destroy()


class SLAIntakeWindow(tk.Toplevel):
    """
    SLA Intake UI (consistent with other intake windows).
    """

    def __init__(self, master, *, app_dir: Path, vault_root: Path | None, on_submit, resume: bool = False):
        super().__init__(master)
        self.app_dir = app_dir
        self.vault_root = vault_root
        self.on_submit = on_submit

        self.title("SLA Intake")
        self.configure(bg="#111111")
        self.geometry("900x720")
        self.minsize(780, 600)
        self.resizable(True, True)

        self._build_widgets()

        if resume:
            self._apply_draft(load_sla_draft(self.app_dir))
        else:
            # prefill created/updated
            now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            self.vars["created"].set(now)
            self.vars["updated"].set(now)

        self.transient(master)
        self.grab_set()
        self.focus_force()

    def _build_widgets(self):
        cfg = getattr(self.master, "cfg", {}) or {}
        analyst_default = (cfg.get("analyst_profile", {}) or {}).get("name", "")

        container = tk.Frame(self, bg="#111111")
        container.pack(fill="both", expand=True)

        canvas = tk.Canvas(container, bg="#111111", highlightthickness=0)
        vbar = tk.Scrollbar(container, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vbar.set)
        vbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        content = tk.Frame(canvas, bg="#111111")
        win_id = canvas.create_window((0, 0), window=content, anchor="nw")

        def _sync(_=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfigure(win_id, width=canvas.winfo_width())

        content.bind("<Configure>", _sync)
        canvas.bind("<Configure>", _sync)
        canvas.bind_all("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1 * (e.delta / 120)), "units"))

        hdr = tk.Frame(content, bg="#111111")
        hdr.pack(fill="x", padx=16, pady=(14, 10))
        tk.Label(hdr, text="Service Level Agreement Intake", font=("Segoe UI", 16, "bold"), fg="#FFFFFF", bg="#111111").pack(anchor="w")
        tk.Label(hdr, text="Capture SLA metadata and generate an Obsidian note.", font=("Segoe UI", 10),
                 fg="#BBBBBB", bg="#111111").pack(anchor="w")

        form = tk.Frame(content, bg="#111111")
        form.pack(fill="both", expand=True, padx=16)
        left = tk.Frame(form, bg="#111111")
        right = tk.Frame(form, bg="#111111")
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        right.pack(side="right", fill="both", expand=True, padx=(10, 0))

        self.vars = {
            "sla_id": tk.StringVar(),
            "sla_title": tk.StringVar(),
            "sla_status": tk.StringVar(value="Draft"),
            "sla_type": tk.StringVar(value="Internal"),
            "sla_owner_primary": tk.StringVar(value=analyst_default),
            "sla_owning_team": tk.StringVar(),
            "sla_provider_contact": tk.StringVar(),
            "sla_customer_contact": tk.StringVar(),
            "sla_effective_date": tk.StringVar(),
            "sla_expiration_date": tk.StringVar(),
            "tags": tk.StringVar(),
            "tlp_classification": tk.StringVar(value="TLP:CLEAR"),
            "created": tk.StringVar(),
            "updated": tk.StringVar(),
        }

        self._date_pickers = {}

        def add_date_picker(parent, label: str, key: str, var: tk.StringVar):
            """Date-only picker; stores as YYYY-MM-DD 00:00:00."""
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            tk.Label(row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w")
            picker_row = tk.Frame(row, bg="#111111")
            picker_row.pack(fill="x")
            value = (var.get() or "").strip()
            d, _, _, _ = _parse_datetime_for_picker(value if value else "")
            date_entry = None
            fallback_entry = None
            if DateEntry is not None:
                try:
                    style = ttk.Style(self)
                    style.configure(
                        "SLA.DateEntry",
                        fieldbackground="#1B1B1B",
                        background="#2A2A2A",
                        foreground="#FFFFFF",
                        arrowcolor="#FFFFFF",
                    )
                    date_entry = DateEntry(
                        picker_row, date_pattern="y-mm-dd", year=d.year, month=d.month, day=d.day,
                        style="SLA.DateEntry",
                    )
                except Exception:
                    date_entry = DateEntry(
                        picker_row, date_pattern="y-mm-dd", year=d.year, month=d.month, day=d.day,
                    )
                date_entry.pack(side="left", padx=(0, 8))
            else:
                fallback_entry = tk.Entry(picker_row, bg="#1B1B1B", fg="#FFFFFF", relief="flat", width=12)
                fallback_entry.insert(0, d.isoformat())
                fallback_entry.pack(side="left", padx=(0, 8))
            self._date_pickers[key] = {"date_entry": date_entry, "fallback_entry": fallback_entry}
            return date_entry, fallback_entry

        def _sync_expiration_from_effective():
            """Set Expiration date to Effective date + 1 year."""
            eff_p = self._date_pickers.get("sla_effective_date")
            exp_p = self._date_pickers.get("sla_expiration_date")
            if not eff_p or not exp_p:
                return
            try:
                if eff_p["date_entry"] is not None:
                    d = eff_p["date_entry"].get_date()
                elif eff_p.get("fallback_entry"):
                    raw = eff_p["fallback_entry"].get().strip()
                    if not raw:
                        return
                    d = datetime.strptime(raw, "%Y-%m-%d").date()
                else:
                    return
                exp_d = d + timedelta(days=365)
                if exp_p["date_entry"] is not None:
                    exp_p["date_entry"].set_date(exp_d)
                elif exp_p.get("fallback_entry"):
                    exp_p["fallback_entry"].delete(0, "end")
                    exp_p["fallback_entry"].insert(0, exp_d.isoformat())
            except (ValueError, TypeError):
                pass

        def add_entry(parent, label, var, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            ent = tk.Entry(row, textvariable=var, bg="#1B1B1B", fg="#FFFFFF",
                           insertbackground="#FFFFFF", relief="flat")
            ent.pack(fill="x", ipady=6)
            return ent

        def add_combo(parent, label, var, values, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        def add_labeled_combo(parent, label, var, values, required=False):
            row = tk.Frame(parent, bg="#111111")
            row.pack(fill="x", pady=6)
            lbl_row = tk.Frame(row, bg="#111111")
            lbl_row.pack(anchor="w")
            tk.Label(lbl_row, text=label, fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            if required:
                tk.Label(lbl_row, text=" *", fg="#E05A5A", bg="#111111", font=("Segoe UI", 10)).pack(side="left")
            cb = ttk.Combobox(row, textvariable=var, values=values, state="readonly")
            cb.pack(fill="x", ipady=3)
            return cb

        add_entry(left, "SLA ID (5 digits, numerals only)", self.vars["sla_id"])
        add_entry(left, "SLA Title", self.vars["sla_title"], required=True)
        add_combo(left, "Status", self.vars["sla_status"], ["Draft", "Active", "Suspended", "Retired"])
        add_combo(left, "SLA Type", self.vars["sla_type"], ["Internal", "Customer", "Vendor", "Regulatory", "Other"])
        eff_de, eff_fb = add_date_picker(left, "Effective Date", "sla_effective_date", self.vars["sla_effective_date"])
        add_date_picker(right, "Expiration Date", "sla_expiration_date", self.vars["sla_expiration_date"])
        _sync_expiration_from_effective()
        if eff_de is not None:
            eff_de.bind("<<DateEntrySelected>>", lambda e: _sync_expiration_from_effective())
        elif eff_fb is not None:
            eff_fb.bind("<FocusOut>", lambda e: _sync_expiration_from_effective())
        add_entry(right, "SLA Owner", self.vars["sla_owner_primary"])
        add_entry(right, "Team Ownership", self.vars["sla_owning_team"])
        add_entry(right, "SLA Contact Name", self.vars["sla_provider_contact"])
        add_entry(right, "SLA Customer Name", self.vars["sla_customer_contact"])
        add_entry(left, "Tags", self.vars["tags"])
        add_combo(left, "TLP Classification", self.vars["tlp_classification"], ["TLP:CLEAR", "TLP:GREEN", "TLP:AMBER", "TLP:AMBER+STRICT", "TLP:RED"])
        add_entry(right, "Created Date", self.vars["created"])
        add_entry(right, "Updated Date", self.vars["updated"])

        tk.Label(content, text="Description", fg="#DDDDDD", bg="#111111", font=("Segoe UI", 10)).pack(anchor="w", padx=16, pady=(12, 0))
        self.service_description_txt = tk.Text(content, height=8, bg="#1B1B1B", fg="#FFFFFF", insertbackground="#FFFFFF", relief="flat", wrap="word")
        self.service_description_txt.pack(fill="x", padx=16, pady=(4, 10))

        bar = tk.Frame(self, bg="#111111")
        bar.pack(side="bottom", fill="x", padx=16, pady=(6, 14))

        self.open_note_var = tk.BooleanVar(value=True)  # Hidden; always True

        def btn(t, cmd):
            return tk.Button(bar, text=t, command=cmd, bg="#2A2A2A", fg="#FFFFFF", relief="flat", padx=14, pady=8)

        btn("Cancel", self.destroy).pack(side="right", padx=(8, 0))
        btn("Save & Continue", self._save_continue).pack(side="right")

    def _get_date_picker_value(self, key: str) -> str:
        """Get date from picker as YYYY-MM-DD 00:00:00."""
        p = self._date_pickers.get(key)
        if not p:
            return self.vars.get(key, tk.StringVar()).get().strip()
        try:
            if p["date_entry"] is not None:
                d = p["date_entry"].get_date()
            elif p.get("fallback_entry"):
                d = datetime.strptime(p["fallback_entry"].get().strip(), "%Y-%m-%d").date()
            else:
                return self.vars.get(key, tk.StringVar()).get().strip()
            return f"{d.isoformat()} 00:00:00"
        except (ValueError, TypeError):
            return self.vars.get(key, tk.StringVar()).get().strip()

    def _set_date_picker_value(self, key: str, value: str) -> None:
        p = self._date_pickers.get(key)
        if not p:
            self.vars[key].set(value)
            return
        d, _, _, _ = _parse_datetime_for_picker(value)
        if p["date_entry"] is not None:
            p["date_entry"].set_date(d)
        elif p.get("fallback_entry"):
            p["fallback_entry"].delete(0, "end")
            p["fallback_entry"].insert(0, d.isoformat())

    def _apply_draft(self, draft: dict):
        if not draft:
            return
        for k, v in draft.items():
            if k in self.vars and isinstance(v, str):
                self.vars[k].set(v)
        for key in ("sla_effective_date", "sla_expiration_date"):
            if key in draft and key in self._date_pickers:
                self._set_date_picker_value(key, str(draft[key]))
        if "sla_service_description" in draft:
            self.service_description_txt.delete("1.0", "end")
            self.service_description_txt.insert("1.0", draft.get("sla_service_description", ""))

    def _collect(self) -> dict:
        data = {k: v.get().strip() for k, v in self.vars.items()}
        for key in ("sla_effective_date", "sla_expiration_date"):
            if key in self._date_pickers:
                data[key] = self._get_date_picker_value(key)
        data["sla_service_description"] = self.service_description_txt.get("1.0", "end").strip()
        return data

    def _save_draft(self):
        save_sla_draft(self.app_dir, self._collect())
        messagebox.showinfo("SLA Intake", "Draft saved.", parent=self)

    def _clear_draft(self):
        clear_sla_draft(self.app_dir)
        for v in self.vars.values():
            v.set("")
        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        self.vars["created"].set(now)
        self.vars["updated"].set(now)
        self.service_description_txt.delete("1.0", "end")
        for key in ("sla_effective_date", "sla_expiration_date"):
            if key in self._date_pickers:
                self._set_date_picker_value(key, "")
        messagebox.showinfo("SLA Intake", "Draft cleared.", parent=self)

    def _save_continue(self):
        data = self._collect()
        if not data.get("sla_title"):
            messagebox.showwarning("SLA Intake", "SLA Title is required.", parent=self)
            return
        sla_id_raw = (data.get("sla_id") or "").strip()
        if sla_id_raw and not SLA_ID_RE.match(sla_id_raw):
            messagebox.showwarning(
                "SLA Intake",
                "SLA ID must be exactly 5 digits (numerals only).",
                parent=self,
            )
            return

        eff = (data.get("sla_effective_date") or "").strip().split()[0] if (data.get("sla_effective_date") or "").strip() else ""
        exp = (data.get("sla_expiration_date") or "").strip().split()[0] if (data.get("sla_expiration_date") or "").strip() else ""
        if eff and exp:
            try:
                eff_d = datetime.strptime(eff, "%Y-%m-%d").date()
                exp_d = datetime.strptime(exp, "%Y-%m-%d").date()
                if eff_d >= exp_d:
                    messagebox.showwarning(
                        "SLA Intake",
                        "Effective Date must be earlier than Expiration Date.",
                        parent=self,
                    )
                    return
            except ValueError:
                pass

        # Normalize timestamps
        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not data.get("created"):
            data["created"] = now
        if len(data.get("created", "")) == 10:
            data["created"] = data["created"] + " 00:00:00"
        data["updated"] = now  # always stamp updated

        save_sla_draft(self.app_dir, data)

        note_path = None
        if self.vault_root is not None:
            try:
                note_path = create_obsidian_sla_note(self.vault_root, data)
            except Exception as e:
                messagebox.showwarning("SLA Intake", f"Saved draft, but failed to create SLA note:\n{e}", parent=self)
            try:
                _log_shift_entry_from_launcher(self.master, self.vault_root, "sla", data, note_path)
            except Exception:
                pass

        if note_path is not None and self.open_note_var.get():
            try:
                vault_name = (getattr(getattr(self, "master", None), "cfg", {}) or {}).get("paths", {}).get("OBSIDIAN_VAULT", "")
                open_note_in_obsidian(vault_name, self.vault_root, note_path)
            except Exception:
                pass

        try:
            self.on_submit({**data, **({"obsidian_note_path": str(note_path)} if note_path else {})})
        finally:
            self.destroy()



def main():
    app_dir = get_app_dir()
    try:
        cfg = load_config(app_dir)
    except Exception as e:
        messagebox.showerror("Configuration error", str(e))
        return

    tokens = build_tokens(app_dir, cfg)

    state_path = get_state_path("SCOUT_Launcher")
    app = LauncherApp(cfg, tokens, state_path)
    app.mainloop()


if __name__ == "__main__":
    main()